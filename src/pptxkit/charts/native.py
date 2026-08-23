"""The chart renderer: a real OOXML chart part with an embedded worksheet."""

from __future__ import annotations

from typing import TYPE_CHECKING, cast

from pptx.chart.chart import Chart
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.chart.point import Point
from pptx.chart.series import LineSeries, RadarSeries, XySeries
from pptx.dml.color import RGBColor
from pptx.shapes.graphfrm import GraphicFrame
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from pptxkit.charts._effects import apply_shadow
from pptxkit.charts._native_types import (
    _AXIS_CHART_TYPES,
    _CHART_TYPES,
    _CONNECTED_CHART_TYPES,
    _GAP_WIDTH_CHART_TYPES,
    _HORIZONTAL_BAR_CHART_TYPES,
    _LABEL_POSITIONS,
    _MARKER_CHART_TYPES,
    _MARKER_STYLES,
    _NO_DATA_LABEL_CHART_TYPES,
    _PERCENT_AXIS_CHART_TYPES,
    _PIE_FAMILY_CHART_TYPES,
    _SERIES_FILL_CHART_TYPES,
    _STROKE_CHART_TYPES,
    _STRUCTURAL_GRIDLINE_CHART_TYPES,
)
from pptxkit.charts._shared import label_number_format, lighten
from pptxkit.charts.model import _BUBBLE_CHART_TYPES, _XY_CHART_TYPES, ChartSpec
from pptxkit.errors import LayoutError, ThemeError
from pptxkit.theme.chartstyle import ChartStyle
from pptxkit.theme.model import Rect
from pptxkit.utils.text import text_em

if TYPE_CHECKING:
    from pptxkit.layouts.registry import SlideCtx


def add_native_chart(ctx: SlideCtx, spec: ChartSpec, rect: Rect) -> GraphicFrame:
    """Build a chart part for ``spec`` inside ``rect`` on ``ctx.slide``, styled from the theme.

    Args:
        ctx: The slide context; supplies the slide to draw on and the theme to style from.
        spec: The validated chart body to render.
        rect: Where the chart frame lands, in inches.

    Returns:
        The ``GraphicFrame`` shape containing the chart; ``.chart`` reaches the chart itself.

    Raises:
        LayoutError: ``spec.type`` has no creatable native chart type.
    """
    try:
        chart_type = _CHART_TYPES[spec.type]
    except KeyError:
        raise LayoutError(
            f"slide {ctx.spec.index}: chart type {spec.type!r} has no native chart mapping; "
            f"known types: {', '.join(_CHART_TYPES)}"
        ) from None

    chart_data = _build_chart_data(spec)

    frame = ctx.slide.shapes.add_chart(
        chart_type,
        Inches(rect.left),
        Inches(rect.top),
        Inches(rect.width),
        Inches(rect.height),
        chart_data,
    )
    chart = frame.chart
    _style_series(ctx, chart, spec)
    if spec.type not in _NO_DATA_LABEL_CHART_TYPES:
        _style_data_labels(ctx, chart, spec)
    if spec.type in _AXIS_CHART_TYPES:
        _style_axes(ctx, chart, spec)
    if spec.type in _GAP_WIDTH_CHART_TYPES:
        chart.plots[0].gap_width = ctx.theme.chart.gap_width
    if spec.type in _SIDE_LABEL_CHART_TYPES:
        _reserve_label_column(
            chart, spec, rect=rect, size_pt=ctx.style("caption").size, face=ctx.theme.face
        )
    chart.has_legend = len(spec.series) > 1
    _style_text(ctx, chart)
    return frame


def _build_chart_data(spec: ChartSpec) -> CategoryChartData | XyChartData | BubbleChartData:
    """The embedded worksheet, shaped to match the spec: one value per category,
    or (x, y)/(x, y, size) points."""
    if spec.type in _XY_CHART_TYPES:
        xy_data = XyChartData()
        for series in spec.series:
            xy_series = xy_data.add_series(series.name)
            for x, y in cast("tuple[tuple[float, float], ...]", series.points):
                xy_series.add_data_point(x, y)
        return xy_data
    if spec.type in _BUBBLE_CHART_TYPES:
        bubble_data = BubbleChartData()
        for series in spec.series:
            bubble_series = bubble_data.add_series(series.name)
            for x, y, size in cast("tuple[tuple[float, float, float], ...]", series.points):
                bubble_series.add_data_point(x, y, size)
        return bubble_data
    category_data = CategoryChartData()
    category_data.categories = spec.categories
    for series in spec.series:
        category_data.add_series(series.name, series.values)
    return category_data


def _series_colors(ctx: SlideCtx) -> tuple[RGBColor, ...]:
    """The palette's accent ramp, resolved to colours in cycle order.

    Raises:
        ThemeError: the palette declares no accent roles, so there is nothing to
            colour a series with.
    """
    palette = ctx.theme.palette
    accents = tuple(RGBColor.from_string(palette.role(name)) for name in palette.accents)
    if not accents:
        raise ThemeError(f"theme {ctx.theme.name!r} declares no accent roles")
    return accents


def _highlight_color(ctx: SlideCtx) -> RGBColor:
    """The colour that marks ``highlight:`` — the second accent.

    Raises:
        ThemeError: the palette has fewer than two accents, so the marked point
            would be painted the same colour as its neighbours.
    """
    accents = _series_colors(ctx)
    if len(accents) < 2:
        raise ThemeError(
            f"theme {ctx.theme.name!r} declares {len(accents)} accent role(s); "
            f"'highlight' marks a point with the second accent, so a palette with "
            f"fewer cannot show one"
        )
    return accents[1]


def _style_series(ctx: SlideCtx, chart: Chart, spec: ChartSpec) -> None:
    """Cycle the theme's palette by point (pie family) or by series (else); highlight wins.

    Series-level fill stays solid regardless of ``theme.chart.gradient``: it is what a
    legend swatch reads, and a swatch has no point to gradient.
    """
    palette = _series_colors(ctx)
    highlight = _highlight_color(ctx) if spec.highlight is not None else None
    style = ctx.theme.chart
    solid_only = spec.type in _STROKE_CHART_TYPES
    angle = (
        (style.gradient_angle - 90) % 360
        if spec.type in _HORIZONTAL_BAR_CHART_TYPES
        else style.gradient_angle
    )

    if spec.type in _PIE_FAMILY_CHART_TYPES:
        for index, point in enumerate(chart.series[0].points):
            colour = (
                highlight
                if highlight is not None and index == spec.highlight
                else palette[index % len(palette)]
            )
            _fill_point(point, colour, style, angle=angle, solid_only=solid_only)
        return

    for series_index, series in enumerate(chart.series):
        colour = palette[series_index % len(palette)]
        if spec.type in _SERIES_FILL_CHART_TYPES:
            _fill_series(series, colour, style, angle=angle)
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colour
        if spec.type in _CONNECTED_CHART_TYPES:
            series.format.line.color.rgb = colour
            series.format.line.width = Pt(ctx.theme.line_weight)
        if spec.type in _MARKER_CHART_TYPES:
            _style_marker(series, colour, style)
        if spec.type in _STROKE_CHART_TYPES:
            # A point here has no fillable shape, and a bare `c:dPt` fill on a stroke
            # series is what LibreOffice misassigns to the neighbouring series' marks.
            if highlight is not None and spec.type in _MARKER_CHART_TYPES:
                for index, point in enumerate(series.points):
                    if index == spec.highlight:
                        point.marker.format.fill.solid()
                        point.marker.format.fill.fore_color.rgb = highlight
            continue
        # On one series, emphasis is isolation: the marked point keeps the accent and
        # everything else goes muted, so it is the only saturated thing on the plot. A
        # second hue reads as a second category. Multi-series keeps the second accent —
        # muting other points there would erase the series distinction.
        mute = ctx.color("muted") if highlight is not None and len(spec.series) == 1 else None
        for index, point in enumerate(series.points):
            if mute is not None:
                point_colour = colour if index == spec.highlight else mute
            else:
                point_colour = (
                    highlight if highlight is not None and index == spec.highlight else colour
                )
            _fill_point(point, point_colour, style, angle=angle, solid_only=solid_only)


def _style_marker(
    series: LineSeries | RadarSeries | XySeries, colour: RGBColor, style: ChartStyle
) -> None:
    """Theme a line/radar/xy-scatter series' marker to match its own line colour."""
    marker = series.marker
    marker.style = _MARKER_STYLES[style.marker_style]
    marker.size = style.marker_size
    marker.format.fill.solid()
    marker.format.fill.fore_color.rgb = colour


def _fill_series(series, colour: RGBColor, style: ChartStyle, *, angle: float) -> None:
    """Gradient or solid fill on the series itself, plus its shadow."""
    fill = series.format.fill
    if style.gradient:
        fill.gradient()
        stops = fill.gradient_stops
        stops[0].color.rgb = colour
        stops[1].color.rgb = lighten(colour)
        fill.gradient_angle = angle
    else:
        fill.solid()
        fill.fore_color.rgb = colour
    if style.shadow:
        apply_shadow(series.format.element.get_or_add_spPr(), style)


def _fill_point(
    point: Point, colour: RGBColor, style: ChartStyle, *, angle: float, solid_only: bool = False
) -> None:
    """Solid or gradient fill for one data point, plus its drop shadow.

    Every point's fill is set explicitly, even when unchanged: an unset point's fill
    raises on read rather than reading back its series colour.
    """
    fill = point.format.fill
    if style.gradient and not solid_only:
        fill.gradient()
        stops = fill.gradient_stops
        stops[0].color.rgb = colour
        stops[1].color.rgb = lighten(colour)
        fill.gradient_angle = angle
    else:
        fill.solid()
        fill.fore_color.rgb = colour
    if style.shadow:
        apply_shadow(point.format.element.get_or_add_spPr(), style)


def _style_text(ctx: SlideCtx, chart: Chart) -> None:
    """Chart-wide text defaults, and the legend.

    The legend and the title PowerPoint generates for a single-series chart carry no
    colour of their own, so without this they take the presentation theme's dark ink
    and vanish on a dark slide.
    """
    chart.font.name = ctx.theme.face
    chart.font.color.rgb = ctx.fg()
    if not chart.has_legend:
        return
    chart.legend.font.name = ctx.theme.face
    chart.legend.font.size = Pt(ctx.style("caption").size)
    chart.legend.font.color.rgb = ctx.fg()


def _style_data_labels(ctx: SlideCtx, chart: Chart, spec: ChartSpec) -> None:
    style = ctx.theme.chart
    plot = chart.plots[0]
    if style.label_position == "none":
        plot.has_data_labels = False
        return

    plot.has_data_labels = True
    # The value is the point of the chart, so it outranks the axis scale that frames it.
    value_style = ctx.style("kicker")
    labels = plot.data_labels
    labels.font.size = Pt(value_style.size)
    labels.font.bold = True
    labels.font.italic = value_style.italic
    labels.font.name = ctx.theme.face
    labels.font.color.rgb = ctx.fg()
    if style.label_position in _LABEL_POSITIONS:
        labels.position = _LABEL_POSITIONS[style.label_position]
    # A 100% type is already a percentage, so a series unit of "%" would print a second sign.
    unit = None if spec.type in _PERCENT_AXIS_CHART_TYPES else spec.series[0].unit
    number_format = label_number_format(unit, thousands_sep=style.thousands_sep)
    if number_format is not None:
        labels.number_format = number_format
    if spec.type in _PIE_FAMILY_CHART_TYPES:
        # Colour alone no longer names a wedge, and a pie/doughnut's single series
        # never gets a legend — the label is the only thing that can name a slice.
        labels.show_category_name = True


def _style_axes(ctx: SlideCtx, chart: Chart, spec: ChartSpec) -> None:
    rule = ctx.color("line")
    show_grid = ctx.theme.chart.grid == "horizontal"
    caption = ctx.style("caption")
    for axis in (chart.category_axis, chart.value_axis):
        axis.format.line.color.rgb = rule
        axis.has_minor_gridlines = False
        # Untouched, tick labels inherit the template's size and dominate the values.
        axis.tick_labels.font.size = Pt(caption.size)
        axis.tick_labels.font.name = ctx.theme.face
        axis.tick_labels.font.color.rgb = ctx.dim()
    if spec.type not in _STRUCTURAL_GRIDLINE_CHART_TYPES:
        chart.category_axis.has_major_gridlines = False
        if show_grid:
            # .major_gridlines.format adds the element as a side effect (python-pptx's
            # get_or_add) — has_major_gridlines is never assigned True explicitly.
            chart.value_axis.major_gridlines.format.line.color.rgb = rule
        else:
            chart.value_axis.has_major_gridlines = False
    if spec.type in _PERCENT_AXIS_CHART_TYPES:
        chart.value_axis.tick_labels.number_format = "0%"
        chart.value_axis.tick_labels.number_format_is_linked = False
    if spec.y_min is not None:
        chart.value_axis.minimum_scale = spec.y_min
    if spec.y_max is not None:
        chart.value_axis.maximum_scale = spec.y_max


# Nothing in the file says how much room left-hand category labels get, so renderers
# disagree — Keynote runs them off the slide. A manual plot-area layout settles it.
_SIDE_LABEL_CHART_TYPES = frozenset({"bar", "bar-stacked", "bar-stacked-100"})
_LABEL_PAD_EM = 1.2  # breathing room beyond the longest label
_MIN_PLOT_FRACTION = 0.45  # never give the labels more than this much of the frame
_PLOT_TOP, _PLOT_BOTTOM_PAD = 0.04, 0.16  # room for the value axis beneath the plot


def _reserve_label_column(
    chart, spec: ChartSpec, *, rect: Rect, size_pt: float, face: str | None
) -> None:
    """Pin the plot area so the bar family's left-hand category labels fit inside the frame."""
    longest = max((str(c) for c in spec.categories), key=lambda c: text_em(c, face), default="")
    if not longest:
        return
    needed_in = (text_em(longest, face) + _LABEL_PAD_EM) * size_pt / 72
    x = min(needed_in / rect.width, 1.0 - _MIN_PLOT_FRACTION)
    if x <= 0:
        return
    plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    layout = plot_area.find(qn("c:layout"))
    if layout is None:
        layout = parse_xml(
            '<c:layout xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"/>'
        )
        plot_area.insert(0, layout)
    layout.append(
        parse_xml(
            '<c:manualLayout xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">'
            '<c:layoutTarget val="inner"/>'
            '<c:xMode val="edge"/><c:yMode val="edge"/>'
            f'<c:x val="{x:.4f}"/><c:y val="{_PLOT_TOP:.4f}"/>'
            f'<c:w val="{1.0 - x - 0.02:.4f}"/>'
            f'<c:h val="{1.0 - _PLOT_TOP - _PLOT_BOTTOM_PAD:.4f}"/>'
            "</c:manualLayout>"
        )
    )
