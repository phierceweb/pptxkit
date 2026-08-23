"""Choose the slide layout generated slides are composed on."""

from __future__ import annotations

from typing import Any

from lxml import etree
from pptx.enum.shapes import PP_PLACEHOLDER

from pptxkit.errors import ThemeError

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
# python-pptx clones every layout placeholder onto a new slide except these three,
# matched on type — idx is arbitrary, so a picture placeholder at idx 10 is cloned.
_LATENT_TYPES = frozenset({PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER})


def _emptiness(layout) -> tuple[int, int, int]:
    cloned = sum(1 for ph in layout.placeholders if ph.placeholder_format.type not in _LATENT_TYPES)
    return cloned, len(layout.placeholders), len(layout.shapes)


def _inherited(layout, master_index: int) -> tuple[int, str, bytes]:
    """What a slide would inherit from this layout beyond its own shapes."""
    element = layout.element
    bg = element.find(f"{{{_P}}}cSld/{{{_P}}}bg")
    return (
        master_index,
        str(element.get("showMasterSp", "1")),
        etree.tostring(bg) if bg is not None else b"",
    )


def pick_compose_layout(prs, *, prefer: str | None = None) -> Any:
    """The emptiest slide layout in ``prs``, searched across every master.

    ``prs.slide_layouts`` exposes only the first master's layouts, so a usable layout on
    a later master is otherwise invisible. Cloneable placeholders dominate the ranking
    because python-pptx copies them onto every slide it creates. Equally empty layouts
    that would compose *differently* are an ambiguity, not a free pick; ``prefer``
    settles it.

    Raises:
        ThemeError: the template has no layouts on any master, or its emptiest
            layouts disagree on what a slide would inherit.
    """
    candidates = [
        (_emptiness(layout), master_index, layout)
        for master_index, master in enumerate(prs.slide_masters)
        for layout in master.slide_layouts
    ]
    if not candidates:
        raise ThemeError("template defines no slide layouts on any master")

    if prefer is not None:
        named = [(mi, layout) for _, mi, layout in candidates if layout.name == prefer]
        if not named:
            available = ", ".join(sorted({layout.name for _, _, layout in candidates}))
            raise ThemeError(
                f"compose_layout names {prefer!r}, which the template does not "
                f"define; it has: {available}"
            )
        if len(named) > 1:
            where = ", ".join(f"master {mi}" for mi, _ in named)
            raise ThemeError(
                f"compose_layout {prefer!r} is ambiguous — {len(named)} layouts carry "
                f"that name ({where}); rename one in the template"
            )
        return named[0][1]

    best = min(rank for rank, _, _ in candidates)
    tied = [(master_index, layout) for rank, master_index, layout in candidates if rank == best]
    variants = {_inherited(layout, master_index) for master_index, layout in tied}
    if len(variants) > 1:
        listed = ", ".join(
            f"{layout.name!r} (master {master_index})" for master_index, layout in tied
        )
        raise ThemeError(
            f"template's emptiest slide layouts disagree on what a slide would "
            f"inherit — {listed}. Composing on either would rebrand the deck, so "
            f"name the one to use with 'compose_layout:' in the theme."
        )
    return tied[0][1]
