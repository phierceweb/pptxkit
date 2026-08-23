"""A native PowerPoint table: an optional header band over ruled body rows.

The band is a declared pair, so its ink is checked against its own fill; body cells are
left unfilled unless a pair is named, which lets one declaration read on ``page`` and on
``inverse``. Every cell is recorded on its own, or QA sees the table as one rectangle and
never reports a lost column.
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

from pptxkit.errors import LayoutError
from pptxkit.layouts.components import BodyResult, component
from pptxkit.layouts.registry import SlideCtx
from pptxkit.theme.model import Rect, TypeStyle
from pptxkit.theme.palette import Pair, lum
from pptxkit.utils.color import relative_luminance
from pptxkit.utils.shapes import ALIGN, ANCHOR, ANCHORS, para

from pptxkit.components._shape import known_fields, stroke, weight_pt
from pptxkit.components._tablegeom import aligns, heights, widths
from pptxkit.components._tablespec import Cell, Placed, Row, read_rows

_FIELDS = (
    "rows",
    "header",
    "total",
    "align",
    "valign",
    "widths",
    "head_pair",
    "body_pair",
    "color",
    "weight",
    "banding",
    "rules",
    "density",
)
_HEAD_PAIR_DEFAULT = "surface"
_BAND_PAIR_DEFAULT = "surface"
_RULE_ROLE_DEFAULT = "line"
# Every real table in the corpus centres its cells; the placement's own ``anchor:``
# says where the block sits in its box, which is a different question.
_VALIGN_DEFAULT = "middle"
RULES = ("rows", "header", "grid", "none")
"""What the hairlines mark: every row, the header alone, the whole grid, or nothing."""
# Office's "No Style, No Grid": python-pptx stamps a banded blue style on every new
# table, and it would repaint every fill and rule this component sets.
_NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
# Cell padding above and below the type, in ems of the body size.
_PAD_Y_EM = 0.35
# How far a banded row's fill moves from the surface it bands. Enough to follow a row
# across a wide table, not enough to read as a second colour.
_BAND_SHIFT = 0.045


@dataclass(frozen=True)
class _Look:
    """What every cell of one table is drawn with."""

    style: TypeStyle
    pad_x: float
    pad_y: float
    valign: str
    rules: str
    rule: RGBColor
    rule_pt: float


@dataclass(frozen=True)
class _CellBox:
    """A cell's geometry in the shape of what the recorder reads off a shape.

    A cell is not a shape: no id, no position of its own.
    """

    shape_id: int
    name: str
    left: int
    top: int
    width: int
    height: int


@component("table")
def table(ctx: SlideCtx) -> BodyResult:
    """Draw the placement as a table and record every cell it sets."""
    known_fields(ctx, _FIELDS)
    rows, columns, placed = read_rows(ctx)
    columns_align = aligns(ctx, columns)
    rect = ctx.body_rect
    columns_in = widths(ctx, columns, total=rect.width)

    style = ctx.style("body")
    density, valign, ruling = _density(ctx), _valign(ctx), _rules(ctx)
    pad_x = ctx.grid.gutter * density
    pad_y = style.size * _PAD_Y_EM / 72 * density
    if min(columns_in) - 2 * pad_x <= 0:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'table'): {columns} columns leave "
            f"{min(columns_in):.2f}in each, less than the {2 * pad_x:.2f}in the theme's "
            f"gutter pads a cell by — drop a column, widen the placement, or tighten "
            f"it with 'density:'"
        )

    rows_in = heights(
        rows,
        placed,
        columns_in,
        size_pt=style.size,
        pad_x=pad_x,
        pad_y=pad_y,
        limit=rect.height,
        face=ctx.theme.font_for(style),
    )
    extent = sum(rows_in)
    if extent > rect.height:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'table'): {len(rows)} rows need "
            f"{extent:.2f}in of height but the placement is only {rect.height:.2f}in — "
            f"split the table or shorten the cells"
        )

    frame = ctx.slide.shapes.add_table(
        len(rows),
        columns,
        Inches(rect.left),
        Inches(rect.top),
        Inches(rect.width),
        Inches(extent),
    )
    _unstyle(frame.table)
    for index, width in enumerate(columns_in):
        frame.table.columns[index].width = Inches(width)
    for index, depth in enumerate(rows_in):
        frame.table.rows[index].height = Inches(depth)
    ctx.manifest.record(frame)

    look = _Look(
        style=style,
        pad_x=pad_x,
        pad_y=pad_y,
        valign=valign,
        rules=ruling,
        rule=stroke(ctx, _RULE_ROLE_DEFAULT),
        rule_pt=weight_pt(ctx),
    )
    _draw(ctx, frame, rows, placed, rows_in, columns_in, columns_align, look, rect)
    return BodyResult(groups=[[frame.shape_id]], height=extent)


def _draw(
    ctx: SlideCtx,
    frame,
    rows: list[Row],
    placed: list[Placed],
    rows_in,
    columns_in,
    columns_align,
    look: _Look,
    rect: Rect,
) -> None:
    """Paint and record every cell, merging the ones that reach."""
    fills = _row_fills(ctx, rows)
    columns = len(columns_in)
    for spot in placed:
        cell = spot.cell
        box = Rect(
            rect.left + sum(columns_in[: spot.col]),
            rect.top + sum(rows_in[: spot.row]),
            sum(columns_in[spot.col : spot.col + cell.across]),
            sum(rows_in[spot.row : spot.row + cell.down]),
        )
        target = frame.table.cell(spot.row, spot.col)
        if cell.across > 1 or cell.down > 1:
            target.merge(frame.table.cell(spot.last_row, spot.last_col))
        row = rows[spot.row]
        pair = ctx.theme.palette.pair(cell.pair) if cell.pair else fills[spot.row]
        ink = _cell(
            ctx,
            target,
            cell,
            look=look,
            pair=pair,
            row=row,
            align=cell.align or columns_align[spot.col],
            edges=_edges(rows, spot, look.rules, columns=columns),
        )
        _record(
            ctx,
            frame,
            box,
            cell.text,
            ink=ink,
            size_pt=look.style.size,
            behind=pair.bg if pair is not None else ctx.behind(box),
            at=(spot.row, spot.col),
        )


def _edges(rows: list[Row], spot: Placed, rules: str, *, columns: int) -> tuple[bool, bool, bool]:
    """Which of this cell's edges are ruled: below, above, to the right.

    A cell's underside is the underside of everything it covers, so a cell reaching
    down carries the rule where its span ends and leaves the boundaries it crosses
    unruled *in its own column only* — which is what parting a merged label from the
    rows beside it looks like.
    """
    if rules == "none":
        return False, False, False
    last = spot.last_row
    below = rows[last].kind == "head" if rules == "header" else last < len(rows) - 1
    right = rules == "grid" and spot.last_col < columns - 1
    return below, rows[spot.row].kind == "total", right


def _row_fills(ctx: SlideCtx, rows: list[Row]) -> list[Pair | None]:
    """The pair painted behind each row, banding counted over body rows alone."""
    head_pair = ctx.theme.palette.pair(str(ctx.body.get("head_pair", _HEAD_PAIR_DEFAULT)))
    named = ctx.body.get("body_pair")
    body_pair = ctx.theme.palette.pair(str(named)) if named is not None else None
    bands = _band_fill(ctx, body_pair) if ctx.body.get("banding") else None

    out: list[Pair | None] = []
    body_index = 0
    for row in rows:
        if row.kind == "head":
            out.append(head_pair)
            continue
        if row.kind == "body":
            banded = bands is not None and body_index % 2 == 1
            body_index += 1
            out.append(bands if banded else body_pair)
            continue
        out.append(body_pair)
    return out


def _band_fill(ctx: SlideCtx, body: Pair | None) -> Pair:
    """The alternate row's colour, taken from the palette rather than from the slide.

    Nudging whatever the slide happens to show puts a mid-grey under the type the moment
    the master is a photograph. With no body fill the band is the ``surface`` pair; with
    one, that pair's own colour shifted toward its ink, near enough to keep its contrast.
    """
    if body is None:
        return ctx.theme.palette.pair(_BAND_PAIR_DEFAULT)
    toward = 1 if relative_luminance(body.fg) > relative_luminance(body.bg) else -1
    return Pair(body.fg, lum(body.bg, 1.0, _BAND_SHIFT * toward))


def _valign(ctx: SlideCtx) -> str:
    value = ctx.body.get("valign", _VALIGN_DEFAULT)
    if str(value) not in ANCHORS:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'table'): 'valign' is one of "
            f"{', '.join(ANCHORS)}, got {value!r}"
        )
    return str(value)


def _rules(ctx: SlideCtx) -> str:
    value = ctx.body.get("rules", RULES[0])
    if str(value) not in RULES:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'table'): 'rules' is one of "
            f"{', '.join(RULES)}, got {value!r}"
        )
    return str(value)


def _density(ctx: SlideCtx) -> float:
    """The multiplier on both cell paddings. Below 1 is tighter, above 1 airier."""
    value = ctx.body.get("density", 1.0)
    try:
        factor = 0.0 if isinstance(value, bool) else float(value)
    except (TypeError, ValueError):
        factor = 0.0
    if factor <= 0:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'table'): 'density' scales the cell "
            f"padding, so it is a positive number — 0.6 for a dense table, 1.4 for an "
            f"airy one; got {value!r}"
        )
    return factor


def _unstyle(tbl) -> None:
    """Drop the banded Office style python-pptx stamps on a new table."""
    tbl.first_row = False
    tbl.horz_banding = False
    style_id = tbl._tbl.get_or_add_tblPr().find(qn("a:tableStyleId"))
    if style_id is not None:
        style_id.text = _NO_STYLE


def _cell(
    ctx: SlideCtx,
    cell_shape,
    cell: Cell,
    *,
    look: _Look,
    pair: Pair | None,
    row: Row,
    align: str,
    edges: tuple[bool, bool, bool],
) -> str:
    """Paint one cell and set its type; returns the ink it was written in."""
    if pair is None:
        cell_shape.fill.background()
        ink = ctx.pair.fg
    else:
        cell_shape.fill.solid()
        cell_shape.fill.fore_color.rgb = ctx.rgb(pair.bg)
        ink = pair.fg
    cell_shape.margin_left = cell_shape.margin_right = Inches(look.pad_x)
    cell_shape.margin_top = cell_shape.margin_bottom = Inches(look.pad_y)
    cell_shape.vertical_anchor = ANCHOR[cell.valign or look.valign]
    tf = cell_shape.text_frame
    tf.word_wrap = True
    bold = row.kind in ("head", "total") or cell.emphasis
    if cell.text:
        para(
            tf,
            cell.text,
            look.style.size,
            ctx.rgb(ink),
            bold=bold,
            align=ALIGN[align],
            first=True,
            space_after=0,
            font=ctx.theme.font_for(look.style),
        )
    below, above, right = edges
    # Each rule goes to the front of a:tcPr, so they are written back to front: the
    # schema's order is lnL, lnR, lnT, lnB, and PowerPoint refuses a file out of it.
    if below:
        _rule(cell_shape, "lnB", look.rule, look.rule_pt)
    if above:
        _rule(cell_shape, "lnT", look.rule, look.rule_pt * 2)
    if right:
        _rule(cell_shape, "lnR", look.rule, look.rule_pt)
    return ink


def _rule(cell_shape, edge: str, colour: RGBColor, weight_pt_: float) -> None:
    """Rule one edge of the cell. python-pptx exposes no cell-border API at all."""
    line = parse_xml(
        f'<a:{edge} {nsdecls("a")} w="{Pt(weight_pt_).emu}" cap="flat" cmpd="sng" '
        f'algn="ctr"><a:solidFill><a:srgbClr val="{colour}"/></a:solidFill></a:{edge}>'
    )
    # Every border tag precedes the fill in the a:tcPr sequence, so the front is in order.
    cell_shape._tc.get_or_add_tcPr().insert(0, line)


def _record(
    ctx: SlideCtx,
    frame,
    rect: Rect,
    text: str,
    *,
    ink: str,
    size_pt: float,
    behind: str,
    at: tuple[int, int],
) -> None:
    row, col = at
    box = _CellBox(
        shape_id=int(frame.shape_id),
        name=f"{frame.name} r{row + 1}c{col + 1}",
        left=Inches(rect.left),
        top=Inches(rect.top),
        width=Inches(rect.width),
        height=Inches(rect.height),
    )
    cell = f"r{row + 1}c{col + 1}"
    if not text:
        ctx.manifest.record(box, part=cell)
        return
    ctx.manifest.record(box, part=cell, lines=[text], font_pt=size_pt, fg=ink, bg=behind)
