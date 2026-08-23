"""A monospace listing on a themed plate.

The only other way to put code on a slide is `document:`, which rasterizes through a
browser — this draws real text, so it is selectable, themed, and in the manifest as
lines QA can measure.
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from pptxkit.errors import LayoutError
from pptxkit.layouts.components import BodyResult, RevealItem, component
from pptxkit.layouts.registry import SlideCtx
from pptxkit.utils.shapes import para, rrect

from pptxkit.components._shape import known_fields, pair_named
from pptxkit.components._shared import head as head_line

_FIELDS = ("lines", "text", "heading", "pair", "accent", "wrap", "size")
_PAIR_DEFAULT = "surface"
# The plate's own padding and the space a heading takes, in inches.
_PAD_X = 0.26
_PAD_Y = 0.2
_HEADING_H = 0.42
# Line advance as a multiple of the type size, and the rung a listing is set at.
_LINE_ADVANCE = 1.55
_SIZE_RUNG = "caption"


@component("code")
def code(ctx: SlideCtx) -> BodyResult:
    """One plate, one line of monospace per entry."""
    known_fields(ctx, _FIELDS)
    lines = _lines(ctx)
    pair = pair_named(ctx, _PAIR_DEFAULT)
    accents = tuple(str(a) for a in (ctx.body.get("accent") or ()))
    heading = str(ctx.body.get("heading", "")) or None
    size = _size(ctx)

    rect = ctx.body_rect
    top = rect.top
    groups: list[list[RevealItem]] = []
    if heading is not None:
        frame = _heading(ctx, heading, top)
        groups.append([frame.shape_id])
        top += _HEADING_H

    advance = Inches(size * _LINE_ADVANCE / 72).inches
    plate_h = advance * len(lines) + 2 * _PAD_Y
    if top + plate_h > rect.bottom:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'code'): {len(lines)} lines need "
            f"{plate_h:.2f}in but only {rect.bottom - top:.2f}in is left — shorten the "
            f"listing, split the slide, or lower 'size'"
        )

    plate = rrect(ctx.slide, rect.left, top, rect.width, plate_h, ctx.rgb(pair.bg), radius=0.04)
    frame = plate.text_frame
    frame.word_wrap = bool(ctx.body.get("wrap", False))
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = frame.margin_right = Inches(_PAD_X)
    frame.margin_top = frame.margin_bottom = Inches(_PAD_Y)

    ink = ctx.rgb(ctx.ink_on(pair.bg))
    lit = ctx.rgb(ctx.accent_on(pair.bg, size_pt=size))
    drawn: list[str] = []
    for index, line in enumerate(lines):
        emphasised = any(line.lstrip().startswith(a) for a in accents)
        para(
            frame,
            line or " ",
            size,
            lit if emphasised else ink,
            bold=emphasised,
            first=(index == 0),
            space_after=0,
            font=ctx.theme.mono,
        )
        drawn.append(line or " ")
    ctx.manifest.record(plate, lines=drawn, font_pt=size, bg=pair.bg, plate=True)
    groups.append([plate.shape_id])
    ctx.panels.append((_plate_rect(rect, top, plate_h), pair.bg))
    return BodyResult(groups=groups, height=(top - rect.top) + plate_h)


def _lines(ctx: SlideCtx) -> list[str]:
    """``lines:`` as written, or ``text:`` split — one of the two is required."""
    raw, text = ctx.body.get("lines"), ctx.body.get("text")
    if raw is not None and text is not None:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'code'): give 'lines' or 'text', "
            f"not both — 'text' is a block scalar, 'lines' a list"
        )
    if text is not None:
        return str(text).rstrip("\n").split("\n")
    if not isinstance(raw, list) or not raw:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'code'): needs 'lines' (a list) or "
            f"'text' (a block scalar) — a listing with nothing in it draws an empty plate"
        )
    return [str(line) for line in raw]


def _size(ctx: SlideCtx) -> float:
    raw = ctx.body.get("size")
    if raw is None:
        return ctx.style(_SIZE_RUNG).size
    try:
        value = float(raw)
    except (TypeError, ValueError):
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'code'): 'size' is a point size, got {raw!r}"
        ) from None
    if value < ctx.theme.min_pt:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'code'): 'size' {value}pt is below the "
            f"theme's {ctx.theme.min_pt:.1f}pt minimum"
        )
    return value


def _heading(ctx: SlideCtx, text: str, top: float):
    from pptxkit.utils.shapes import textbox

    frame = textbox(ctx.slide, ctx.body_rect.left, top, ctx.body_rect.width, _HEADING_H)
    head_line(ctx, frame, text, first=True)
    ctx.manifest.record(frame._parent, text=text)
    return frame._parent


def _plate_rect(rect, top: float, height: float):
    from pptxkit.theme.model import Rect

    return Rect(left=rect.left, top=top, width=rect.width, height=height)
