"""Fill, stroke, effect and field-coercion helpers shared by the shape primitives.

A fill comes from a declared pair, already contrast-checked. A stroke has no pair, so it
takes the quietest role clearing WCAG's 3:1 non-text minimum against the paper it lands on.
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

from pptxkit.charts._effects import apply_shadow
from pptxkit.errors import LayoutError
from pptxkit.utils.keys import unknown_field
from pptxkit.layouts.registry import SlideCtx
from pptxkit.theme.model import Rect
from pptxkit.theme.palette import Pair
from pptxkit.utils.color import AA_LARGE, contrast_ratio
from pptxkit.utils.shapes import solid

ARROWS = ("none", "end", "both")
"""``arrow:`` values a connector accepts."""

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_ARROW_TAGS = {"none": (), "end": ("tailEnd",), "both": ("headEnd", "tailEnd")}
# Roles tried after the one the author named, before the slide pair's own ink.
_STROKE_FALLBACK_ROLES = ("muted",)
# Below this ratio two colours cannot be told apart: a fill vanishes into the paper,
# and a line drawn in it is not there at all.
_MIN_VISIBLE_RATIO = 1.2
_WEIGHT_RANGE = (0.1, 8.0)


def known_fields(ctx: SlideCtx, fields: tuple[str, ...]) -> None:
    """Reject an unknown component field, naming the ones this primitive reads."""
    unknown = sorted(set(ctx.body) - set(fields))
    if unknown:
        raise LayoutError(unknown_field(unknown[0], fields, where=_where(ctx)))


def known_item_fields(
    ctx: SlideCtx,
    item: dict,
    fields: frozenset[str],
    *,
    index: int | None = None,
    noun: str = "item",
) -> None:
    """Reject an unknown key inside one entry of a component's item list.

    ``known_fields`` sees only the component's own keys, so without this an unknown item
    key is silently dropped — commonest cause being an unquoted comma in a label, which
    YAML has already split into two keys.
    """
    unknown = sorted(set(item) - fields)
    if not unknown:
        return
    # An entry in a list is "item 3"; a mapping the component names outright — a versus
    # side, say — is just 'left', and numbering it reads as nonsense.
    if index is None:
        where, subject = f"{noun!r}", "a side"
    else:
        where = f"{noun} {index}"
        subject = f"{'an' if noun[0] in 'aeiou' else 'a'} {noun}"
    raise LayoutError(
        unknown_field(
            unknown[0],
            sorted(fields),
            where=_where(ctx),
            lead=f"{where} has the unknown field",
            label=f"{subject} reads",
        )
    )


def pair_named(ctx: SlideCtx, default: str) -> Pair:
    """The palette pair the component's ``pair:`` names."""
    return ctx.theme.palette.pair(str(ctx.body.get("pair", default)))


def visible_edge(ctx: SlideCtx, pair: Pair) -> RGBColor | None:
    """The theme's rule colour when ``pair``'s fill would vanish into the slide's own."""
    if contrast_ratio(pair.bg, ctx.pair.bg) >= _MIN_VISIBLE_RATIO:
        return None
    return ctx.color("line")


def stroke(ctx: SlideCtx, default_role: str) -> RGBColor:
    """The colour a line is drawn in.

    A role the author named is refused only when invisible on this slide's paper. The
    *default* has no author behind it, so it must clear WCAG's 3:1 non-text minimum or
    give way to a role that does.
    """
    named = ctx.body.get("color")
    if named is not None:
        colour = ctx.theme.palette.role(str(named))
        ratio = contrast_ratio(colour, ctx.pair.bg)
        if ratio < _MIN_VISIBLE_RATIO:
            raise LayoutError(
                f"{_where(ctx)}: color {named!r} is {colour} against this slide's "
                f"{ctx.pair.bg}, {ratio:.2f}:1 — a line that close to the paper cannot "
                f"be seen; name a role that stands off it"
            )
        return ctx.rgb(colour)
    roles = (default_role, *_STROKE_FALLBACK_ROLES)
    for candidate in (ctx.theme.palette.role(r) for r in roles):
        if contrast_ratio(candidate, ctx.pair.bg) >= AA_LARGE:
            return ctx.rgb(candidate)
    return ctx.fg()


def weight_pt(ctx: SlideCtx) -> float:
    """Stroke weight in points: ``weight:`` as a multiple of the theme's line weight."""
    lo, hi = _WEIGHT_RANGE
    multiple = number(ctx, "weight", default=1.0)
    if not lo <= multiple <= hi:
        raise LayoutError(
            f"{_where(ctx)}: 'weight' is a multiple of the theme's line weight "
            f"({ctx.theme.line_weight:.2f}pt), not a point size — it runs {lo} to {hi}, "
            f"got {multiple}"
        )
    return ctx.theme.line_weight * multiple


def fraction(ctx: SlideCtx, key: str, *, default: float, what: str) -> float:
    """A field that must be a fraction in ``0 < value <= 1``."""
    value = number(ctx, key, default=default)
    if not 0.0 < value <= 1.0:
        raise LayoutError(
            f"{_where(ctx)}: {key!r} is {what}, so it runs above 0 and up to 1; got {value}"
        )
    return value


def number(ctx: SlideCtx, key: str, *, default: float) -> float:
    """Coerce a component field to ``float``, failing with the slide named."""
    raw = ctx.body.get(key)
    if raw is None:
        return default
    try:
        return float(raw)
    except (TypeError, ValueError):
        raise LayoutError(f"{_where(ctx)}: {key!r} must be a number, got {raw!r}") from None


def flag(ctx: SlideCtx, key: str) -> bool:
    """A boolean component field; anything but a real bool is refused."""
    value = ctx.body.get(key, False)
    if not isinstance(value, bool):
        raise LayoutError(f"{_where(ctx)}: {key!r} must be true or false, got {value!r}")
    return value


def choice(ctx: SlideCtx, key: str, options: tuple[str, ...], *, default: str) -> str:
    """A component field restricted to a fixed vocabulary."""
    value = ctx.body.get(key)
    if value is None:
        return default
    if str(value) not in options:
        raise LayoutError(
            f"{_where(ctx)}: {key!r} must be one of {', '.join(options)}, got {value!r}"
        )
    return str(value)


def oval(ctx: SlideCtx, rect: Rect, fill: RGBColor, *, line: RGBColor | None = None):
    """Add an ellipse filling ``rect``."""
    shape = ctx.slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(rect.left),
        Inches(rect.top),
        Inches(rect.width),
        Inches(rect.height),
    )
    solid(shape, fill, line)
    return shape


def line_shape(
    ctx: SlideCtx, kind, start: tuple[float, float], end: tuple[float, float], *, default_role: str
):
    """Add a line between two points, coloured and weighted from the theme.

    A connector arrives carrying a ``<p:style>`` pointing at the template's own line and
    effect styles, so the inherited effect is dropped.
    """
    shape = ctx.slide.shapes.add_connector(
        kind, Inches(start[0]), Inches(start[1]), Inches(end[0]), Inches(end[1])
    )
    shape.line.color.rgb = stroke(ctx, default_role)
    shape.line.width = Pt(weight_pt(ctx))
    shape.shadow.inherit = False
    return shape


def shadow(ctx: SlideCtx, shape) -> None:
    """Drop the theme's declared shadow behind ``shape`` when the component asked.

    The geometry is the theme's ``chart:`` shadow, the only one a theme declares.
    """
    if not flag(ctx, "shadow"):
        return
    apply_shadow(shape._element.spPr, ctx.theme.chart)


def arrow(shape, kind: str) -> None:
    """Put arrowheads on a line's ends. python-pptx exposes no API for either end."""
    if kind == "none":
        return
    ln = shape.line._get_or_add_ln()
    for tag in _ARROW_TAGS[kind]:
        ln.append(parse_xml(f'<a:{tag} xmlns:a="{_A}" type="triangle" w="med" len="med"/>'))


def anchored(rect: Rect, width: float, height: float, *, align: str, anchor: str) -> Rect:
    """Place a ``width`` x ``height`` shape inside ``rect`` per the placement's keys."""
    left = {
        "left": rect.left,
        "center": rect.left + (rect.width - width) / 2,
        "right": rect.right - width,
    }[align]
    top = {
        "top": rect.top,
        "middle": rect.top + (rect.height - height) / 2,
        "bottom": rect.bottom - height,
    }[anchor]
    return Rect(left, top, width, height)


def canvas_point(ctx: SlideCtx, value: Any, *, key: str) -> tuple[float, float]:
    """A ``[x, y]`` pair of canvas fractions, in inches."""
    ok = (
        isinstance(value, (list, tuple))
        and len(value) == 2
        and all(isinstance(n, (int, float)) and not isinstance(n, bool) for n in value)
    )
    if not ok:
        raise LayoutError(
            # No square brackets in the literal text: the CLI prints through a rich
            # console, which reads a bare bracketed run as markup and drops it.
            f"{_where(ctx)}: {key!r} must be a placement id, or two numbers, x then y, "
            f"as fractions of the canvas; got {value!r}"
        )
    x, y = (float(n) for n in value)
    return ctx.grid.scale.x(x), ctx.grid.scale.y(y)


def _where(ctx: SlideCtx) -> str:
    return f"slide {ctx.spec.index} (component {ctx.component!r})"
