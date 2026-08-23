"""One call on the left, the work it sets off on the right, joined by a bus.

A fan, not a sequence: every branch leaves the same point at once. Use `flow` where each
step follows the one before.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from pptxkit.errors import LayoutError
from pptxkit.icons.draw import place_icon
from pptxkit.layouts.components import BodyResult, RevealItem, component
from pptxkit.layouts.registry import SlideCtx
from pptxkit.theme.model import Rect
from pptxkit.utils.shapes import para, rect, rrect, solid, textbox

from pptxkit.components._shape import known_fields, known_item_fields

_FIELDS = ("source", "items", "weight")
_ITEM_FIELDS = frozenset({"text", "icon"})

_MIN_ITEMS = 2
_SOURCE_FRACTION = 0.29
_GAP_FRACTION = 0.11
_STROKE = 0.022
_WEIGHT_RANGE = (0.5, 4.0)
_ICON_SIDE = 0.28
_HEAD = 0.15
_SOURCE_MAX_H = 1.15
_CHIP_MAX_H = 0.48


def _weight(ctx: SlideCtx) -> float:
    raw = ctx.body.get("weight", 1.0)
    try:
        value = float(raw)
    except (TypeError, ValueError):
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'fanout'): 'weight' scales the bus "
            f"stroke, got {raw!r}"
        ) from None
    low, high = _WEIGHT_RANGE
    if not low <= value <= high:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'fanout'): 'weight' scales the bus "
            f"stroke, {low} to {high}; got {value:g}"
        )
    return value


def _items(ctx: SlideCtx) -> list[dict]:
    raw = ctx.body.get("items")
    if not isinstance(raw, list) or len(raw) < _MIN_ITEMS:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'fanout'): a fanout needs at least "
            f"{_MIN_ITEMS} items — one consequence is the 'connector' component"
        )
    for index, item in enumerate(raw, start=1):
        if not isinstance(item, dict) or not str(item.get("text", "")).strip():
            raise LayoutError(
                f"slide {ctx.spec.index} (component 'fanout'): item {index} needs a 'text'"
            )
        known_item_fields(ctx, item, _ITEM_FIELDS, index=index)
    return raw


@component("fanout")
def fanout(ctx: SlideCtx) -> BodyResult:
    """The source plate, then one branch per consequence; one reveal group each.

    The plate, the trunk and the spine arrive together as the first group, so the shape
    is established before anything hangs off it.
    """
    known_fields(ctx, _FIELDS)
    source = str(ctx.body.get("source", "")).strip()
    if not source:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'fanout'): 'source' is the call the "
            f"branches leave from"
        )
    items = _items(ctx)
    stroke = _STROKE * _weight(ctx)
    r = ctx.body_rect

    src_w = r.width * _SOURCE_FRACTION
    gap = r.width * _GAP_FRACTION
    spine_x = r.left + src_w + gap * 0.42
    icon_x = r.left + src_w + gap * 1.02
    chip_x = icon_x + _ICON_SIDE + 0.16
    chip_w = r.right - chip_x
    if chip_w < 1.0:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'fanout'): the source plate and the bus "
            f"leave {chip_w:.2f}in for the consequences — widen the placement"
        )
    lane = r.height / len(items)
    chip_h = min(lane * 0.8, _CHIP_MAX_H)

    body, caption = ctx.style("body"), ctx.style("caption")
    bus = ctx.color("accent-2")
    accent_hex = ctx.theme.palette.role("accent-1")
    groups: list[list[RevealItem]] = []

    src_h = min(r.height * 0.32, _SOURCE_MAX_H)
    src_y = r.top + (r.height - src_h) / 2
    plate = rrect(ctx.slide, r.left, src_y, src_w, src_h, ctx.color("accent-1"), radius=0.12)
    ctx.manifest.record(plate)
    tf = textbox(ctx.slide, r.left + 0.16, src_y, src_w - 0.32, src_h, anchor=MSO_ANCHOR.MIDDLE)
    para(
        tf,
        source,
        caption.size * 1.3,
        ctx.rgb(ctx.ink_on(accent_hex)),
        bold=True,
        align=PP_ALIGN.CENTER,
        first=True,
        space_after=0,
        font=ctx.theme.mono,
    )
    ctx.manifest.record(
        tf._parent,
        lines=[source],
        font_pt=caption.size * 1.3,
        fg=ctx.ink_on(accent_hex),
        bg=accent_hex,
    )

    first_y = r.top + lane / 2
    last_y = r.top + (len(items) - 1) * lane + lane / 2
    trunk = rect(
        ctx.slide,
        r.left + src_w,
        src_y + src_h / 2 - stroke / 2,
        spine_x - (r.left + src_w),
        stroke,
        bus,
    )
    spine = rect(ctx.slide, spine_x - stroke / 2, first_y, stroke, last_y - first_y, bus)
    ctx.manifest.record(trunk)
    ctx.manifest.record(spine)
    groups.append([plate.shape_id, tf._parent.shape_id, trunk.shape_id, spine.shape_id])

    for index, item in enumerate(items):
        mid = r.top + index * lane + lane / 2
        ids: list[RevealItem] = []

        branch = rect(
            ctx.slide, spine_x, mid - stroke / 2, icon_x - spine_x - _HEAD - 0.04, stroke, bus
        )
        ctx.manifest.record(branch)
        ids.append(branch.shape_id)

        # A triangle turned a quarter is the arrowhead; a rectangle cannot carry one.
        tip = ctx.slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(icon_x - _HEAD - 0.04),
            Inches(mid - _HEAD * 0.60),
            Inches(_HEAD),
            Inches(_HEAD * 1.2),
        )
        tip.rotation = 90
        tip.line.fill.background()
        solid(tip, bus)
        ctx.manifest.record(tip)
        ids.append(tip.shape_id)

        if item.get("icon"):
            # Into the row's group: a shape outside every group is unanimated, which
            # means on screen from the first beat — every mark ahead of its own text.
            mark = place_icon(
                ctx.slide,
                str(item["icon"]),
                Rect(icon_x, mid - _ICON_SIDE / 2, _ICON_SIDE, _ICON_SIDE),
                fill=accent_hex,
                theme=ctx.theme,
            )
            ctx.manifest.record(mark)
            ids.append(mark.shape_id)

        text = str(item["text"])
        chip = textbox(
            ctx.slide, chip_x, mid - chip_h / 2, chip_w, chip_h, anchor=MSO_ANCHOR.MIDDLE
        )
        para(chip, text, body.size, ctx.fg(), first=True, space_after=0, font=ctx.theme.face)
        ctx.manifest.record(
            chip._parent, lines=[text], font_pt=body.size, fg=ctx.pair.fg, bg=ctx.pair.bg
        )
        ids.append(chip._parent.shape_id)
        groups.append(ids)

    return BodyResult(groups=groups, height=r.height)
