"""A rounded block carrying a heading, a line of copy and an optional icon.

The plate's colour is a declared pair, so its ink is contrast-checked against its own
fill rather than against the slide's.
"""

from __future__ import annotations

import math

from pptx.util import Inches

from pptxkit.errors import LayoutError
from pptxkit.icons.draw import place_icon
from pptxkit.layouts.components import BodyResult, component
from pptxkit.layouts.registry import SlideCtx
from pptxkit.theme.media import resolve_media
from pptxkit.theme.model import Rect
from pptxkit.utils.shapes import ANCHOR, para, rrect, textbox
from pptxkit.utils.text import LINE_HEIGHT, wrapped_lines

from pptxkit.components._shape import (
    known_fields,
    number,
    pair_named,
    shadow,
    visible_edge,
)
from pptxkit.components._shared import HEAD_SPACE_AFTER_PT, mark_colour

_FIELDS = ("pair", "heading", "body", "icon", "radius", "shadow")
_PAIR_DEFAULT = "surface"
_RADIUS_DEFAULT = 0.06
_MAX_RADIUS = 0.5
# The icon is a mark beside the copy, not an illustration: two heading line-heights square.
_ICON_LINES = 2.0


@component("card")
def card(ctx: SlideCtx) -> BodyResult:
    """Paint the placement as a plate and set the card's contents on it."""
    known_fields(ctx, _FIELDS)
    heading = _text(ctx, "heading")
    copy = _text(ctx, "body")
    icon = ctx.body.get("icon")
    if not (heading or copy or icon):
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'card'): a card needs a 'heading', a "
            f"'body' or an 'icon' — an empty plate is the 'panel' component"
        )
    pair = pair_named(ctx, _PAIR_DEFAULT)
    rect = ctx.body_rect
    radius = _radius(ctx)
    plate = rrect(
        ctx.slide,
        rect.left,
        rect.top,
        rect.width,
        rect.height,
        ctx.rgb(pair.bg),
        line=visible_edge(ctx, pair),
        radius=radius,
    )
    shadow(ctx, plate)
    ctx.manifest.record(plate, bg=pair.bg)
    # Chrome drawn over this plate reads its fill, not the slide's surface.
    ctx.panels.append((rect, pair.bg))
    shapes = [plate.shape_id]

    inset = ctx.grid.gutter
    side = inset + _corner_reach(rect, radius, depth=inset)
    area = Rect(rect.left + side, rect.top + inset, rect.width - 2 * side, rect.height - 2 * inset)
    if area.width <= 0 or area.height <= 0:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'card'): the placement is smaller than "
            f"the {inset:.2f}in the theme's gutter insets a card by on every side"
        )
    if icon:
        area = _icon(ctx, area, str(icon), shapes)
    if heading or copy:
        shapes.append(_copy(ctx, area, heading=heading, copy=copy, ink=pair.fg, bg=pair.bg))
    return BodyResult(groups=[shapes], height=rect.height)


def copy_height(ctx: SlideCtx, *, width: float, heading: str, copy: str) -> float:
    """The height a card's type needs set to ``width``, in inches."""
    head_style, body_style = ctx.style("head"), ctx.style("body")
    needed = 0.0
    if heading:
        needed += (
            wrapped_lines(
                heading,
                width_in=width,
                size_pt=head_style.size,
                face=ctx.theme.font_for(head_style),
            )
            * head_style.size
            * LINE_HEIGHT
            + HEAD_SPACE_AFTER_PT
        ) / 72
    if copy:
        needed += (
            wrapped_lines(
                copy, width_in=width, size_pt=body_style.size, face=ctx.theme.font_for(body_style)
            )
            * body_style.size
            * LINE_HEIGHT
            / 72
        )
    return needed


def icon_side(ctx: SlideCtx) -> float:
    """The side of a card's icon: two heading line-heights, in inches."""
    return ctx.style("head").size * LINE_HEIGHT / 72 * _ICON_LINES


def plate_height(
    ctx: SlideCtx,
    *,
    width: float,
    heading: str,
    copy: str,
    icon: bool = False,
    radius: float = _RADIUS_DEFAULT,
) -> float:
    """The plate a card of this content fills at ``width``: its parts plus the inset.

    ``icon`` counts, or a caller sizes its steps by their copy and then has no room for
    the mark above it.
    """
    inset = ctx.grid.gutter
    mark = icon_side(ctx) + inset if icon else 0.0

    def depth(side: float) -> float:
        return (
            copy_height(ctx, width=width - 2 * side, heading=heading, copy=copy) + mark + 2 * inset
        )

    # The corner reach needs the plate's height, so size once against the square corner
    # and again against the curve that height earns — the same rule `card()` insets by.
    height = depth(inset)
    side = inset + _corner_reach(Rect(0.0, 0.0, width, height), radius, depth=inset)
    return height if side == inset else depth(side)


def _text(ctx: SlideCtx, key: str) -> str:
    value = ctx.body.get(key)
    return "" if value is None else str(value)


def _radius(ctx: SlideCtx) -> float:
    value = number(ctx, "radius", default=_RADIUS_DEFAULT)
    if not 0.0 <= value <= _MAX_RADIUS:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'card'): 'radius' is a fraction of the "
            f"plate's short side, 0..{_MAX_RADIUS} ({_MAX_RADIUS} is a stadium); "
            f"got {value}"
        )
    return value


def _corner_reach(rect: Rect, radius: float, *, depth: float) -> float:
    """How far the corner arc cuts in horizontally, ``depth`` below the plate's top edge.

    The corner is a quarter-circle of radius ``r = radius * min(w, h)`` centred at
    ``(r, r)``, so its edge at height ``d`` sits at ``r - sqrt(r^2 - (r - d)^2)``.
    """
    r = radius * min(rect.width, rect.height)
    if depth >= r:
        return 0.0
    return r - math.sqrt(r * r - (r - depth) * (r - depth))


def _icon(ctx: SlideCtx, area: Rect, name: str, shapes: list[int]) -> Rect:
    """Place the icon at the top of the card's inner area; return what is left below it.

    A bare name is a glyph, drawn as vector and painted from the palette; anything
    carrying a suffix or a path separator is a picture file, placed as it is.
    """
    side = icon_side(ctx)
    if side >= area.height:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'card'): the icon alone wants "
            f"{side:.2f}in of the {area.height:.2f}in inside the plate — grow the "
            f"placement or drop the icon"
        )
    box = Rect(area.left, area.top, side, side)
    if _is_file(name):
        media = resolve_media(name, template=ctx.theme.template, roots=ctx.media_roots)
        shape = ctx.slide.shapes.add_picture(
            str(media), Inches(box.left), Inches(box.top), Inches(side), Inches(side)
        )
        ctx.manifest.record(shape, rendered="image")
    else:
        fill = mark_colour(ctx, box)
        shape = place_icon(ctx.slide, name, box, fill=fill, theme=ctx.theme)
        ctx.manifest.record(shape, fg=fill, bg=ctx.behind(box, ink=fill))
    shapes.append(shape.shape_id)
    gap = ctx.grid.gutter
    return Rect(area.left, area.top + side + gap, area.width, area.height - side - gap)


def _is_file(name: str) -> bool:
    return "." in name or "/" in name or "\\" in name


def _copy(ctx: SlideCtx, area: Rect, *, heading: str, copy: str, ink: str, bg: str) -> int:
    """Set the card's type in one frame, refusing copy that will not fit the plate."""
    head_style, body_style = ctx.style("head"), ctx.style("body")
    needed = copy_height(ctx, width=area.width, heading=heading, copy=copy)
    if needed > area.height:
        raise LayoutError(
            f"slide {ctx.spec.index} (component 'card'): the card's type wants "
            f"{needed:.2f}in but only {area.height:.2f}in is left inside the plate — "
            f"shorten the copy or grow the placement"
        )
    tf = textbox(ctx.slide, area.left, area.top, area.width, area.height, anchor=ANCHOR["top"])
    lines: list[str] = []
    if heading:
        para(
            tf,
            heading,
            head_style.size,
            ctx.rgb(ink),
            bold=head_style.bold,
            italic=head_style.italic,
            align=ctx.text_align(),
            first=True,
            space_after=HEAD_SPACE_AFTER_PT,
            font=ctx.theme.font_for(head_style),
        )
        lines.append(heading)
    if copy:
        para(
            tf,
            copy,
            body_style.size,
            ctx.rgb(ink),
            italic=body_style.italic,
            align=ctx.text_align(),
            first=not heading,
            space_after=0,
            font=ctx.theme.font_for(body_style),
        )
        lines.append(copy)
    ctx.manifest.record(
        tf._parent,
        lines=lines,
        font_pt=head_style.size if heading else body_style.size,
        line_pt=([head_style.size] if heading else []) + ([body_style.size] if copy else []),
        fg=ink,
        bg=bg,
    )
    return tf._parent.shape_id
