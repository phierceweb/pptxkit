"""Read a brand template and write the theme YAML that drives it.

A starting point for a human to edit, never a build-time path: the library specializes
into a template it is *told* about.
"""

from __future__ import annotations

import collections
from pathlib import Path
from typing import Any, cast

from pptx.util import Emu

from PIL import Image

from pptxkit.imagery.fit import ImageFit
from pptxkit.imagery.sample import cells, load, weakest
from pptxkit.layouts.chrome import CHROME_ORDER
from pptxkit.layouts.resolve import pick_compose_layout
from pptxkit.theme.clrscheme import parse_color_scheme, read_theme_xml
from pptxkit.theme.defaults import DEFAULT_RAMP, default_grid
from pptxkit.theme.media import resolve_media
from pptxkit.theme.model import Rect
from pptxkit.theme.palette import lum
from pptxkit.theme.scale import Scale
from pptxkit.theme.stock import is_stock_accent
from pptxkit.theme.surface import inherited_surface
from pptxkit.utils.color import AA_NORMAL, contrast_ratio, relative_luminance
from pptxkit.utils.deck import open_presentation

_DARK_SLOTS = ("dk1", "dk2")
_LIGHT_SLOTS = ("lt1", "lt2")
# How far a measured rung may sit from a built-in one and still be worth restating.
_RUNG_TOLERANCE = 0.15
# How far secondary text and plate fills move from the ink toward the page.
_MUTED_TOWARD_PAGE = 0.35
# Backed off in this step until the blend is readable on the page it sits on.
_MUTED_STEP = 0.05
_SURFACE_TOWARD_INK = 0.12
# Stepped away from the ink until the panel's own text reads on it.
_SURFACE_STEP = 0.04
_SURFACE_LIMIT = 0.60
# Below this contrast the master's paint and the page slot are the same surface, and
# restating every role as a literal would only churn the theme it writes.
_SAME_PAGE = 1.2
# The grid the background is probed on for artwork, and how far a column's pixels may
# spread before it counts as decorated rather than free.
_PROBE_COLS, _PROBE_ROWS = 48, 32
_UNIFORM_SPREAD = 12
# A run narrower than this much of the canvas is not worth steering a title into.
_MIN_CLEAR_RUN = 0.3
# Below this, a colour is a neutral: black, white, or a grey with no brand in it.
_MIN_CHROMA = 0.08
# The vertical slice of the canvas the chrome stack occupies.
_CHROME_BAND = (0.0, 0.3)


def _shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from _shapes(shape.shapes)
        else:
            yield shape


def _typography(prs, slide_h: float) -> tuple[str | None, dict[str, float]]:
    """The face the template's slides are really set in, and the rungs they use.

    Never the ``fontScheme``: a template routinely declares one face while every slide
    in it uses another.
    """
    faces: collections.Counter[str] = collections.Counter()
    rungs: collections.Counter[float] = collections.Counter()
    for slide in prs.slides:
        for shape in _shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name:
                        faces[run.font.name] += len(run.text)
                    if run.font.size:
                        rungs[round(run.font.size.pt / slide_h, 3)] += 1
    face = faces.most_common(1)[0][0] if faces else None
    ramp = {
        name: measured
        for name, default in DEFAULT_RAMP.items()
        for measured in [min(rungs, key=lambda r: abs(r - default), default=None)]
        if measured is not None and abs(measured - default) <= _RUNG_TOLERANCE
    }
    return face, ramp


def _surfaces(scheme: dict[str, str]) -> tuple[str, str]:
    """The slots to bind ``page`` and ``ink`` to: the lightest and the darkest.

    Not ``lt1``/``dk1`` by position — a template's real dark may sit in ``dk2``.
    """
    slots = [s for s in (*_LIGHT_SLOTS, *_DARK_SLOTS) if s in scheme]
    by_light = sorted(slots, key=lambda s: relative_luminance(scheme[s]))
    return by_light[-1], by_light[0]


def _painted(surface, art: Path | None, canvas: Rect, scheme: dict[str, str]) -> str | None:
    """The colour the template's own master really shows, if it paints one at all.

    A picture is sampled: the returned colour is the spot reading *worst* against the
    ink that will end up on it, so the palette's stated contrast is the real floor.
    """
    if surface is None:
        return None
    if art is None:
        return min(surface.fills, key=relative_luminance, default=None)
    sampled = cells(art, ImageFit(canvas).window, base="FFFFFF")
    candidates = [scheme[s] for s in (*_LIGHT_SLOTS, *_DARK_SLOTS) if s in scheme]
    ink = max(candidates, key=lambda c: contrast_ratio(c, weakest(sampled, ink=c)))
    return weakest(sampled, ink=ink)


def _muted(ink: str, painted: str) -> str:
    """Ink softened toward the page, but never past readable on it."""
    amount = _MUTED_TOWARD_PAGE
    while amount > 0:
        candidate = _mix(ink, painted, amount)
        if contrast_ratio(candidate, painted) >= AA_NORMAL:
            return candidate
        amount -= _MUTED_STEP
    return ink


def _surface(painted: str, ink: str, *, away: int) -> str:
    """A panel shifted off the page, far enough that ``ink`` reads on it."""
    # A panel is normally a slight recess; on a mid-tone page that runs into its own ink.
    for direction in (away, -away):
        offset = _SURFACE_TOWARD_INK
        while offset <= _SURFACE_LIMIT:
            candidate = lum(painted, 1.0, offset * direction)
            if contrast_ratio(ink, candidate) >= AA_NORMAL:
                return candidate
            offset += _SURFACE_STEP
    return lum(painted, 1.0, _SURFACE_TOWARD_INK * away)


def _over(painted: str, scheme: dict[str, str]) -> dict[str, str]:
    """Every surface role, re-derived to read on a page the template itself paints.

    The built-in roles assume a light page, so they are restated as literals relative
    to the ink actually chosen — a ``clrScheme`` holds no slot for "secondary text on
    this template's photograph".
    """
    slots = [s for s in (*_LIGHT_SLOTS, *_DARK_SLOTS) if s in scheme]
    ink = max(slots, key=lambda s: contrast_ratio(scheme[s], painted))
    toward = relative_luminance(painted) > relative_luminance(scheme[ink])
    return {
        "page": painted,
        "ink": ink,
        "muted": _muted(scheme[ink], painted),
        "surface": _surface(painted, scheme[ink], away=(-1 if toward else 1)),
        "surface-ink": scheme[ink],
    }


def _clear_run(path: Path, band: tuple[float, float]) -> tuple[float, float] | None:
    """The widest horizontal run of ``band`` a background picture leaves uniform.

    A column counts as decorated when it differs from the *typical* column of the
    band, not when it varies down its own length — a banded background varies
    everywhere and none of that is artwork.

    Returns:
        ``(start, end)`` as fractions of the width, or None when the whole band is
        plain (nothing to avoid) or none of it is (nowhere to go).
    """
    top, bottom = band
    image = load(path).convert("RGB").resize((_PROBE_COLS, _PROBE_ROWS), Image.Resampling.BOX)
    rows = range(
        max(0, int(top * _PROBE_ROWS)), max(1, min(_PROBE_ROWS, int(bottom * _PROBE_ROWS) + 1))
    )
    columns = [
        [cast("tuple[int, int, int]", image.getpixel((x, y))) for y in rows]
        for x in range(_PROBE_COLS)
    ]
    typical = [
        tuple(sorted(channel)[len(channel) // 2] for channel in zip(*pixels, strict=True))
        for pixels in zip(*columns, strict=True)
    ]
    quiet = [
        max(
            abs(a - b)
            for pixel, ref in zip(column, typical, strict=True)
            for a, b in zip(pixel, ref, strict=True)
        )
        <= _UNIFORM_SPREAD
        for column in columns
    ]
    if all(quiet) or not any(quiet):
        return None
    best: tuple[int, int] | None = None
    start = 0
    run = 0
    for x, free in (*enumerate(quiet), (_PROBE_COLS, False)):
        if free:
            start = x if run == 0 else start
            run += 1
            continue
        if best is None or run > best[1] - best[0]:
            best = (start, x) if run else best
        run = 0
    if best is None or (best[1] - best[0]) < _PROBE_COLS * _MIN_CLEAR_RUN:
        return None
    return best[0] / _PROBE_COLS, best[1] / _PROBE_COLS


def _columns(run: tuple[float, float]) -> tuple[int, int]:
    """A canvas-fraction run as a pair of grid column indices, clamped to the grid."""
    grid = default_grid(Scale(slide_w=1.0, slide_h=1.0))
    span = 1.0 - grid.left_frac - grid.right_frac
    edges = [
        min(grid.columns, max(0, round((frac - grid.left_frac) / span * grid.columns)))
        for frac in run
    ]
    return edges[0], max(edges[0] + 1, edges[1])


def _chroma(hex_colour: str) -> float:
    """How far a colour is from neutral, 0 (grey) to 1 (fully saturated)."""
    channels = [int(hex_colour[i : i + 2], 16) for i in (0, 2, 4)]
    return (max(channels) - min(channels)) / 255


def _mix(colour: str, other: str, amount: float) -> str:
    """``colour`` moved ``amount`` of the way to ``other``, per channel."""
    return "".join(
        f"{round(int(colour[i : i + 2], 16) * (1 - amount) + int(other[i : i + 2], 16) * amount):02X}"
        for i in (0, 2, 4)
    )


def derive(template: str | Path, *, prefer: str | None = None) -> dict[str, Any]:
    """The theme YAML for ``template``, as a mapping ready to dump.

    Raises:
        ThemeError: the template is not a readable .pptx, defines no usable layout,
            or carries no theme part.
    """
    template = Path(template)
    prs = open_presentation(template)
    slide_h = Emu(prs.slide_height).inches
    canvas = Rect(0.0, 0.0, Emu(prs.slide_width).inches, slide_h)
    scheme = parse_color_scheme(
        read_theme_xml(pick_compose_layout(prs, prefer=prefer).slide_master)
    )

    surface = inherited_surface(pick_compose_layout(prs, prefer=prefer))
    art = (
        resolve_media(surface.media, template=template)
        if surface is not None and surface.media is not None
        else None
    )

    page_slot, ink_slot = _surfaces(scheme)
    bind: dict[str, str] = {"page": page_slot, "ink": ink_slot}
    # The master's own paint outranks the scheme: it is what a slide will show.
    painted = _painted(surface, art, canvas, scheme)
    if painted is not None and contrast_ratio(painted, scheme[page_slot]) > _SAME_PAGE:
        bind.update(_over(painted, scheme))
    # An accent still holding Microsoft's shipped value says nothing about the brand.
    real = [
        s
        for s in (f"accent{i}" for i in range(1, 7))
        if s in scheme and not is_stock_accent(scheme[s])
    ]
    for i, slot in enumerate(real, start=1):
        bind[f"accent-{i}"] = slot
    slots = [s for s in (*_DARK_SLOTS, *_LIGHT_SLOTS) if s in scheme]
    dark = [s for s in _DARK_SLOTS if s in scheme]
    if dark:
        # A dark carrying hue is this brand's dark; pure black is every brand's black.
        hued = [s for s in dark if _chroma(scheme[s]) >= _MIN_CHROMA]
        chosen = min(hued or dark, key=lambda s: relative_luminance(scheme[s]))
        lightest = max(slots, key=lambda s: relative_luminance(scheme[s]))
        # Some schemes set every dark slot light — there is then nothing to reverse out
        # of, and binding one gives a pair whose ink and ground are the same colour.
        if contrast_ratio(scheme[lightest], scheme[chosen]) < AA_NORMAL:
            chosen = max(slots, key=lambda s: contrast_ratio(scheme[lightest], scheme[s]))
        bind["inverse"] = chosen

    face, ramp = _typography(prs, slide_h)
    theme: dict[str, Any] = {
        "name": template.stem[:40],
        "template": template.name,
        "drop_template_slides": True,
        "bind": bind,
    }
    clear = None if art is None else _clear_run(art, _CHROME_BAND)
    if clear is not None:
        start, end = _columns(clear)
        theme["chrome"] = {
            name: {"at": {"cols": {"from": start, "to": end}}} for name in CHROME_ORDER
        }
    type_block: dict[str, Any] = {}
    if face:
        type_block["face"] = face
        type_block["heading_face"] = face
    if ramp:
        # Measured in points off the template's own slides, so points are what is
        # written; reference_height is the canvas they were measured on.
        type_block["reference_height"] = round(slide_h, 3)
        type_block["ramp"] = {
            name: {"pt": round(rung * slide_h, 1)} for name, rung in sorted(ramp.items())
        }
    if type_block:
        theme["type"] = type_block
    if prefer:
        # Without it the written theme cannot resolve the layout it was derived against.
        theme["compose_layout"] = prefer
    return theme


def notes(template: str | Path, *, prefer: str | None = None) -> list[str]:
    """What a reader of the derived theme should know about this template."""
    prs = open_presentation(template)
    scheme = parse_color_scheme(
        read_theme_xml(pick_compose_layout(prs, prefer=prefer).slide_master)
    )
    page_slot, ink_slot = _surfaces(scheme)
    out = [
        f"canvas {Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f}in",
        f"composes on {pick_compose_layout(prs, prefer=prefer).name!r} "
        f"across {len(prs.slide_masters)} master(s)",
        f"page {page_slot}={scheme[page_slot]}, ink {ink_slot}={scheme[ink_slot]} "
        f"({contrast_ratio(scheme[ink_slot], scheme[page_slot]):.1f}:1)",
    ]
    stock = [
        s for s in (f"accent{i}" for i in range(1, 7)) if s in scheme and is_stock_accent(scheme[s])
    ]
    if stock:
        out.append(f"ignored {len(stock)} unedited stock accent(s): {', '.join(stock)}")
    if ink_slot != "dk1":
        out.append(f"ink came from {ink_slot}, not dk1 — dk1 is not this template's darkest")
    return out
