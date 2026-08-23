"""Render a panel once, optionally cut it into regions, and place it on a slide."""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from PIL import Image
from pptx.util import Inches

from pf_core.log import get_logger
from pf_core.utils.env import resolve_int
from pf_core.utils.io import atomic_write_bytes

from pptxkit.errors import LayoutError
from pptxkit.panels.cache import PanelRenderer, cached_png
from pptxkit.panels.model import Panel

logger = get_logger(__name__)

_SCALE_DEFAULT = 2


def _scale() -> int:
    return resolve_int(None, "PPTXKIT_SHOT_SCALE", default=_SCALE_DEFAULT)


def _save_atomically(image: Image.Image, path: Path) -> None:
    """Write the crop so a build reading a same-named one never sees a torn PNG."""
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    atomic_write_bytes(path, buffer.getvalue())


def place_panel(
    ctx,
    panel: Panel,
    *,
    left: float,
    top: float,
    width: float | None = None,
    height: float | None = None,
    slice_by: str | None = None,
    max_height: float | None = None,
    render: PanelRenderer,
) -> dict[str, Any]:
    """Place ``panel`` at (``left``, ``top``) inches; return its pictures by region name.

    Pass exactly one of ``width`` / ``height`` to scale while preserving aspect ratio.
    ``slice_by`` cuts the one rendered PNG at region boundaries into a picture each, so
    the parts can animate independently. ``max_height`` raises rather than silently
    placing a picture taller than that budget.
    """
    if (width is None) == (height is None):
        raise LayoutError("place_panel needs exactly one of width or height")
    if slice_by and not panel.regions:
        raise LayoutError(f"panel has no regions to slice by {slice_by!r}")

    scale = _scale()
    png = cached_png(panel, scale=scale, theme_hash=ctx.theme.hash, render=render)
    image = Image.open(png)
    full_w, full_h = image.size
    if width is not None:
        inches_w, inches_h = width, width * full_h / full_w
    elif height is not None:
        inches_w, inches_h = height * full_w / full_h, height
    else:  # unreachable — the guard above requires exactly one
        raise LayoutError("place_panel needs exactly one of width or height")
    if max_height is not None and inches_h > max_height:
        raise LayoutError(
            f"slide {ctx.spec.index}: panel placed at {inches_h:.2f}in tall, "
            f"exceeds the {max_height:.2f}in budget"
        )
    per_px = inches_w / full_w

    if not slice_by:
        picture = ctx.slide.shapes.add_picture(
            str(png), Inches(left), Inches(top), width=Inches(inches_w)
        )
        ctx.manifest.record(picture, rendered="image")
        logger.info("panel_placed", regions=1, sliced=False)
        return {"": picture}

    placed: dict[str, Any] = {}
    for region in panel.regions:
        # Region coords are CSS px; the render is scale x larger, so the crop box scales too.
        box = (
            region.left * scale,
            region.top * scale,
            (region.left + region.width) * scale,
            (region.top + region.height) * scale,
        )
        # Image.crop() never raises on an out-of-bounds box — it zero-pads instead — so a
        # stale region literal or an autocropped render would ship a black-padded slice.
        if box[0] < 0 or box[1] < 0 or box[2] > full_w or box[3] > full_h:
            raise LayoutError(
                f"region {region.name!r} ({region.left},{region.top})-"
                f"({region.left + region.width},{region.top + region.height}) at scale "
                f"{scale} does not fit the rendered panel ({full_w}x{full_h})"
            )
        part_path = png.with_name(f"{png.stem}-{region.name}.png")
        _save_atomically(image.crop(box), part_path)
        picture = ctx.slide.shapes.add_picture(
            str(part_path),
            Inches(left + region.left * scale * per_px),
            Inches(top + region.top * scale * per_px),
            width=Inches(region.width * scale * per_px),
        )
        ctx.manifest.record(picture, rendered="image")
        placed[region.name] = picture
    logger.info("panel_placed", regions=len(placed), sliced=True)
    return placed
