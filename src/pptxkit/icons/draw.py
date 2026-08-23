"""Put an icon on a slide — a preset shape where the name is one, else a freeform."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml  # not lxml's: python-pptx keys its shape classes off this
from pptx.util import Emu, Inches

from pptxkit.icons.load import Glyph, load, supplied
from pptxkit.icons.path import UNITS
from pptxkit.icons.shapes import SHAPES
from pptxkit.theme.model import Rect

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

_SHAPE = """
<p:sp xmlns:p="{p}" xmlns:a="{a}">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect l="0" t="0" r="r" b="b"/>
      <a:pathLst><a:path w="{units}" h="{units}">{path}</a:path></a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>
"""


def place_icon(slide, name: str, box: Rect, *, fill: str, theme=None):
    """Draw the icon called ``name``, whichever kind it is. Returns the new shape.

    A configured directory wins outright, so a brand shipping its own ``circle.svg``
    gets that and not the preset.
    """
    if name in SHAPES and supplied(name, theme) is None:
        return place_preset(slide, name, box, fill=fill)
    return place_glyph(slide, load(name, theme=theme), box, fill=fill)


def place_preset(slide, name: str, box: Rect, *, fill: str):
    """Draw a plain shape as DrawingML preset geometry, squared off like a glyph."""
    side = min(box.width, box.height)
    shape = slide.shapes.add_shape(
        SHAPES[name],
        Inches(box.left + (box.width - side) / 2),
        Inches(box.top + (box.height - side) / 2),
        Inches(side),
        Inches(side),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.name = f"Icon {name}"
    return shape


def place_glyph(slide, glyph: Glyph, box: Rect, *, fill: str):
    """Draw ``glyph`` filling ``box``, painted ``fill``. Returns the new shape.

    The box is squared off and centred first: stretching a square viewBox into a wide
    placement is the one distortion nothing downstream can undo.
    """
    side = min(box.width, box.height)
    left = box.left + (box.width - side) / 2
    top = box.top + (box.height - side) / 2
    tree = slide.shapes._spTree
    element = parse_xml(
        _SHAPE.format(
            p=_P,
            a=_A,
            sid=_next_id(tree),
            name=f"Icon {glyph.name}",
            x=Emu(int(left * 914400)),
            y=Emu(int(top * 914400)),
            cx=Emu(int(side * 914400)),
            cy=Emu(int(side * 914400)),
            units=UNITS,
            path=glyph.drawingml(),
            fill=fill,
        )
    )
    tree.append(element)
    return slide.shapes[-1]


def _next_id(tree) -> int:
    """One past the highest shape id in the tree — a duplicate id opens as corrupt."""
    used = [
        int(el.get("id", 0))
        for el in tree.iter(f"{{{_P}}}cNvPr", f"{{{_P}}}cNvGrpSpPr")
        if el.get("id")
    ]
    return max(used, default=1) + 1
