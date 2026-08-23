"""Put a fitted picture and its scrim on a slide, in the OOXML those need.

python-pptx has no API for a masked picture or a translucent fill, so both are written
as elements: ``a:prstGeom`` masks the picture, and ``a:alpha`` inside a fill's colour is
what makes a scrim a scrim rather than a lid.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches
from lxml import etree

from pptxkit.errors import LayoutError
from pptxkit.imagery.fit import ImageFit
from pptxkit.imagery.scrim import Scrim
from pptxkit.theme.model import Rect

_PRST = {"none": "rect", "circle": "ellipse", "rounded": "roundRect"}
_PER_UNIT = 100000
# a:lin angle for a fill running down the shape, in 60000ths of a degree.
_DOWNWARD = 5400000


def place_picture(slide, path: str, fit: ImageFit, *, mask: str = "none", radius: float = 0.0):
    """Add ``path`` at ``fit.dest``, showing ``fit.window``, masked if asked.

    The trim becomes ``a:srcRect``, so the source file itself is never rewritten.
    """
    if mask not in _PRST:
        raise LayoutError(f"unknown mask {mask!r}; expected one of {', '.join(_PRST)}")
    dest = fit.dest
    pic = slide.shapes.add_picture(
        path, Inches(dest.left), Inches(dest.top), Inches(dest.width), Inches(dest.height)
    )
    left, top, right, bottom = fit.trim
    pic.crop_left, pic.crop_top = left, top
    pic.crop_right, pic.crop_bottom = right, bottom
    if mask != "none":
        _mask(pic, mask, radius)
    return pic


def paint_scrim(slide, rect: Rect, scrim: Scrim):
    """Lay ``scrim`` over ``rect``. Returns the shape, or None when it is invisible."""
    if scrim.opacity <= 0:
        return None
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(rect.left),
        Inches(rect.top),
        Inches(rect.width),
        Inches(rect.height),
    )
    shape.line.fill.background()
    shape.shadow.inherit = False
    spPr = shape._element.spPr
    _clear_fill(spPr)
    fill = (
        _solid(scrim.colour, scrim.opacity)
        if scrim.gradient == "none"
        else _gradient(scrim.colour, scrim.opacity, scrim.gradient)
    )
    spPr.insert(list(spPr).index(spPr.find(qn("a:prstGeom"))) + 1, fill)
    return shape


def _mask(pic, mask: str, radius: float) -> None:
    geom = pic._element.spPr.find(qn("a:prstGeom"))
    geom.set("prst", _PRST[mask])
    avLst = geom.find(qn("a:avLst"))
    for child in list(avLst):
        avLst.remove(child)
    if mask != "rounded":
        return
    gd = etree.SubElement(avLst, qn("a:gd"))
    gd.set("name", "adj")
    gd.set("fmla", f"val {round(radius * _PER_UNIT)}")


def _clear_fill(spPr) -> None:
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        found = spPr.find(qn(tag))
        if found is not None:
            spPr.remove(found)


def _solid(colour: str, opacity: float):
    fill = etree.Element(qn("a:solidFill"))
    fill.append(_colour(colour, opacity))
    return fill


def _gradient(colour: str, opacity: float, direction: str):
    """A one-sided fade: fully ``opacity`` at ``direction``'s edge, clear at the other."""
    fill = etree.Element(qn("a:gradFill"))
    stops = etree.SubElement(fill, qn("a:gsLst"))
    alphas = (opacity, 0.0) if direction == "top" else (0.0, opacity)
    for position, alpha in zip((0, _PER_UNIT), alphas, strict=True):
        stop = etree.SubElement(stops, qn("a:gs"))
        stop.set("pos", str(position))
        stop.append(_colour(colour, alpha))
    line = etree.SubElement(fill, qn("a:lin"))
    line.set("ang", str(_DOWNWARD))
    line.set("scaled", "0")
    return fill


def _colour(colour: str, opacity: float):
    srgb = etree.Element(qn("a:srgbClr"))
    srgb.set("val", colour)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(round(opacity * _PER_UNIT)))
    return srgb
