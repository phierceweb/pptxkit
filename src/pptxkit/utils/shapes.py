"""Styled autoshape and text primitives for python-pptx.

Colors are ``RGBColor``; positions and sizes are in inches.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_FONT = "Helvetica"

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
"""Spec ``align:`` values to paragraph alignment."""

ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
"""Spec ``anchor:`` values to text-frame vertical anchoring."""

ALIGNS = tuple(ALIGN)
ANCHORS = tuple(ANCHOR)


def solid(shape, color: RGBColor, line: RGBColor | None = None) -> None:
    """Give ``shape`` a solid fill (optional 1pt line, else none) and drop the theme shadow."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False


def para(
    tf,
    text: str,
    size: float,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    first: bool = False,
    space_after: float = 6,
    font: str = DEFAULT_FONT,
):
    """Append a styled single-run paragraph to text frame ``tf``.

    When ``first`` is set and the frame's first paragraph is still empty, that
    paragraph is reused instead of adding a new one. Returns the paragraph.
    """
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return p


def textbox(
    slide, x: float, y: float, w: float, h: float, anchor=MSO_ANCHOR.TOP, wrap: bool = True
):
    """Add a zero-margin textbox at (x, y) sized w×h inches; return its text frame."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def rrect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius: float = 0.08,
):
    """Add a rounded rectangle (inches) with a solid fill; ``radius`` is the corner adjustment."""
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    solid(sp, fill, line)
    try:
        sp.adjustments[0] = radius
    except Exception:  # noqa: BLE001 — some shapes expose no adjustment handle
        pass
    return sp


def rect(
    slide, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None
):
    """Add a plain rectangle (inches) with a solid fill."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    solid(sp, fill, line)
    return sp


def notes(slide, text: str) -> None:
    """Set the slide's speaker-notes text."""
    slide.notes_slide.notes_text_frame.text = text


def bring_to_front(shape) -> None:
    """Move ``shape`` to the top of its slide's z-order — python-pptx exposes no z-order API."""
    el = shape._element
    el.getparent().append(el)
