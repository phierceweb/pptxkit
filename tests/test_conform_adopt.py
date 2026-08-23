"""What `--adopt` is allowed to write, and what it must refuse to write over. Adoption is the one
thing `conform` does outside `out/`: it writes into the theme directory, beside the template it
binds to, where a file may be somebody's hand-edited theme and nothing would recover it."""

from __future__ import annotations

import pytest
import yaml
from pptx import Presentation
from typer.testing import CliRunner

from pptxkit.cli import app
from pptxkit.compile.build import build_deck
from pptxkit.conform import conform, install, plan
from pptxkit.errors import ThemeError

runner = CliRunner()

# One any template can draw, and one no placement can hold — 60 columns leave less
# width than the theme's gutter pads a cell by, which is a LayoutError.
FINE = {
    "title": "A slide",
    "place": [{"at": {"cols": "full"}, "bullets": {"items": ["One", "Two"]}}],
}
DOOMED = {
    "title": "Too many columns",
    "place": [{"at": {"cols": "full"}, "table": {"rows": [["c"] * 60]}}],
}


def _template(path, slides: int = 1):
    """A stock Office deck is a usable template: it carries a theme part to derive from.
    ``slides`` gives two templates of the same name different bytes."""
    prs = Presentation()
    for _ in range(slides):
        prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    return path


@pytest.fixture
def themes(tmp_path, monkeypatch):
    """Adoption writes where ``build`` resolves theme names from, so it reads the
    same env var — point both at a scratch directory."""
    root = tmp_path / "themes"
    root.mkdir()
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(root))
    return root


@pytest.fixture
def template(themes):
    """A template is adopted where it lives, so it starts in the theme directory."""
    return _template(themes / "Brand.pptx")


def test_an_adopted_theme_builds_a_deck_that_names_it(themes, template, tmp_path):
    """The end of onboarding: `theme: brand` and nothing else supplied."""
    conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")

    spec = tmp_path / "d.deck.yaml"
    spec.write_text("theme: brand\nout: d.pptx\n---\ntitle: T\n")
    result = build_deck(spec, out=tmp_path / "d.pptx")

    assert result.deck.is_file()
    assert (themes / "brand.theme.yaml").is_file()
    assert not (themes / "assets").exists(), "adoption must not copy the template"


def test_the_adopted_theme_carries_the_adopted_name(themes, template, tmp_path):
    """The derivation names the theme after the template stem. Adopting renames it,
    so the file, its name and the `theme:` a deck writes are one word."""
    conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")

    assert yaml.safe_load((themes / "brand.theme.yaml").read_text())["name"] == "brand"


def test_install_repoints_the_theme_at_the_template_beside_it(themes, template, tmp_path):
    """The installed theme names its template as a bare filename resolved beside itself — the
    source path is a fact about where it was derived, not about where it now lives."""
    derived = tmp_path / "derived.theme.yaml"
    derived.write_text(
        yaml.safe_dump({"name": "derived", "template": str(tmp_path / "elsewhere.pptx")})
    )

    installed = install(plan("brand", template), derived)

    assert yaml.safe_load(installed.read_text())["template"] == "Brand.pptx"
    assert installed.parent == template.parent


def test_re_adopting_the_same_template_keeps_the_hand_edits(themes, template, tmp_path):
    """The theme beside the template IS this run's sidecar, so a refresh carries the edits
    through verbatim. Georgia is the tell — no derivation from a stock template makes it."""
    hand_edited = themes / "brand.theme.yaml"
    hand_edited.write_text(
        "name: brand\ntemplate: Brand.pptx\ntype:\n  face: Georgia\n  heading_face: Georgia\n"
    )

    conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")

    assert yaml.safe_load(hand_edited.read_text())["type"]["face"] == "Georgia"


def test_force_re_derives_and_discards_the_hand_edits(themes, template, tmp_path):
    """`--force` is the only way to throw tuning away: it ignores the sidecar that would
    otherwise restore it."""
    (themes / "brand.theme.yaml").write_text(
        "name: brand\ntemplate: Brand.pptx\ntype:\n  face: Georgia\n"
    )

    conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand", force=True)

    assert (
        yaml.safe_load((themes / "brand.theme.yaml").read_text()).get("type", {}).get("face")
        != "Georgia"
    )


def test_a_template_outside_the_theme_directory_is_refused(themes, tmp_path):
    """The theme binds its template by bare filename resolved beside itself, so adopting one
    that lives elsewhere would write a theme pointing at nothing. One directory, one copy."""
    stray = _template(tmp_path / "Downloads-Brand.pptx")

    with pytest.raises(ThemeError, match="a template is adopted where it lives"):
        conform(stray, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")

    assert not (themes / "brand.theme.yaml").exists()


def test_a_name_already_bound_to_a_different_template_is_refused(themes, template, tmp_path):
    """Two brands can both want `brand`. The second must not silently repoint the first."""
    other = _template(themes / "Other.pptx", slides=3)
    (themes / "brand.theme.yaml").write_text(
        yaml.safe_dump({"name": "brand", "template": other.name})
    )

    with pytest.raises(ThemeError, match="binds"):
        conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")

    assert yaml.safe_load((themes / "brand.theme.yaml").read_text())["template"] == "Other.pptx"


def test_the_same_template_adopted_twice_is_not_a_collision(themes, template, tmp_path):
    """Adopting one template under two names is ordinary — a light and a dark theme
    over the same brand deck. Only *differing* bytes are a collision."""
    conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")
    conform(template, tmp_path / "out2", exercises={"fine": FINE}, adopt="brand-dark")

    assert (themes / "brand-dark.theme.yaml").is_file()


def test_a_name_that_escapes_the_theme_directory_is_refused(themes, template, tmp_path):
    """`themes` is a directory below `tmp_path`, so `../escaped` names a real file
    outside it — the write a bare-name check exists to stop."""
    with pytest.raises(ThemeError, match="bare theme name"):
        conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="../escaped")

    assert not (tmp_path / "escaped.yaml").exists()


def test_the_name_is_vetted_before_a_single_exercise_is_built(themes, template, tmp_path):
    """A collision found at the end of a two-minute run has already cost the run it
    exists to prevent. Nothing may be written before the target is known good."""
    other = _template(themes / "Other.pptx", slides=3)
    (themes / "brand.theme.yaml").write_text(
        yaml.safe_dump({"name": "brand", "template": other.name})
    )
    out = tmp_path / "out"

    with pytest.raises(ThemeError):
        conform(template, out, exercises={"fine": FINE}, adopt="brand")

    assert not out.exists()


def test_a_failing_exercise_does_not_block_adoption(themes, template, tmp_path):
    """A FAIL says this template cannot carry that shape, and editing the theme is the answer —
    which needs the theme somewhere that survives, not the disposable directory."""
    result = conform(
        template, tmp_path / "out", exercises={"fine": FINE, "doomed": DOOMED}, adopt="brand"
    )

    assert [name for name, _ in result.failed] == ["doomed"]
    assert result.adopted == themes / "brand.theme.yaml"
    assert (themes / "brand.theme.yaml").is_file()


def test_nothing_is_adopted_when_nothing_built(themes, template, tmp_path):
    """Every exercise failing means the theme does not describe the template at all.
    Installing that under a project name would bless an artefact that builds nothing."""
    result = conform(template, tmp_path / "out", exercises={"doomed": DOOMED}, adopt="brand")

    assert not (themes / "brand.theme.yaml").exists()
    assert result.adopted is None
    assert "no exercise built" in result.refused
    assert "not adopted:" in result.report()


def test_a_missing_template_names_itself(themes, tmp_path):
    """python-pptx raises `PackageNotFoundError` for this, which reaches the CLI as a
    traceback rather than as a message with the path in it."""
    with pytest.raises(ThemeError, match="template not found"):
        conform(tmp_path / "absent.pptx", tmp_path / "out", exercises={"fine": FINE})


def test_a_template_that_will_not_open_names_itself(themes, template, tmp_path):
    """Present but unreadable — a truncated download, or a .key someone renamed.
    `notes()` is the first thing a run opens, so it is the site that has to say so."""
    template.write_bytes(b"PK\x03\x04 truncated")

    with pytest.raises(ThemeError, match=r"template .*Brand\.pptx is not a readable \.pptx"):
        conform(template, tmp_path / "out", exercises={"fine": FINE})


def test_the_cli_passes_adopt_and_force_through(themes, template, tmp_path, monkeypatch):
    """The flags are the whole feature. One exercise stands in for the registry the real
    command runs — this is about the two flags, not the corpus."""
    monkeypatch.setattr("pptxkit.conform.run.EXERCISE", {"fine": FINE})
    other = _template(themes / "Other.pptx", slides=3)
    (themes / "brand.theme.yaml").write_text(
        yaml.safe_dump({"name": "brand", "template": other.name})
    )

    refused = runner.invoke(
        app, ["conform", str(template), "-o", str(tmp_path / "a"), "--adopt", "brand"]
    )
    assert refused.exit_code != 0
    assert yaml.safe_load((themes / "brand.theme.yaml").read_text())["template"] == "Other.pptx"

    forced = runner.invoke(
        app, ["conform", str(template), "-o", str(tmp_path / "b"), "--adopt", "brand", "--force"]
    )
    assert forced.exit_code == 0, forced.output
    assert "adopted -> " in forced.output
    assert yaml.safe_load((themes / "brand.theme.yaml").read_text())["template"] == "Brand.pptx"


def test_a_name_that_starts_legal_and_then_escapes_is_refused(themes, template, tmp_path):
    """`../escaped` fails on its first character, so it cannot tell `match` from `fullmatch`;
    `brand/../../escaped` is the one that needs the end anchored."""
    with pytest.raises(ThemeError, match="bare theme name"):
        conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand/../../escaped")

    assert not (tmp_path / "escaped.yaml").exists()
    assert not (themes / "brand.theme.yaml").exists(), "a legal prefix was written as a theme"


def test_a_sidecar_theme_beside_the_template_is_installed_instead_of_the_derivation(
    themes, template, tmp_path
):
    """A tuned theme kept beside the binary survives the round trip verbatim. Georgia is the tell
    — no derivation from a stock template produces it."""
    sidecar = template.with_name("brand.theme.yaml")
    sidecar.write_text(
        "name: whatever\ntemplate: old.pptx\ntype:\n  face: Georgia\n  heading_face: Georgia\n"
    )

    result = conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")

    assert result.adopted is not None
    installed = yaml.safe_load((themes / "brand.theme.yaml").read_text())
    assert installed["type"]["face"] == "Georgia"
    assert installed["name"] == "brand"
    assert installed["template"] == template.name
    assert any("sidecar" in n for n in result.notes), result.notes


def test_without_a_sidecar_the_derivation_is_what_installs(themes, template, tmp_path):
    """The sidecar is an override, not a requirement — a first meeting still derives."""
    conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="fresh")
    installed = yaml.safe_load((themes / "fresh.theme.yaml").read_text())
    assert installed.get("type", {}).get("face") != "Georgia"


def test_adoption_writes_the_theme_and_nothing_else(themes, template, tmp_path):
    """One directory, one copy. Adoption adds a theme file beside the template and
    changes nothing else about the directory — put a `shutil.copy` back into `install`
    and this reddens on the extra entry."""
    before = {p.name for p in themes.iterdir()}

    conform(template, tmp_path / "out", exercises={"fine": FINE}, adopt="brand")

    added = {p.name for p in themes.iterdir()} - before
    assert added == {"brand.theme.yaml"}, added


def test_no_module_copies_a_template(themes, template, tmp_path):
    """The gate on the whole design. A second copy is a second filename, and a second
    filename drifts from the first.

    `out/` is regenerated per run and may hold whatever it likes; nothing that outlives
    a run may copy a brand binary.
    """
    import pathlib
    import re

    src = pathlib.Path(__file__).resolve().parents[1] / "src/pptxkit"
    offenders = [
        f"{path.relative_to(src)}:{i}"
        for path in src.rglob("*.py")
        for i, line in enumerate(path.read_text().splitlines(), 1)
        if re.search(r"shutil\.copy", line)
    ]
    assert offenders == [], (
        "a template is adopted where it lives and is never copied: " + ", ".join(offenders)
    )
