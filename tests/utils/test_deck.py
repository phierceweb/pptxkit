"""Package-level fixups python-pptx leaves undone."""

from __future__ import annotations

import re
import textwrap
import zipfile

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

from pptxkit.compile import build_deck
from pptxkit.theme import blank_presentation
from pptxkit.utils.deck import register_notes_master


def _with_a_note():
    prs = blank_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "a speaker note"
    return prs


def test_a_notes_master_relationship_is_declared_on_the_presentation():
    """python-pptx adds the relationship and never declares it; Keynote then
    refuses to read the package at all."""
    prs = _with_a_note()
    register_notes_master(prs)

    entries = prs._element.find(qn("p:notesMasterIdLst"))
    assert entries is not None, "no p:notesMasterIdLst was written"
    declared = [e.get(qn("r:id")) for e in entries]
    related = [r.rId for r in prs.part.rels.values() if r.reltype == RT.NOTES_MASTER]
    assert declared == related


def test_the_declaration_sits_where_the_schema_puts_it():
    """p:notesMasterIdLst follows p:sldMasterIdLst; anywhere else is invalid."""
    prs = _with_a_note()
    register_notes_master(prs)
    order = [el.tag.split("}")[-1] for el in prs._element]
    assert order[:2] == ["sldMasterIdLst", "notesMasterIdLst"]


def test_a_deck_with_no_notes_declares_no_notes_master():
    """An empty p:notesMasterIdLst would name a relationship that is not there."""
    prs = blank_presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    register_notes_master(prs)
    assert prs._element.find(qn("p:notesMasterIdLst")) is None


def test_calling_it_twice_does_not_declare_the_master_twice():
    prs = _with_a_note()
    register_notes_master(prs)
    register_notes_master(prs)
    assert len(prs._element.findall(qn("p:notesMasterIdLst"))) == 1


def test_a_built_deck_carrying_notes_declares_its_notes_master(tmp_path):
    """The end-to-end guard: build_deck must run the fixup before it saves."""
    spec = tmp_path / "d.deck.yaml"
    spec.write_text(
        textwrap.dedent("""
        theme: base
        title: Notes
        out: Notes.pptx
        ---
        title: A slide that carries a speaker note
        notes: Say this part out loud.
    """).lstrip()
    )
    result = build_deck(spec)

    with zipfile.ZipFile(result.deck) as z:
        presentation = z.read("ppt/presentation.xml").decode()
        rels = z.read("ppt/_rels/presentation.xml.rels").decode()
        assert re.search(r"ppt/notesSlides/notesSlide1\.xml", " ".join(z.namelist()))
    rid = re.search(r'<Relationship Id="([^"]+)" Type="[^"]*notesMaster"', rels)
    assert rid, "no notes master relationship was written"
    assert f'<p:notesMasterId r:id="{rid.group(1)}"/>' in presentation
