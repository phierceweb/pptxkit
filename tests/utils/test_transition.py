"""Slide transitions — the ``<p:transition>`` element, not the timing tree. Every direction
vocabulary is per element: the "orientation" effects take ``dir`` not ``orient``, ``strips`` takes
corners only, and the child order survives an animation being written first."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from pptxkit.errors import LayoutError
from pptxkit.motion import add_click_build, add_transition, transition_xml
from pptxkit.motion.transition import EFFECTS

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_the_base_schema_offers_exactly_twenty_one_effects():
    assert len(EFFECTS) == 21


def test_a_transition_names_its_effect_as_a_child_element():
    assert transition_xml("fade") == (
        f'<p:transition xmlns:p="{_P}" spd="fast"><p:fade/></p:transition>'
    )


def test_speed_is_written_as_spd():
    assert 'spd="slow"' in transition_xml("wipe", speed="slow")


def test_an_orientation_effect_takes_dir_not_orient():
    """Every readable summary of this schema calls these "orientation" transitions and
    says `orient`. The XSD rejects `orient` on all four of them."""
    assert '<p:blinds dir="vert"/>' in transition_xml("blinds", direction="vert")


def test_strips_refuses_an_edge_direction():
    """`strips` takes corners only — a shared l/u/r/d list is invalid here alone."""
    with pytest.raises(LayoutError, match=r"'strips' has no direction 'l'.*lu, ru, ld, rd"):
        transition_xml("strips", direction="l")


def test_push_refuses_a_corner_direction():
    with pytest.raises(LayoutError, match=r"'push' has no direction 'lu'"):
        transition_xml("push", direction="lu")


def test_an_effect_with_no_direction_refuses_one():
    with pytest.raises(LayoutError, match=r"'fade' takes no direction, got 'l'"):
        transition_xml("fade", direction="l")


def test_an_unknown_effect_lists_the_ones_that_exist():
    with pytest.raises(LayoutError, match=r"unknown transition 'ripple'; known"):
        transition_xml("ripple")


def test_an_unknown_speed_is_refused():
    with pytest.raises(LayoutError, match=r"speed must be one of slow, med, fast"):
        transition_xml("fade", speed="quick")


def test_transition_precedes_timing_even_when_the_animation_was_written_first():
    """``CT_Slide`` is an xsd:sequence. LibreOffice repairs a wrong order on import,
    so no render or round trip can show this being broken."""
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_click_build(slide, [box.shape_id])

    add_transition(slide, "fade")

    tags = [el.tag.split("}")[1] for el in slide._element]
    assert tags == ["cSld", "clrMapOvr", "transition", "timing"]


def test_a_second_transition_on_one_slide_is_refused():
    prs = Presentation()
    slide = _blank_slide(prs)
    add_transition(slide, "fade")
    with pytest.raises(LayoutError, match="already carries a transition"):
        add_transition(slide, "wipe")
    assert len(slide._element.findall(qn("p:transition"))) == 1
