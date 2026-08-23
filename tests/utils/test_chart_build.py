"""``add_chart_build`` — the ``<p:bldGraphic>``/``<a:bldChart>`` chart build, as distinct from the
text-shape ``<p:bldP>`` builds ``motion/builds.py`` emits. Malformed timing parses fine in
LibreOffice and lxml but makes PowerPoint offer to repair the file, so these assert the structural
facts known to cause it: the namespace, the placement inside ``<p:bldLst>``, and that ``spid``
names a real shape."""

from __future__ import annotations

import re

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from pptxkit.errors import LayoutError
from pptxkit.motion import add_chart_build

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _timing(slide):
    return slide._element.find(qn("p:timing"))


def test_it_appends_a_timing_tree():
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, box.shape_id, "category")
    assert _timing(slide) is not None


def test_bld_graphic_sits_inside_bld_lst_not_beside_it():
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, box.shape_id, "category")
    timing = _timing(slide)

    bld_lst = timing.find(qn("p:bldLst"))
    assert bld_lst is not None
    bld_graphic = bld_lst.find(qn("p:bldGraphic"))
    assert bld_graphic is not None

    # Not a sibling of bldLst, and not floating elsewhere in the tree.
    assert timing.find(qn("p:bldGraphic")) is None
    assert list(bld_lst) == [bld_graphic]


def test_the_namespaces_are_presentationml_and_drawingml():
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, box.shape_id, "category")
    timing = _timing(slide)

    bld_graphic = timing.find(f"{{{_P}}}bldLst/{{{_P}}}bldGraphic")
    assert bld_graphic is not None
    bld_chart = bld_graphic.find(f"{{{_P}}}bldSub/{{{_A}}}bldChart")
    assert bld_chart is not None
    # A namespace mixup (p:bldChart, or a:bldGraphic) parses but does nothing —
    # confirm the tags actually resolved to the namespaces above, not by name alone.
    assert bld_graphic.tag == f"{{{_P}}}bldGraphic"
    assert bld_chart.tag == f"{{{_A}}}bldChart"


def test_spid_matches_the_real_shape_id():
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, box.shape_id, "category")
    timing = _timing(slide)
    bld_graphic = timing.find(f"{{{_P}}}bldLst/{{{_P}}}bldGraphic")
    assert bld_graphic.get("spid") == str(box.shape_id)


@pytest.mark.parametrize(
    "by, bld",
    [
        ("category", "category"),
        ("series", "series"),
        ("all", "allAtOnce"),
    ],
)
def test_by_maps_to_the_ooxml_bld_value(by, bld):
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, box.shape_id, by)
    timing = _timing(slide)
    bld_chart = timing.find(f"{{{_P}}}bldLst/{{{_P}}}bldGraphic/{{{_P}}}bldSub/{{{_A}}}bldChart")
    assert bld_chart.get("bld") == bld


def test_an_unknown_by_raises_naming_the_valid_values():
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    with pytest.raises(LayoutError, match=r"category.*series.*element.*all|by"):
        add_chart_build(slide, box.shape_id, "sideways")


def test_grp_id_is_zero():
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, box.shape_id, "category")
    timing = _timing(slide)
    bld_graphic = timing.find(f"{{{_P}}}bldLst/{{{_P}}}bldGraphic")
    assert bld_graphic.get("grpId") == "0"


def test_survives_a_save_reload_round_trip(tmp_path):
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, box.shape_id, "category")
    path = tmp_path / "t.pptx"
    prs.save(str(path))

    reloaded = Presentation(str(path))
    timing = _timing(reloaded.slides[0])
    assert timing is not None
    bld_chart = timing.find(f"{{{_P}}}bldLst/{{{_P}}}bldGraphic/{{{_P}}}bldSub/{{{_A}}}bldChart")
    assert bld_chart is not None
    assert bld_chart.get("bld") == "category"


def _chart_slide():
    """A real graphicFrame, so spid assertions are against an actual chart."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    slide = _blank_slide(prs)
    data = CategoryChartData()
    data.categories = ["A", "B", "C", "D"]
    data.add_series("S", (1, 2, 3, 4))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), data
    )
    return slide, frame.shape_id


def _timing_xml(slide) -> str:
    from lxml import etree

    return etree.tostring(_timing(slide)).decode()


def test_each_category_gets_its_own_click():
    """<a:bldChart> only declares the build type — without one animation node per
    element PowerPoint fades the whole chart at once."""
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "category", parts=4)
    xml = _timing_xml(slide)

    assert xml.count('nodeType="clickEffect"') == 5  # background + 4 categories
    assert re.findall(r'categoryIdx="(-?\d+)"', xml)[::2] == ["-1", "0", "1", "2", "3"]


def test_a_build_with_no_parts_falls_back_to_one_click():
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "category")
    xml = _timing_xml(slide)
    assert xml.count('nodeType="clickEffect"') == 1
    assert "<a:chart" not in xml


def test_every_click_effect_has_a_matching_build_group():
    """A click effect with no bldGraphic of that grpId is what PowerPoint repairs."""
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "category", parts=3)
    xml = _timing_xml(slide)

    assert xml.count('nodeType="clickEffect"') == 4
    assert re.findall(r'<p:bldGraphic spid="\d+" grpId="(\d+)"', xml) == ["0", "1", "2", "3"]


def test_a_series_build_staggers_by_series_not_category():
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "series", parts=2)
    xml = _timing_xml(slide)
    assert re.findall(r'seriesIdx="(-?\d+)"', xml)[::2] == ["-1", "0", "1"]


def test_bars_grow_out_of_the_axis_while_the_grid_fades():
    """PowerPoint's wipe filters are inverted — "grow up" is subtype 1 / wipe(down).
    The background fades because wiping gridlines upward reads as a glitch."""
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "category", parts=4)
    xml = _timing_xml(slide)

    assert re.findall(r'filter="([^"]+)"', xml) == ["fade"] + ["wipe(down)"] * 4
    assert re.findall(r'presetID="(\d+)"', xml) == ["10"] + ["22"] * 4


def test_the_entrance_kind_is_selectable():
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "category", parts=2, kind="fade")
    assert re.findall(r'filter="([^"]+)"', _timing_xml(slide)) == ["fade"] * 3


def test_every_chart_target_declares_its_build_step():
    """``bldStep`` is ``use="required"`` on ``CT_AnimationChartElement``, and omitting it is the one
    thing in this module's output that no render or round trip shows."""
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "category", parts=3)
    xml = _timing_xml(slide)

    assert xml.count("<a:chart ") == len(re.findall(r'bldStep="', xml))


def test_the_build_step_names_the_axis_and_the_background_names_the_grid():
    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "category", parts=3)
    steps = re.findall(r'bldStep="([^"]+)"', _timing_xml(slide))
    assert steps[::2] == ["gridLegend", "category", "category", "category"]

    slide, spid = _chart_slide()
    add_chart_build(slide, spid, "series", parts=2)
    steps = re.findall(r'bldStep="([^"]+)"', _timing_xml(slide))
    assert steps[::2] == ["gridLegend", "series", "series"]
