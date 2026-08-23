import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxkit.errors import SpecError
from pptxkit.qa.inspect import inspect_deck


def test_inspect_deck_reports_real_ids_and_inch_boxes(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.0), Inches(0.5))
    deck = tmp_path / "d.pptx"
    prs.save(str(deck))

    slides = inspect_deck(deck)
    assert len(slides) == 1
    assert slides[0]["index"] == 1
    assert slides[0]["layout"] == slide.slide_layout.name

    shapes = slides[0]["shapes"]
    assert len(shapes) == 1
    assert shapes[0]["shape_id"] == shape.shape_id
    assert shapes[0]["name"] == shape.name
    assert shapes[0]["box"] == {"x": 1.0, "y": 2.0, "w": 3.0, "h": 0.5}


def test_inspect_deck_returns_one_entry_per_slide_in_order(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    deck = tmp_path / "d.pptx"
    prs.save(str(deck))

    slides = inspect_deck(deck)
    assert [s["index"] for s in slides] == [1, 2]


def test_a_deck_that_will_not_open_is_refused(tmp_path):
    """Inspection is handed a file by a human, not by `build`, so it meets decks no
    compiler wrote — and a bad one is the caller's input, not a theme problem."""
    deck = tmp_path / "d.pptx"
    deck.write_text("not a package")

    with pytest.raises(SpecError, match=r"deck .*d\.pptx is not a readable \.pptx"):
        inspect_deck(deck)
