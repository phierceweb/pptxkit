"""Will PowerPoint open it? Every case is built by corrupting a real saved deck, because the point
is what a *package* holds. LibreOffice renders all of these without complaint."""

from __future__ import annotations

import re
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxkit.qa.model import Severity
from pptxkit.qa.package import check_package


@pytest.fixture
def deck(tmp_path):
    """A saved two-slide deck with a picture and a shape on it."""
    png = tmp_path / "dot.png"
    png.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
            "890000000d4944415478da63f8cfc0500f0004000180ff9fd6a4f10000000049454e44ae426082"
        )
    )
    prs = Presentation()
    first = prs.slides.add_slide(prs.slide_layouts[6])
    first.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "A"
    first.shapes.add_picture(str(png), Inches(1), Inches(3), Inches(2), Inches(2))
    second = prs.slides.add_slide(prs.slide_layouts[6])
    second.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "B"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def corrupt(deck: Path, out: Path, transform, *, part: str | None = None) -> Path:
    """Rewrite one slide part of ``deck`` through ``transform``."""
    with zipfile.ZipFile(deck) as zin, zipfile.ZipFile(out, "w") as zo:
        target = part or next(
            n for n in sorted(zin.namelist()) if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)
        )
        for item in zin.infolist():
            data = zin.read(item.filename)
            zo.writestr(item, transform(data) if item.filename == target else data)
    return out


def test_a_sound_deck_reports_nothing(deck):
    assert check_package(deck) == []


def test_a_duplicate_shape_id_is_an_error(deck, tmp_path):
    """PowerPoint repairs the slide and drops content; nothing else here would see it."""
    bad = corrupt(
        deck, tmp_path / "dup.pptx", lambda d: re.sub(rb'(<p:cNvPr id=")\d+(")', rb"\g<1>9\g<2>", d)
    )
    findings = check_package(bad)
    assert findings and findings[0].severity is Severity.ERROR
    assert "shape id 9 is used by both" in findings[0].detail


def test_a_zero_shape_id_is_an_error(deck, tmp_path):
    bad = corrupt(
        deck, tmp_path / "zero.pptx", lambda d: d.replace(b'<p:cNvPr id="2"', b'<p:cNvPr id="0"', 1)
    )
    assert any("outside 1.." in f.detail for f in check_package(bad))


def test_an_animation_targeting_a_missing_shape_is_an_error(deck, tmp_path):
    """Timing is the one tree written as raw XML, so its ids are not kept in step."""
    timing = (
        b'<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/'
        b'2006/main"><p:tnLst><p:par><p:cTn><p:childTnLst><p:set><p:cBhvr>'
        b'<p:tgtEl><p:spTgt spid="9999"/></p:tgtEl></p:cBhvr></p:set>'
        b"</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing></p:sld>"
    )
    bad = corrupt(deck, tmp_path / "anim.pptx", lambda d: d.replace(b"</p:sld>", timing))
    findings = check_package(bad)
    assert findings and "targets shape id 9999" in findings[0].detail


def test_a_relationship_reference_nothing_declares_is_an_error(deck, tmp_path):
    bad = corrupt(
        deck, tmp_path / "rel.pptx", lambda d: d.replace(b'r:embed="rId', b'r:embed="rIdNope', 1)
    )
    assert any("declared in no relationship part" in f.detail for f in check_package(bad))


def test_a_relationship_pointing_at_a_missing_part_is_an_error(deck, tmp_path):
    """The reference and its declaration agree; the file they name is simply gone."""
    out = tmp_path / "gone.pptx"
    with zipfile.ZipFile(deck) as zin, zipfile.ZipFile(out, "w") as zo:
        for item in zin.infolist():
            if item.filename.startswith("ppt/media/"):
                continue
            zo.writestr(item, zin.read(item.filename))
    assert any("the package does not contain" in f.detail for f in check_package(out))


def test_malformed_slide_xml_is_an_error(deck, tmp_path):
    bad = corrupt(deck, tmp_path / "broken.pptx", lambda d: d[:-20])
    findings = check_package(bad)
    assert findings and "not well-formed XML" in findings[0].detail


def test_a_file_that_is_not_a_pptx_is_an_error(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_text("this is not a presentation")
    findings = check_package(path)
    assert findings and "not a readable .pptx" in findings[0].detail


def test_the_finding_names_the_slide_it_is_on(deck, tmp_path):
    """A 67-slide deck needs the number, not just 'somewhere in here'."""
    with zipfile.ZipFile(deck) as z:
        second = sorted(n for n in z.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n))[1]
    bad = corrupt(
        deck,
        tmp_path / "second.pptx",
        lambda d: d.replace(b'<p:cNvPr id="2"', b'<p:cNvPr id="0"', 1),
        part=second,
    )
    findings = check_package(bad)
    assert findings and findings[0].slide == int(re.search(r"(\d+)", second).group(1))


# --- charts the render cannot verify ----------------------------------------


def _charted(tmp_path, chart_type, values, name="c.pptx"):
    """A saved deck holding one chart of ``chart_type``."""
    from pptx.chart.data import CategoryChartData

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["Up", "Down"]
    data.add_series("s", values)
    slide.shapes.add_chart(chart_type, Inches(1), Inches(1), Inches(6), Inches(4), data)
    path = tmp_path / name
    prs.save(str(path))
    return path


def test_a_bar_chart_with_a_negative_value_is_flagged_unverifiable(tmp_path):
    """The file is correct; the render this suite and `pptxkit render` both go through plots the
    absolute value, so nothing automated can check the chart. Reproduced with bare python-pptx."""
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _charted(tmp_path, XL_CHART_TYPE.COLUMN_CLUSTERED, (271, -146))
    findings = [f for f in check_package(deck) if f.check == "chart-negative"]
    assert len(findings) == 1
    assert findings[0].severity is Severity.WARN
    assert findings[0].slide == 1
    assert "diverge" in findings[0].detail


def test_an_all_positive_bar_chart_is_clean(tmp_path):
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _charted(tmp_path, XL_CHART_TYPE.COLUMN_CLUSTERED, (271, 146))
    assert [f for f in check_package(deck) if f.check == "chart-negative"] == []


def test_a_line_chart_with_a_negative_value_is_not_flagged(tmp_path):
    """Line series render their negatives correctly, so warning about them is noise."""
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _charted(tmp_path, XL_CHART_TYPE.LINE, (271, -146))
    assert [f for f in check_package(deck) if f.check == "chart-negative"] == []


def test_two_shapes_sharing_a_name_are_flagged(deck, tmp_path):
    """Legal OOXML that silently costs the deck the traceability naming exists for: two distinct
    textboxes named `s2.chrome` write a package PowerPoint opens and an author cannot navigate."""
    bad = corrupt(
        deck,
        tmp_path / "dupname.pptx",
        lambda d: d.replace(b'name="TextBox 1"', b'name="Picture 2"', 1),
    )
    findings = [f for f in check_package(bad) if f.check == "shape-name"]
    assert len(findings) == 1
    assert findings[0].severity is Severity.WARN
    assert "Picture 2" in findings[0].detail


def test_distinct_shape_names_are_clean(deck):
    assert [f for f in check_package(deck) if f.check == "shape-name"] == []
