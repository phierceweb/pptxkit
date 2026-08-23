"""Drive every brand template you have through every capability.

The set is whatever `.pptx` files live in the theme directory — the same ones your decks
name with `theme:`. That makes this guard exactly as strong as the diversity of what is
there, which is a fact to read at runtime rather than a number to write down here.

`templates/` is gitignored, so this module skips where it is empty, and it never runs in
CI."""

from __future__ import annotations

import yaml

import pathlib

import pytest

from pptxkit.conform import conform
from pptxkit.qa.geometry import (
    check_bounds,
    check_contrast,
    check_placement_fit,
    check_reserved,
)
from pptxkit.conform.sample import is_sample
from pptxkit.qa.package import check_package
from pptxkit.theme import load_theme

SAMPLES = pathlib.Path(__file__).resolve().parents[1] / "templates"
# A -4-3 file is the same design at another aspect, so one of each pair is enough; a generated
# `pptxkit sample` is refused by its own docProps mark, since it varies in nothing.
TEMPLATES = (
    sorted(p for p in SAMPLES.glob("*.pptx") if not p.stem.endswith("-4-3") and not is_sample(p))
    if SAMPLES.is_dir()
    else []
)

#: Release gate. Unset, any number of templates is accepted; set, a thinner set fails
#: rather than passing quietly.
_MIN_ENV_VAR = "PPTXKIT_TEMPLATES_MIN"


def _required_templates() -> int:
    from pf_core.utils.env import resolve_int

    n: int = resolve_int(None, _MIN_ENV_VAR, default=0)
    return max(n, 0)


def test_the_template_set_is_as_large_as_this_run_requires():
    """Skipped along with the rest when there are no templates at all — the runbook says to
    run this on a machine holding them. What it catches is a set that is present but thin."""
    required = _required_templates()
    if not required:
        pytest.skip(f"{_MIN_ENV_VAR} unset — any number of templates is accepted")
    assert len(TEMPLATES) >= required, (
        f"{_MIN_ENV_VAR}={required} but {SAMPLES} holds {len(TEMPLATES)} usable "
        f"template(s): {[p.name for p in TEMPLATES]}. This is the primary guard; "
        f"a narrow set is a weaker release, not a smaller one."
    )


pytestmark = pytest.mark.skipif(not TEMPLATES, reason="no brand templates are present")


def _sidecar_for(template: pathlib.Path) -> pathlib.Path | None:
    """The adopted theme bound to this template, if it has one."""
    for path in sorted(SAMPLES.glob("*.theme.yaml")):
        declared = (yaml.safe_load(path.read_text()) or {}).get("template")
        if declared and pathlib.Path(declared).name == template.name:
            return path
    return None


@pytest.fixture(scope="module", params=[p.name for p in TEMPLATES])
def built(request, tmp_path_factory):
    """Every capability built against one template, once, through its adopted theme where
    it has one — that is the theme its decks resolve."""
    template = SAMPLES / request.param
    out = tmp_path_factory.mktemp("conform") / template.stem[:40]
    return conform(template, out, theme=_sidecar_for(template))


def _derived(built) -> dict:
    """The theme YAML `conform` wrote — what `derive` actually decided. Read as YAML, never through
    `load_theme`, which fills every unbound role from pptxkit's own defaults."""
    import yaml

    return yaml.safe_load(pathlib.Path(built.theme).read_text())


def _scheme_of(built) -> dict[str, str]:
    """The template's own clrScheme, so a bound slot can be resolved to a colour."""
    from pptx import Presentation

    from pptxkit.theme.clrscheme import parse_color_scheme, read_theme_xml

    # Through the theme's own pointer, which is how the loader reaches it — nothing
    # copies the binary any more, so a hardcoded layout here would be a second guess.
    theme_path = pathlib.Path(built.theme)
    declared = yaml.safe_load(theme_path.read_text())["template"]
    master = Presentation(str((theme_path.parent / declared).resolve())).slide_masters[0]
    return parse_color_scheme(read_theme_xml(master))


def _brand_accents_in_scheme(built) -> dict[str, str]:
    """Accent slots holding a colour Microsoft does not ship — what a bind could take."""
    return {
        slot: hex_
        for slot, hex_ in _scheme_of(built).items()
        if slot.startswith("accent") and hex_.upper() not in _MICROSOFT_ACCENTS
    }


def _bound_accents(built) -> dict[str, str]:
    """`{accent-N: hex}` for the accents the derived theme really bound."""
    bind = _derived(built).get("bind", {})
    scheme = _scheme_of(built)
    return {
        role: scheme[slot].upper()
        for role, slot in sorted(bind.items())
        if role.startswith("accent-") and slot in scheme
    }


def test_every_capability_builds_on_the_template(built):
    assert built.failed == [], f"{built.template}: {built.failed}"


def test_the_deck_it_produced_can_be_read_back(built):
    from pptx import Presentation

    assert built.deck is not None
    assert len(Presentation(str(built.deck)).slides) == len(built.passed)


def test_powerpoint_would_open_what_the_build_saved(built):
    """Timing is the one tree written as raw XML, so an id it names is checked by nothing the
    manifest holds — PowerPoint repairs the file and drops the build."""
    broken = check_package(built.deck)
    assert broken == [], f"{built.template}: {[f.detail for f in broken]}"


def test_no_text_is_unreadable_on_what_was_really_painted_behind_it(built):
    """Ink lost on its own background. Errors only: below the 3:1 floor is an error, and a
    warning is a brand accent a percent under AA rather than a defect."""
    import json

    from pptxkit.qa.model import Severity

    manifest = json.loads(built.deck.with_suffix(".manifest.json").read_text())
    theme = load_theme(built.theme)
    unreadable = [f for f in check_contrast(manifest, theme) if f.severity is Severity.ERROR]
    assert unreadable == [], f"{built.template}: {[f.detail for f in unreadable]}"


def test_a_template_that_paints_a_picture_has_its_pixels_read_not_its_page_role(built):
    """A photograph is not one colour, so a deck on one cannot record only one. Without this,
    dropping the inherited-backdrop path leaves the check comparing the palette against itself."""
    import json

    theme = load_theme(built.theme)
    if theme.surface is None or theme.surface.media is None:
        pytest.skip(f"{built.template} paints no picture behind its slides")
    manifest = json.loads(built.deck.with_suffix(".manifest.json").read_text())
    # Only slides carrying no picture of their own: a placed photograph is sampled by a different
    # path, and its colours would keep this green with the inherited backdrop deleted.
    # `.get("shapes", [])` — a slide carrying only `notes:` has no shapes key at all.
    bare = [
        slide
        for slide in manifest["slides"]
        if slide.get("background") == "page"
        and not any(s.get("rendered") == "picture" for s in slide.get("shapes", []))
    ]
    assert bare, f"{built.template}: no slide composes on the template's own surface"
    papers = {
        str(shape["bg"]).upper()
        for slide in bare
        for shape in slide.get("shapes", [])
        if shape.get("bg")
    }
    # The samples are compared with each other, never with the palette: `derive` seeds `page` from
    # the same flattening of the same artwork, so the better the code agrees with itself the more
    # certainly every sample is also a declared role.
    assert len(papers) > 1, (
        f"{built.template}: every slide recorded the same background {papers} — the "
        f"master's picture was never sampled, only its page role read"
    )


def test_nothing_lands_outside_the_slide(built):
    import json

    manifest = json.loads(built.deck.with_suffix(".manifest.json").read_text())
    theme = load_theme(built.theme)
    assert check_bounds(manifest, theme) == []


def test_nothing_escapes_the_rect_it_was_placed_in(built):
    """`bounds` above only asks whether a shape is on the slide, so a component that
    overruns its own placement and lands on its neighbour passes it."""
    import json

    manifest = json.loads(built.deck.with_suffix(".manifest.json").read_text())
    theme = load_theme(built.theme)
    assert check_placement_fit(manifest, theme) == []


def test_nothing_intrudes_on_a_region_the_brand_reserves(built):
    """Only meaningful for a template onboarded with a tuned sidecar: `derive` emits no `reserve:`,
    and `check_reserved` returns `[]` for a theme declaring none. Skipped rather than left green so
    the gap is visible; the check itself is covered in `tests/qa/test_geometry_bounds.py`."""
    import json

    theme = load_theme(built.theme)
    if not theme.reserve:
        pytest.skip(f"{built.template}: the derived theme declares no reserved region")
    manifest = json.loads(built.deck.with_suffix(".manifest.json").read_text())
    assert check_reserved(manifest, theme) == []


def test_every_brand_accent_the_template_owns_is_bound(built):
    """Skipped for a template that never edited the Office palette: it has no brand to bind, which
    is a fact about that template, not a defect. Asserted against `derive`'s own `bind` —
    `Theme.palette.accents` is filled from pptxkit's defaults when nothing was bound."""
    if not _brand_accents_in_scheme(built):
        pytest.skip(f"{built.template}: every accent slot holds a stock Office colour")
    assert _bound_accents(built), (
        f"{built.template} owns brand accents but bound none — the binding path stopped running"
    )


# Spelled out here rather than imported: a test that asks the code under test whether
# a colour is stock passes just as happily when that check is gutted.
_MICROSOFT_ACCENTS = frozenset(
    {
        "4F81BD",
        "C0504D",
        "9BBB59",
        "8064A2",
        "4BACC6",
        "F79646",  # Office 2007
        "4472C4",
        "ED7D31",
        "A5A5A5",
        "FFC000",
        "5B9BD5",
        "70AD47",  # Office 2013-2021
        "156082",
        "E97132",
        "196B24",
        "0F9ED5",
        "A02B93",
        "4EA72E",  # Office 2024
    }
)


def test_no_unedited_microsoft_accent_is_adopted_as_a_brand_colour(built):
    """Some templates ship stock accents they never use; adopting one is a lie. Over the
    colours actually bound, so an empty result means "nothing stock got in", never "nothing bound".
    """
    if not _brand_accents_in_scheme(built):
        pytest.skip(f"{built.template}: every accent slot holds a stock Office colour")
    accents = _bound_accents(built)
    assert accents, "nothing bound — the assertion below would hold over an empty set"
    adopted = sorted((role, hex_) for role, hex_ in accents.items() if hex_ in _MICROSOFT_ACCENTS)
    assert adopted == [], f"{built.template} adopted Microsoft's own colours: {adopted}"
