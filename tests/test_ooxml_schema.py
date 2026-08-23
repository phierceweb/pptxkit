"""Validate the raw OOXML `pptxkit.motion` writes against ISO/IEC 29500-4:2016 — nothing else in the
project can. LibreOffice converts schema-invalid timing to PDF without complaint and `pptxkit qa`
renders only a slide's final state; schema-valid is itself only a floor, since real PowerPoint is
what says whether a file opens without a repair prompt. See `docs/pptx-deck-building.md`."""

from __future__ import annotations

import io
import pathlib
import re
import zipfile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptxkit.motion import (
    add_chart_build,
    add_click_build,
    add_click_reveals,
    add_click_sequence,
    add_transition,
)
from pptxkit.motion.transition import EFFECTS, SPEEDS

SCHEMA = pathlib.Path(__file__).parent / "schemas" / "ooxml" / "pml.xsd"
_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"


@pytest.fixture(scope="module")
def schema():
    return etree.XMLSchema(etree.parse(str(SCHEMA)))


def _deck():
    """A slide carrying the three shape kinds the builds are pointed at."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    box.text_frame.text = "text"
    blank = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    data.add_series("S1", (1, 2, 3))
    data.add_series("S2", (3, 2, 1))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(4), Inches(1), Inches(4), Inches(3), data
    )
    return prs, slide, box.shape_id, blank.shape_id, frame.shape_id


def _slide_xml(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return zipfile.ZipFile(buf).read("ppt/slides/slide1.xml")


def _validate(schema, prs) -> list[str]:
    doc = etree.fromstring(_slide_xml(prs))
    if schema.validate(doc):
        return []
    return [e.message for e in schema.error_log]


def test_the_vendored_schema_is_present():
    """A skipped schema gate proves nothing; this fails rather than skips."""
    assert SCHEMA.is_file(), f"missing vendored schema at {SCHEMA}"


def test_a_click_build_validates(schema):
    prs, slide, text, blank, chart = _deck()
    add_click_build(slide, [text, blank])
    assert _validate(schema, prs) == []


def test_a_staggered_click_build_validates(schema):
    prs, slide, text, blank, chart = _deck()
    add_click_build(slide, [text, blank], 80)
    assert _validate(schema, prs) == []


@pytest.mark.parametrize("kind", ["fade", "wipeup", "wiperight"])
def test_every_click_sequence_effect_kind_validates(schema, kind):
    prs, slide, text, blank, chart = _deck()
    add_click_sequence(slide, [[(text, kind)], [(blank, kind)]], 60)
    assert _validate(schema, prs) == []


@pytest.mark.parametrize("by", ["category", "series", "element", "all"])
def test_every_chart_build_validates(schema, by):
    """`bldStep` is required on `<a:chart>`; omitting it is the defect this caught."""
    prs, slide, text, blank, chart = _deck()
    add_chart_build(slide, chart, by, parts=3)
    assert _validate(schema, prs) == []


@pytest.mark.parametrize("kind", ["fade", "wipeup", "wiperight"])
def test_every_after_previous_chain_validates(schema, kind):
    """`afterEffect` group starts and their beat gate, the auto-advance path."""
    prs, slide, text, blank, chart = _deck()
    add_click_sequence(slide, [[(text, kind)], [(blank, kind)]], beat_ms=250)
    assert _validate(schema, prs) == []


def test_a_build_over_only_text_free_shapes_validates(schema):
    """No bldP is legal for these, so the build list is omitted — and `<p:bldLst/>`
    with no children would be invalid."""
    prs, slide, text, blank, chart = _deck()
    add_click_build(slide, [blank])
    assert _validate(schema, prs) == []


def test_a_group_holding_both_shape_kinds_validates(schema):
    """A worded shape and a text-free one on one click: the first entrance node carries `grpId` and
    the second does not. The XSD accepts either — it is `use="optional"` on
    `CT_TLCommonTimeNodeData` — so this gates only spelling and position;
    `tests/test_pptx_helpers.py` gates which node."""
    prs, slide, text, blank, chart = _deck()
    add_click_sequence(slide, [[text, blank]])
    assert _validate(schema, prs) == []


def test_an_interactive_reveal_validates(schema):
    prs, slide, text, blank, chart = _deck()
    add_click_reveals(slide, [(text, blank)])
    assert _validate(schema, prs) == []


def test_two_triggers_onto_one_target_stay_siblings_under_the_root(schema):
    """Every shape a trigger placement drew listens, so a two-shape trigger writes two `<p:seq>`.
    Both hang off the tmRoot: burying the second inside the first validates, but it would only arm
    once the outer sequence had run."""
    prs, slide, text, blank, chart = _deck()
    add_click_reveals(slide, [(text, blank), (chart, blank)])

    assert _validate(schema, prs) == []
    seqs = list(etree.fromstring(_slide_xml(prs)).iter(f"{{{_PML}}}seq"))
    assert len(seqs) == 2
    assert {s.getparent().getparent().get("nodeType") for s in seqs} == {"tmRoot"}


@pytest.mark.parametrize("kind", sorted(EFFECTS))
def test_every_transition_validates(schema, kind):
    prs, slide, text, blank, chart = _deck()
    add_transition(slide, kind)
    assert _validate(schema, prs) == []


@pytest.mark.parametrize(
    "kind, direction",
    [(k, d) for k, dirs in sorted(EFFECTS.items()) for d in dirs],
)
def test_every_transition_direction_validates(schema, kind, direction):
    """Direction vocabularies are per element — a shared list is invalid at `strips`."""
    prs, slide, text, blank, chart = _deck()
    add_transition(slide, kind, direction=direction)
    assert _validate(schema, prs) == []


@pytest.mark.parametrize("speed", SPEEDS)
def test_every_transition_speed_validates(schema, speed):
    prs, slide, text, blank, chart = _deck()
    add_transition(slide, "wipe", speed=speed)
    assert _validate(schema, prs) == []


def test_a_transition_beside_an_animation_validates_and_keeps_child_order(schema):
    """`CT_Slide` is an xsd:sequence. LibreOffice repairs a wrong order on import, so
    this is the only mechanical check that the two writers do not collide."""
    prs, slide, text, blank, chart = _deck()
    add_click_sequence(slide, [[text], [blank]])
    add_transition(slide, "push", direction="u")

    assert _validate(schema, prs) == []
    tags = [el.tag.split("}")[1] for el in slide._element]
    assert tags.index("transition") < tags.index("timing")


def test_the_gate_catches_a_missing_required_attribute(schema):
    """The negative control. Without it, a validator that silently passed everything
    would look identical to one that works."""
    prs, slide, text, blank, chart = _deck()
    add_chart_build(slide, chart, "category", parts=3)
    sabotaged = re.sub(rb' bldStep="[^"]+"', b"", _slide_xml(prs))

    assert not schema.validate(etree.fromstring(sabotaged))
    assert "bldStep" in schema.error_log[0].message
