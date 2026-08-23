import dataclasses

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxkit.compile.manifest import ManifestRecorder
from pptxkit.errors import LayoutError
from pptxkit.layouts.components import component
from pptxkit.layouts.compose import render_slide
from pptxkit.layouts.place import Reserved
from pptxkit.layouts.registry import SlideCtx
from pptxkit.spec.model import Placement, SlideSpec

_SEEN: dict[str, object] = {}


@component("geo-probe")
def _probe(ctx):
    _SEEN[ctx.body["tag"]] = ctx.rect
    return []


@pytest.fixture
def drawn():
    """The rect each probe placement was handed, keyed by its tag."""
    _SEEN.clear()
    return _SEEN


def _ctx(theme, spec):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    manifest = ManifestRecorder(deck="d", theme="t")
    manifest.begin_slide(1, background="page")
    return SlideCtx(slide=slide, theme=theme, spec=spec, manifest=manifest)


def _place(tag, at, **kw):
    return Placement(at=at, component="geo-probe", body={"tag": tag}, **kw)


def test_two_placements_side_by_side_get_adjacent_rects(theme, drawn):
    spec = SlideSpec(
        index=1,
        title="T",
        place=(_place("l", {"cols": "left-half"}), _place("r", {"cols": "right-half"})),
    )
    render_slide(_ctx(theme, spec))
    assert drawn["r"].left >= drawn["l"].right


def test_two_placements_sharing_a_column_are_rejected_naming_both(theme, drawn):
    spec = SlideSpec(
        index=1, place=(_place("a", {"cols": (0, 7)}), _place("b", {"cols": "right-half"}))
    )
    with pytest.raises(
        LayoutError,
        match=r"slide 1 placement 1 \(geo-probe\) overlaps "
        r"slide 1 placement 2 \(geo-probe\)",
    ):
        render_slide(_ctx(theme, spec))


def test_a_named_placement_publishes_its_rect(theme, drawn):
    spec = SlideSpec(index=1, place=(_place("a", {"cols": "left-half"}, id="left"),))
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    assert ctx.placements["left"] == drawn["a"]


def test_a_placement_under_short_chrome_starts_exactly_at_the_body_top(theme, drawn):
    """Chrome that fits above body_top must not move the band the design system set."""
    spec = SlideSpec(index=1, kicker="K", title="T", place=(_place("a", {"cols": "full"}),))
    render_slide(_ctx(theme, spec))
    assert drawn["a"].top == pytest.approx(theme.grid.body_top)


def test_a_bleeding_placement_may_cover_the_whole_canvas(theme, drawn):
    spec = SlideSpec(index=1, place=(_place("a", {"box": (0, 0, 1, 1)}, bleed=True),))
    render_slide(_ctx(theme, spec))
    assert (drawn["a"].left, drawn["a"].width) == pytest.approx((0.0, 13.333))


WRAPS = "Retrieval-augmented generation cut our support ticket backlog in half"
SUBTITLE = "Q3 results across all four regions"


def test_the_chrome_lines_share_one_frame_starting_at_the_top_margin(theme):
    """One frame, so a title that wraps wider than the estimate pushes the subtitle
    down the frame rather than drawing the next line through it."""
    ctx = _ctx(theme, SlideSpec(index=1, kicker="K", title="T", subtitle="S"))
    render_slide(ctx)
    (frame,) = ctx.slide.shapes
    assert frame.top == Inches(theme.grid.top)
    assert frame.text_frame.text == "K\nT\nS"


def _chrome_height(theme, title):
    ctx = _ctx(theme, SlideSpec(index=1, title=title, subtitle=SUBTITLE))
    render_slide(ctx)
    return ctx.slide.shapes[0].height


def test_a_wrapping_title_reserves_a_second_line_before_the_subtitle(theme):
    one_line = _chrome_height(theme, "A short title")
    two_lines = _chrome_height(theme, WRAPS)
    assert two_lines - one_line == pytest.approx(
        Inches(theme.style("title").size * 1.2 / 72), abs=2
    )


def test_a_wrapping_title_leaves_the_first_placement_below_the_chrome(theme, drawn):
    spec = SlideSpec(
        index=1, title=WRAPS, subtitle=SUBTITLE, place=(_place("a", {"cols": "full"}),)
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    frame = ctx.slide.shapes[0]
    assert Inches(drawn["a"].top) >= frame.top + frame.height


def test_a_full_width_placement_is_narrowed_to_clear_a_reserved_wedge(theme, drawn):
    """cols: full is the commonest placement there is, so a reserved logo corner
    must narrow it rather than turn it into a build failure."""
    wedged = dataclasses.replace(
        theme,
        reserve=(Reserved(name="logo-wedge", poly=((1.0, 0.7227), (1.0, 1.0), (0.825, 1.0))),),
    )
    spec = SlideSpec(index=1, place=(_place("a", {"cols": "full"}),))
    render_slide(_ctx(wedged, spec))
    assert drawn["a"].right < theme.grid.right_edge


def test_a_placement_bounded_above_the_wedge_keeps_the_full_content_width(theme, drawn):
    wedged = dataclasses.replace(
        theme,
        reserve=(Reserved(name="logo-wedge", poly=((1.0, 0.7227), (1.0, 1.0), (0.825, 1.0))),),
    )
    spec = SlideSpec(index=1, place=(_place("a", {"cols": "full", "rows": "top-half"}),))
    render_slide(_ctx(wedged, spec))
    assert drawn["a"].width == pytest.approx(theme.grid.content_w)


def test_an_exact_box_over_a_reserved_region_is_rejected_not_narrowed(theme, drawn):
    """A box: is geometry the author stated outright; moving it would be a lie."""
    wedged = dataclasses.replace(
        theme,
        reserve=(Reserved(name="logo-wedge", poly=((1.0, 0.7227), (1.0, 1.0), (0.825, 1.0))),),
    )
    spec = SlideSpec(index=1, place=(_place("a", {"box": (0.85, 0.8, 0.1, 0.1)}),))
    with pytest.raises(LayoutError, match=r"overlaps the reserved region 'logo-wedge'"):
        render_slide(_ctx(wedged, spec))


def test_a_box_placement_reaches_above_the_content_band(theme, drawn):
    """Through the composer, not by hand: check_placements only exempts a box it was
    *told* is one, so a composer that stops saying so puts the ceiling back."""
    spec = SlideSpec(
        index=1,
        title="A title in the usual place",
        place=(_place("high", {"box": (0.06, 0.03, 0.4, 0.08)}),),
    )
    render_slide(_ctx(theme, spec))
    assert drawn["high"].top == pytest.approx(0.03 * 7.5, abs=1e-3)


def test_a_box_placement_off_the_canvas_is_still_refused(theme):
    spec = SlideSpec(index=1, place=(_place("gone", {"box": (0.85, 0.4, 0.4, 0.1)}),))
    with pytest.raises(LayoutError, match="falls outside the canvas"):
        render_slide(_ctx(theme, spec))


@component("geo-short")
def _short(ctx):
    """Draws one plate a third of its placement tall, at the top — the shape every
    content-sized component makes."""
    from pptxkit.utils.shapes import rect as fill_rect

    r = ctx.body_rect
    shape = fill_rect(ctx.slide, r.left, r.top, r.width, r.height / 3, ctx.color("line"))
    ctx.manifest.record(shape)
    return []


def _short_place(anchor):
    return Placement(at={"cols": "full"}, component="geo-short", body={}, anchor=anchor)


@pytest.mark.parametrize("anchor,share", [("top", 0.0), ("middle", 0.5), ("bottom", 1.0)])
def test_render_slide_settles_a_short_component_per_its_anchor(theme, anchor, share):
    """The wiring, not the arithmetic: `_settle` is unit-tested, and nothing else notices
    when `render_slide` stops calling it."""
    ctx = _ctx(theme, SlideSpec(index=1, title="T", place=(_short_place(anchor),)))
    render_slide(ctx)
    plate = next(s for s in ctx.slide.shapes if s.name.startswith("s1.p1.geo-short"))
    rect = ctx.rect
    slack = rect.height - plate.height / 914400
    assert plate.top / 914400 == pytest.approx(rect.top + slack * share, abs=0.02)


def test_a_bleeding_placement_is_left_where_it_was_drawn(theme):
    """A bleed is a declared overrun; settling it would undo what the author asked for."""
    ctx = _ctx(
        theme,
        SlideSpec(
            index=1,
            title="T",
            place=(
                Placement(
                    at={"cols": "full"}, component="geo-short", body={}, anchor="middle", bleed=True
                ),
            ),
        ),
    )
    render_slide(ctx)
    plate = next(s for s in ctx.slide.shapes if s.name.startswith("s1.p1.geo-short"))
    assert plate.top / 914400 == pytest.approx(ctx.rect.top, abs=0.02)
