import re
import dataclasses

import pytest
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

import pptxkit.components  # noqa: F401 — registers the built-ins
from pptxkit.compile.manifest import ManifestRecorder
from pptxkit.errors import LayoutError, ThemeError
from pptxkit.layouts.chrome import ChromeField
from pptxkit.layouts.components import component
from pptxkit.layouts.compose import render_slide
from pptxkit.layouts.registry import SlideCtx
from pptxkit.spec.model import Background, Placement, SlideSpec
from pptxkit.theme import default_theme
from pptxkit.utils.color import AA_NORMAL, contrast_ratio
from pptxkit.utils.shapes import ALIGN, ANCHOR, para, textbox

_EMU = 914400
_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
_AML = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _ctx(theme, spec):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    manifest = ManifestRecorder(deck="d", theme="t")
    manifest.begin_slide(spec.index, background=spec.background.pair)
    return SlideCtx(slide=slide, theme=theme, spec=spec, manifest=manifest)


def _texts(slide):
    return [s.text_frame.text for s in slide.shapes if s.has_text_frame]


def _probe(name, seen):
    @component(name)
    def _draw(ctx):
        seen.append((ctx.component, dict(ctx.body)))
        tf = textbox(ctx.slide, ctx.body_rect.left, ctx.body_rect.top, 2.0, 0.4)
        para(tf, str(ctx.body.get("label", name)), 12, ctx.fg(), first=True)
        return [[tf._parent.shape_id]]

    return _draw


def test_a_title_is_written_as_chrome(theme):
    ctx = _ctx(theme, SlideSpec(index=1, title="Revenue up 40 percent"))
    render_slide(ctx)
    assert "Revenue up 40 percent" in "\n".join(_texts(ctx.slide))


def test_the_kicker_title_and_subtitle_share_one_frame_in_chrome_order(theme):
    ctx = _ctx(theme, SlideSpec(index=1, kicker="Q3", title="First", subtitle="Sub"))
    render_slide(ctx)
    assert _texts(ctx.slide) == ["Q3\nFirst\nSub"]


def test_a_slide_with_no_chrome_draws_nothing(theme):
    ctx = _ctx(theme, SlideSpec(index=1))
    render_slide(ctx)
    assert len(ctx.slide.shapes) == 0


def test_an_inverse_background_paints_the_whole_slide(theme):
    ctx = _ctx(theme, SlideSpec(index=1, background=Background(kind="inverse")))
    render_slide(ctx)
    painted = ctx.slide.shapes[0]
    assert painted.left == 0 and painted.top == 0
    assert painted.width / _EMU == pytest.approx(13.333, abs=0.01)
    assert painted.height / _EMU == pytest.approx(7.5, abs=0.01)


def test_the_painted_surface_is_the_pairs_own_background_colour(theme):
    ctx = _ctx(theme, SlideSpec(index=1, background=Background(kind="inverse")))
    render_slide(ctx)
    assert str(ctx.slide.shapes[0].fill.fore_color.rgb) == theme.palette.pair("inverse").bg


def test_a_theme_with_no_template_still_paints_an_inverse_slide_legibly():
    """The original bug: a themeless deck put white chrome on an unpainted white slide."""
    themeless = default_theme()
    assert themeless.marks == {}
    ctx = _ctx(themeless, SlideSpec(index=1, title="T", background=Background(kind="inverse")))
    render_slide(ctx)
    surface = str(ctx.slide.shapes[0].fill.fore_color.rgb)
    written = [r for r in ctx.manifest.slides[0].shapes if r.text == "T"][0]
    assert surface == written.bg
    assert contrast_ratio(written.fg, surface) >= AA_NORMAL


def test_a_theme_mark_lays_its_art_over_the_painted_surface(theme, tmp_path):
    Image.new("RGB", (64, 36), "black").save(tmp_path / "art.png")
    themed = dataclasses.replace(
        theme, template=tmp_path / "t.pptx", marks={"inverse": {"media": "art.png"}}
    )
    ctx = _ctx(themed, SlideSpec(index=1, background=Background(kind="inverse")))
    render_slide(ctx)
    assert [r.rendered for r in ctx.manifest.slides[0].shapes] == ["native", "picture"]


def test_a_theme_mark_without_media_is_reported_as_a_theme_error(theme):
    themed = dataclasses.replace(theme, marks={"inverse": "art.png"})
    ctx = _ctx(themed, SlideSpec(index=1, background=Background(kind="inverse")))
    with pytest.raises(ThemeError, match=r"theme mark 'inverse' needs a 'media:'"):
        render_slide(ctx)


def test_chrome_on_an_inverse_background_is_recorded_against_the_inverse_pair(theme):
    ctx = _ctx(theme, SlideSpec(index=1, title="T", background=Background(kind="inverse")))
    render_slide(ctx)
    written = [r for r in ctx.manifest.slides[0].shapes if r.text == "T"][0]
    pair = theme.palette.pair("inverse")
    assert (written.fg, written.bg) == (pair.fg, pair.bg)


def test_chrome_on_the_page_is_recorded_against_the_page_pair(theme):
    ctx = _ctx(theme, SlideSpec(index=1, title="T"))
    render_slide(ctx)
    written = [r for r in ctx.manifest.slides[0].shapes if r.text == "T"][0]
    pair = theme.palette.pair("page")
    assert (written.fg, written.bg) == (pair.fg, pair.bg)


def test_an_image_background_is_placed_full_bleed(theme, tmp_path):
    Image.new("RGB", (64, 36), "black").save(tmp_path / "bg.png")
    themed = dataclasses.replace(theme, template=tmp_path / "t.pptx")
    ctx = _ctx(themed, SlideSpec(index=1, background=Background(kind="image", image="bg.png")))
    render_slide(ctx)
    picture = ctx.slide.shapes[1]
    assert picture.width / _EMU == pytest.approx(13.333, abs=0.01)
    assert ctx.manifest.slides[0].shapes[1].rendered == "picture"


def test_an_image_background_paints_the_surface_under_the_image(theme, tmp_path):
    Image.new("RGB", (64, 36), "black").save(tmp_path / "bg.png")
    themed = dataclasses.replace(theme, template=tmp_path / "t.pptx")
    ctx = _ctx(themed, SlideSpec(index=1, background=Background(kind="image", image="bg.png")))
    render_slide(ctx)
    assert str(ctx.slide.shapes[0].fill.fore_color.rgb) == theme.palette.pair("inverse").bg


def test_a_placement_hands_its_component_its_own_mapping(theme):
    seen = []
    _probe("t-solo", seen)
    spec = SlideSpec(
        index=1, place=(Placement(at={"cols": (0, 6)}, component="t-solo", body={"label": "L"}),)
    )
    render_slide(_ctx(theme, spec))
    assert seen == [("t-solo", {"label": "L"})]


def test_every_placement_is_drawn_in_spec_order(theme):
    seen = []
    _probe("t-left", seen)
    _probe("t-right", seen)
    spec = SlideSpec(
        index=1,
        place=(
            Placement(at={"cols": (0, 6)}, component="t-left", body={"label": "L"}),
            Placement(at={"cols": (6, 12)}, component="t-right", body={"label": "R"}),
        ),
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    assert [name for name, _ in seen] == ["t-left", "t-right"]
    assert {"L", "R"} <= set(_texts(ctx.slide))


def test_a_named_placement_publishes_its_rect_on_the_context(theme):
    _probe("t-named", [])
    spec = SlideSpec(
        index=1, place=(Placement(at={"cols": (0, 6)}, component="t-named", id="left"),)
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    assert ctx.placements["left"].left == pytest.approx(theme.grid.col_x(0))


def test_an_anonymous_placement_publishes_nothing(theme):
    _probe("t-anon", [])
    spec = SlideSpec(index=1, place=(Placement(at={"cols": (0, 6)}, component="t-anon"),))
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    assert ctx.placements == {}


def test_a_component_on_an_inverse_background_draws_the_inverse_ink(theme):
    _probe("t-oninverse", [])
    spec = SlideSpec(
        index=1,
        background=Background(kind="inverse"),
        place=(Placement(at={"cols": (0, 12)}, component="t-oninverse"),),
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    drawn = [
        s for s in ctx.slide.shapes if s.has_text_frame and s.text_frame.text == "t-oninverse"
    ][0]
    run = drawn.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == theme.palette.pair("inverse").fg


def test_an_unknown_component_is_reported_by_name(theme):
    spec = SlideSpec(index=1, place=(Placement(at={"cols": (0, 6)}, component="t-absent"),))
    with pytest.raises(LayoutError, match=r"unknown body component 't-absent'"):
        render_slide(_ctx(theme, spec))


def test_an_unknown_animate_names_the_slide(theme):
    spec = SlideSpec(index=4, title="T", animate="wobble")
    with pytest.raises(LayoutError, match=r"slide 4: unknown animate 'wobble'"):
        render_slide(_ctx(theme, spec))


def test_one_at_a_time_builds_a_click_sequence_across_placements(theme):
    seen = []
    _probe("t-a", seen)
    _probe("t-b", seen)
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(
            Placement(at={"cols": (0, 6)}, component="t-a"),
            Placement(at={"cols": (6, 12)}, component="t-b"),
        ),
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    animations = ctx.manifest.slides[0].animations
    assert animations[0]["kind"] == "click_sequence"
    assert len(animations[0]["steps"]) == 2


def test_together_builds_one_click_build(theme):
    _probe("t-together", [])
    spec = SlideSpec(
        index=1,
        animate="together",
        place=(Placement(at={"cols": (0, 12)}, component="t-together"),),
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    assert ctx.manifest.slides[0].animations[0]["kind"] == "click_build"


def test_animate_none_builds_nothing(theme):
    _probe("t-still", [])
    spec = SlideSpec(
        index=1, animate="none", place=(Placement(at={"cols": (0, 12)}, component="t-still"),)
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    assert ctx.manifest.slides[0].animations == []


def test_by_series_survives_the_slide_level_animate_whitelist(theme):
    """The only end-to-end drive of ``animate: by_series``: drop it from the chart vocabulary and
    compose rejects it before the component's own tests ever see it."""
    spec = SlideSpec(
        index=1,
        animate="by_series",
        place=(
            Placement(
                at={"cols": (0, 12), "rows": (0, 8)},
                component="chart",
                body={
                    "kind": "column",
                    "data": [
                        {"category": "Q1", "values": {"North": 12, "South": 5}},
                        {"category": "Q2", "values": {"North": 34, "South": 21}},
                    ],
                },
            ),
        ),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    timing = ctx.slide._element.find(qn("p:timing"))
    builds = timing.findall(f"{{{_PML}}}bldLst/{{{_PML}}}bldGraphic")
    assert [b.find(f"{{{_PML}}}bldSub/{{{_AML}}}bldChart").get("bld") for b in builds] == [
        "series",
        "series",
        "series",
    ]  # the plot background, then one per series
    assert ctx.manifest.slides[0].animations[0]["kind"] == "chart_build"


def _two_shape_component(name):
    @component(name)
    def _draw(ctx):
        a = textbox(ctx.slide, 1.0, 1.0, 2.0, 0.4)
        b = textbox(ctx.slide, 1.0, 2.0, 2.0, 0.4)
        para(a, "a", 12, ctx.fg(), first=True)
        para(b, "b", 12, ctx.fg(), first=True)
        return [[a._parent.shape_id, b._parent.shape_id]]

    return _draw


def test_the_themes_stagger_reaches_the_emitted_delay(theme):
    """The theme's pacing has to travel to the timing tree, and only this reads it back — the
    corpus builds the path but asserts nothing about the delay."""
    _two_shape_component("t-two")
    staggered = dataclasses.replace(theme, motion=dataclasses.replace(theme.motion, stagger_ms=70))
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(Placement(at={"cols": (0, 12)}, component="t-two"),),
    )
    ctx = _ctx(staggered, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert re.findall(r'delay="([^"]+)"', xml) == ["indefinite", "0", "0", "0", "70", "0", "0", "0"]


def test_together_cascades_the_whole_slide_on_its_one_click(theme):
    """``together`` is where a stagger reads as a cascade: every shape on the slide
    shares one click, so the offset runs across components rather than inside one."""
    _two_shape_component("t-two-together")
    staggered = dataclasses.replace(theme, motion=dataclasses.replace(theme.motion, stagger_ms=60))
    spec = SlideSpec(
        index=1,
        animate="together",
        place=(Placement(at={"cols": (0, 12)}, component="t-two-together"),),
    )
    ctx = _ctx(staggered, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert re.findall(r'delay="([^"]+)"', xml) == ["indefinite", "0", "0", "0", "60", "0", "0", "0"]


def test_a_theme_that_does_not_stagger_leaves_every_item_on_the_click(theme):
    _two_shape_component("t-two-flat")
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(Placement(at={"cols": (0, 12)}, component="t-two-flat"),),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert re.findall(r'delay="([^"]+)"', xml) == ["indefinite", "0", "0", "0", "0", "0", "0", "0"]


def _line_component(name):
    @component(name)
    def _draw(ctx):
        tf = textbox(ctx.slide, 1.0, 1.0, 2.0, 0.4)
        para(tf, "x", 12, ctx.fg(), first=True)
        return [[(tf._parent.shape_id, "line")]]

    return _draw


def test_a_components_motion_role_resolves_to_the_themes_entrance(theme):
    """The component says "I am a line"; the theme says lines wipe. Neither names a
    preset — this is the only path by which `wiperight` reaches a built deck."""
    _line_component("t-line")
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(Placement(at={"cols": (0, 12)}, component="t-line"),),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert 'filter="wipe(left)"' in xml
    assert 'presetID="22"' in xml


def test_a_theme_may_rebind_a_role_without_touching_the_spec(theme):
    _line_component("t-line-rebound")
    rebound = dataclasses.replace(
        theme,
        motion=dataclasses.replace(theme.motion, roles={**theme.motion.roles, "line": "fade"}),
    )
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(Placement(at={"cols": (0, 12)}, component="t-line-rebound"),),
    )
    ctx = _ctx(rebound, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert 'filter="fade"' in xml
    assert 'presetID="22"' not in xml


def test_a_role_the_theme_does_not_bind_is_reported_against_the_component(theme):
    @component("t-badrole")
    def _draw(ctx):
        tf = textbox(ctx.slide, 1.0, 1.0, 2.0, 0.4)
        para(tf, "x", 12, ctx.fg(), first=True)
        return [[(tf._parent.shape_id, "squiggle")]]

    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(Placement(at={"cols": (0, 12)}, component="t-badrole"),),
    )
    with pytest.raises(LayoutError, match=r"motion role 'squiggle'.*known roles"):
        render_slide(_ctx(theme, spec))


def test_after_previous_chains_the_groups_onto_one_click(theme):
    """`one_at_a_time` normally spends a click per group; this spends one in total."""
    _two_shape_component("t-chain")
    chained = dataclasses.replace(
        theme, motion=dataclasses.replace(theme.motion, advance="after_previous", beat_ms=250)
    )
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(
            Placement(at={"cols": (0, 6)}, component="t-chain"),
            Placement(at={"cols": (6, 12)}, component="t-chain"),
        ),
    )
    ctx = _ctx(chained, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert re.findall(r'nodeType="(clickEffect|afterEffect)"', xml) == [
        "clickEffect",
        "afterEffect",
    ]
    assert 'delay="250"' in xml


def test_on_click_spends_a_click_per_group(theme):
    _two_shape_component("t-clicky")
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(
            Placement(at={"cols": (0, 6)}, component="t-clicky"),
            Placement(at={"cols": (6, 12)}, component="t-clicky"),
        ),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert re.findall(r'nodeType="(clickEffect|afterEffect)"', xml) == [
        "clickEffect",
        "clickEffect",
    ]


def _themed_transition(theme, **kw):
    from pptxkit.theme.model import Transition

    return dataclasses.replace(
        theme, motion=dataclasses.replace(theme.motion, transition=Transition(**kw))
    )


def test_the_themes_transition_is_written_onto_every_slide(theme):
    ctx = _ctx(
        _themed_transition(theme, kind="push", direction="u", speed="slow"),
        SlideSpec(index=1, title="T"),
    )
    render_slide(ctx)
    el = ctx.slide._element.find(qn("p:transition"))
    assert el.get("spd") == "slow"
    assert [c.tag.split("}")[1] for c in el] == ["push"]
    assert el[0].get("dir") == "u"


def test_a_theme_that_names_no_transition_writes_none(theme):
    ctx = _ctx(theme, SlideSpec(index=1, title="T"))
    render_slide(ctx)
    assert ctx.slide._element.find(qn("p:transition")) is None


def test_a_slide_may_refuse_the_themes_transition(theme):
    ctx = _ctx(
        _themed_transition(theme, kind="fade"), SlideSpec(index=1, title="T", transition="none")
    )
    render_slide(ctx)
    assert ctx.slide._element.find(qn("p:transition")) is None


def test_a_slide_naming_a_transition_is_rejected_and_points_at_the_theme(theme):
    """Which transition is a look, and a look is the theme's."""
    ctx = _ctx(
        _themed_transition(theme, kind="fade"), SlideSpec(index=1, title="T", transition="push")
    )
    with pytest.raises(LayoutError, match=r"slide 1: transition 'push'.*only say 'none'"):
        render_slide(ctx)


def test_the_transition_precedes_the_timing_on_an_animated_slide(theme):
    _probe("t-trans", [])
    ctx = _ctx(
        _themed_transition(theme, kind="fade"),
        SlideSpec(
            index=1,
            animate="together",
            place=(Placement(at={"cols": (0, 6)}, component="t-trans"),),
        ),
    )
    render_slide(ctx)
    tags = [el.tag.split("}")[1] for el in ctx.slide._element]
    assert tags.index("transition") < tags.index("timing")


def test_a_chart_animate_on_a_non_chart_placement_is_rejected(theme):
    _probe("t-notchart", [])
    spec = SlideSpec(
        index=1,
        animate="by_category",
        place=(Placement(at={"cols": (0, 12)}, component="t-notchart"),),
    )
    with pytest.raises(
        LayoutError, match=r"slide 1: animate 'by_category' only applies to a native chart"
    ):
        render_slide(_ctx(theme, spec))


# --- composed chrome -------------------------------------------------------


def test_a_placed_chrome_line_gets_a_frame_of_its_own(theme):
    spec = SlideSpec(
        index=1,
        kicker="Q3",
        title="First",
        subtitle="Sub",
        chrome={"title": ChromeField(at={"box": (0.1, 0.6, 0.5, 0.1)})},
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    assert sorted(_texts(ctx.slide)) == ["First", "Q3\nSub"]


def test_a_placed_chrome_line_lands_where_its_box_says(theme):
    spec = SlideSpec(
        index=1, title="Low", chrome={"title": ChromeField(at={"box": (0.1, 0.63, 0.5, 0.1)})}
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    box = next(s for s in ctx.slide.shapes if s.has_text_frame)
    assert box.top / _EMU == pytest.approx(7.5 * 0.63)
    assert box.left / _EMU == pytest.approx(13.333 * 0.1)


def test_a_centred_chrome_line_is_written_centred(theme):
    spec = SlideSpec(
        index=1,
        title="Middle",
        chrome={"title": ChromeField(at={"box": (0, 0.05, 1, 0.1)}, align="center")},
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    box = next(s for s in ctx.slide.shapes if s.has_text_frame)
    assert box.text_frame.paragraphs[0].alignment == ALIGN["center"]


def test_a_bottom_anchored_chrome_line_grows_upward_from_its_baseline(theme):
    spec = SlideSpec(
        index=1,
        title="Hero",
        chrome={"title": ChromeField(at={"box": (0.05, 0.4, 0.8, 0.2)}, anchor="bottom")},
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    box = next(s for s in ctx.slide.shapes if s.has_text_frame)
    assert box.text_frame.vertical_anchor == ANCHOR["bottom"]


def test_a_chrome_line_can_take_a_rung_other_than_its_own_name(theme):
    spec = SlideSpec(
        index=1,
        title="Hero",
        chrome={"title": ChromeField(at={"box": (0.05, 0.4, 0.8, 0.2)}, rung="display")},
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    run = next(s for s in ctx.slide.shapes if s.has_text_frame).text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == pytest.approx(theme.style("display").size)


def test_a_chrome_line_takes_its_ink_from_the_pair_it_names(theme):
    pair = theme.palette.pair("accent-1")
    spec = SlideSpec(
        index=1,
        title="On a panel",
        chrome={"title": ChromeField(at={"box": (0.05, 0.4, 0.4, 0.2)}, pair="accent-1")},
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    run = next(s for s in ctx.slide.shapes if s.has_text_frame).text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == pair.fg
    recorded = ctx.manifest.slides[0].shapes[0]
    # The pair lends its ink; the paper stays whatever the slide really painted, so a
    # contrast check cannot be fooled by a background nothing laid down.
    assert (recorded.fg, recorded.bg) == (pair.fg, theme.palette.pair("page").bg)


def test_an_ink_that_vanishes_on_the_paper_the_slide_paints_is_rejected(theme):
    """A named colour that cannot be read is the white-on-white failure, so it raises."""
    spec = SlideSpec(
        index=1,
        kicker="INVISIBLE",
        chrome={"kicker": ChromeField(at={"box": (0.05, 0.4, 0.4, 0.1)}, ink="page")},
    )
    # A kicker is small type, so the threshold is the 4.5:1 QA would hold it to.
    with pytest.raises(LayoutError, match="below 4.5:1"):
        render_slide(_ctx(theme, spec))


def test_an_ink_over_a_panel_reads_against_the_panel_not_the_slide(theme):
    """The panel is painted first, so text on it is measured against the panel."""
    spec = SlideSpec(
        index=1,
        kicker="ON THE PANEL",
        chrome={"kicker": ChromeField(at={"box": (0.05, 0.05, 0.3, 0.1)}, ink="page")},
        place=[
            Placement(
                at={"box": (0, 0, 0.45, 1)}, bleed=True, component="panel", body={"pair": "inverse"}
            )
        ],
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    kicker = next(r for r in ctx.manifest.slides[0].shapes if r.text == "ON THE PANEL")
    assert kicker.bg == theme.palette.pair("inverse").bg


def test_chrome_is_drawn_over_the_placements_it_overlaps(theme):
    """A title told to sit on a painted panel has to land on top of it."""
    spec = SlideSpec(
        index=1,
        title="Reversed",
        chrome={"title": ChromeField(at={"box": (0.05, 0.3, 0.3, 0.1)}, pair="accent-1")},
        place=(
            Placement(
                at={"box": (0.0, 0.0, 0.45, 1.0)},
                component="panel",
                body={"pair": "accent-1"},
                bleed=True,
            ),
        ),
    )
    ctx = _ctx(theme, spec)
    render_slide(ctx)
    order = [s.shape_type for s in ctx.slide.shapes]
    last = list(ctx.slide.shapes)[-1]
    assert len(order) == 2
    assert last.has_text_frame and last.text_frame.text == "Reversed"


# --- interactive click-to-reveal --------------------------------------------


def test_a_reveals_placement_wires_an_interactive_trigger(theme):
    """`add_click_reveals` had exactly one caller in the repo — a test — until this."""
    _probe("t-trigger", [])
    _probe("t-target", [])
    spec = SlideSpec(
        index=1,
        place=(
            Placement(at={"cols": (0, 6)}, component="t-trigger", id="question"),
            Placement(at={"cols": (6, 12)}, component="t-target", reveals="question"),
        ),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    xml = ctx.slide._element.find(qn("p:timing")).xml
    assert 'nodeType="interactiveSeq"' in xml
    assert 'evt="onClick"' in xml
    assert ctx.manifest.slides[0].animations[0]["kind"] == "click_reveals"


def test_reveals_names_the_trigger_shape_and_the_target_shape(theme):
    _probe("t-trig2", [])
    _probe("t-targ2", [])
    spec = SlideSpec(
        index=1,
        place=(
            Placement(at={"cols": (0, 6)}, component="t-trig2", id="q"),
            Placement(at={"cols": (6, 12)}, component="t-targ2", reveals="q"),
        ),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    trigger, target = (s.shape_id for s in ctx.slide.shapes)
    timing = ctx.slide._element.find(qn("p:timing"))
    seq = [s for s in timing.iter(qn("p:cTn")) if s.get("nodeType") == "interactiveSeq"]
    assert len(seq) == 1
    # The trigger is what the sequence listens to; the target is what it makes visible.
    listens = seq[0].find(qn("p:stCondLst")).find(qn("p:cond"))
    assert listens.get("evt") == "onClick"
    assert listens.find(f".//{{{_PML}}}spTgt").get("spid") == str(trigger)
    revealed = {s.get("spid") for s in seq[0].iter(qn("p:spTgt"))} - {str(trigger)}
    assert revealed == {str(target)}


def test_every_shape_the_trigger_placement_drew_listens_for_the_click(theme):
    """A card's plate is drawn first and its words last, and only the shape clicked
    fires. Wiring one of the two leaves the other half of the trigger dead."""
    _two_shape_component("t-trigpair")
    _probe("t-targ5", [])
    spec = SlideSpec(
        index=1,
        place=(
            Placement(at={"cols": (0, 6)}, component="t-trigpair", id="q"),
            Placement(at={"cols": (6, 12)}, component="t-targ5", reveals="q"),
        ),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    plate, words, target = (s.shape_id for s in ctx.slide.shapes)
    timing = ctx.slide._element.find(qn("p:timing"))
    seqs = [s for s in timing.iter(qn("p:cTn")) if s.get("nodeType") == "interactiveSeq"]
    assert len(seqs) == 2
    listens = {
        s.find(qn("p:stCondLst")).find(qn("p:cond")).find(f".//{{{_PML}}}spTgt").get("spid")
        for s in seqs
    }
    assert listens == {str(plate), str(words)}
    revealed = {t.get("spid") for s in seqs for t in s.iter(qn("p:spTgt"))} - listens
    assert revealed == {str(target)}


def test_a_two_shape_trigger_still_records_one_revealed_group(theme):
    """Two interactions, one shape revealed. A group per pair would tell QA the slide
    reveals the same shape twice."""
    _two_shape_component("t-trigpair2")
    _probe("t-targ6", [])
    spec = SlideSpec(
        index=1,
        place=(
            Placement(at={"cols": (0, 6)}, component="t-trigpair2", id="q"),
            Placement(at={"cols": (6, 12)}, component="t-targ6", reveals="q"),
        ),
    )
    ctx = _ctx(theme, spec)

    render_slide(ctx)

    steps = ctx.manifest.slides[0].animations[0]["steps"]
    # One step, not one per trigger shape. The probe draws without recording, so the
    # step keeps the id — see test_a_reveal_target_that_was_never_recorded_keeps_its_id.
    assert steps == [[f"shape {list(ctx.slide.shapes)[-1].shape_id}"]]


def test_a_reveals_naming_no_placement_lists_the_ids_that_exist(theme):
    _probe("t-trig3", [])
    _probe("t-targ3", [])
    spec = SlideSpec(
        index=1,
        place=(
            Placement(at={"cols": (0, 6)}, component="t-trig3", id="question"),
            Placement(at={"cols": (6, 12)}, component="t-targ3", reveals="typo"),
        ),
    )
    with pytest.raises(LayoutError, match=r"'reveals: typo' names no placement.*question"):
        render_slide(_ctx(theme, spec))


def test_a_reveals_naming_itself_is_refused(theme):
    _probe("t-self", [])
    spec = SlideSpec(
        index=1, place=(Placement(at={"cols": (0, 6)}, component="t-self", id="me", reveals="me"),)
    )
    with pytest.raises(LayoutError, match="names itself"):
        render_slide(_ctx(theme, spec))


def test_reveals_and_animate_cannot_share_a_slide(theme):
    """Both want the slide's one timing tree, and they are different kinds."""
    _probe("t-trig4", [])
    _probe("t-targ4", [])
    spec = SlideSpec(
        index=1,
        animate="one_at_a_time",
        place=(
            Placement(at={"cols": (0, 6)}, component="t-trig4", id="q"),
            Placement(at={"cols": (6, 12)}, component="t-targ4", reveals="q"),
        ),
    )
    with pytest.raises(LayoutError, match="cannot share a slide"):
        render_slide(_ctx(theme, spec))
