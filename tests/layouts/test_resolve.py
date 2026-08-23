import pathlib
import types

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from pptxkit.errors import ThemeError
from pptxkit.layouts.resolve import pick_compose_layout

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

SAMPLES = pathlib.Path(__file__).resolve().parents[2] / "templates"


def _picture_heavy_sample() -> pathlib.Path | None:
    """The first corpus template most of whose layouts hold a picture placeholder, chosen by that
    property rather than by filename. Read with python-pptx, independent of the resolver."""
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    for path in sorted(SAMPLES.glob("*.pptx")) if SAMPLES.is_dir() else []:
        layouts = Presentation(str(path)).slide_layouts
        pictured = sum(
            1
            for layout in layouts
            if any(
                ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE for ph in layout.placeholders
            )
        )
        if len(layouts) >= 6 and pictured > len(layouts) // 2:
            return path
    return None


PICTURE_LAYOUTS = _picture_heavy_sample()


def _element(*, bg: str | None = None, show_master: str = "1"):
    inner = f"<p:bg><p:bgPr>{bg}</p:bgPr></p:bg>" if bg else ""
    return etree.fromstring(
        f'<p:sldLayout xmlns:p="{_P}" xmlns:a="{_A}" showMasterSp="{show_master}">'
        f"<p:cSld>{inner}</p:cSld></p:sldLayout>"
    )


_LATENT = (PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER)


def _ph(idx, ph_type):
    return types.SimpleNamespace(placeholder_format=types.SimpleNamespace(idx=idx, type=ph_type))


class _Layout:
    """A stand-in layout. ``latent`` placeholders sit at idx >= 10, like the real ones."""

    def __init__(self, name, *, content=0, latent=0, picture=0, decor=0, **element_kwargs):
        self.name = name
        self.placeholders = [_ph(i, PP_PLACEHOLDER.BODY) for i in range(content)]
        self.placeholders += [_ph(10 + i, PP_PLACEHOLDER.PICTURE) for i in range(picture)]
        self.placeholders += [
            _ph(10 + picture + i, _LATENT[i % len(_LATENT)]) for i in range(latent)
        ]
        self.shapes = list(self.placeholders) + [object()] * decor
        self.element = _element(**element_kwargs)


class _Master:
    def __init__(self, *layouts):
        self.slide_layouts = list(layouts)


class _Prs:
    def __init__(self, *masters):
        self.slide_masters = list(masters)


def test_the_emptiest_layout_wins():
    picked = pick_compose_layout(
        _Prs(_Master(_Layout("busy", content=5), _Layout("bare", content=1)))
    )
    assert picked.name == "bare"


def test_a_content_placeholder_outweighs_three_latent_ones():
    """Date, footer and slide-number placeholders are never cloned onto a slide."""
    picked = pick_compose_layout(
        _Prs(_Master(_Layout("titled", content=1), _Layout("blank", latent=3)))
    )
    assert picked.name == "blank"


def test_a_picture_placeholder_at_a_latent_index_is_still_cloneable():
    """python-pptx clones by placeholder type, so idx says nothing about latency. Count
    a lone picture placeholder as empty and every generated slide carries its frame."""
    picked = pick_compose_layout(
        _Prs(_Master(_Layout("pictured", picture=1), _Layout("blank", latent=3)))
    )
    assert picked.name == "blank"


def test_decoration_breaks_a_placeholder_tie():
    picked = pick_compose_layout(
        _Prs(_Master(_Layout("decorated", latent=1, decor=4), _Layout("plain", latent=1)))
    )
    assert picked.name == "plain"


def test_a_layout_on_a_later_master_is_reachable():
    """prs.slide_layouts exposes only the first master's layouts."""
    prs = _Prs(
        _Master(_Layout("m0-a", content=4), _Layout("m0-b", content=3)),
        _Master(_Layout("m1-bare", content=1)),
    )
    assert pick_compose_layout(prs).name == "m1-bare"


def test_layouts_that_would_compose_identically_break_by_document_order():
    prs = _Prs(_Master(_Layout("first", latent=1), _Layout("second", latent=1)))
    assert pick_compose_layout(prs).name == "first"


def test_equally_empty_layouts_on_two_masters_are_rejected():
    """Two masters means two colour schemes; picking one silently rebrands the deck."""
    prs = _Prs(_Master(_Layout("m0-bare", latent=1)), _Master(_Layout("m1-bare", latent=1)))
    with pytest.raises(ThemeError, match="disagree on what a slide would inherit"):
        pick_compose_layout(prs)


def test_equally_empty_layouts_with_different_backgrounds_are_rejected():
    prs = _Prs(_Master(_Layout("plain", latent=1), _Layout("navy", latent=1, bg="<a:solidFill/>")))
    with pytest.raises(ThemeError, match="'navy'"):
        pick_compose_layout(prs)


def test_a_template_with_no_layouts_at_all_is_rejected():
    with pytest.raises(ThemeError, match="no slide layouts"):
        pick_compose_layout(_Prs(_Master()))


def test_it_picks_the_blank_layout_of_a_real_template(synthetic_template):
    """Blank carries only date/footer/slide-number; every other stock layout carries
    more. Chosen by emptiness — the name is only the proof."""
    assert pick_compose_layout(Presentation(str(synthetic_template))).name == "Blank"


@pytest.mark.skipif(PICTURE_LAYOUTS is None, reason="no picture-heavy corpus sample")
def test_a_corpus_template_of_picture_layouts_composes_a_slide_with_no_frames():
    """An idx-based ranking would prefer a picture frame over the body of every slide. The
    contract is the outcome: whatever layout is picked, the slide starts with no frames."""
    prs = Presentation(str(PICTURE_LAYOUTS))
    picked = pick_compose_layout(prs)
    assert list(prs.slides.add_slide(picked).placeholders) == []


def test_compose_layout_names_a_layout_outright(synthetic_template):
    """The escape hatch for a template whose emptiest layouts cannot be ranked apart."""
    from pptx import Presentation

    prs = Presentation(str(synthetic_template))
    assert pick_compose_layout(prs, prefer="Title Only").name == "Title Only"


def test_compose_layout_naming_an_absent_layout_lists_what_there_is(synthetic_template):
    from pptx import Presentation

    prs = Presentation(str(synthetic_template))
    with pytest.raises(ThemeError, match="which the template does not define"):
        pick_compose_layout(prs, prefer="No Such Layout")
