"""The icon component: what it actually puts on the slide. A corpus build cannot tell a glyph
from a coloured rectangle, so these read the shape's own geometry."""

from __future__ import annotations

import pytest

import pptxkit.components  # noqa: F401 — registers the built-in components
from pptxkit.errors import SpecError
from pptxkit.layouts.components import get_component

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _draw(ctx):
    get_component("icon")(ctx)
    return list(ctx.slide.shapes)[-1]


def test_an_icon_is_drawn_as_path_geometry_not_as_a_box(ctx_factory):
    """The whole point of the SVG conversion: a real outline, not a placeholder."""
    xml = _draw(ctx_factory({"icon": {"name": "check"}}))._element.xml
    assert f'<a:custGeom xmlns:a="{_A}"' in xml or "<a:custGeom>" in xml
    assert "<a:moveTo>" in xml and "<a:pt " in xml
    assert "<a:prstGeom" not in xml


def test_every_legacy_name_lands_as_real_geometry(ctx_factory, legacy_glyphs):
    """A path that parsed to nothing draws an invisible shape and reports no error. A glyph
    carries an outline, a plain shape a preset; neither may carry no geometry at all."""
    from pptxkit.icons.shapes import SHAPES

    for name in legacy_glyphs:
        xml = _draw(ctx_factory({"icon": {"name": name}}))._element.xml
        if name in SHAPES:
            assert "<a:prstGeom" in xml, f"{name} lost its preset"
            assert "<a:custGeom" not in xml, f"{name} drew art instead of geometry"
        else:
            assert xml.count("<a:pt ") >= 3, name


def test_a_plain_shape_is_a_preset_and_not_the_icon_sets_drawing_of_one(ctx_factory):
    """No rhombus exists in the vendored set, so `diamond` resolved to a faceted gem. A square
    is a square — the geometry is the meaning, and a preset states it exactly."""
    from pptx.enum.shapes import MSO_SHAPE

    for name, preset in (
        ("diamond", MSO_SHAPE.DIAMOND),
        ("circle", MSO_SHAPE.OVAL),
        ("square", MSO_SHAPE.RECTANGLE),
        ("ring", MSO_SHAPE.DONUT),
        ("triangle", MSO_SHAPE.ISOSCELES_TRIANGLE),
    ):
        assert _draw(ctx_factory({"icon": {"name": name}})).auto_shape_type == preset, name


def test_a_configured_directory_still_beats_the_preset(ctx_factory, tmp_path, monkeypatch):
    """A brand shipping its own diamond.svg gets its diamond. The preset is the answer
    when nobody has a better one, not a name the library takes away."""
    (tmp_path / "diamond.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24">'
        '<path d="M2 2 H22 V22 H2 Z"/></svg>'
    )
    monkeypatch.setenv("PPTXKIT_ICON_DIR", str(tmp_path))
    xml = _draw(ctx_factory({"icon": {"name": "diamond"}}))._element.xml
    assert "<a:custGeom" in xml and "<a:prstGeom" not in xml


def test_a_box_no_colour_reads_across_gets_a_plate_like_a_chrome_line(ctx_factory):
    """Half the box white, half the fixture ink: nothing reads across it, so the mark must
    plate rather than go invisible."""
    from pptxkit.theme.model import Rect

    ctx = ctx_factory({"icon": {"name": "check"}})
    rect = ctx.body_rect
    # The glyph squares off to the placement's short side, anchored at the left.
    side = min(rect.width, rect.height)
    ctx.panels.append((Rect(rect.left, rect.top, side / 2, rect.height), "FFFFFF"))
    ctx.panels.append((Rect(rect.left + side / 2, rect.top, side / 2, rect.height), "2D0937"))
    before = len(ctx.slide.shapes)
    get_component("icon")(ctx)
    assert len(ctx.slide.shapes) == before + 2  # the plate, then the glyph
    recorded = ctx.manifest.slides[0].shapes[-1]
    assert (recorded.fg, recorded.bg) == ("2D0937", "FFFFFF")


def test_an_icon_is_squared_off_inside_a_wide_placement(ctx_factory):
    """A glyph stretched to a wide box is the one distortion nothing downstream undoes."""
    shape = _draw(ctx_factory({"icon": {"name": "square"}}))
    assert shape.width == shape.height


def test_size_scales_the_glyph_within_its_placement(ctx_factory):
    full = _draw(ctx_factory({"icon": {"name": "square"}})).width
    half = _draw(ctx_factory({"icon": {"name": "square", "size": 0.5}})).width
    assert half == pytest.approx(full / 2, rel=0.02)


def test_an_ink_role_is_painted_verbatim(ctx_factory):
    """A named colour is the author's instruction, not a suggestion to contrast-check."""
    ctx = ctx_factory({"icon": {"name": "check", "ink": "accent-2"}})
    expected = ctx.theme.palette.role("accent-2")
    assert f'<a:srgbClr val="{expected}"/>' in _draw(ctx)._element.xml


def test_an_unnamed_icon_points_at_the_catalogue(ctx_factory):
    with pytest.raises(SpecError, match="needs a 'name:'.*docs/glyphs.md"):
        _draw(ctx_factory({"icon": {}}))


def test_an_unknown_field_is_rejected(ctx_factory):
    from pptxkit.errors import LayoutError

    with pytest.raises(LayoutError, match="unknown field 'colour'"):
        _draw(ctx_factory({"icon": {"name": "check", "colour": "red"}}))


def test_a_card_draws_a_bare_icon_name_as_geometry(ctx_factory):
    """`icon: rocket` is a glyph; `icon: rocket.png` stays a picture file."""
    ctx = ctx_factory({"card": {"icon": "star", "heading": "H", "body": "B"}})
    get_component("card")(ctx)
    assert any("<a:custGeom" in s._element.xml for s in ctx.slide.shapes)


def test_shape_ids_stay_unique_when_several_glyphs_land(ctx_factory):
    """A duplicate id opens as a repair prompt, which no check downstream would see."""
    ctx = ctx_factory({"icon": {"name": "check"}})
    for name in ("check", "star", "globe", "gear"):
        ctx.body = {"name": name}
        get_component("icon")(ctx)
    ids = [s.shape_id for s in ctx.slide.shapes]
    assert len(ids) == len(set(ids)), ids
