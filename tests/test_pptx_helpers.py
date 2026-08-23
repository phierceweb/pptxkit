"""Tests for the generic deck / animation helpers."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import re

import pytest

from pptxkit.errors import LayoutError
from pptxkit.motion import (
    add_chart_build,
    add_click_build,
    add_click_reveals,
    add_click_sequence,
)
from pptxkit.utils.deck import delete_slide
from pptxkit.utils.shapes import bring_to_front


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_delete_slide_removes_one():
    prs = Presentation()
    _blank_slide(prs)
    _blank_slide(prs)
    before = len(list(prs.slides._sldIdLst))
    delete_slide(prs, 0)
    assert len(list(prs.slides._sldIdLst)) == before - 1


def test_add_click_reveals_appends_timing():
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    add_click_reveals(slide, [(box.shape_id, box.shape_id)])
    assert slide._element.find(qn("p:timing")) is not None


def test_a_reveal_with_no_pairs_is_refused_before_anything_is_written():
    """No pairs leaves the tmRoot's `<p:childTnLst/>` empty, and `CT_TimeNodeList`
    requires a child — a tree that reveals nothing is not a tree with nothing in it."""
    prs = Presentation()
    slide = _blank_slide(prs)

    with pytest.raises(LayoutError, match="needs at least one"):
        add_click_reveals(slide, [])
    assert slide._element.find(qn("p:timing")) is None


def test_a_second_timing_tree_on_one_slide_is_refused():
    """``CT_Slide`` allows one ``<p:timing>``; two charts both asked to build append two — an
    invalid file that LibreOffice converts happily and ``check_package`` reports clean."""
    prs = Presentation()
    slide = _blank_slide(prs)
    one = slide.shapes.add_textbox(0, 0, 100, 100)
    two = slide.shapes.add_textbox(0, 0, 100, 100)
    add_chart_build(slide, one.shape_id, "category", parts=2)

    with pytest.raises(LayoutError, match="already carries an animation"):
        add_chart_build(slide, two.shape_id, "category", parts=2)
    assert len(slide._element.findall(qn("p:timing"))) == 1


def test_timing_is_inserted_before_ext_lst_not_appended_after_it():
    """``CT_Slide`` is an xsd:sequence — ``timing`` precedes ``extLst``, and a bare append lands
    after it. LibreOffice silently repairs the order on import, so no render shows the corruption.
    """
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    slide._element.append(slide._element.makeelement(qn("p:extLst"), {}))

    add_click_build(slide, [box.shape_id])

    tags = [el.tag.split("}")[1] for el in slide._element]
    assert tags.index("timing") < tags.index("extLst")


def _texty(slide, words: str):
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text_frame.text = words
    return box


def _wordless(slide):
    return slide.shapes.add_shape(MSO_SHAPE.OVAL, 0, 0, 100, 100)


def _grp_ids(slide) -> dict[str, str | None]:
    """The `grpId` on each entrance node, keyed by the spid its `spTgt` names."""
    ids: dict[str, str | None] = {}
    for node in slide._element.find(qn("p:timing")).iter(qn("p:cTn")):
        if node.get("nodeType") not in ("clickEffect", "withEffect", "afterEffect"):
            continue
        (spid,) = {t.get("spid") for t in node.iter(qn("p:spTgt"))}
        ids[spid] = node.get("grpId")
    return ids


def test_bld_p_is_written_only_for_shapes_that_carry_text():
    """[MS-OI29500] 19.5.16(c): a bldP's spid must name an sp holding a `t` element, so naming a
    picture, connector or text-free icon is a claim the shape cannot honour."""
    prs = Presentation()
    slide = _blank_slide(prs)
    worded = _texty(slide, "a heading")
    empty = slide.shapes.add_textbox(0, 0, 100, 100)  # an sp with no text
    add_click_build(slide, [worded.shape_id, empty.shape_id])

    xml = slide._element.find(qn("p:timing")).xml
    assert re.findall(r'<p:bldP spid="(\d+)"', xml) == [str(worded.shape_id)]
    # Both still animate — only the build-list entry is withheld.
    assert xml.count('<p:spTgt spid="%d"' % empty.shape_id) == 2


def test_a_slide_of_only_text_free_shapes_writes_no_build_list_and_no_grp_id():
    """`CT_BuildList` requires a child, so an empty `<p:bldLst/>` is invalid. With no
    build list there is nothing for `grpId` to name, and the schema makes it optional."""
    prs = Presentation()
    slide = _blank_slide(prs)
    one = slide.shapes.add_textbox(0, 0, 100, 100)
    two = slide.shapes.add_textbox(0, 0, 100, 100)
    add_click_build(slide, [one.shape_id, two.shape_id])

    xml = slide._element.find(qn("p:timing")).xml
    assert "<p:bldLst" not in xml
    assert "grpId" not in xml


def test_grp_id_is_decided_per_shape_not_per_slide():
    """`grpId` names a build-list entry, and on a mixed slide only the worded shape has one —
    deciding it once for the whole slide writes it onto the icon's node, where it names nothing."""
    prs = Presentation()
    slide = _blank_slide(prs)
    worded = _texty(slide, "a heading")
    icon = _wordless(slide)
    add_click_build(slide, [worded.shape_id, icon.shape_id])

    assert _grp_ids(slide) == {str(worded.shape_id): "0", str(icon.shape_id): None}
    timing = slide._element.find(qn("p:timing"))
    built = timing.findall(f"{qn('p:bldLst')}/{qn('p:bldP')}")
    assert [b.get("spid") for b in built] == [str(worded.shape_id)]


def test_a_click_sequence_withholds_bld_p_for_text_free_shapes_too():
    prs = Presentation()
    slide = _blank_slide(prs)
    worded = _texty(slide, "step one")
    icon = slide.shapes.add_textbox(0, 0, 100, 100)
    add_click_sequence(slide, [[worded.shape_id], [icon.shape_id]])

    xml = slide._element.find(qn("p:timing")).xml
    assert re.findall(r'<p:bldP spid="(\d+)"', xml) == [str(worded.shape_id)]
    assert xml.count('nodeType="clickEffect"') == 2


def test_a_click_sequence_decides_grp_id_inside_the_group_too():
    """One click, one group, two shapes of different kinds: the choice is per item, so
    a group with a worded shape in it does not lend its `grpId` to the icon beside it."""
    prs = Presentation()
    slide = _blank_slide(prs)
    worded = _texty(slide, "step one")
    icon = _wordless(slide)
    add_click_sequence(slide, [[worded.shape_id, icon.shape_id]])

    assert _grp_ids(slide) == {str(worded.shape_id): "0", str(icon.shape_id): None}


def test_a_click_reveal_writes_no_grp_id_because_it_writes_no_build_list():
    """An interactive tree carries no `<p:bldLst>` at all, so every `grpId` in it names
    an entry that was never written."""
    prs = Presentation()
    slide = _blank_slide(prs)
    trigger = _texty(slide, "the question")
    answer = _texty(slide, "the answer")
    add_click_reveals(slide, [(trigger.shape_id, answer.shape_id)])

    timing = slide._element.find(qn("p:timing"))
    assert timing.find(qn("p:bldLst")) is None
    assert {n.get("grpId") for n in timing.iter(qn("p:cTn"))} == {None}


def test_bring_to_front_moves_shape_last():
    prs = Presentation()
    slide = _blank_slide(prs)
    first = slide.shapes.add_textbox(0, 0, 100, 100)
    slide.shapes.add_textbox(0, 0, 100, 100)  # second shape, currently drawn on top
    bring_to_front(first)
    assert first._element.getparent()[-1] is first._element
