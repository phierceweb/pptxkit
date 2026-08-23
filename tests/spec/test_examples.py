"""Every deck shipped under examples/ must parse — nothing else in the suite reads them,
so a broken demo deck otherwise ships with the suite green."""

import pathlib

import pytest

from pptxkit.spec import parse_deck

EXAMPLES = pathlib.Path(__file__).resolve().parents[2] / "examples"
DECKS = sorted(EXAMPLES.glob("*.deck.yaml"))


def test_there_are_example_decks_to_check():
    assert DECKS, f"no *.deck.yaml under {EXAMPLES}"


@pytest.mark.parametrize("deck", DECKS, ids=lambda p: p.name)
def test_a_shipped_example_deck_parses(deck):
    assert parse_deck(deck).slides


TREATMENTS = EXAMPLES / "title-treatments.deck.yaml"


@pytest.fixture
def treatments(tmp_path, monkeypatch, synthetic_template):
    """The shipped title-treatment deck, compiled: parsing catches a bad field, only a build
    catches a title drawn *under* its panel. The theme is a stand-in for the brand's own."""
    from pptx import Presentation

    from pptxkit.compile import build_deck

    (tmp_path / "assets").mkdir()
    (tmp_path / "assets" / "t.pptx").write_bytes(synthetic_template.read_bytes())
    (tmp_path / "base.yaml").write_text(
        "name: base\ntemplate: assets/t.pptx\n"
        "bind: {page: lt1, ink: dk1, inverse: dk1, line: lt2}\n"
    )
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path))
    result = build_deck(TREATMENTS, out=tmp_path / "out.pptx")
    return result, Presentation(str(result.deck))


def test_the_title_treatment_deck_compiles_end_to_end(treatments):
    result, _ = treatments
    assert result.slides == len(parse_deck(TREATMENTS).slides)


def test_the_title_on_a_panel_is_drawn_over_the_panel(treatments):
    _, prs = treatments
    title_text = "Reversed out of a panel"
    slide = next(
        s
        for s in prs.slides
        if title_text in "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
    )
    order = [sh.text_frame.text if sh.has_text_frame else "" for sh in slide.shapes]
    # The panel carries no text of its own, so it is the empty entry.
    assert order.index("") < order.index(title_text)
