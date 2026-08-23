"""What a conformance run is allowed to leave behind. A failed exercise skips the cleanup that
removes a passing one's spec, deck and manifest, so its files have to be in scratch rather than in
the directory a reader opens — entirely a property of the raising path, which no passing run takes."""

from __future__ import annotations

import pytest
from pptx import Presentation

from pptxkit.conform import conform
from pptxkit.paths import SCRATCH

# One that any template can draw, and one that no placement can hold — 60 columns
# leave less width than the theme's gutter pads a cell by, which is a LayoutError.
EXERCISES = {
    "fine": {
        "title": "A slide",
        "place": [{"at": {"cols": "full"}, "bullets": {"items": ["One", "Two"]}}],
    },
    "doomed": {
        "title": "Too many columns",
        "place": [{"at": {"cols": "full"}, "table": {"rows": [["c"] * 60]}}],
    },
}


@pytest.fixture
def template(tmp_path):
    """A stock Office deck is a usable template: it carries a theme part to derive from."""
    path = tmp_path / "stock.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    return path


def test_a_failing_exercise_leaves_nothing_in_the_directory_a_reader_opens(template, tmp_path):
    out = tmp_path / "conformed"
    result = conform(template, out, exercises=EXERCISES)

    assert result.passed == ["fine"], result.report()
    assert [name for name, _ in result.failed] == ["doomed"], result.report()
    strays = sorted(p.name for p in out.iterdir() if p.name.startswith("_"))
    assert strays == [], f"per-exercise scratch left in the open: {strays}"


def test_the_failed_exercise_is_kept_in_scratch_so_it_can_be_re_run(template, tmp_path):
    """Tidy is not the same as gone: the spec that produced a FAIL is what you want
    when you go to diagnose it."""
    out = tmp_path / "conformed"
    conform(template, out, exercises=EXERCISES)

    kept = sorted(p.name for p in (out / SCRATCH).iterdir())
    assert "_doomed.deck.yaml" in kept, kept
    assert "_fine.deck.yaml" not in kept, "a passing exercise should clean up after itself"


def test_the_generated_inputs_are_out_of_sight_too(template, tmp_path):
    """The stand-in photographs and the notes file are build inputs, not outputs."""
    out = tmp_path / "conformed"
    conform(template, out, exercises=EXERCISES)

    assert not list(out.glob("photo-*.png")), "photographs sit beside the deck"
    assert not (out / "NOTES.md").exists(), "the document component's source sits beside the deck"
    assert (out / SCRATCH / "photo-wide.png").is_file()


def test_the_scratch_directory_is_hidden():
    """Pinned as a literal because every other test here reads the constant, and so follows it
    wherever it is renamed."""
    assert SCRATCH == ".build"
