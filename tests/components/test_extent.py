"""`BodyResult.height` must cover what the component drew.

`place:` gives the next placement everything below that number, so under-reporting is
silent: nothing raises, and only `pptxkit qa` on a built deck sees the overhang. A shape
that exactly fills the placement is excluded.
"""

from __future__ import annotations

import pytest

import pptxkit.components  # noqa: F401 — registers the built-ins
from pptxkit.layouts.components import as_body_result, get_component

_EMU_PER_INCH = 914400
_TOLERANCE = 0.01

# One minimal body per component that needs no external file, browser or template. Optional
# fields that change geometry are set: an unset caption exercises no caption arithmetic.
BODIES = {
    "bullets": {"items": ["One", "Two", "Three"]},
    "callouts": {"items": [{"head": "A", "body": "First."}, {"head": "B", "body": "Second."}]},
    "card": {"heading": "A card", "body": "Some copy."},
    "code": {"heading": "spec.yaml", "lines": ["place:", "  - at: {cols: full}"]},
    "grid": {"caption": "12 columns"},
    "panel": {},
    "prose": {
        "paragraphs": [
            "One paragraph of real copy, long enough to wrap at the capped measure.",
            "And a second.",
        ],
        "cite": "A speaker",
    },
    "rule": {},
    "stats": {"items": [{"value": "12", "label": "one"}, {"value": "34", "label": "two"}]},
    "swatches": {
        "roles": ["ink", "accent-1", "line"],
        "caption": "A caption long enough to wrap onto a second line in the box "
        "the component sizes for it, which is where both defects were.",
    },
    "table": {"rows": [["a", "b"], ["c", "d"]]},
}


def _measured_bottom(ctx) -> float | None:
    """Lowest edge drawn, or None when every shape simply fills the placement."""
    rect = ctx.body_rect
    bottoms = []
    for shape in ctx.slide.shapes:
        top, height = shape.top / _EMU_PER_INCH, shape.height / _EMU_PER_INCH
        fills = abs(top - rect.top) < _TOLERANCE and abs(height - rect.height) < _TOLERANCE
        if not fills:
            bottoms.append(top + height)
    return max(bottoms) if bottoms else None


@pytest.mark.parametrize("name", sorted(BODIES))
def test_a_component_draws_no_lower_than_the_height_it_reported(name, ctx_factory):
    ctx = ctx_factory({"title": "T", name: BODIES[name]})
    reported = as_body_result(get_component(name)(ctx)).height
    assert ctx.slide.shapes, f"{name} drew nothing, so this proves nothing"

    bottom = _measured_bottom(ctx)
    if bottom is None:
        pytest.skip(f"{name} draws only a placement-sized frame — no extent to compare")
    drawn = bottom - ctx.body_rect.top
    assert drawn <= reported + _TOLERANCE, (
        f"{name} reported {reported:.3f}in but drew down to {drawn:.3f}in — "
        f"{drawn - reported:.3f}in of overhang the next placement is laid over"
    )


@pytest.mark.parametrize("name", sorted(BODIES))
def test_a_component_draws_nothing_above_its_placement(name, ctx_factory):
    """The other half of the same contract, and what `qa` reports as a bounds error."""
    ctx = ctx_factory({"title": "T", name: BODIES[name]})
    get_component(name)(ctx)
    top = min(s.top / _EMU_PER_INCH for s in ctx.slide.shapes)
    assert top >= ctx.body_rect.top - _TOLERANCE, (
        f"{name} drew {ctx.body_rect.top - top:.3f}in above its placement"
    )


@pytest.mark.parametrize("anchor,where", [("top", 0.0), ("middle", 0.5), ("bottom", 1.0)])
def test_a_placement_settles_its_content_where_anchor_says(anchor, where, ctx_factory):
    """`anchor:` moves the component's own extent, not only the type inside its frames.

    Without this a component that sizes to its content leaves every slide's slack piled at
    the bottom, which is the difference between airy and unfinished.
    """
    from pptxkit.layouts.compose import _settle

    ctx = ctx_factory({"title": "T", "stats": BODIES["stats"]})
    get_component("stats")(ctx)
    drawn = list(ctx.slide.shapes)
    rect = ctx.body_rect
    extent = (max(s.top + s.height for s in drawn) - min(s.top for s in drawn)) / _EMU_PER_INCH

    _settle(ctx, rect, anchor, drawn)
    top = min(s.top for s in drawn) / _EMU_PER_INCH
    expected = rect.top + (rect.height - extent) * where
    assert top == pytest.approx(expected, abs=0.02)


def test_settling_never_pushes_content_out_of_its_placement(ctx_factory):
    """The slack is what is distributed; a component filling its rect must not move."""
    from pptxkit.layouts.compose import _settle

    ctx = ctx_factory({"title": "T", "bullets": BODIES["bullets"]})
    get_component("bullets")(ctx)
    drawn = list(ctx.slide.shapes)
    before = [s.top for s in drawn]
    _settle(ctx, ctx.body_rect, "middle", drawn)
    assert [s.top for s in drawn] == before


@pytest.mark.parametrize(
    "name,body,n_frames",
    [
        ("callouts", {"items": [{"head": "A"}, {"head": "B", "body": "Longer copy here."}]}, 2),
        ("stats", {"items": [{"value": "1", "label": "one"}]}, 1),
        ("card", {"heading": "H", "body": "Copy."}, 1),
    ],
)
def test_content_sized_frames_ignore_the_placement_anchor(name, body, n_frames, ctx_factory):
    """`anchor:` positions the block via settling; re-anchoring type inside a padded frame
    drags it off the marker or sibling beside it."""
    from pptx.enum.text import MSO_ANCHOR

    ctx = ctx_factory({"title": "T", name: body})
    ctx.anchor = "middle"
    get_component(name)(ctx)
    frames = [s for s in ctx.slide.shapes if s.has_text_frame and s.text_frame.text]
    assert len(frames) >= n_frames
    assert all(f.text_frame.vertical_anchor in (MSO_ANCHOR.TOP, None) for f in frames)


def test_a_full_rect_frame_still_honours_the_placement_anchor(ctx_factory):
    """`bullets` draws one frame filling its placement, so frame anchoring IS the
    component's position — settling has no slack to act on there."""
    from pptx.enum.text import MSO_ANCHOR

    ctx = ctx_factory({"title": "T", "bullets": {"items": ["one", "two"]}})
    ctx.anchor = "middle"
    get_component("bullets")(ctx)
    frame = next(s for s in ctx.slide.shapes if s.has_text_frame)
    assert frame.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
