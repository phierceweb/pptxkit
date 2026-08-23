import base64
import dataclasses

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

import pptxkit.components  # noqa: F401 — registers the built-in components
from pptxkit.components.card import _ICON_LINES
from pptxkit.errors import LayoutError, ThemeError
from pptxkit.layouts.components import get_component, registered_components
from pptxkit.theme.model import Rect
from pptxkit.utils.text import LINE_HEIGHT

EMU = 914400

_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAE"
    "hQGAhKmMIQAAAABJRU5ErkJggg=="
)


@pytest.fixture
def icon_theme(theme, tmp_path, synthetic_template):
    """A theme whose template directory holds one icon file beside the template."""
    (tmp_path / "icon.png").write_bytes(_PNG_1X1)
    return dataclasses.replace(theme, template=synthetic_template)


def _draw(ctx):
    """Draw the card; return its shapes, the plate first."""
    get_component("card")(ctx)
    return list(ctx.slide.shapes)


def test_a_file_icon_is_found_beside_the_deck_spec(ctx_factory, theme, tmp_path):
    """A card's file icon resolves like any deck image — the spec's own directory first, not
    only the template's."""
    spec_dir = tmp_path / "deckdir"
    spec_dir.mkdir()
    (spec_dir / "logo.png").write_bytes(_PNG_1X1)
    ctx = ctx_factory({"card": {"icon": "logo.png", "heading": "Discovery"}}, base=spec_dir)
    shapes = _draw(ctx)
    assert any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes)


def _card(ctx):
    return _draw(ctx)[0]


def _paragraphs(ctx):
    return _draw(ctx)[-1].text_frame.paragraphs


def test_the_plate_fills_the_placement(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    plate = _card(ctx)
    rect = ctx.body_rect
    assert plate.left / EMU == pytest.approx(rect.left)
    assert plate.width / EMU == pytest.approx(rect.width)
    assert plate.height / EMU == pytest.approx(rect.height)


def test_the_plate_is_painted_in_the_background_of_the_pair_it_names(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery", "pair": "accent-3"}})
    assert str(_card(ctx).fill.fore_color.rgb) == ctx.theme.palette.pair("accent-3").bg


def test_the_default_plate_is_the_surface_pair(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    assert str(_card(ctx).fill.fore_color.rgb) == ctx.theme.palette.pair("surface").bg


def test_the_type_takes_the_plates_own_ink_not_the_slides(ctx_factory):
    """A light card on a dark slide needs dark type, which ctx.fg() is not."""
    ctx = ctx_factory({"card": {"heading": "Discovery", "body": "one line"}}, background="inverse")
    paragraphs = _paragraphs(ctx)
    expected = ctx.theme.palette.pair("surface").fg
    assert [str(p.runs[0].font.color.rgb) for p in paragraphs] == [expected, expected]
    assert expected != str(ctx.fg())


def test_the_heading_is_set_at_the_head_rung_and_the_body_at_the_body_rung(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery", "body": "one line"}})
    head, copy = _paragraphs(ctx)
    assert head.runs[0].font.size.pt == pytest.approx(ctx.theme.style("head").size)
    assert copy.runs[0].font.size.pt == pytest.approx(ctx.theme.style("body").size)


def test_a_body_only_card_starts_in_the_frames_first_paragraph(ctx_factory):
    ctx = ctx_factory({"card": {"body": "just the copy"}})
    paragraphs = _paragraphs(ctx)
    assert len(paragraphs) == 1
    assert paragraphs[0].runs[0].text == "just the copy"


def test_align_sets_the_cards_type_without_moving_the_plate(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    ctx.align = "center"
    plate, text = _draw(ctx)
    assert text.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    assert plate.left / EMU == pytest.approx(ctx.body_rect.left)


def test_the_type_is_inset_from_the_plate_by_the_themes_gutter(ctx_factory):
    """A 2.5in short side keeps the default radius's corner inside the 0.18in gutter,
    so the plate's straight sides are what the type is measured off."""
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    ctx.rect = Rect(1.0, 1.0, 4.0, 2.5)
    text = _draw(ctx)[-1]
    assert text.left / EMU == pytest.approx(ctx.body_rect.left + ctx.grid.gutter)
    assert text.top / EMU == pytest.approx(ctx.body_rect.top + ctx.grid.gutter)


def test_the_plate_is_rounded_by_default(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    assert _card(ctx).adjustments[0] > 0


def test_a_radius_of_zero_squares_the_plate(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery", "radius": 0}})
    assert _card(ctx).adjustments[0] == pytest.approx(0.0)


def test_a_radius_outside_the_fraction_range_is_refused(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery", "radius": 0.9}})
    with pytest.raises(LayoutError, match="'radius' is a fraction of the plate"):
        get_component("card")(ctx)


def test_a_stadiums_curve_insets_the_type_past_the_plates_square_corner(ctx_factory):
    """A stadium's corner is a circle of radius 2.65in, which at the 0.18in-deep top of the
    text frame has reached 1.69in in — so the frame starts at 1.87in, its top unmoved."""
    ctx = ctx_factory({"card": {"heading": "radius 0.5", "radius": 0.5}})
    text = _draw(ctx)[-1]
    assert text.left / EMU == pytest.approx(ctx.body_rect.left + 1.87)
    assert text.top / EMU == pytest.approx(ctx.body_rect.top + 0.18)


def test_the_default_radius_insets_the_type_by_the_curve_a_tall_plate_earns(ctx_factory):
    """Guards the small-radius end against a threshold: the same 5.30in-tall plate at
    the default 0.06 has a 0.318in corner, reaching 0.0315in in at the frame's top."""
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    text = _draw(ctx)[-1]
    assert text.left / EMU == pytest.approx(ctx.body_rect.left + 0.2115039)


def test_a_fill_too_close_to_the_paper_to_see_is_given_an_edge(ctx_factory):
    from pptx.enum.dml import MSO_FILL

    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    plate = _card(ctx)
    assert plate.line.fill.type == MSO_FILL.SOLID
    assert str(plate.line.color.rgb) == ctx.theme.palette.role("line")


def test_copy_that_will_not_fit_the_plate_is_refused_with_both_measurements(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    ctx.rect = Rect(1.0, 1.0, 2.0, 0.5)
    with pytest.raises(LayoutError, match="the card's type wants .*in but only .*in is left"):
        get_component("card")(ctx)


def test_a_plate_smaller_than_its_own_inset_is_refused(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    ctx.rect = Rect(1.0, 1.0, 0.2, 0.2)
    with pytest.raises(LayoutError, match="smaller than the .*in the theme's gutter insets"):
        get_component("card")(ctx)


def test_a_card_with_nothing_on_it_is_refused_and_names_the_panel_component(ctx_factory):
    ctx = ctx_factory({"card": {}})
    with pytest.raises(LayoutError, match="an empty plate is the 'panel' component"):
        get_component("card")(ctx)


def test_the_icon_is_placed_at_the_top_of_the_plate_sized_off_the_head_rung(
    ctx_factory, icon_theme
):
    ctx = ctx_factory(
        {"card": {"icon": "icon.png", "heading": "Discovery"}}, theme_override=icon_theme
    )
    ctx.rect = Rect(1.0, 1.0, 4.0, 2.5)
    get_component("card")(ctx)
    picture = ctx.slide.shapes[1]
    side = ctx.theme.style("head").size * LINE_HEIGHT / 72 * _ICON_LINES
    assert picture.width / EMU == pytest.approx(side)
    assert picture.height / EMU == pytest.approx(side)
    assert picture.left / EMU == pytest.approx(ctx.body_rect.left + ctx.grid.gutter)


def test_the_cards_type_starts_below_the_icon(ctx_factory, icon_theme):
    ctx = ctx_factory(
        {"card": {"icon": "icon.png", "heading": "Discovery"}}, theme_override=icon_theme
    )
    get_component("card")(ctx)
    picture, text = ctx.slide.shapes[1], ctx.slide.shapes[2]
    assert text.top >= picture.top + picture.height


def test_an_icon_only_card_needs_no_text(ctx_factory, icon_theme):
    ctx = ctx_factory({"card": {"icon": "icon.png"}}, theme_override=icon_theme)
    get_component("card")(ctx)
    assert len(ctx.slide.shapes) == 2


def test_an_icon_the_theme_cannot_resolve_is_refused_by_name(ctx_factory, icon_theme):
    ctx = ctx_factory({"card": {"icon": "missing.png"}}, theme_override=icon_theme)
    with pytest.raises(ThemeError, match="missing.png"):
        get_component("card")(ctx)


def test_an_icon_taller_than_the_plate_is_refused(ctx_factory, icon_theme):
    ctx = ctx_factory({"card": {"icon": "icon.png"}}, theme_override=icon_theme)
    ctx.rect = Rect(1.0, 1.0, 2.0, 0.6)
    with pytest.raises(LayoutError, match="the icon alone wants"):
        get_component("card")(ctx)


def test_the_shadow_the_theme_declares_is_dropped_behind_the_plate(ctx_factory, theme):
    shadowed = dataclasses.replace(theme, chart=dataclasses.replace(theme.chart, shadow=True))
    ctx = ctx_factory({"card": {"heading": "Discovery", "shadow": True}}, theme_override=shadowed)
    assert "outerShdw" in _card(ctx)._element.spPr.xml


def test_no_shadow_is_drawn_unless_the_card_asks(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery"}})
    assert "outerShdw" not in _card(ctx)._element.spPr.xml


def test_an_unknown_field_lists_the_ones_the_card_reads(ctx_factory):
    ctx = ctx_factory({"card": {"title": "Discovery"}})
    with pytest.raises(
        LayoutError, match="known fields: pair, heading, body, icon, radius, shadow"
    ):
        get_component("card")(ctx)


def test_the_card_is_recorded_for_qa_with_its_lines_on_its_own_fill(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery", "body": "one line"}})
    get_component("card")(ctx)
    record = ctx.manifest.slides[0].shapes[-1]
    pair = ctx.theme.palette.pair("surface")
    assert record.lines == ["Discovery", "one line"]
    assert (record.fg, record.bg) == (pair.fg, pair.bg)


def test_the_whole_card_arrives_as_one_reveal_group(ctx_factory):
    ctx = ctx_factory({"card": {"heading": "Discovery", "body": "one line"}})
    result = get_component("card")(ctx)
    assert result.groups == [[s.shape_id for s in ctx.slide.shapes]]


def test_the_card_is_registered():
    assert "card" in registered_components()


def test_a_bleeding_placements_shapes_are_recorded_as_bleeding(ctx_factory):
    """The intent has to survive into the manifest, or QA can only see the geometry."""
    from pptxkit.layouts.compose import _draw
    from pptxkit.spec.model import Placement

    ctx = ctx_factory({"card": {"heading": "H", "body": "B"}})
    _draw(
        ctx,
        Placement(
            at={"cols": "left-half"},
            component="card",
            body={"heading": "H", "body": "B"},
            bleed=True,
        ),
        ctx.body_rect,
        origin="s1.p1.card",
    )
    recorded = ctx.manifest.to_dict()["slides"][0]["shapes"]
    assert recorded and all(s["bleed"] for s in recorded), recorded


def test_the_flag_does_not_leak_to_the_next_placement(ctx_factory):
    """Set per placement, so the one after a bleed is not marked as one too."""
    from pptxkit.layouts.compose import _draw
    from pptxkit.spec.model import Placement

    ctx = ctx_factory({"card": {"heading": "H", "body": "B"}})
    body = {"heading": "H", "body": "B"}
    _draw(
        ctx,
        Placement(at={"cols": "left-half"}, component="card", body=body, bleed=True),
        ctx.body_rect,
        origin="s1.p1.card",
    )
    _draw(
        ctx,
        Placement(at={"cols": "left-half"}, component="card", body=body, bleed=False),
        ctx.body_rect,
        origin="s1.p2.card",
    )
    recorded = ctx.manifest.to_dict()["slides"][0]["shapes"]
    # Marked only where declared: a default-valued key is not recorded at all.
    marked = [s for s in recorded if "bleed" in s]
    assert marked and len(marked) < len(recorded)


def test_a_component_that_raises_mid_bleed_does_not_leave_the_flag_set(ctx_factory):
    """The next placement would otherwise inherit an exemption it never asked for —
    which the happy-path test above cannot see, because it clears the flag on entry."""
    from pptxkit.layouts.compose import _draw
    from pptxkit.spec.model import Placement

    ctx = ctx_factory({"card": {"heading": "H", "body": "B"}})
    tiny = Rect(0.0, 0.0, 0.2, 0.2)
    with pytest.raises(LayoutError):
        _draw(
            ctx,
            Placement(
                at={"cols": (0, 1)},
                component="card",
                body={"heading": "H", "body": "B"},
                bleed=True,
            ),
            tiny,
            origin="s1.p1.card",
        )
    assert ctx.manifest.bleeding is False
    assert ctx.manifest.origin is None


def test_plate_height_insets_by_the_same_curve_the_card_does(ctx_factory):
    """The helper sizes the plate and the component fills it; disagree about how far the
    corner cuts in and the copy runs past the bottom. A stadium is where the two must meet."""
    from pptxkit.components.card import plate_height

    ctx = ctx_factory({"card": {"heading": "H", "body": "b"}})
    # 2.5in and 24 words: narrow enough that the arc's 0.24in a side changes the line count.
    # At 4in the copy wraps identically and a radius-blind helper would pass.
    words = " ".join(["word"] * 24)
    square = plate_height(ctx, width=2.5, heading="Heading", copy=words, radius=0.0)
    stadium = plate_height(ctx, width=2.5, heading="Heading", copy=words, radius=0.5)
    assert stadium > square, f"a stadium eats measure: {square} vs {stadium}"


def test_plate_height_is_unchanged_for_the_radius_every_card_uses(ctx_factory):
    """The second pass must not move the common case: a flow sizes every step through
    this helper, and shifting them all would be a re-flow nobody asked for."""
    from pptxkit.components.card import _RADIUS_DEFAULT, plate_height

    ctx = ctx_factory({"card": {"heading": "H", "body": "b"}})
    words = " ".join(["word"] * 12)
    assert plate_height(
        ctx, width=2.4, heading="H", copy=words, radius=_RADIUS_DEFAULT
    ) == plate_height(ctx, width=2.4, heading="H", copy=words, radius=0.0)
