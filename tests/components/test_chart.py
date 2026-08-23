"""Tests for the ``chart`` body component — the consumer half of ChartSpec plus
the native renderer."""

from __future__ import annotations

import pytest
from pptx.oxml.ns import qn

import pptxkit.components  # noqa: F401 — registers the built-ins
from pptxkit.errors import LayoutError
from pptxkit.layouts.components import as_body_result, get_component, registered_components

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

_CHART = {
    "kind": "column",
    "data": [
        {"category": "Q1", "value": 12},
        {"category": "Q2", "value": 34},
        {"category": "Q3", "value": 58},
        {"category": "Q4", "value": 91, "highlight": True},
    ],
}


def _ctx(ctx_factory, *, animate=None, **overrides):
    chart = {**_CHART, **overrides}
    return ctx_factory({"chart": chart}, animate=animate)


def test_a_native_chart_is_recorded_so_its_text_is_checkable(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("chart")(ctx)
    assert [s.rendered for s in ctx.manifest.slides[0].shapes] == ["native"]


def test_the_native_path_places_a_chart_not_a_picture(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("chart")(ctx)
    assert ctx.slide.shapes[0].has_chart


def test_it_reports_the_frames_height_not_its_width(ctx_factory):
    """``height > 0`` is satisfied by the frame's width just as well, and the body rect is
    always wider than tall — so the wrong dimension passes the whole suite."""
    ctx = _ctx(ctx_factory)
    result = as_body_result(get_component("chart")(ctx))
    frame = ctx.slide.shapes[0]
    assert frame.width != frame.height  # otherwise the assertion below cannot discriminate
    assert result.height == pytest.approx(frame.height / 914400)


def test_chart_is_registered():
    assert "chart" in registered_components()


def test_an_empty_chart_names_the_slide(ctx_factory):
    ctx = ctx_factory({"chart": {}})
    with pytest.raises(LayoutError, match=r"slide 1.*'kind' must be one of"):
        get_component("chart")(ctx)


def test_one_reveal_group_holds_the_native_chart(ctx_factory):
    ctx = _ctx(ctx_factory)
    result = as_body_result(get_component("chart")(ctx))
    ids = {s.shape_id for s in ctx.slide.shapes}
    assert len(result.groups) == 1 and result.groups[0][0] in ids


def test_the_native_chart_stays_inside_the_body_rect(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("chart")(ctx)
    rect = ctx.body_rect
    frame = ctx.slide.shapes[0]
    assert frame.left / 914400 == pytest.approx(rect.left)
    assert frame.width / 914400 == pytest.approx(rect.width)


# Chart build animation — animate: by_category / by_series on a native chart.


def test_a_by_category_animate_emits_a_chart_build_on_the_real_frame(ctx_factory):
    ctx = _ctx(ctx_factory, animate="by_category")
    get_component("chart")(ctx)
    frame = ctx.slide.shapes[0]

    timing = ctx.slide._element.find(qn("p:timing"))
    assert timing is not None
    bld_graphic = timing.find(f"{{{_P}}}bldLst/{{{_P}}}bldGraphic")
    assert bld_graphic is not None
    assert bld_graphic.get("spid") == str(frame.shape_id)
    bld_chart = bld_graphic.find(f"{{{_P}}}bldSub/{{{_A}}}bldChart")
    assert bld_chart.get("bld") == "category"


def test_a_by_series_animate_maps_to_the_series_bld_value(ctx_factory):
    ctx = _ctx(ctx_factory, animate="by_series")
    get_component("chart")(ctx)
    timing = ctx.slide._element.find(qn("p:timing"))
    bld_chart = timing.find(f"{{{_P}}}bldLst/{{{_P}}}bldGraphic/{{{_P}}}bldSub/{{{_A}}}bldChart")
    assert bld_chart.get("bld") == "series"


def test_a_by_category_animate_reports_no_generic_reveal_groups(ctx_factory):
    """The chart build is already emitted directly — a second, generic click
    build over the same shape would either duplicate or conflict with it."""
    ctx = _ctx(ctx_factory, animate="by_category")
    result = as_body_result(get_component("chart")(ctx))
    assert result.groups == []


def test_a_plain_one_at_a_time_animate_on_a_native_chart_is_unaffected(ctx_factory):
    """Only by_category/by_series are chart-build animations; one_at_a_time still
    falls through to the generic single-shape click build."""
    ctx = _ctx(ctx_factory, animate="one_at_a_time")
    result = as_body_result(get_component("chart")(ctx))
    frame = ctx.slide.shapes[0]
    assert result.groups == [[frame.shape_id]]
    assert ctx.slide._element.find(qn("p:timing")) is None  # chart.py itself emits none


def test_a_category_build_staggers_one_part_per_category(ctx_factory):
    """A zero part count collapses the stagger to one click while the XML stays valid."""
    from lxml import etree

    def targets(n):
        ctx = _ctx(
            ctx_factory,
            animate="by_category",
            data=[{"category": f"Q{i}", "value": i} for i in range(1, n + 1)],
        )
        get_component("chart")(ctx)
        return etree.tostring(ctx.slide.element).decode().count("categoryIdx=")

    three, five = targets(3), targets(5)
    assert five > three, "the build does not follow the category count"
    assert five - three == (5 - 3) * (three // 4)


# --- a build the chart cannot perform -----------------------------------------

_RADAR_ROWS = [
    {"category": "Speed", "value": 8},
    {"category": "Fidelity", "value": 6},
    {"category": "Reuse", "value": 7},
]


def test_a_radar_cannot_build_by_category(ctx_factory):
    """Its categories are vertices of one outline, so the build is a click per
    category with nothing to show — observed in Keynote: six clicks, no change."""
    ctx = ctx_factory(
        {"chart": {"kind": "radar-filled", "data": _RADAR_ROWS}}, animate="by_category"
    )
    with pytest.raises(LayoutError, match="cannot build by category"):
        get_component("chart")(ctx)


def test_the_refusal_names_a_kind_that_does_build(ctx_factory):
    ctx = ctx_factory({"chart": {"kind": "radar", "data": _RADAR_ROWS}}, animate="by_category")
    with pytest.raises(LayoutError, match="column"):
        get_component("chart")(ctx)


def test_a_column_chart_still_builds_by_category(ctx_factory):
    """The exclusion is per kind, not a retreat from category builds."""
    ctx = ctx_factory({"chart": {"kind": "column", "data": _RADAR_ROWS}}, animate="by_category")
    get_component("chart")(ctx)
    assert ctx.manifest.to_dict()["slides"][0]["animations"]


def test_a_radar_may_still_animate_together(ctx_factory):
    ctx = ctx_factory({"chart": {"kind": "radar-filled", "data": _RADAR_ROWS}}, animate="together")
    get_component("chart")(ctx)


# --- room for the labels, stated in the file ----------------------------------


def _plot_x(ctx):
    import re

    frame = next(s for s in ctx.slide.shapes if s.has_chart)
    m = re.search(r'<c:manualLayout>.*?<c:x val="([\d.]+)"/>', frame.chart._chartSpace.xml, re.S)
    return float(m.group(1)) if m else None


_LONG = [
    {"category": "Unknown field on a slide", "value": 34},
    {"category": "Placement outside the band", "value": 27},
]
_SHORT = [{"category": "Q1", "value": 34}, {"category": "Q2", "value": 27}]


def test_a_bar_chart_reserves_a_column_for_its_category_labels(ctx_factory):
    """Nothing else states the reservation, so each renderer decides for itself:
    LibreOffice shrinks the plot, Keynote runs the labels off the slide."""
    ctx = ctx_factory({"chart": {"kind": "bar", "data": _LONG}})
    get_component("chart")(ctx)
    assert _plot_x(ctx) is not None


def test_the_reservation_grows_with_the_longest_label(ctx_factory):
    long_ctx = ctx_factory({"chart": {"kind": "bar", "data": _LONG}})
    get_component("chart")(long_ctx)
    short_ctx = ctx_factory({"chart": {"kind": "bar", "data": _SHORT}})
    get_component("chart")(short_ctx)
    assert _plot_x(long_ctx) > _plot_x(short_ctx)


def test_the_reservation_never_swallows_the_plot(ctx_factory):
    """A category name longer than the frame must not leave the bars no room."""
    ctx = ctx_factory({"chart": {"kind": "bar", "data": [{"category": "x" * 400, "value": 3}]}})
    get_component("chart")(ctx)
    assert _plot_x(ctx) <= 0.56


def test_a_column_chart_reserves_nothing(ctx_factory):
    """Its labels sit under the plot, in width they already have."""
    ctx = ctx_factory({"chart": {"kind": "column", "data": _LONG}})
    get_component("chart")(ctx)
    assert _plot_x(ctx) is None
