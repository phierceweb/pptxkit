import pytest
from pptx.util import Inches

import pptxkit.components  # noqa: F401
from pptxkit.errors import LayoutError
from pptxkit.layouts.components import get_component

ITEMS = [
    {"value": "20", "label": "integrations"},
    {"value": "200+", "label": "database tables"},
    {"value": "260+", "label": "data models"},
]


def _texts(slide):
    return "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def _ctx(ctx_factory, **body):
    return ctx_factory({"stats": {"items": ITEMS, **body}})


def test_each_tile_shows_its_value_and_label(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("stats")(ctx)
    text = _texts(ctx.slide)
    assert "200+" in text and "database tables" in text


def test_recorded_lines_match_the_drawn_paragraphs(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("stats")(ctx)
    by_id = {s.shape_id: s for s in ctx.slide.shapes}
    recorded = [r for r in ctx.manifest.slides[0].shapes if r.lines]
    # Value and label are separate lines in the manifest because qa's overflow
    # check looks for each of them in the rendered text.
    assert [r.lines for r in recorded] == [[i["value"], i["label"]] for i in ITEMS]
    for record in recorded:
        drawn = [p.text for p in by_id[record.shape_id].text_frame.paragraphs]
        assert drawn == record.lines


def test_one_reveal_group_per_tile(ctx_factory):
    ctx = _ctx(ctx_factory)
    assert len(get_component("stats")(ctx).groups) == len(ITEMS)


def test_tiles_are_laid_out_left_to_right_without_overlap(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("stats")(ctx)
    tiles = sorted((s.left, s.width) for s in ctx.slide.shapes if s.has_text_frame)
    for (left_a, width_a), (left_b, _) in zip(tiles, tiles[1:], strict=False):
        assert left_a + width_a <= left_b + 1


def test_a_gutter_separates_the_tiles_in_a_row(ctx_factory):
    """The horizontal pitch carries the grid's gutter, not just the tile width — the overlap
    test above is satisfied by tiles that merely touch."""
    ctx = _ctx(ctx_factory)
    get_component("stats")(ctx)
    tiles = sorted((s.left, s.width) for s in ctx.slide.shapes if s.has_text_frame)
    gaps = [nxt - (left + width) for (left, width), (nxt, _) in zip(tiles, tiles[1:], strict=False)]
    assert len(gaps) == len(ITEMS) - 1
    for gap in gaps:
        assert gap == pytest.approx(Inches(ctx.grid.gutter), abs=Inches(0.01))


def test_an_optional_caption_is_rendered(ctx_factory):
    ctx = _ctx(ctx_factory, caption="Six dimensions of scale.")
    get_component("stats")(ctx)
    assert "Six dimensions of scale." in _texts(ctx.slide)


def test_the_caption_joins_the_last_reveal_group(ctx_factory):
    plain = len(get_component("stats")(_ctx(ctx_factory)).groups[-1])
    with_caption = len(get_component("stats")(_ctx(ctx_factory, caption="x")).groups[-1])
    assert with_caption == plain + 1


def test_items_is_required(ctx_factory):
    ctx = ctx_factory({"stats": {}})
    with pytest.raises(LayoutError, match=r"slide 1 .*'items'"):
        get_component("stats")(ctx)


def test_an_item_without_a_value_names_the_slide(ctx_factory):
    ctx = ctx_factory({"stats": {"items": [{"label": "no value"}]}})
    with pytest.raises(LayoutError, match=r"slide 1 .*item 1.*'value'"):
        get_component("stats")(ctx)


def test_content_stays_inside_the_body_rect(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("stats")(ctx)
    rect = ctx.body_rect
    for shape in ctx.slide.shapes:
        assert shape.left / 914400 >= rect.left - 0.01
        assert (shape.left + shape.width) / 914400 <= rect.right + 0.01
        assert shape.top / 914400 >= rect.top - 0.01
        assert (shape.top + shape.height) / 914400 <= rect.bottom + 0.01


def test_every_returned_id_is_a_real_shape(ctx_factory):
    ctx = _ctx(ctx_factory)
    groups = get_component("stats")(ctx).groups
    ids = {s.shape_id for s in ctx.slide.shapes}
    assert all(spid in ids for group in groups for spid in group)


def test_the_reported_height_is_the_extent_of_the_tiles_it_drew(ctx_factory):
    """``0 < h <= body_rect.height`` was true by construction: the tile height derives
    from the body rect, so returning ``rect.height`` outright passed it."""
    ctx = _ctx(ctx_factory)
    result = get_component("stats")(ctx)
    bottom = max((s.top + s.height) / 914400 for s in ctx.slide.shapes)
    assert result.height == pytest.approx(bottom - ctx.body_rect.top, abs=0.01)
    assert result.height < ctx.body_rect.height  # else `rect.height` would also pass


def test_items_wrap_onto_a_second_row(ctx_factory):
    items = [{"value": str(i), "label": f"m{i}"} for i in range(6)]
    ctx = ctx_factory({"stats": {"items": items}})
    get_component("stats")(ctx)
    tops = {s.top for s in ctx.slide.shapes}
    assert len(tops) == 2


def test_a_wrapped_layout_still_fits_the_body_rect(ctx_factory):
    items = [{"value": str(i), "label": f"m{i}"} for i in range(6)]
    ctx = ctx_factory({"stats": {"items": items, "caption": "c"}})
    get_component("stats")(ctx)
    rect = ctx.body_rect
    for shape in ctx.slide.shapes:
        assert (shape.top + shape.height) / 914400 <= rect.bottom + 0.01


def test_too_many_rows_for_the_body_rect_is_rejected(ctx_factory):
    items = [{"value": str(i), "label": f"m{i}"} for i in range(13)]
    ctx = ctx_factory({"stats": {"items": items}})
    with pytest.raises(LayoutError, match="of height but the body rect is"):
        get_component("stats")(ctx)


def test_a_non_numeric_columns_raises_a_layout_error(ctx_factory):
    ctx = _ctx(ctx_factory, columns="two")
    with pytest.raises(LayoutError, match=r"slide 1 .*'columns'.*got 'two'"):
        get_component("stats")(ctx)


def test_stats_is_registered():
    from pptxkit.layouts.components import registered_components

    assert "stats" in registered_components()


def test_a_tile_records_the_fill_its_text_actually_sits_on(ctx_factory):
    """The manifest is what the contrast check reads. Recording the page colour where
    the text sits on a tile is the intent-versus-reality bug that hid white-on-white."""
    ctx = ctx_factory({"stats": {"items": [{"value": "42", "label": "things"}]}})
    get_component("stats")(ctx)
    tile = next(s for s in ctx.slide.shapes if s.has_text_frame)
    recorded = ctx.manifest.slides[0].shapes[0]
    assert recorded.bg == str(tile.fill.fore_color.rgb)
    assert recorded.bg != ctx.pair.bg, "the tile is not the page; recording page hides it"


def test_a_caption_is_counted_in_the_height_the_tiles_need(ctx_factory):
    """Dropping it from extent lets a caption overflow the body without a word."""
    items = [{"value": str(n), "label": "x"} for n in range(1, 4)]
    tall = get_component("stats")(ctx_factory({"stats": {"items": items}})).height
    with_caption = get_component("stats")(
        ctx_factory({"stats": {"items": items, "caption": "a note"}})
    ).height
    assert with_caption > tall


def test_one_column_puts_every_tile_at_the_same_left_edge(ctx_factory):
    """A floor of two would split these into a row, which only the x tells you."""
    ctx = ctx_factory(
        {
            "stats": {
                "columns": 1,
                "items": [{"value": "1", "label": "a"}, {"value": "2", "label": "b"}],
            }
        }
    )
    get_component("stats")(ctx)
    lefts = {s.left for s in ctx.slide.shapes if s.has_text_frame}
    assert len(lefts) == 1


# --- a mark above the number -------------------------------------------------

MARKED = [
    {"value": "42", "label": "first", "icon": "users"},
    {"value": "68", "label": "second", "icon": "globe"},
]
PLAIN = [{"value": "42", "label": "first"}, {"value": "68", "label": "second"}]


def _tiles(ctx):
    return [s for s in ctx.slide.shapes if "roundRect" in s._element.xml]


def test_a_marked_tile_is_taller_than_a_plain_one(ctx_factory):
    """Sizing the tile by its type alone crops either the number or the mark."""
    plain = ctx_factory({"stats": {"items": PLAIN}})
    get_component("stats")(plain)
    marked = ctx_factory({"stats": {"items": MARKED}})
    get_component("stats")(marked)
    assert _tiles(marked)[0].height > _tiles(plain)[0].height


def test_the_number_is_pushed_below_the_mark(ctx_factory):
    """Same frame, so the text's top margin has to clear the glyph drawn over it."""
    plain = ctx_factory({"stats": {"items": PLAIN}})
    get_component("stats")(plain)
    marked = ctx_factory({"stats": {"items": MARKED}})
    get_component("stats")(marked)
    assert _tiles(marked)[0].text_frame.margin_top > _tiles(plain)[0].text_frame.margin_top


def test_each_tile_carries_its_own_glyph(ctx_factory):
    ctx = ctx_factory({"stats": {"items": MARKED}})
    get_component("stats")(ctx)
    assert sum("<a:custGeom" in s._element.xml for s in ctx.slide.shapes) == 2


def test_the_glyph_reveals_with_its_tile(ctx_factory):
    ctx = ctx_factory({"stats": {"items": MARKED}})
    result = get_component("stats")(ctx)
    assert all(len(group) == 2 for group in result.groups)


def test_a_mark_is_measured_against_the_tile_not_the_slide(ctx_factory):
    """The tile is its own surface; a glyph checked against the page picks the wrong ink."""
    import json

    ctx = ctx_factory({"stats": {"items": MARKED}})
    get_component("stats")(ctx)
    recorded = json.loads(json.dumps(ctx.manifest.to_dict()))["slides"][0]["shapes"]
    fill = ctx.theme.palette.role("line")
    marks = [s for s in recorded if s.get("fg") and s.get("bg") == fill and not s.get("lines")]
    assert marks, [s.get("bg") for s in recorded]
