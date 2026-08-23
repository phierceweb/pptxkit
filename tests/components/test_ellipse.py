import dataclasses

import pytest
from pptx.enum.dml import MSO_FILL
from pptx.enum.text import PP_ALIGN

import pptxkit.components  # noqa: F401 — registers the built-in components
from pptxkit.components.ellipse import _INSCRIBED
from pptxkit.errors import LayoutError, ThemeError
from pptxkit.layouts.components import get_component, registered_components
from pptxkit.theme.model import Rect
from pptxkit.utils.text import LINE_HEIGHT, wrapped_lines

EMU = 914400


def _disc(ctx):
    get_component("ellipse")(ctx)
    return ctx.slide.shapes[0]


def test_the_disc_is_a_circle_of_the_placements_short_side(ctx_factory):
    ctx = ctx_factory({"ellipse": {}})
    shape = _disc(ctx)
    short = min(ctx.body_rect.width, ctx.body_rect.height)
    assert shape.width == shape.height
    assert shape.width / EMU == pytest.approx(short)


def test_size_scales_the_diameter_off_the_placement(ctx_factory):
    ctx = ctx_factory({"ellipse": {"size": 0.25}})
    short = min(ctx.body_rect.width, ctx.body_rect.height)
    assert _disc(ctx).width / EMU == pytest.approx(short * 0.25)


def test_a_disc_smaller_than_its_placement_sits_where_align_and_anchor_say(ctx_factory):
    ctx = ctx_factory({"ellipse": {"size": 0.2}})
    ctx.align, ctx.anchor = "right", "bottom"
    shape = _disc(ctx)
    rect = ctx.body_rect
    assert shape.left / EMU == pytest.approx(rect.right - shape.width / EMU)
    assert shape.top / EMU == pytest.approx(rect.bottom - shape.height / EMU)


def test_the_default_disc_sits_at_the_top_left_of_its_placement(ctx_factory):
    ctx = ctx_factory({"ellipse": {"size": 0.2}})
    shape = _disc(ctx)
    assert shape.left / EMU == pytest.approx(ctx.body_rect.left)
    assert shape.top / EMU == pytest.approx(ctx.body_rect.top)


def test_the_disc_is_filled_with_the_background_of_the_pair_it_names(ctx_factory):
    ctx = ctx_factory({"ellipse": {"pair": "accent-3"}})
    assert str(_disc(ctx).fill.fore_color.rgb) == ctx.theme.palette.pair("accent-3").bg


def test_the_default_pair_is_the_first_accent(ctx_factory):
    ctx = ctx_factory({"ellipse": {}})
    assert str(_disc(ctx).fill.fore_color.rgb) == ctx.theme.palette.pair("accent-1").bg


def test_a_fill_too_close_to_the_paper_to_see_is_given_an_edge(ctx_factory):
    ctx = ctx_factory({"ellipse": {"pair": "surface"}})
    shape = _disc(ctx)
    assert shape.line.fill.type == MSO_FILL.SOLID
    assert str(shape.line.color.rgb) == ctx.theme.palette.role("line")


def test_the_label_is_centred_in_the_disc_in_the_pairs_own_ink(ctx_factory):
    ctx = ctx_factory({"ellipse": {"label": "3", "pair": "accent-3"}})
    run = _disc(ctx).text_frame.paragraphs[0].runs[0]
    assert run.text == "3"
    assert str(run.font.color.rgb) == ctx.theme.palette.pair("accent-3").fg
    assert _disc(ctx).text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER


def test_the_label_is_set_at_the_rung_it_names(ctx_factory):
    ctx = ctx_factory({"ellipse": {"label": "7", "rung": "stat"}})
    run = _disc(ctx).text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == pytest.approx(ctx.theme.style("stat").size)


def test_a_label_wider_than_the_disc_is_refused_naming_the_measure(ctx_factory):
    ctx = ctx_factory({"ellipse": {"label": "a step in the process", "size": 0.05}})
    with pytest.raises(LayoutError, match="needs more than the .* across the middle"):
        get_component("ellipse")(ctx)


def test_a_label_that_fits_across_the_disc_but_not_down_it_is_refused(ctx_factory):
    """The label's line height binds before its measure does on any short label."""
    ctx = ctx_factory({"ellipse": {"label": "1"}})
    size_pt = ctx.theme.style("caption").size
    ctx.rect = Rect(1.0, 1.0, 0.15 / _INSCRIBED, 0.15 / _INSCRIBED)
    assert wrapped_lines("1", width_in=0.15, size_pt=size_pt) == 1
    assert size_pt * LINE_HEIGHT / 72 > 0.15
    with pytest.raises(LayoutError, match="needs more than"):
        get_component("ellipse")(ctx)


def test_a_bare_disc_carries_no_text_frame_content(ctx_factory):
    ctx = ctx_factory({"ellipse": {}})
    assert _disc(ctx).text_frame.text == ""


def test_the_shadow_the_theme_declares_is_dropped_behind_the_disc(ctx_factory, theme):
    shadowed = dataclasses.replace(theme, chart=dataclasses.replace(theme.chart, shadow=True))
    ctx = ctx_factory({"ellipse": {"shadow": True}}, theme_override=shadowed)
    xml = _disc(ctx)._element.spPr.xml
    assert "outerShdw" in xml
    assert f'blurRad="{round(theme.chart.shadow_blur_pt * 12700)}"' in xml


def test_no_shadow_is_drawn_unless_the_component_asks(ctx_factory):
    ctx = ctx_factory({"ellipse": {}})
    assert "outerShdw" not in _disc(ctx)._element.spPr.xml


def test_a_non_boolean_shadow_is_refused(ctx_factory):
    ctx = ctx_factory({"ellipse": {"shadow": "yes"}})
    with pytest.raises(LayoutError, match="'shadow' must be true or false"):
        get_component("ellipse")(ctx)


def test_a_size_outside_the_fraction_range_is_refused(ctx_factory):
    ctx = ctx_factory({"ellipse": {"size": 1.5}})
    with pytest.raises(LayoutError, match="'size' is the diameter as a fraction"):
        get_component("ellipse")(ctx)


def test_a_non_numeric_size_is_refused(ctx_factory):
    ctx = ctx_factory({"ellipse": {"size": "big"}})
    with pytest.raises(LayoutError, match="'size' must be a number"):
        get_component("ellipse")(ctx)


def test_an_undeclared_pair_is_refused_by_name(ctx_factory):
    ctx = ctx_factory({"ellipse": {"pair": "accent-9"}})
    with pytest.raises(ThemeError, match="no colour pair 'accent-9'"):
        get_component("ellipse")(ctx)


def test_an_unknown_field_lists_the_ones_the_disc_reads(ctx_factory):
    ctx = ctx_factory({"ellipse": {"colour": "red"}})
    with pytest.raises(LayoutError, match="known fields: pair, label, rung, size, shadow"):
        get_component("ellipse")(ctx)


def test_the_disc_is_recorded_for_qa_with_its_label_and_its_own_fill(ctx_factory):
    ctx = ctx_factory({"ellipse": {"label": "9", "pair": "accent-3"}})
    get_component("ellipse")(ctx)
    record = ctx.manifest.slides[0].shapes[-1]
    pair = ctx.theme.palette.pair("accent-3")
    assert (record.text, record.fg, record.bg) == ("9", pair.fg, pair.bg)


def test_the_reveal_group_is_the_disc_itself(ctx_factory):
    ctx = ctx_factory({"ellipse": {}})
    result = get_component("ellipse")(ctx)
    assert result.groups == [[ctx.slide.shapes[0].shape_id]]


def test_the_label_takes_the_discs_own_ink_not_the_slides(ctx_factory):
    """A white-filled badge on a dark slide needs dark type, which ctx.fg() is not."""
    ctx = ctx_factory({"ellipse": {"label": "1", "pair": "page"}}, background="inverse")
    ink = str(_disc(ctx).text_frame.paragraphs[0].runs[0].font.color.rgb)
    assert ink == ctx.theme.palette.pair("page").fg
    assert ink != str(ctx.fg())


def test_a_disc_filled_in_the_slides_own_paper_is_given_an_edge(ctx_factory):
    ctx = ctx_factory({"ellipse": {"pair": "inverse"}}, background="inverse")
    shape = _disc(ctx)
    assert shape.line.fill.type == MSO_FILL.SOLID
    assert str(shape.line.color.rgb) == ctx.theme.palette.role("line")


def test_the_disc_is_registered():
    assert "ellipse" in registered_components()
