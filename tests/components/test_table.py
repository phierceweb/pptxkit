import pytest
from pptx.enum.dml import MSO_FILL
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

import pptxkit.components  # noqa: F401
from pptxkit.errors import LayoutError
from pptxkit.layouts.components import get_component
from pptxkit.qa.geometry import check_contrast
from pptxkit.utils.color import contrast_ratio
from pptxkit.theme.model import Rect

HEADER = ["Item", "Where", "Count"]
ROWS = [["The first thing", "On the left", "12"], ["The second thing", "In the middle", "34"]]


def _ctx(ctx_factory, **body):
    return ctx_factory({"table": {"header": HEADER, "rows": ROWS, **body}})


def _table(ctx):
    return next(s.table for s in ctx.slide.shapes if s.has_table)


# The manifest's box granularity. Half of it is not enough: 6.6715 rounds to
# 6.671 (half-to-even), leaving a gap of exactly 0.0005 plus float noise.
_BOX_TOLERANCE = 0.001


def _cells(ctx):
    return [r for r in ctx.manifest.slides[0].shapes if r.lines]


# --- what a successful build never reaches: validation ---------------------------


def test_rows_is_required(ctx_factory):
    ctx = ctx_factory({"table": {"header": HEADER}})
    with pytest.raises(LayoutError, match=r"slide 1 .*'rows' must be a non-empty list"):
        get_component("table")(ctx)


def test_a_row_that_is_not_a_list_names_its_position(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [["a", "b"], "c, d"]}})
    with pytest.raises(LayoutError, match=r"row 2 must be a non-empty list of cells"):
        get_component("table")(ctx)


def test_a_row_shorter_than_the_header_is_refused(ctx_factory):
    ctx = ctx_factory({"table": {"header": HEADER, "rows": [["only", "two"]]}})
    with pytest.raises(LayoutError, match=r"row 1 covers 2 column\(s\) but the table is 3 wide"):
        get_component("table")(ctx)


def test_rows_of_differing_widths_are_refused_without_a_header(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [["a", "b"], ["c"]]}})
    with pytest.raises(LayoutError, match=r"row 2 covers 1 column\(s\) but the table is 2 wide"):
        get_component("table")(ctx)


def test_align_must_have_one_entry_per_column(ctx_factory):
    ctx = _ctx(ctx_factory, align=["left", "right"])
    with pytest.raises(LayoutError, match=r"'align' is one entry per column, so it needs 3"):
        get_component("table")(ctx)


def test_an_unknown_alignment_lists_the_vocabulary(ctx_factory):
    ctx = _ctx(ctx_factory, align=["left", "middle", "right"])
    with pytest.raises(LayoutError, match=r"'align' entries must be one of left, center, right"):
        get_component("table")(ctx)


def test_widths_must_have_one_entry_per_column(ctx_factory):
    ctx = _ctx(ctx_factory, widths=[1, 1])
    with pytest.raises(LayoutError, match=r"'widths' is one entry per column, so it needs 3"):
        get_component("table")(ctx)


def test_a_zero_width_column_is_refused(ctx_factory):
    ctx = _ctx(ctx_factory, widths=[1, 0, 1])
    with pytest.raises(LayoutError, match=r"'widths' entries must be positive numbers, got 0"):
        get_component("table")(ctx)


def test_a_non_numeric_width_is_refused(ctx_factory):
    ctx = _ctx(ctx_factory, widths=[1, "wide", 1])
    with pytest.raises(LayoutError, match=r"'widths' entries must be positive numbers, got 'wide'"):
        get_component("table")(ctx)


def test_an_unknown_field_names_the_component(ctx_factory):
    ctx = _ctx(ctx_factory, caption="not a table field")
    with pytest.raises(LayoutError, match=r"component 'table'.*unknown field 'caption'"):
        get_component("table")(ctx)


def test_more_rows_than_the_placement_can_hold_is_refused(ctx_factory):
    rows = [[f"row {n}", "b", "c"] for n in range(40)]
    ctx = ctx_factory({"table": {"header": HEADER, "rows": rows}})
    with pytest.raises(LayoutError, match=r"41 rows need .*in of height but the placement"):
        get_component("table")(ctx)


def test_columns_narrower_than_the_cell_padding_are_refused(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [["a"] * 60]}})
    with pytest.raises(LayoutError, match=r"60 columns leave .*less than the .*gutter pads"):
        get_component("table")(ctx)


# --- wiring the built output hides ------------------------------------------------


def test_the_banded_office_table_style_is_replaced_by_no_style_no_grid(ctx_factory):
    """python-pptx stamps Medium Style 2 Accent 1 on every new table, and PowerPoint
    would repaint every fill and rule this component sets with it."""
    ctx = _ctx(ctx_factory)
    get_component("table")(ctx)
    tbl = _table(ctx)._tbl
    style = tbl.find(qn("a:tblPr")).find(qn("a:tableStyleId"))
    assert style.text == "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    assert tbl.find(qn("a:tblPr")).get("firstRow") in (None, "0")
    assert tbl.find(qn("a:tblPr")).get("bandRow") in (None, "0")


def test_a_rule_is_drawn_under_every_row_but_the_last(ctx_factory):
    """Also pins the `rules:` default. Without the `lnR` line, moving that default to
    `grid` gives every table a column rule it never had and nothing goes red."""
    ctx = _ctx(ctx_factory)
    get_component("table")(ctx)
    rules = _table(ctx)._tbl.findall(f".//{qn('a:lnB')}")
    # Header plus two body rows; the last row's underside is the table's own edge.
    assert len(rules) == 3 * 2
    assert {r.find(f".//{qn('a:srgbClr')}").get("val") for r in rules} == {"573C65"}
    assert _table(ctx)._tbl.findall(f".//{qn('a:lnR')}") == []


def test_the_header_sits_on_the_surface_pair_and_the_body_on_nothing(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("table")(ctx)
    tbl = _table(ctx)
    assert str(tbl.cell(0, 0).fill.fore_color.rgb) == "F5F6F8"
    assert tbl.cell(1, 0).fill.type == MSO_FILL.BACKGROUND


def test_a_named_body_pair_fills_the_body_cells(ctx_factory):
    ctx = _ctx(ctx_factory, body_pair="inverse")
    get_component("table")(ctx)
    assert str(_table(ctx).cell(1, 0).fill.fore_color.rgb) == "2D0937"


def test_each_column_takes_its_own_alignment(ctx_factory):
    ctx = _ctx(ctx_factory, align=["left", "center", "right"])
    get_component("table")(ctx)
    tbl = _table(ctx)
    drawn = [tbl.cell(1, c).text_frame.paragraphs[0].alignment for c in range(3)]
    assert drawn == [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT]


def test_columns_with_no_align_follow_the_placements_own(ctx_factory):
    ctx = _ctx(ctx_factory)
    ctx.align = "right"
    get_component("table")(ctx)
    tbl = _table(ctx)
    assert [tbl.cell(1, c).text_frame.paragraphs[0].alignment for c in range(3)] == [
        PP_ALIGN.RIGHT
    ] * 3


def test_widths_are_relative_weights_of_the_placement(ctx_factory):
    ctx = _ctx(ctx_factory, widths=[2, 1, 1])
    get_component("table")(ctx)
    widths = [c.width / 914400 for c in _table(ctx).columns]
    total = ctx.body_rect.width
    assert widths == pytest.approx([total / 2, total / 4, total / 4])


# --- the manifest QA reads --------------------------------------------------------


def test_every_cell_is_recorded_with_its_own_box(ctx_factory):
    """A table recorded as one rectangle is a table whose clipped column QA cannot see."""
    ctx = _ctx(ctx_factory)
    get_component("table")(ctx)
    cells = _cells(ctx)
    assert [c.lines for c in cells] == [[t] for row in [HEADER, *ROWS] for t in row]
    assert len({tuple(c.box) for c in cells}) == 9


def test_a_recorded_cell_box_matches_the_column_it_was_drawn_in(ctx_factory):
    ctx = _ctx(ctx_factory, widths=[2, 1, 1])
    get_component("table")(ctx)
    rect = ctx.body_rect
    lefts = [c.box.x for c in _cells(ctx)[:3]]
    # A box is recorded to 0.001in; it meets unrounded geometry within that.
    assert lefts == pytest.approx(
        [rect.left, rect.left + rect.width / 2, rect.left + rect.width * 3 / 4], abs=_BOX_TOLERANCE
    )


def test_a_cell_records_the_colour_really_behind_it_not_the_page(ctx_factory):
    """Recording the page where a cell sits on a panel is the intent-versus-reality
    bug that hid white-on-white."""
    ctx = _ctx(ctx_factory)
    ctx.panels.append((Rect(0.0, 0.0, 13.333, 7.5), "2D0937"))
    get_component("table")(ctx)
    body = [c for c in _cells(ctx) if c.lines[0] in ("12", "34")]
    assert [c.bg for c in body] == ["2D0937", "2D0937"]


def test_qa_sees_a_body_cell_lost_on_what_is_painted_behind_it(ctx_factory, theme):
    ctx = _ctx(ctx_factory)
    ctx.panels.append((Rect(0.0, 0.0, 13.333, 7.5), "2D0937"))
    get_component("table")(ctx)
    findings = check_contrast(ctx.manifest.to_dict(), theme)
    assert [f.shape for f in findings] == [f"Table 1 r{r}c{c}" for r in (2, 3) for c in (1, 2, 3)]


def test_an_empty_cell_records_its_box_but_claims_no_ink(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [["filled", ""]]}})
    get_component("table")(ctx)
    blank = ctx.manifest.slides[0].shapes[-1]
    assert blank.lines == [] and blank.fg is None and blank.font_pt is None
    assert blank.box.w == pytest.approx(ctx.body_rect.width / 2, abs=_BOX_TOLERANCE)


# --- geometry ---------------------------------------------------------------------


def test_a_wrapping_cell_makes_its_row_taller(ctx_factory):
    short = ctx_factory({"table": {"rows": [["a", "b"]]}})
    long = ctx_factory({"table": {"rows": [["a", " ".join(["word"] * 40)]]}})
    assert get_component("table")(long).height > get_component("table")(short).height


def test_the_table_starts_at_the_placement_and_reports_the_height_it_used(ctx_factory):
    ctx = _ctx(ctx_factory)
    result = get_component("table")(ctx)
    frame = next(s for s in ctx.slide.shapes if s.has_table)
    assert frame.left / 914400 == pytest.approx(ctx.body_rect.left)
    assert frame.top / 914400 == pytest.approx(ctx.body_rect.top)
    assert result.height == pytest.approx(frame.height / 914400, abs=0.001)
    assert result.height < ctx.body_rect.height  # else `rect.height` would also pass


def test_the_whole_table_is_one_reveal_group(ctx_factory):
    ctx = _ctx(ctx_factory)
    groups = get_component("table")(ctx).groups
    ids = {s.shape_id for s in ctx.slide.shapes}
    assert groups == [[next(iter(ids))]]


def test_table_is_registered():
    from pptxkit.layouts.components import registered_components

    assert "table" in registered_components()


# --- cells that span, rows that total, rows that band -------------------------


def _rows(ctx):
    """Recorded cells as {(row, col): (text, box)}. ``box`` is the manifest's own mapping —
    ``["w"]``, ``["h"]``, never an index."""
    out = {}
    for sh in ctx.manifest.to_dict()["slides"][0]["shapes"]:
        name = sh.get("name", "")
        if " r" not in name:
            continue
        tag = name.split(" r")[-1]
        row, col = (int(n) for n in tag.split("c"))
        out[(row, col)] = ((sh.get("lines") or [""])[0], sh["box"])
    return out


def test_a_cell_may_span_its_neighbours(ctx_factory):
    """Merged in the file, not merely recorded as merged: the manifest is written from
    the spans, so recording alone would agree with a table whose cells never joined."""
    ctx = ctx_factory(
        {
            "table": {
                "header": [{"text": "Effort", "across": 2}, "Owner"],
                "rows": [["12", "8", "Compiler"]],
            }
        }
    )
    get_component("table")(ctx)
    frame = next(s for s in ctx.slide.shapes if s.has_table)
    assert 'gridSpan="2"' in frame._element.xml, "the cells were never merged"
    cells = _rows(ctx)
    assert cells[(1, 1)][0] == "Effort"
    assert (1, 2) not in cells
    assert cells[(1, 3)][0] == "Owner"


def test_a_spanned_cell_records_the_box_it_really_covers(ctx_factory):
    """Recording only the first column would have QA measure a third of the cell."""
    ctx = ctx_factory(
        {
            "table": {
                "header": [{"text": "Effort", "across": 2}, "Owner"],
                "rows": [["12", "8", "Compiler"]],
            }
        }
    )
    get_component("table")(ctx)
    cells = _rows(ctx)
    spanned = cells[(1, 1)][1]["w"]
    single = cells[(2, 1)][1]["w"]
    assert spanned == pytest.approx(single * 2, rel=0.02)


def test_a_row_whose_spans_do_not_reach_the_width_is_refused(ctx_factory):
    ctx = ctx_factory(
        {"table": {"header": ["A", "B", "C"], "rows": [[{"text": "wide", "across": 2}]]}}
    )
    with pytest.raises(LayoutError, match=r"row 1 covers 2 column\(s\) but the table is 3"):
        get_component("table")(ctx)


def test_a_span_counts_toward_the_header_width(ctx_factory):
    """Two header cells spanning 2 and 1 make a three-column table, not a two."""
    ctx = ctx_factory(
        {"table": {"header": [{"text": "Pair", "across": 2}, "Solo"], "rows": [["a", "b", "c"]]}}
    )
    get_component("table")(ctx)
    assert len(_rows(ctx)) >= 5


def test_a_total_row_is_ruled_above_and_set_apart(ctx_factory):
    ctx = ctx_factory(
        {
            "table": {
                "header": ["Item", "Cost"],
                "rows": [["One", "12"], ["Two", "30"]],
                "total": ["Total", "42"],
            }
        }
    )
    get_component("table")(ctx)
    frame = next(s for s in ctx.slide.shapes if s.has_table)
    xml = frame._element.xml
    assert "<a:lnT" in xml, "the total carries no rule above it"
    cells = _rows(ctx)
    assert cells[(4, 1)][0] == "Total"


def test_a_table_with_no_total_rules_no_top_edge(ctx_factory):
    ctx = ctx_factory({"table": {"header": ["Item", "Cost"], "rows": [["One", "12"]]}})
    get_component("table")(ctx)
    frame = next(s for s in ctx.slide.shapes if s.has_table)
    assert "<a:lnT" not in frame._element.xml


def test_banding_fills_alternate_body_rows_only(ctx_factory):
    """The header keeps its own pair; banding is a body reading aid. Asserted on the fills,
    since banding cannot change how many rows a table has."""
    ctx = ctx_factory(
        {
            "table": {
                "banding": True,
                "header": ["A", "B"],
                "rows": [["1", "2"], ["3", "4"], ["5", "6"]],
            }
        }
    )
    get_component("table")(ctx)
    fills = [s[1] for (r, c), s in sorted(_rows(ctx).items()) if c == 1]
    assert len(fills) == 4
    head, first, second, third = (
        next(
            sh["bg"]
            for sh in ctx.manifest.to_dict()["slides"][0]["shapes"]
            if sh.get("name", "").endswith(f"r{n}c1")
        )
        for n in (1, 2, 3, 4)
    )
    assert head != first, "the header band must not be the body's fill"
    assert first != second, "alternate body rows must differ"
    assert first == third, "and the third body row returns to the first's fill"


def test_without_banding_every_body_row_shares_a_surface(ctx_factory):
    ctx = ctx_factory({"table": {"header": ["A", "B"], "rows": [["1", "2"], ["3", "4"]]}})
    get_component("table")(ctx)
    shapes = {sh["name"]: sh for sh in ctx.manifest.to_dict()["slides"][0]["shapes"]}
    first = next(s for n, s in shapes.items() if n.endswith("r2c1"))
    second = next(s for n, s in shapes.items() if n.endswith("r3c1"))
    assert first["bg"] == second["bg"]


def test_a_cell_may_take_a_pair_of_its_own(ctx_factory):
    ctx = ctx_factory(
        {
            "table": {
                "header": ["A", "Verdict"],
                "rows": [["1", {"text": "Yes", "pair": "accent-1"}]],
            }
        }
    )
    get_component("table")(ctx)
    shapes = {sh["name"]: sh for sh in ctx.manifest.to_dict()["slides"][0]["shapes"]}
    marked = next(s for n, s in shapes.items() if n.endswith("r2c2"))
    assert marked["bg"] == ctx.theme.palette.pair("accent-1").bg


def test_a_cell_may_align_unlike_its_column(ctx_factory):
    """Per-column align is the rule; a cell overriding it is how a total label sits left
    in a right-aligned column."""
    ctx = ctx_factory(
        {
            "table": {
                "align": ["right", "right"],
                "rows": [[{"text": "Total", "align": "left"}, "42"]],
            }
        }
    )
    get_component("table")(ctx)
    frame = next(s for s in ctx.slide.shapes if s.has_table)
    first = frame.table.cell(0, 0).text_frame.paragraphs[0]
    second = frame.table.cell(0, 1).text_frame.paragraphs[0]
    assert first.alignment != second.alignment


def test_an_unknown_cell_key_is_refused(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [[{"text": "x", "colour": "red"}]]}})
    with pytest.raises(LayoutError, match="has no key 'colour'"):
        get_component("table")(ctx)


def test_a_span_of_zero_is_refused(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [[{"text": "x", "across": 0}]]}})
    with pytest.raises(LayoutError, match="sets 'across' to 0"):
        get_component("table")(ctx)


def test_merging_reclaims_the_padding_between_the_columns_it_swallowed():
    """Two merged columns hold more text than their two measures added up. On the arithmetic:
    whether the extra measure changes the line count depends on where the text falls."""
    from pptxkit.components._tablegeom import measure

    widths, pad = [2.0, 2.0, 2.0], 0.2
    one = measure(widths, 0, 1, pad)
    two = measure(widths, 0, 2, pad)
    assert one == pytest.approx(1.6)  # 2.0 less a margin either side
    assert two == pytest.approx(3.6)  # 4.0 less *one* pair, not two
    # The 0.4 that was padding between the two columns is measure now.
    assert two == pytest.approx(2 * one + 2 * pad)


def test_a_band_never_derives_its_colour_from_the_slides_own_surface(ctx_factory):
    """A banded row is painted, so it needs a colour something reads on. Nudging whatever the
    slide shows puts a mid-grey under the type the moment the master is a photograph."""
    ctx = ctx_factory(
        {"table": {"banding": True, "header": ["A", "B"], "rows": [["1", "2"], ["3", "4"]]}}
    )
    get_component("table")(ctx)
    shapes = {sh["name"]: sh for sh in ctx.manifest.to_dict()["slides"][0]["shapes"]}
    banded = next(s for n, s in shapes.items() if n.endswith("r3c1"))
    surface = ctx.theme.palette.pair("surface")
    assert banded["bg"] == surface.bg
    assert banded["fg"] == surface.fg


def test_a_banded_row_over_a_body_pair_keeps_that_pairs_ink(ctx_factory):
    """Shifting the fill must not strand the ink the pair was validated with."""
    ctx = ctx_factory(
        {
            "table": {
                "banding": True,
                "body_pair": "inverse",
                "header": ["A", "B"],
                "rows": [["1", "2"], ["3", "4"]],
            }
        }
    )
    get_component("table")(ctx)
    shapes = {sh["name"]: sh for sh in ctx.manifest.to_dict()["slides"][0]["shapes"]}
    plain = next(s for n, s in shapes.items() if n.endswith("r2c1"))
    banded = next(s for n, s in shapes.items() if n.endswith("r3c1"))
    inverse = ctx.theme.palette.pair("inverse")
    assert plain["bg"] == inverse.bg
    assert banded["bg"] != inverse.bg  # it is visibly a band
    assert banded["fg"] == inverse.fg  # in the ink the pair was checked with
    assert contrast_ratio(banded["fg"], banded["bg"]) >= 4.5


# --- cells that reach down --------------------------------------------------------


def test_a_cell_may_reach_down_over_the_rows_below(ctx_factory):
    """Merged in the file, not merely recorded as merged: a table whose vertical spans
    never joined records exactly the same manifest."""
    ctx = ctx_factory(
        {
            "table": {
                "header": ["Phase", "Task"],
                "rows": [[{"text": "Build", "down": 2}, "Compile"], ["Render"]],
            }
        }
    )
    get_component("table")(ctx)
    frame = next(s for s in ctx.slide.shapes if s.has_table)
    assert 'rowSpan="2"' in frame._element.xml, "the rows were never merged"
    cells = _rows(ctx)
    assert cells[(2, 1)][0] == "Build"
    assert (3, 1) not in cells  # row 3 column 1 belongs to Build
    assert cells[(3, 2)][0] == "Render"


def test_a_spanning_cell_records_the_depth_of_every_row_it_covers(ctx_factory):
    """Recording only its first row would have QA measure half the cell."""
    ctx = ctx_factory({"table": {"rows": [[{"text": "Build", "down": 2}, "Compile"], ["Render"]]}})
    get_component("table")(ctx)
    cells = _rows(ctx)
    assert cells[(1, 1)][1]["h"] == pytest.approx(
        cells[(1, 2)][1]["h"] + cells[(2, 2)][1]["h"], rel=0.01
    )


def test_a_row_beneath_a_reaching_cell_is_written_with_the_columns_left_to_it(ctx_factory):
    """The column above is spoken for, so the author writes two cells, not three."""
    ctx = ctx_factory(
        {
            "table": {
                "header": ["Phase", "Task", "Owner"],
                "rows": [[{"text": "Build", "down": 2}, "Compile", "Ana"], ["Render", "Bo"]],
            }
        }
    )
    get_component("table")(ctx)
    assert _rows(ctx)[(3, 2)][0] == "Render"


def test_a_reach_past_the_last_row_is_refused(ctx_factory):
    ctx = ctx_factory({"table": {"header": ["A", "B"], "rows": [[{"text": "x", "down": 4}, "y"]]}})
    with pytest.raises(LayoutError, match=r"sets 'down' to 4, reaching past the last"):
        get_component("table")(ctx)


def test_a_cell_landing_on_one_that_reaches_down_is_refused(ctx_factory):
    """Four columns, and the second row's spanning cell walks into the covered one."""
    ctx = ctx_factory(
        {
            "table": {
                "header": ["A", "B", "C", "D"],
                "rows": [
                    ["a", {"text": "b", "down": 2}, "c", "d"],
                    [{"text": "wide", "across": 3}],
                ],
            }
        }
    )
    with pytest.raises(LayoutError, match=r"reaches across a column that a 'down:' cell"):
        get_component("table")(ctx)


def test_a_row_that_forgets_its_inherited_column_is_refused(ctx_factory):
    ctx = ctx_factory(
        {
            "table": {
                "header": ["A", "B", "C"],
                "rows": [[{"text": "a", "down": 2}, "b", "c"], ["x", "y", "z"]],
            }
        }
    )
    with pytest.raises(LayoutError, match=r"row 2 covers 4 column\(s\) but the table is 3"):
        get_component("table")(ctx)


def test_a_reach_of_zero_is_refused(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [[{"text": "x", "down": 0}]]}})
    with pytest.raises(LayoutError, match=r"sets 'down' to 0"):
        get_component("table")(ctx)


def test_a_spanning_cell_that_wraps_deepens_the_rows_it_covers(ctx_factory):
    """A cell reaching down is skipped by the pass that sizes rows from their own
    cells, so nothing but the second pass can pay for its wrap."""
    short = ctx_factory({"table": {"rows": [[{"text": "Build", "down": 2}, "a"], ["b"]]}})
    long = ctx_factory(
        {"table": {"rows": [[{"text": " ".join(["word"] * 60), "down": 2}, "a"], ["b"]]}}
    )
    assert get_component("table")(long).height > get_component("table")(short).height


def test_a_spanning_cells_extra_depth_is_shared_by_the_rows_it_covers(ctx_factory):
    """Sizing a spanning label as though it lived in its first row alone builds a table just as
    tall, so the total height proves nothing and the split has to be read."""
    plain = ctx_factory({"table": {"rows": [[{"text": "x", "down": 2}, "a"], ["b"]]}})
    tall = ctx_factory(
        {"table": {"rows": [[{"text": " ".join(["word"] * 60), "down": 2}, "a"], ["b"]]}}
    )
    get_component("table")(plain)
    get_component("table")(tall)
    first, second = _rows(tall)[(1, 2)][1]["h"], _rows(tall)[(2, 2)][1]["h"]
    assert first == pytest.approx(second, rel=0.01)
    assert first > _rows(plain)[(1, 2)][1]["h"]  # the shortfall was real


def test_a_span_is_credited_the_depth_its_rows_already_hold(ctx_factory):
    """A span that ignores what its rows already hold pays for its whole height again — so it
    stops being shorter than the same text spelled out over two separate rows."""
    words = " ".join(["word"] * 60)
    spanned = ctx_factory({"table": {"rows": [[{"text": words, "down": 2}, "a"], ["b"]]}})
    apart = ctx_factory({"table": {"rows": [[words, "a"], ["", "b"]]}})
    assert get_component("table")(spanned).height < get_component("table")(apart).height


def test_a_cell_is_ruled_under_where_its_span_ends_not_under_where_it_starts(ctx_factory):
    """A cell reaching the bottom of the table has the table's own edge beneath it, so
    it takes no rule — even though the row it *starts* in is not the last one."""
    plain = ctx_factory({"table": {"rows": [["a", "b"], ["c", "d"], ["e", "f"]]}})
    spanned = ctx_factory({"table": {"rows": [["a", "b"], [{"text": "c", "down": 2}, "d"], ["f"]]}})
    get_component("table")(plain)
    get_component("table")(spanned)
    assert len(_table(plain)._tbl.findall(f".//{qn('a:lnB')}")) == 4
    assert len(_table(spanned)._tbl.findall(f".//{qn('a:lnB')}")) == 3


# --- how a cell's text sits in it -------------------------------------------------


def test_cells_are_middle_anchored_and_do_not_follow_the_placements_anchor(ctx_factory):
    """The placement's ``anchor:`` says where the block sits in its box — a different question —
    and reading it here top-sets every table."""
    ctx = _ctx(ctx_factory)
    assert ctx.anchor == "top"
    get_component("table")(ctx)
    assert _table(ctx).cell(1, 0).vertical_anchor == MSO_ANCHOR.MIDDLE


def test_valign_sets_the_table_and_a_cell_may_override_it(ctx_factory):
    ctx = ctx_factory(
        {"table": {"valign": "top", "rows": [["a", {"text": "b", "valign": "bottom"}]]}}
    )
    get_component("table")(ctx)
    tbl = _table(ctx)
    assert tbl.cell(0, 0).vertical_anchor == MSO_ANCHOR.TOP
    assert tbl.cell(0, 1).vertical_anchor == MSO_ANCHOR.BOTTOM


def test_an_unknown_valign_lists_the_vocabulary(ctx_factory):
    ctx = _ctx(ctx_factory, valign="centre")
    with pytest.raises(LayoutError, match=r"'valign' is one of top, middle, bottom"):
        get_component("table")(ctx)


def test_an_unknown_cell_valign_lists_the_vocabulary(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [[{"text": "x", "valign": "centre"}]]}})
    with pytest.raises(LayoutError, match=r"sets valign 'centre'; known: top, middle"):
        get_component("table")(ctx)


# --- what the hairlines mark ------------------------------------------------------


def _lines(ctx, edge):
    return len(_table(ctx)._tbl.findall(f".//{qn(f'a:{edge}')}"))


def test_rules_none_draws_no_line_at_all(ctx_factory):
    ctx = _ctx(ctx_factory, rules="none", total=["T", "1", "2"])
    get_component("table")(ctx)
    assert (_lines(ctx, "lnB"), _lines(ctx, "lnT"), _lines(ctx, "lnR")) == (0, 0, 0)


def test_rules_header_parts_the_band_from_the_body_and_nothing_else(ctx_factory):
    ctx = _ctx(ctx_factory, rules="header")
    get_component("table")(ctx)
    assert _lines(ctx, "lnB") == 3  # the header's three cells, and no body row


def test_rules_header_still_rules_a_total_above_itself(ctx_factory):
    ctx = _ctx(ctx_factory, rules="header", total=["T", "1", "2"])
    get_component("table")(ctx)
    assert _lines(ctx, "lnT") == 3


def test_rules_grid_adds_a_line_between_the_columns(ctx_factory):
    ctx = _ctx(ctx_factory, rules="grid")
    get_component("table")(ctx)
    # Three rows, and every column but the last carries a rule down its right side.
    assert _lines(ctx, "lnR") == 3 * 2
    assert _lines(ctx, "lnB") == 3 * 2  # unchanged from the default


def test_a_grid_rules_the_right_of_where_a_span_ends_not_where_it_starts(ctx_factory):
    """A cell reaching across into the last column has the table's own edge to its
    right, so it takes no rule — though the column it *starts* in is not the last."""
    ctx = ctx_factory(
        {
            "table": {
                "rules": "grid",
                "header": ["A", "B", "C"],
                "rows": [["a", {"text": "wide", "across": 2}]],
            }
        }
    )
    get_component("table")(ctx)
    # Header: two of three cells ruled. Body: 'a' is ruled, the spanning cell is not.
    assert _lines(ctx, "lnR") == 3


def test_an_unknown_rules_value_lists_the_vocabulary(ctx_factory):
    ctx = _ctx(ctx_factory, rules="dotted")
    with pytest.raises(LayoutError, match=r"'rules' is one of rows, header, grid, none"):
        get_component("table")(ctx)


# --- density ----------------------------------------------------------------------


def test_density_scales_both_cell_paddings(ctx_factory):
    plain = _ctx(ctx_factory)
    tight = _ctx(ctx_factory, density=0.5)
    get_component("table")(plain)
    get_component("table")(tight)
    wide, narrow = _table(plain).cell(0, 0), _table(tight).cell(0, 0)
    assert narrow.margin_left == pytest.approx(wide.margin_left / 2, rel=0.01)
    assert narrow.margin_top == pytest.approx(wide.margin_top / 2, rel=0.01)


def test_a_tighter_density_makes_the_table_shorter(ctx_factory):
    plain = _ctx(ctx_factory)
    tight = _ctx(ctx_factory, density=0.4)
    assert get_component("table")(tight).height < get_component("table")(plain).height


def test_a_density_of_zero_is_refused(ctx_factory):
    ctx = _ctx(ctx_factory, density=0)
    with pytest.raises(LayoutError, match=r"'density' scales the cell padding"):
        get_component("table")(ctx)


def test_a_non_numeric_density_is_refused(ctx_factory):
    ctx = _ctx(ctx_factory, density="tight")
    with pytest.raises(LayoutError, match=r"'density' scales the cell padding.*'tight'"):
        get_component("table")(ctx)


def test_a_tighter_density_lets_a_column_through_that_the_gutter_refused(ctx_factory):
    """The padding guard measures the padding actually used, so tightening it is the
    fix its own message names."""
    columns = [["c"] * 60]
    refused = ctx_factory({"table": {"rows": columns}})
    with pytest.raises(LayoutError, match=r"60 columns leave"):
        get_component("table")(refused)
    allowed = ctx_factory({"table": {"rows": columns, "density": 0.1}})
    assert get_component("table")(allowed).height > 0


def test_a_spanning_cell_wraps_against_the_columns_it_covers(ctx_factory):
    """Pins the call site that hands `measure()` the cell's span. 20 words is ~50em: over one
    column's ~37em and under the span's ~77."""
    words = " ".join(["word"] * 20)
    spanned = ctx_factory(
        {"table": {"header": ["A", "B"], "rows": [[{"text": words, "across": 2}]]}}
    )
    single = ctx_factory({"table": {"header": ["A", "B"], "rows": [[words, ""]]}})
    assert get_component("table")(spanned).height < get_component("table")(single).height


def test_two_spans_sharing_a_row_do_not_refuse_a_table_that_fits(ctx_factory):
    """Sharing a span's shortfall evenly grows a row once per span covering it. 210 words is the
    window where the two disagree: 5.94in even against 4.25in tight, in a 5.30in rect."""
    words = " ".join(["word"] * 210)
    ctx = ctx_factory(
        {
            "table": {
                "rows": [[{"text": words, "down": 2}, "x"], [{"text": words, "down": 2}], ["y"]]
            }
        }
    )
    result = get_component("table")(ctx)
    # `<= rect.height` was true by construction: the even split *raises* rather than
    # overrunning, so reaching this line at all satisfied it.
    frame = next(s for s in ctx.slide.shapes if s.has_table)
    assert result.height == pytest.approx(frame.height / 914400, abs=0.001)
    assert result.height < ctx.body_rect.height
    cells = _rows(ctx)
    for at in ((1, 1), (2, 2)):  # both spanning cells, two rows each
        assert cells[at][1]["h"] > 0


def test_the_even_split_is_still_what_a_table_with_room_gets(ctx_factory):
    """The tight packing pays a span off in its last row, which would leave the first
    row of every two-row label at the floor. It is the fallback, not the rule."""
    ctx = ctx_factory(
        {"table": {"rows": [[{"text": " ".join(["word"] * 60), "down": 2}, "a"], ["b"]]}}
    )
    get_component("table")(ctx)
    cells = _rows(ctx)
    assert cells[(1, 2)][1]["h"] == pytest.approx(cells[(2, 2)][1]["h"], rel=0.01)


def test_a_row_every_column_of_which_is_covered_from_above_is_refused(ctx_factory):
    """A wholly-vMerged row is height and nothing else, and LibreOffice drops the last row of
    any table containing one."""
    ctx = ctx_factory(
        {
            "table": {
                "header": ["Phase", "Detail"],
                "rows": [[{"text": "Discovery", "down": 2}, {"text": "Interviews", "down": 2}], []],
            }
        }
    )
    with pytest.raises(LayoutError, match=r"row 2 has no cells of its own"):
        get_component("table")(ctx)


def test_a_reach_beside_a_column_that_does_not_reach_is_fine(ctx_factory):
    """The refusal above is about the *whole* row being covered, not about stacking
    reaches: two of three columns spanning leaves the row a cell, and builds."""
    ctx = ctx_factory(
        {
            "table": {
                "header": ["Phase", "Detail", "Who"],
                "rows": [
                    [{"text": "Discovery", "down": 2}, {"text": "Interviews", "down": 2}, "Ana"],
                    ["Bea"],
                ],
            }
        }
    )
    get_component("table")(ctx)
    assert sorted(_rows(ctx)) == [
        (1, 1),
        (1, 2),
        (1, 3),  # the header
        (2, 1),
        (2, 2),
        (2, 3),  # two reaching, one not
        (3, 3),
    ]  # all row 3 has left


def test_an_empty_first_row_with_no_header_names_what_fixes_the_width(ctx_factory):
    ctx = ctx_factory({"table": {"rows": [[], ["a", "b"]]}})
    with pytest.raises(LayoutError, match=r"the first row fixes the table's width"):
        get_component("table")(ctx)


def test_a_row_with_a_cell_too_many_is_told_so_rather_than_told_to_add_one(ctx_factory):
    ctx = ctx_factory(
        {"table": {"header": ["A", "B"], "rows": [[{"text": "x", "down": 2}, "y"], ["z", "w"]]}}
    )
    with pytest.raises(LayoutError, match=r"row 2 covers 3 .*this row has a cell too many"):
        get_component("table")(ctx)


def test_the_tight_packing_pays_each_claim_off_in_the_row_where_it_ends():
    """Taking the claims in order of where they end is the whole of what makes the packing
    minimal: a long span settled first leaves depth past the reach of the shorter claims."""
    from pptxkit.components._tablegeom import _tight
    from pptxkit.components._tablespec import Cell, Placed, Row

    rows = [Row(cells=()), Row(cells=()), Row(cells=())]
    deep = Placed(Cell(text="a", down=3), row=0, col=0)  # rows 0-2, wants 4.0
    pair = Placed(Cell(text="b", down=2), row=0, col=1)  # rows 0-1, wants 2.0
    last = Placed(Cell(text="c"), row=2, col=1)  # row 2, wants the floor
    needs = {id(deep): 4.0, id(pair): 2.0, id(last): 0.5}

    packed = _tight(rows, [deep, pair, last], needs, floor=0.5)
    # The pair settles in row 1, and the deep span then only owes what row 2 must add.
    assert packed == [0.5, 1.5, 2.0]
    assert sum(packed) == 4.0  # the deep span's own need, nothing over


def test_an_empty_row_that_nothing_covers_is_told_it_is_short_not_degenerate(ctx_factory):
    """A stray `- []` in an ordinary table is a missing row of cells, not the covered
    row the message about height and nothing else is for."""
    ctx = ctx_factory({"table": {"header": ["A", "B"], "rows": [[], ["a", "b"]]}})
    with pytest.raises(LayoutError, match=r"row 1 covers 0 column\(s\).*add a cell"):
        get_component("table")(ctx)


def test_an_over_long_row_with_nothing_reaching_into_it_blames_the_comma(ctx_factory):
    """The commonest way a table row grows a cell is a comma YAML already split, and
    that is a different fix from a row standing under a reaching cell."""
    ctx = ctx_factory({"table": {"header": ["A", "B"], "rows": [["a", "b", "c"]]}})
    with pytest.raises(LayoutError, match=r"cell holding an unquoted comma"):
        get_component("table")(ctx)


def test_a_band_shifts_toward_the_ink_it_will_be_read_against(ctx_factory):
    """Which way the shift goes is the whole of whether banding is visible — away from the ink
    and a light table's band goes lighter still, invisible on the page it sits on."""
    from pptxkit.utils.color import relative_luminance

    def band_bg(pair_name):
        ctx = ctx_factory(
            {
                "table": {
                    "banding": True,
                    "body_pair": pair_name,
                    "header": ["A", "B"],
                    "rows": [["1", "2"], ["3", "4"]],
                }
            }
        )
        get_component("table")(ctx)
        shapes = {sh["name"]: sh for sh in ctx.manifest.to_dict()["slides"][0]["shapes"]}
        return next(s for n, s in shapes.items() if n.endswith("r3c1"))["bg"]

    light = ctx_factory({"table": {"rows": [["x"]]}}).theme.palette.pair("surface")
    dark = ctx_factory({"table": {"rows": [["x"]]}}).theme.palette.pair("inverse")
    assert relative_luminance(band_bg("surface")) < relative_luminance(light.bg), (
        "a light table's band must darken, or it disappears into its own fill"
    )
    assert relative_luminance(band_bg("inverse")) > relative_luminance(dark.bg), (
        "a dark table's band must lighten"
    )


def test_a_total_is_ruled_at_twice_the_weight_of_a_row_rule(ctx_factory):
    """The doubled weight is what makes a total read as a summary rather than one more row —
    an `lnT` at the row weight erases the distinction."""
    ctx = ctx_factory(
        {"table": {"header": ["Item", "Cost"], "rows": [["One", "12"]], "total": ["Total", "12"]}}
    )
    get_component("table")(ctx)
    tbl = _table(ctx)._tbl
    under = {int(e.get("w")) for e in tbl.findall(f".//{qn('a:lnB')}")}
    above = {int(e.get("w")) for e in tbl.findall(f".//{qn('a:lnT')}")}
    assert len(under) == 1 and len(above) == 1, (under, above)
    assert above.pop() == 2 * under.pop()
