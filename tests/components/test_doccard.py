import pytest
from PIL import Image

import pptxkit.components  # noqa: F401 — registers the built-ins
from pptxkit.errors import LayoutError
from pptxkit.layouts.components import as_body_result, get_component


@pytest.fixture
def fake_render(tmp_path, monkeypatch):
    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path))

    def render(html, path, *, width, scale):
        # Wide enough an aspect that the card fits the fixture theme's body rect.
        Image.new("RGB", (width * scale, 400 * scale), "white").save(path)
        return str(path)

    return render


@pytest.fixture
def doc(tmp_path):
    path = tmp_path / "SSV.md"
    path.write_text("# Self-service verification\n\nHow a customer verifies.\n")
    return path


def _ctx(ctx_factory, doc, **body):
    ctx = ctx_factory({"title": "T", "document": {"source": str(doc), **body}})
    return ctx


def test_the_card_is_placed_as_a_picture(ctx_factory, doc, fake_render, monkeypatch):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    ctx = _ctx(ctx_factory, doc)
    get_component("document")(ctx)
    assert len(ctx.slide.shapes) == 1


def test_it_is_recorded_as_an_image_not_native(ctx_factory, doc, fake_render, monkeypatch):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    ctx = _ctx(ctx_factory, doc)
    get_component("document")(ctx)
    assert [s.rendered for s in ctx.manifest.slides[0].shapes] == ["image"]


def test_it_reports_the_cards_height_not_its_width(ctx_factory, doc, fake_render, monkeypatch):
    """``height > 0`` is satisfied by the card's width, or by the whole body rect, just as well."""
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    ctx = _ctx(ctx_factory, doc)
    result = as_body_result(get_component("document")(ctx))
    picture = ctx.slide.shapes[0]
    assert picture.width != picture.height  # the fixture card is 2.5:1
    assert result.height == pytest.approx(picture.height / 914400)
    assert result.height < ctx.body_rect.height  # else `rect.height` would also pass


def test_one_reveal_group_holds_the_card(ctx_factory, doc, fake_render, monkeypatch):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    ctx = _ctx(ctx_factory, doc)
    result = as_body_result(get_component("document")(ctx))
    ids = {s.shape_id for s in ctx.slide.shapes}
    assert len(result.groups) == 1 and result.groups[0][0] in ids


def test_the_source_is_found_beside_the_deck_spec(ctx_factory, fake_render, monkeypatch, tmp_path):
    """Spec-relative is the rule `extends:`, `image:` and `background:` already follow, so a
    deck directory stays movable."""
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    spec_dir = tmp_path / "deckdir"
    spec_dir.mkdir()
    (spec_dir / "notes.md").write_text("# Notes\n\nBeside the spec.\n")
    monkeypatch.chdir(tmp_path)  # cwd is the PARENT, so only spec-relative can find it
    ctx = ctx_factory({"title": "T", "document": {"source": "notes.md"}}, base=spec_dir)
    get_component("document")(ctx)
    assert len(ctx.slide.shapes) == 1


def test_a_source_naming_neither_place_says_where_it_looked(ctx_factory, tmp_path):
    ctx = ctx_factory({"title": "T", "document": {"source": "nope.md"}}, base=tmp_path)
    with pytest.raises(LayoutError, match=r"source not found: nope\.md — looked beside"):
        get_component("document")(ctx)


def test_source_is_required(ctx_factory):
    ctx = ctx_factory({"title": "T", "document": {}})
    with pytest.raises(LayoutError, match="'source'"):
        get_component("document")(ctx)


def test_a_missing_source_file_names_the_slide(ctx_factory, tmp_path):
    ctx = ctx_factory({"title": "T", "document": {"source": str(tmp_path / "absent.md")}})
    with pytest.raises(LayoutError, match=r"slide 1 .*not found"):
        get_component("document")(ctx)


def test_a_left_side_card_stays_in_the_left_half(ctx_factory, doc, fake_render, monkeypatch):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    ctx = _ctx(ctx_factory, doc, side="left")
    get_component("document")(ctx)
    rect = ctx.body_rect
    picture = ctx.slide.shapes[0]
    assert (picture.left + picture.width) / 914400 <= rect.left + rect.width / 2 + 0.2


def test_a_right_side_card_stays_in_the_right_half(ctx_factory, doc, fake_render, monkeypatch):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    ctx = _ctx(ctx_factory, doc, side="right")
    get_component("document")(ctx)
    rect = ctx.body_rect
    picture = ctx.slide.shapes[0]
    assert picture.left / 914400 >= rect.left + rect.width / 2 - 0.2


def test_an_unrecognised_side_is_rejected(ctx_factory, doc):
    ctx = _ctx(ctx_factory, doc, side="up")
    with pytest.raises(LayoutError, match="'side'"):
        get_component("document")(ctx)


def test_a_directory_source_says_not_a_file(ctx_factory, tmp_path):
    ctx = ctx_factory({"title": "T", "document": {"source": str(tmp_path)}})
    with pytest.raises(LayoutError, match=r"slide 1 .*not a file"):
        get_component("document")(ctx)


def test_a_non_utf8_source_names_the_slide(ctx_factory, tmp_path):
    bad = tmp_path / "bad.md"
    bad.write_bytes(b"\xff\xfe not valid utf-8")
    ctx = ctx_factory({"title": "T", "document": {"source": str(bad)}})
    with pytest.raises(LayoutError, match=r"slide 1 .*bad\.md"):
        get_component("document")(ctx)


def test_a_non_positive_max_width_names_the_slide(ctx_factory, doc):
    ctx = _ctx(ctx_factory, doc, max_width=0)
    with pytest.raises(LayoutError, match="'max_width'"):
        get_component("document")(ctx)


def test_the_document_content_reaches_the_card(ctx_factory, doc, tmp_path, monkeypatch):
    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path))
    seen = {}

    def render(html, path, *, width, scale):
        seen["html"] = html
        # Aspect 2.5 regardless of width — fits the fixture theme's body rect.
        Image.new("RGB", (width * scale, int(width * scale / 2.5)), "white").save(path)
        return str(path)

    monkeypatch.setattr("pptxkit.components.doccard._render", render)
    get_component("document")(_ctx(ctx_factory, doc))
    assert "Self-service verification" in seen["html"]
    assert "How a customer verifies" in seen["html"]


def test_max_width_feeds_the_renderer(ctx_factory, doc, tmp_path, monkeypatch):
    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path))
    seen = {}

    def render(html, path, *, width, scale):
        seen["width"] = width
        Image.new("RGB", (width * scale, int(width * scale / 2.5)), "white").save(path)
        return str(path)

    monkeypatch.setattr("pptxkit.components.doccard._render", render)
    get_component("document")(_ctx(ctx_factory, doc, max_width=650))
    assert seen["width"] == 650


def test_a_non_numeric_max_width_names_the_slide(ctx_factory, doc):
    ctx = _ctx(ctx_factory, doc, max_width="wide")
    with pytest.raises(LayoutError, match="'max_width'"):
        get_component("document")(ctx)


def test_the_slides_own_theme_colour_reaches_the_cards_rules(
    ctx_factory, doc, tmp_path, monkeypatch, theme
):
    """The card's CSS rules consume ``ctx.theme`` through ``panel_css`` rather than carry an
    inert copy: a card missing ``content_css=panel_css(ctx.theme)`` produces neither."""
    import dataclasses

    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path))
    seen = {}

    def render(html, path, *, width, scale):
        seen["html"] = html
        Image.new("RGB", (width * scale, int(width * scale / 2.5)), "white").save(path)
        return str(path)

    monkeypatch.setattr("pptxkit.components.doccard._render", render)
    odd_theme = dataclasses.replace(
        theme,
        palette=dataclasses.replace(
            theme.palette, roles={**theme.palette.roles, "muted": "AB12CD"}
        ),
    )
    ctx = ctx_factory({"title": "T", "document": {"source": str(doc)}}, theme_override=odd_theme)
    get_component("document")(ctx)
    assert "#AB12CD" in seen["html"]


def test_a_too_tall_card_names_the_slide_height_and_budget(ctx_factory, doc, tmp_path, monkeypatch):
    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path))

    def render(html, path, *, width, scale):
        # Square, so at full width the card is far taller than the body rect.
        Image.new("RGB", (width * scale, width * scale), "white").save(path)
        return str(path)

    monkeypatch.setattr("pptxkit.components.doccard._render", render)
    ctx = _ctx(ctx_factory, doc)
    with pytest.raises(
        LayoutError, match=r"slide 1 .*came out .*in tall but only .*in is available"
    ):
        get_component("document")(ctx)


def test_document_is_registered():
    from pptxkit.layouts.components import registered_components

    assert "document" in registered_components()


def test_an_unspecified_side_fills_the_body_width(ctx_factory, doc, fake_render, monkeypatch):
    """'full' is the default; 'left' would silently halve every card that omits side."""
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    ctx = _ctx(ctx_factory, doc)
    get_component("document")(ctx)
    placed = ctx.slide.shapes[0]
    from pptx.util import Emu

    assert Emu(placed.width).inches == pytest.approx(ctx.body_rect.width, abs=0.01)


def test_the_titlebar_shows_the_filename_the_deck_asked_for(
    ctx_factory, doc, fake_render, monkeypatch
):
    seen = {}

    def capture(html, path, *, width, scale):
        seen["html"] = html
        return fake_render(html, path, width=width, scale=scale)

    monkeypatch.setattr("pptxkit.components.doccard._render", capture)
    get_component("document")(_ctx(ctx_factory, doc, filename="RENAMED.md"))
    assert "RENAMED.md" in seen["html"]
    assert "SSV.md" not in seen["html"]


@pytest.fixture
def numbered(tmp_path):
    """Ten lines whose text names their own number, so an excerpt is unambiguous."""
    path = tmp_path / "numbered.md"
    path.write_text("".join(f"line {n}\n" for n in range(1, 11)))
    return path


@pytest.fixture
def carded(monkeypatch):
    """The markdown that actually reached the card, rather than the picture it became."""
    seen = []

    def capture(md_text, **kwargs):
        seen.append(md_text)
        return "<html></html>"

    monkeypatch.setattr("pptxkit.components.doccard.markdown_card", capture)
    return seen


def test_lines_cards_only_the_named_range(ctx_factory, numbered, carded, fake_render, monkeypatch):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    get_component("document")(_ctx(ctx_factory, numbered, lines="3-5"))
    assert carded == ["line 3\nline 4\nline 5"]


def test_the_range_is_inclusive_at_both_ends(
    ctx_factory, numbered, carded, fake_render, monkeypatch
):
    """Half-open would drop 'line 4'; zero-based would card lines 4 and 5."""
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    get_component("document")(_ctx(ctx_factory, numbered, lines="3-4"))
    assert carded == ["line 3\nline 4"]


def test_an_end_past_the_last_line_stops_at_the_last_line(
    ctx_factory, numbered, carded, fake_render, monkeypatch
):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    get_component("document")(_ctx(ctx_factory, numbered, lines="9-400"))
    assert carded == ["line 9\nline 10"]


def test_without_lines_the_whole_file_is_carded(
    ctx_factory, numbered, carded, fake_render, monkeypatch
):
    """The default must stay 'all of it' — `lines:` is an opt-in narrowing."""
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    get_component("document")(_ctx(ctx_factory, numbered))
    assert carded == ["".join(f"line {n}\n" for n in range(1, 11))]


@pytest.mark.parametrize("bad", ["12", "twelve-40", "12-", "-40", "12:40", ""])
def test_a_range_that_is_not_two_numbers_is_refused(
    ctx_factory, numbered, bad, fake_render, monkeypatch
):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    with pytest.raises(LayoutError, match="must look like"):
        get_component("document")(_ctx(ctx_factory, numbered, lines=bad))


@pytest.mark.parametrize("bad", ["0-4", "5-2"])
def test_a_range_that_cannot_select_anything_is_refused(
    ctx_factory, numbered, bad, fake_render, monkeypatch
):
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    with pytest.raises(LayoutError, match="must start at 1 or more"):
        get_component("document")(_ctx(ctx_factory, numbered, lines=bad))


def test_a_start_past_the_end_of_the_file_is_refused(
    ctx_factory, numbered, fake_render, monkeypatch
):
    """The excerpt was cut out of the source: silently carding nothing hides that."""
    monkeypatch.setattr("pptxkit.components.doccard._render", fake_render)
    with pytest.raises(LayoutError, match="has only 10 line"):
        get_component("document")(_ctx(ctx_factory, numbered, lines="11-20"))
