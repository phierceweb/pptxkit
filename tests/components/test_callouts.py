import dataclasses

import pytest

import pptxkit.components  # noqa: F401
from pptxkit.errors import LayoutError
from pptxkit.layouts.components import get_component
from pptxkit.theme import Scale

ITEMS = [
    {"head": "One map per subsystem", "body": "A doc for each part of the system."},
    {"head": "Plain text it reads natively", "body": "Markdown, not a special format."},
]

TALL = Scale(26.666, 15.0)


def _texts(slide):
    return "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def _ctx(ctx_factory, items=ITEMS):
    return ctx_factory({"callouts": {"items": items}})


def _tall(theme):
    """The fixture theme on twice the canvas — every rung and margin follows the scale."""
    return dataclasses.replace(
        theme,
        scale=TALL,
        grid=dataclasses.replace(theme.grid, scale=TALL),
        ramp={role: dataclasses.replace(style, scale=TALL) for role, style in theme.ramp.items()},
    )


def test_each_item_renders_its_head_and_body(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("callouts")(ctx)
    text = _texts(ctx.slide)
    assert "One map per subsystem" in text
    assert "Markdown, not a special format." in text


def test_one_reveal_group_per_item(ctx_factory):
    ctx = _ctx(ctx_factory)
    assert len(get_component("callouts")(ctx).groups) == len(ITEMS)


def test_each_group_holds_the_dot_and_the_text(ctx_factory):
    ctx = _ctx(ctx_factory)
    groups = get_component("callouts")(ctx).groups
    assert all(len(group) == 2 for group in groups)


def test_every_returned_id_is_a_real_shape(ctx_factory):
    ctx = _ctx(ctx_factory)
    groups = get_component("callouts")(ctx).groups
    ids = {s.shape_id for s in ctx.slide.shapes}
    assert all(spid in ids for group in groups for spid in group)


def test_items_is_required(ctx_factory):
    ctx = ctx_factory({"callouts": {}})
    with pytest.raises(LayoutError, match=r"slide 1 .*'items'"):
        get_component("callouts")(ctx)


def test_an_item_without_a_head_names_the_slide(ctx_factory):
    ctx = _ctx(ctx_factory, items=[{"body": "no head"}])
    with pytest.raises(LayoutError, match=r"slide 1 .*item 1.*'head'"):
        get_component("callouts")(ctx)


def test_a_bare_string_item_names_the_slide(ctx_factory):
    ctx = _ctx(ctx_factory, items=["not a mapping"])
    with pytest.raises(LayoutError, match=r"slide 1 .*item 1.*'head'"):
        get_component("callouts")(ctx)


def test_items_stack_downward_without_overlapping(ctx_factory):
    ctx = _ctx(ctx_factory, items=ITEMS * 2)
    get_component("callouts")(ctx)
    text_boxes = sorted(
        (s.top, s.top + s.height)
        for s in ctx.slide.shapes
        if s.has_text_frame and s.text_frame.text
    )
    for i, ((_, bottom), (next_top, _)) in enumerate(zip(text_boxes, text_boxes[1:], strict=False)):
        assert bottom <= next_top + 0.01, (
            f"callout rows overlap vertically at transition {i}-{i + 1}: "
            f"row {i} bottom {bottom} > row {i + 1} top {next_top}"
        )


def test_content_stays_inside_the_body_rect(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("callouts")(ctx)
    rect = ctx.body_rect
    for shape in ctx.slide.shapes:
        assert shape.left / 914400 >= rect.left - 0.01
        assert (shape.left + shape.width) / 914400 <= rect.right + 0.01
        assert shape.top / 914400 >= rect.top - 0.01
        assert (shape.top + shape.height) / 914400 <= rect.bottom + 0.01


def test_heads_are_recorded_in_the_manifest(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("callouts")(ctx)
    assert "One map per subsystem" in ctx.manifest.slides[0].texts()


def test_recorded_lines_match_the_drawn_paragraphs(ctx_factory):
    ctx = _ctx(ctx_factory)
    get_component("callouts")(ctx)
    by_id = {s.shape_id: s for s in ctx.slide.shapes}
    recorded = [r for r in ctx.manifest.slides[0].shapes if r.lines]
    # Head and body are separate lines in the manifest because qa's overflow
    # check looks for each of them in the rendered text.
    assert [r.lines for r in recorded] == [[i["head"], i["body"]] for i in ITEMS]
    for record in recorded:
        drawn = [p.text for p in by_id[record.shape_id].text_frame.paragraphs]
        assert drawn == record.lines


def test_more_copy_than_the_rect_holds_raises(ctx_factory):
    """Rows are measured, so the limit is the copy's own depth, not the item count."""
    wordy = {"head": "A heading", "body": "sentence that wraps. " * 25}
    ctx = _ctx(ctx_factory, items=[wordy] * 6)
    with pytest.raises(LayoutError, match=r"slide 1.*items need.*of height but the body"):
        get_component("callouts")(ctx)


def test_eight_short_items_now_fit(ctx_factory):
    """Eight one-line items are 8 x a heading line, well inside the rect — measured rows are
    what lets them fit."""
    items = [{"head": f"P{i}", "body": "line"} for i in range(8)]
    assert len(get_component("callouts")(_ctx(ctx_factory, items=items)).groups) == 8


def test_a_row_is_as_deep_as_its_own_copy(ctx_factory):
    """The invariant the equal-lane layout could not hold, and the reason it overflowed."""
    ctx = _ctx(
        ctx_factory,
        items=[
            {"head": "Short", "body": "one line"},
            {"head": "Long", "body": "a sentence that has to wrap several times over. " * 4},
        ],
    )
    get_component("callouts")(ctx)
    short_h, long_h = _frame_heights(ctx.slide)
    body_line = ctx.style("body").size * 1.2 / 72
    # Equal lanes gave both rows the same frame, so any gap at all reddens the old
    # arithmetic; a whole body line proves the wrap was actually counted.
    assert long_h - short_h >= body_line, (short_h, long_h, body_line)


def test_a_reasonable_item_count_builds(ctx_factory):
    ctx = _ctx(ctx_factory, items=[{"head": f"Item {i}", "body": f"Body {i}"} for i in range(4)])
    groups = get_component("callouts")(ctx).groups
    assert len(groups) == 4


def test_the_reported_height_is_the_rows_it_drew(ctx_factory):
    """``0 < h <= body_rect.height`` is true by construction: the row height derives from the
    body rect. Three rows, where the type's own minimum caps them short of the bottom."""
    ctx = _ctx(ctx_factory, items=ITEMS[:1] * 3)
    result = get_component("callouts")(ctx)
    tops, heights = _row_tops(ctx.slide), _frame_heights(ctx.slide)
    drawn = (tops[-1] + heights[-1]) - ctx.body_rect.top
    assert result.height == pytest.approx(drawn, abs=0.01)
    assert result.height < ctx.body_rect.height  # else `rect.height` would also pass


def _row_tops(slide):
    """Top of each row's text box, in inches — excludes the dot (also a
    text-frame-bearing shape, but with no text) so gaps reflect row spacing."""
    return sorted(s.top / 914400 for s in slide.shapes if s.has_text_frame and s.text_frame.text)


def _frame_heights(slide):
    """Height of each row's text box, in inches — the dot's own frame carries no words."""
    return [s.height / 914400 for s in slide.shapes if s.has_text_frame and s.text_frame.text]


def test_short_lists_do_not_drift_apart(ctx_factory):
    ctx = _ctx(ctx_factory, items=ITEMS[:1] * 3)
    get_component("callouts")(ctx)
    tops = _row_tops(ctx.slide)
    gaps = [b - a for a, b in zip(tops, tops[1:], strict=False)]
    assert all(g < ctx.body_rect.height / 3 for g in gaps)


def test_a_long_list_still_uses_the_available_height(ctx_factory):
    ctx = _ctx(ctx_factory, items=ITEMS[:1] * 5)
    get_component("callouts")(ctx)
    tops = _row_tops(ctx.slide)
    assert (tops[1] - tops[0]) < 1.42


def test_callouts_is_registered():
    from pptxkit.layouts.components import registered_components

    assert "callouts" in registered_components()


def test_align_is_refused_so_a_row_cannot_drift_from_its_dot(ctx_factory):
    ctx = ctx_factory({"callouts": {"items": [{"head": "H"}]}})
    ctx.align = "center"
    with pytest.raises(LayoutError, match="would pull each row's text away from its dot"):
        get_component("callouts")(ctx)


def test_the_dot_is_the_brand_accent_and_sits_clear_of_the_text(ctx_factory):
    """Named colour, and a geometric relation — not a re-read of the constants."""
    from pptx.util import Emu

    ctx = ctx_factory({"callouts": {"items": [{"head": "H", "body": "b"}]}})
    get_component("callouts")(ctx)
    dot = next(s for s in ctx.slide.shapes if not s.has_text_frame or not s.text_frame.text)
    text = next(s for s in ctx.slide.shapes if s.has_text_frame and "H" in s.text_frame.text)

    assert str(dot.fill.fore_color.rgb) == ctx.theme.palette.role("accent-1")
    assert dot.left + dot.width <= text.left, "the dot overlaps the text column"
    assert Emu(dot.width).inches < 0.3, "a dot this large is a disc, not a bullet"


# --- a mark instead of a dot -------------------------------------------------

MARKED = [
    {"head": "Secure", "body": "Checked.", "icon": "lock"},
    {"head": "Fast", "body": "Quick.", "icon": "bolt"},
]


def test_an_item_with_an_icon_draws_a_glyph_where_the_dot_would_be(ctx_factory):
    """The built deck hides this: a dot and a glyph are both a small shape on the rail."""
    ctx = _ctx(ctx_factory, MARKED)
    get_component("callouts")(ctx)
    marks = [s for s in ctx.slide.shapes if "<a:custGeom" in s._element.xml]
    assert len(marks) == 2
    assert not any("roundRect" in s._element.xml for s in ctx.slide.shapes)


def test_a_list_can_mix_dots_and_glyphs(ctx_factory):
    ctx = _ctx(ctx_factory, [MARKED[0], {"head": "Plain", "body": "No mark named."}])
    get_component("callouts")(ctx)
    xml = [s._element.xml for s in ctx.slide.shapes]
    assert sum("<a:custGeom" in x for x in xml) == 1
    assert sum("roundRect" in x for x in xml) == 1


def test_the_text_clears_the_glyph_it_sits_beside(ctx_factory):
    """A glyph is wider than the dot, so a fixed indent would run the head through it."""
    dotted = _ctx(ctx_factory)
    get_component("callouts")(dotted)
    marked = _ctx(ctx_factory, MARKED)
    get_component("callouts")(marked)

    def first_text_left(ctx):
        # A glyph shape carries an empty text frame of its own, so having one is not
        # what makes a shape text — having words in it is.
        return min(
            s.left for s in ctx.slide.shapes if s.has_text_frame and s.text_frame.text.strip()
        )

    def first_mark_right(ctx):
        return min(
            s.left + s.width
            for s in ctx.slide.shapes
            if not (s.has_text_frame and s.text_frame.text.strip())
        )

    assert first_text_left(marked) > first_text_left(dotted)
    assert first_text_left(marked) >= first_mark_right(marked)


def test_the_glyph_and_its_text_reveal_together(ctx_factory):
    """The mark belongs to its row, so a click-build must not leave it behind."""
    ctx = _ctx(ctx_factory, MARKED)
    result = get_component("callouts")(ctx)
    assert all(len(group) == 2 for group in result.groups)
