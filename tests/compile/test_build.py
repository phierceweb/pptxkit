import hashlib
import json
import textwrap
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from pptxkit.compile import build_deck
from pptxkit.compile.build import _drop_template_slides
from pptxkit.errors import SpecError, ThemeError
from pptxkit.theme.chartstyle import ChartStyle
from pptxkit.theme import Grid, Scale
from pptxkit.theme.defaults import DEFAULT_PAIRS
from pptxkit.theme.model import Theme, TypeStyle
from pptxkit.theme.palette import build_palette

FIXTURE_ROLES = {
    "page": "FFFFFF",
    "ink": "2D0937",
    "muted": "573C65",
    "line": "EDEDED",
    "surface": "F5F6F8",
    "surface-ink": "2D0937",
    "inverse": "2D0937",
    "inverse-ink": "FFFFFF",
    "accent-1": "27B94C",
    "accent-2": "18CEDA",
    "accent-3": "4DB6AC",
    "accent-4": "A78BD0",
}


def _theme_stub(*, drop: bool) -> Theme:
    """A minimal Theme for exercising build.py's own logic — mirrors conftest's testtheme.yaml."""
    scale = Scale(13.333, 7.5)
    ramp = {
        "kicker": TypeStyle(14 / 7.5, scale, bold=True),
        "caption": TypeStyle(13 / 7.5, scale),
        "body": TypeStyle(11 / 7.5, scale),
        "lead": TypeStyle(19 / 7.5, scale),
        "subtitle": TypeStyle(15 / 7.5, scale),
        "title": TypeStyle(32 / 7.5, scale, bold=True),
        "display": TypeStyle(46 / 7.5, scale, bold=True),
    }
    grid = Grid(
        scale=scale,
        top_frac=0.30 / 7.5,
        right_frac=0.61 / 13.333,
        bottom_frac=0.5 / 7.5,
        left_frac=0.62 / 13.333,
        columns=12,
        rows=12,
        gutter_frac=0.18 / 13.333,
        body_top_frac=1.7 / 7.5,
    )
    return Theme(
        name="testtheme",
        template=Path("assets/t.pptx"),
        drop_template_slides=drop,
        palette=build_palette(FIXTURE_ROLES, pairs=DEFAULT_PAIRS),
        scale=scale,
        face="Calibri",
        mono="Consolas",
        ramp=ramp,
        min_pt=10.5,
        grid=grid,
        line_weight=2.25,
        chart=ChartStyle(),
    )


def _resolved(recorded: str, manifest) -> Path:
    """A manifest path as the reader resolves it — against the manifest's own directory."""
    path = Path(recorded)
    return path if path.is_absolute() else (Path(manifest).parent / path).resolve()


def test_builds_one_slide_per_spec_document(project):
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    assert result.slides == 3
    assert len(Presentation(str(result.deck)).slides) == 3


def test_writes_the_deck_to_the_spec_out_path(project):
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    assert result.deck == project / "out" / "Demo.pptx"
    assert result.deck.is_file()


def test_out_override_wins_over_the_spec(project, tmp_path):
    dest = tmp_path / "elsewhere" / "Other.pptx"
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml", out=dest)
    assert result.deck == dest and dest.is_file()


def test_a_failed_save_leaves_the_previous_deck_intact(project, monkeypatch):
    """Patched on the stdlib class, not python-pptx's private writer, so a rewrite
    upstream cannot quietly stop this from firing."""
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    good = result.deck.read_bytes()
    listing = sorted(p.name for p in result.deck.parent.iterdir())

    calls = []
    writestr = zipfile.ZipFile.writestr

    def failing(self, *args, **kwargs):
        calls.append(1)
        if len(calls) == 5:
            raise OSError("no space left on device")
        return writestr(self, *args, **kwargs)

    monkeypatch.setattr(zipfile.ZipFile, "writestr", failing)
    with pytest.raises(OSError, match="no space left on device"):
        build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    monkeypatch.undo()

    assert result.deck.read_bytes() == good
    assert sorted(p.name for p in result.deck.parent.iterdir()) == listing


def test_writes_a_manifest_beside_the_deck(project):
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    data = json.loads(result.manifest.read_text())
    assert data["theme"] == "testtheme"
    assert [s["background"] for s in data["slides"]] == ["page", "inverse", "page"]
    texts = [s["text"] for s in data["slides"][0]["shapes"] if s["text"]]
    assert "One" in texts  # the section kicker
    assert "First" in texts  # the slide title


def test_manifest_records_slide_size_and_the_resolved_theme_path(project):
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    data = json.loads(result.manifest.read_text())
    assert data["canvas"] == {"w": 13.333, "h": 7.5, "unit": "in"}
    # Recorded relative to the manifest; what matters is that it resolves back.
    assert _resolved(data["theme_path"], result.manifest) == (project / "testtheme.yaml").resolve()


def test_slide_content_reaches_the_deck(project):
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    prs = Presentation(str(result.deck))
    text = "\n".join(
        sh.text_frame.text for slide in prs.slides for sh in slide.shapes if sh.has_text_frame
    )
    assert "First" in text and "Second" in text and "With a subhead" in text
    assert "The second act." in text


def test_speaker_notes_are_written(project):
    (project / "n.deck.yaml").write_text(
        textwrap.dedent("""
        theme: testtheme
        ---
        title: T
        notes: Remember the framing.
    """)
    )
    result = build_deck(
        project / "n.deck.yaml", theme_path=project / "testtheme.yaml", out=project / "n.pptx"
    )
    prs = Presentation(str(result.deck))
    assert "Remember the framing." in prs.slides[0].notes_slide.notes_text_frame.text


def test_an_unknown_component_names_the_slide_and_the_placement(project):
    (project / "bad.deck.yaml").write_text(
        "theme: testtheme\n---\nplace:\n  - at: {cols: full}\n    nonesuch: {}\n"
    )
    with pytest.raises(SpecError, match=r"slide 1: placement 1: unknown field 'nonesuch'"):
        build_deck(
            project / "bad.deck.yaml",
            theme_path=project / "testtheme.yaml",
            out=project / "bad.pptx",
        )


def test_a_deck_with_no_out_and_no_override_is_rejected(project):
    (project / "noout.deck.yaml").write_text("theme: testtheme\n---\ntitle: x\n")
    with pytest.raises(SpecError, match="no output path"):
        build_deck(project / "noout.deck.yaml", theme_path=project / "testtheme.yaml")


def test_extends_module_components_are_available(project):
    (project / "ext.py").write_text(
        textwrap.dedent("""
        from pptxkit.layouts.components import component
        from pptxkit.utils.shapes import para, textbox

        @component("t-custom-body")
        def custom(ctx):
            tf = textbox(ctx.slide, 1, 2, 5, 1)
            para(tf, "from the escape hatch", 18, ctx.fg(), first=True)
            return []
    """)
    )
    (project / "e.deck.yaml").write_text(
        textwrap.dedent("""
        theme: testtheme
        extends: ext.py
        ---
        place:
          - at: {cols: full}
            t-custom-body: {}
    """)
    )
    result = build_deck(
        project / "e.deck.yaml", theme_path=project / "testtheme.yaml", out=project / "e.pptx"
    )
    prs = Presentation(str(result.deck))
    text = "\n".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert "from the escape hatch" in text


def test_template_sample_slides_are_dropped(synthetic_template):
    prs = Presentation(str(synthetic_template))
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])
    assert len(prs.slides) == 3

    theme = _theme_stub(drop=True)
    _drop_template_slides(prs, theme)
    assert len(prs.slides) == 0


def test_template_slides_are_kept_when_the_theme_says_so(synthetic_template):
    prs = Presentation(str(synthetic_template))
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])
    _drop_template_slides(prs, _theme_stub(drop=False))
    assert len(prs.slides) == 3


def test_theme_dir_env_var_redirects_theme_lookup(project, monkeypatch, tmp_path):
    custom_dir = tmp_path / "custom_themes"
    custom_dir.mkdir(parents=True)
    (custom_dir / "t.pptx").write_bytes((project / "assets" / "t.pptx").read_bytes())
    (custom_dir / "testtheme.theme.yaml").write_text(
        (project / "testtheme.yaml").read_text().replace("assets/t.pptx", "t.pptx")
    )
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(custom_dir))

    result = build_deck(project / "d.deck.yaml", out=project / "env_out.pptx")

    assert result.deck.is_file()
    recorded = json.loads(result.manifest.read_text())["theme_path"]
    assert _resolved(recorded, result.manifest) == (custom_dir / "testtheme.theme.yaml").resolve()


def test_default_theme_path_resolves_under_templates(
    tmp_path, monkeypatch, synthetic_template, theme_yaml
):
    monkeypatch.chdir(tmp_path)
    themes_dir = tmp_path / "templates"
    themes_dir.mkdir(parents=True)
    (themes_dir / "t.pptx").write_bytes(synthetic_template.read_bytes())
    (themes_dir / "testtheme.theme.yaml").write_text(theme_yaml.replace("assets/t.pptx", "t.pptx"))

    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    (deck_dir / "d.deck.yaml").write_text("theme: testtheme\nout: out/D.pptx\n---\ntitle: Solo\n")

    result = build_deck(deck_dir / "d.deck.yaml")
    assert result.deck.is_file()
    assert json.loads(result.manifest.read_text())["theme"] == "testtheme"


_PLACED = """
    theme: testtheme
    title: Demo
    out: out/Placed.pptx
    ---
    title: One
    place:
    - at: {cols: left-half, rows: {from: 1, to: 12}}
      id: keep
      bullets: {items: [A]}
    - at: {cols: right-half, rows: {from: 1, to: 12}}
      bullets: {items: [B]}
"""

_PLACED_WITH_AN_INSERT = """
    theme: testtheme
    title: Demo
    out: out/Placed.pptx
    ---
    title: One
    place:
    - at: {cols: full, rows: {from: 0, to: 1}}
      rule: {}
    - at: {cols: left-half, rows: {from: 1, to: 12}}
      id: keep
      bullets: {items: [A]}
    - at: {cols: right-half, rows: {from: 1, to: 12}}
      bullets: {items: [B]}
"""


def _shape_names(project, spec_text: str) -> list[str]:
    spec = project / "p.deck.yaml"
    spec.write_text(textwrap.dedent(spec_text))
    result = build_deck(spec, theme_path=project / "testtheme.yaml")
    data = json.loads(Path(result.manifest).read_text())
    return [s["name"] for s in data["slides"][0]["shapes"]]


def test_an_id_holds_a_placement_name_when_one_is_inserted_above_it(project):
    """The whole point of naming shapes: an edit above must not rename what is below."""
    before = _shape_names(project, _PLACED)
    after = _shape_names(project, _PLACED_WITH_AN_INSERT)
    assert "s1.keep.bullets#1" in before
    assert "s1.keep.bullets#1" in after
    # The placement with no id is named for its position, so the insert moves it.
    assert "s1.p2.bullets#1" in before
    assert "s1.p2.bullets#1" not in after
    assert "s1.p3.bullets#1" in after


def test_the_manifest_names_the_spec_that_produced_it(project):
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    data = json.loads(Path(result.manifest).read_text())
    # Recorded relative to the manifest, so a delivered file carries no build-machine
    # path — and so the deck directory survives being moved.
    assert data["spec"] == "../d.deck.yaml"
    assert _resolved(data["spec"], result.manifest) == (project / "d.deck.yaml").resolve()
    # Computed here rather than with build's own helper, so this cannot agree with a
    # broken one.
    assert (
        data["spec_hash"] == hashlib.sha256((project / "d.deck.yaml").read_bytes()).hexdigest()[:16]
    )


def test_deck_hash_is_the_deck_that_was_written(project):
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    data = json.loads(Path(result.manifest).read_text())
    assert data["deck_hash"] == hashlib.sha256(result.deck.read_bytes()).hexdigest()[:16]


def test_build_id_is_stable_until_the_spec_changes(project):
    spec, theme = project / "d.deck.yaml", project / "testtheme.yaml"
    first = json.loads(Path(build_deck(spec, theme_path=theme).manifest).read_text())
    again = json.loads(Path(build_deck(spec, theme_path=theme).manifest).read_text())
    assert again["build_id"] == first["build_id"]
    spec.write_text(spec.read_text().replace("title: First", "title: Moved"))
    edited = json.loads(Path(build_deck(spec, theme_path=theme).manifest).read_text())
    assert edited["build_id"] != first["build_id"]


def test_build_id_changes_when_the_theme_changes(project):
    spec, theme = project / "d.deck.yaml", project / "testtheme.yaml"
    first = json.loads(Path(build_deck(spec, theme_path=theme).manifest).read_text())
    theme.write_text(theme.read_text().replace("columns: 12", "columns: 6"))
    after = json.loads(Path(build_deck(spec, theme_path=theme).manifest).read_text())
    assert after["build_id"] != first["build_id"]


def test_the_spec_is_kept_beside_the_deck_it_made(project):
    """An edited spec leaves its earlier builds unregenerable, in a directory labelled
    disposable — the snapshot is what keeps `out/` actually disposable."""
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    kept = result.deck.parent / ".build" / f"{result.deck.stem}.deck.yaml"
    assert kept.is_file()
    assert kept.read_text() == (project / "d.deck.yaml").read_text()


def test_each_version_keeps_its_own_spec(project):
    """The failure this exists for is version 2 overwriting version 1's only copy."""
    first = build_deck(
        project / "d.deck.yaml",
        theme_path=project / "testtheme.yaml",
        out=project / "out" / "Demo v1.pptx",
    )
    (project / "d.deck.yaml").write_text(
        (project / "d.deck.yaml").read_text().replace("Demo", "Demo edited")
    )
    second = build_deck(
        project / "d.deck.yaml",
        theme_path=project / "testtheme.yaml",
        out=project / "out" / "Demo v2.pptx",
    )
    scratch = first.deck.parent / ".build"
    kept = sorted(p.name for p in scratch.glob("*.deck.yaml"))
    assert kept == ["Demo v1.deck.yaml", "Demo v2.deck.yaml"], kept
    assert (scratch / "Demo v1.deck.yaml").read_text() != (
        scratch / "Demo v2.deck.yaml"
    ).read_text()
    assert second.deck.is_file()


def test_recorded_paths_carry_no_build_machine_path(project, tmp_path):
    """A manifest can be handed over beside its deck; the username in an absolute path
    goes with it, and nothing in a delivered manifest needs it."""
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")
    data = json.loads(Path(result.manifest).read_text())
    for key in ("spec", "deck", "theme_path"):
        assert not Path(data[key]).is_absolute(), f"{key} is absolute: {data[key]}"
    assert data["deck"] == Path(result.deck).name


def test_a_path_sharing_only_the_filesystem_root_stays_absolute(project, tmp_path):
    """Nothing about such a build is portable, and a relative form would climb out
    through the same home directory it exists to keep out."""
    out = tmp_path / "elsewhere" / "d.pptx"
    result = build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml", out=out)
    data = json.loads(Path(result.manifest).read_text())
    # `project` and `tmp_path` are siblings under pytest's tmp root, so they do share an
    # ancestor — assert the rule, not the accident: whatever is written resolves back.
    resolved = Path(data["theme_path"])
    if not resolved.is_absolute():
        resolved = (Path(result.manifest).parent / resolved).resolve()
    assert resolved == (project / "testtheme.yaml").resolve()


_NAMED = """
    theme: testtheme
    title: Demo
    out: out/Named.pptx
    ---
    title: One
    place:
    - at: {cols: left-half, rows: {from: 1, to: 12}}
      id: hero
      bullets: {items: [A, B]}
    - at: {cols: right-half, rows: {from: 1, to: 12}}
      stats: {items: [{value: "9", label: nine}]}
"""


def test_the_package_carries_the_shape_names_the_manifest_records(project):
    """Read from the `.pptx` itself, not the manifest: `docs/qa.md` promises the naming
    reaches the package, which is the half a hand-edit in PowerPoint depends on."""
    spec = project / "p.deck.yaml"
    spec.write_text(textwrap.dedent(_NAMED))
    result = build_deck(spec, theme_path=project / "testtheme.yaml")
    data = json.loads(Path(result.manifest).read_text())

    package = {s.name for s in Presentation(str(result.deck)).slides[0].shapes}
    recorded = [s["name"] for s in data["slides"][0]["shapes"]]
    assert recorded, "nothing was recorded, so the comparison below proves nothing"
    assert "s1.hero.bullets#1" in recorded

    for name in recorded:
        # Parts sharing one frame collapse to their origin in the package — one shape
        # cannot carry three names — so the origin is the legal weaker form.
        assert name in package or name.rsplit(".", 1)[0] in package, (
            f"{name!r} is in the manifest but no shape in the .pptx carries it"
        )

    # `qa.md` says *every* shape is named for a spec node. A component that draws without
    # recording leaves python-pptx's own "TextBox 4" — a name that is unique and wrong.
    unnamed = {n for n in package if not n.startswith("s1.")}
    assert unnamed == set(), f"shapes in the .pptx no spec node named: {sorted(unnamed)}"


def test_a_theme_name_missing_from_the_theme_dir_resolves_to_the_packaged_builtin(
    tmp_path, monkeypatch
):
    """An install with no `templates/` dir; `theme: base` must still build."""
    from pptxkit.compile.build import resolve_theme

    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path / "no-such-dir"))
    resolved = resolve_theme("base")
    assert resolved.is_file()
    assert resolved.parent.name == "builtin"


def test_a_theme_dir_file_wins_over_the_packaged_builtin(tmp_path, monkeypatch):
    from pptxkit.compile.build import resolve_theme

    (tmp_path / "base.theme.yaml").write_text("name: local\n")
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path))
    assert resolve_theme("base") == tmp_path / "base.theme.yaml"


def test_an_unknown_theme_name_resolves_to_the_theme_dir_candidate(tmp_path, monkeypatch):
    """The not-found error must name the directory the caller controls, not the package."""
    from pptxkit.compile.build import resolve_theme

    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path))
    assert resolve_theme("nope") == tmp_path / "nope.theme.yaml"


def test_a_deck_naming_base_builds_with_no_theme_dir_at_all(tmp_path, monkeypatch):
    """End to end: the fallback reaches build_deck, not just the resolver."""
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path / "no-such-dir"))
    spec = tmp_path / "d.deck.yaml"
    spec.write_text("theme: base\nout: Demo.pptx\n---\ntitle: Hello\n")
    result = build_deck(spec)
    assert result.slides == 1
    assert (tmp_path / "Demo.pptx").is_file()


def test_an_unknown_spec_theme_name_is_reported_as_a_name(tmp_path, monkeypatch):
    """Resolving to a path before loading turned a bad `theme:` into 'file not found:
    /abs/candidate.theme.yaml', which names neither the name nor the way to fix it."""
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path))
    spec = tmp_path / "d.deck.yaml"
    spec.write_text("theme: acme\nout: Demo.pptx\n---\ntitle: Hello\n")

    with pytest.raises(ThemeError) as excinfo:
        build_deck(spec)
    assert "unknown theme 'acme'" in str(excinfo.value)
    assert "--adopt acme" in str(excinfo.value)
