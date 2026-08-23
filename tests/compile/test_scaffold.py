"""`pptxkit new` — the deck you get before you know the wire format."""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation

from pptxkit.compile.scaffold import new_deck, slug_and_title
from pptxkit.errors import SpecError, ThemeError

import pptxkit.theme as _theme_pkg

_THEMES = Path(_theme_pkg.__file__).parent / "builtin"


@pytest.fixture
def elsewhere(tmp_path, monkeypatch):
    """A working directory of its own. The theme directory is read relative to the process,
    so a scaffold built outside the repo needs it named."""
    monkeypatch.chdir(tmp_path)
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(_THEMES))
    return tmp_path


def test_the_scaffold_builds(elsewhere):
    """The whole promise. A scaffold that stops compiling is worse than none — the
    first thing an author meets would be a build error in code they did not write."""
    made = new_deck("Q4 Review", root=elsewhere / "authoring")
    assert made.built is not None
    assert len(Presentation(str(made.built.deck)).slides) == 6


@pytest.mark.parametrize(
    "title",
    [
        "Q4 Review",
        "Customer Retention Strategy",
        "Annual Security Posture Review",
        "Enterprise Cloud Transformation Roadmap",
        "Platform Reliability And Incident Response Review",
        "Quarterly Business Review For The Northern Region Team",
        "A Rather Long Deck Title About Operational Excellence Program",
    ],
)
def test_the_scaffold_builds_whatever_the_title_is(title, elsewhere):
    """The cover's title box is written once and the title is whatever was typed. Size
    that box for one line and every title that wraps trips the chrome overflow guard —
    which is `pptxkit new` failing on the deck names people actually use."""
    made = new_deck(title, root=elsewhere / "authoring")
    assert made.built is not None
    assert len(Presentation(str(made.built.deck)).slides) == 6


def test_a_title_yaml_would_read_as_syntax_survives_into_the_deck(elsewhere):
    """An unquoted `title: Q4 Review: Growth` is not a mapping value — the scaffold
    would write a spec that cannot be parsed at all."""
    made = new_deck("Q4 Review: Growth And Churn", root=elsewhere / "authoring")
    config = yaml.safe_load(made.spec.read_text().split("\n---\n")[0])
    assert config["title"] == "Q4 Review: Growth And Churn"
    cover = Presentation(str(made.built.deck)).slides[0]
    assert "Q4 Review: Growth And Churn" in [
        s.text_frame.text for s in cover.shapes if s.has_text_frame
    ]


@pytest.mark.parametrize("name", ["../escape hatch", "Q1/Q2 Review", "a/../../b"])
def test_a_deck_name_never_writes_outside_the_root_it_was_given(name, elsewhere):
    root = elsewhere / "authoring"
    made = new_deck(name, root=root, build=False)
    assert made.spec.parent.parent.resolve() == root.resolve()


def test_a_failed_build_frees_the_name_it_just_wrote(elsewhere):
    """Leave the spec behind and the retry meets 'already exists' instead of the
    error that stopped it — one typo and the name is spent."""
    root = elsewhere / "authoring"
    with pytest.raises(ThemeError):
        new_deck("Q4 Review", root=root, theme="no-such-theme")
    assert not (root / "q4-review").exists()
    assert new_deck("Q4 Review", root=root).built is not None


def test_the_spec_lands_under_the_deck_s_own_directory(elsewhere):
    made = new_deck("Q4 Review", root=elsewhere / "authoring", build=False)
    assert made.spec == elsewhere / "authoring" / "q4-review" / "q4-review.deck.yaml"


def test_nothing_is_built_when_the_build_is_declined(elsewhere):
    made = new_deck("Q4 Review", root=elsewhere / "authoring", build=False)
    assert made.built is None
    assert made.spec.is_file()


def test_an_existing_deck_is_never_overwritten(elsewhere):
    """A scaffold is the one command that would destroy work by succeeding."""
    new_deck("Q4 Review", root=elsewhere / "authoring", build=False)
    with pytest.raises(SpecError, match="already exists"):
        new_deck("q4-review", root=elsewhere / "authoring", build=False)


def test_the_theme_reaches_the_spec(elsewhere):
    made = new_deck("D", root=elsewhere / "authoring", theme="motion", build=False)
    assert "theme: motion" in made.spec.read_text()


def test_the_scaffold_says_what_to_change_on_every_slide(elsewhere):
    """Comments are the documentation an author has in hand; a bare spec is a wall."""
    made = new_deck("D", root=elsewhere / "authoring", build=False)
    text = made.spec.read_text()
    slides = text.split("\n---\n")[1:]
    assert len(slides) == 6
    assert sum(1 for s in slides if s.lstrip().startswith("#")) >= 5


@pytest.mark.parametrize(
    "name, slug, title",
    [
        ("Q4 Review", "q4-review", "Q4 Review"),
        ("q4-review", "q4-review", "Q4 Review"),
        ("q4_review", "q4-review", "Q4 Review"),
        ("  spaced  out  ", "spaced-out", "Spaced Out"),
        # str.title() would give 'Ml Pipeline'; an acronym keeps the case it was typed in.
        ("ML-pipeline", "ml-pipeline", "ML Pipeline"),
    ],
)
def test_a_name_gives_the_same_deck_however_it_is_typed(name, slug, title):
    assert slug_and_title(name) == (slug, title)


@pytest.mark.parametrize("name", ["", "   ", "-", "__", "!!!", "//"])
def test_a_nameless_deck_is_refused(name):
    with pytest.raises(SpecError, match="a deck needs a name"):
        slug_and_title(name)
