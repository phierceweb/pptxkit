import io
import re

from PIL import Image
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

from pptxkit.compile.background import _flatten, flatten_master_background

_BG_EMBED = re.compile(rb"<p:bg>.*?r:embed=\"(rId\d+)\"", re.DOTALL)


def _png(mode: str, color, size=(8, 8)) -> bytes:
    buf = io.BytesIO()
    Image.new(mode, size, color).save(buf, format="PNG")
    return buf.getvalue()


def _with_picture_background(tmp_path, template, blob: bytes):
    """A copy of ``template`` whose slide master draws ``blob`` as its background."""
    path = tmp_path / "bg.pptx"
    path.write_bytes(template.read_bytes())
    prs = Presentation(str(path))
    master = prs.slide_masters[0]
    _, rid = master.part.get_or_add_image_part(io.BytesIO(blob))
    master._element.find(qn("p:cSld")).insert(
        0,
        parse_xml(
            f"<p:bg {nsdecls('p', 'a', 'r')}><p:bgPr><a:blipFill>"
            f'<a:blip r:embed="{rid}"/><a:stretch><a:fillRect/></a:stretch>'
            f"</a:blipFill><a:effectLst/></p:bgPr></p:bg>"
        ),
    )
    prs.save(str(path))
    return path


def _background_blob(prs) -> bytes:
    master = prs.slide_masters[0]
    rid = _BG_EMBED.search(master.part.blob).group(1).decode()
    return master.part.related_part(rid).blob


def _alpha_hole(marker) -> bytes:
    """A transparent PNG with one opaque pixel of ``marker`` at (2, 2)."""
    img = Image.new("RGBA", (4, 4), (0, 0, 0, 0))
    img.putpixel((2, 2), (*marker, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def test_opaque_pixels_are_preserved():
    src = Image.new("RGBA", (4, 4), (0, 0, 0, 0))
    src.putpixel((2, 2), (10, 25, 54, 255))
    buf = io.BytesIO()
    src.save(buf, format="PNG")
    flat = Image.open(io.BytesIO(_flatten(buf.getvalue(), (255, 255, 255)))).convert("RGB")
    assert flat.getpixel((2, 2)) == (10, 25, 54)
    assert flat.getpixel((0, 0)) == (255, 255, 255)


def test_the_result_has_no_alpha_channel():
    out = _flatten(_png("RGBA", (0, 0, 0, 0)), (255, 255, 255))
    assert Image.open(io.BytesIO(out)).mode == "RGB"


def test_a_fully_opaque_image_is_left_alone():
    assert _flatten(_png("RGBA", (1, 2, 3, 255)), (255, 255, 255)) is None


def test_an_image_without_alpha_is_left_alone():
    assert _flatten(_png("RGB", (1, 2, 3)), (255, 255, 255)) is None


def test_non_image_bytes_are_left_alone():
    assert _flatten(b"not an image at all", (255, 255, 255)) is None


def test_flattening_composites_onto_the_requested_colour():
    out = _flatten(_png("RGBA", (0, 0, 0, 0)), (10, 25, 54))
    assert Image.open(io.BytesIO(out)).convert("RGB").getpixel((0, 0)) == (10, 25, 54)


def test_a_template_without_a_picture_background_is_unchanged(synthetic_template):
    prs = Presentation(str(synthetic_template))
    assert flatten_master_background(prs, (255, 255, 255)) is False


def test_a_transparent_master_background_is_composited_onto_the_deck_colour(
    tmp_path, synthetic_template
):
    """Keynote composites an alpha master background over black; this is the fix."""
    path = _with_picture_background(tmp_path, synthetic_template, _alpha_hole((200, 30, 40)))
    prs = Presentation(str(path))
    assert flatten_master_background(prs, (10, 25, 54)) is True
    flat = Image.open(io.BytesIO(_background_blob(prs))).convert("RGB")
    assert flat.getpixel((0, 0)) == (10, 25, 54)  # was transparent
    assert flat.getpixel((2, 2)) == (200, 30, 40)  # the artwork itself survives


def test_an_opaque_master_background_is_left_untouched(tmp_path, synthetic_template):
    path = _with_picture_background(tmp_path, synthetic_template, _png("RGB", (1, 2, 3)))
    prs = Presentation(str(path))
    before = _background_blob(prs)
    assert flatten_master_background(prs, (10, 25, 54)) is False
    assert _background_blob(prs) == before
