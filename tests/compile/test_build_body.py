import json
import textwrap

import pytest
from pptx import Presentation

from pptxkit.compile import build_deck
from pptxkit.compile.record import box_of

DECK = """
    theme: testtheme
    sections: [One]
    out: out/Body.pptx
    ---
    section: One
    title: With a body
    animate: one_at_a_time
    place:
      - at: {cols: full}
        callouts:
          items:
            - {head: First point, body: Its supporting line.}
            - {head: Second point, body: Another line.}
"""


def test_a_body_component_renders_through_the_compiler(project):
    (project / "b.deck.yaml").write_text(textwrap.dedent(DECK))
    result = build_deck(project / "b.deck.yaml", theme_path=project / "testtheme.yaml")
    prs = Presentation(str(result.deck))
    text = "\n".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert "First point" in text and "Second point" in text


def test_the_animate_is_recorded_in_the_manifest(project):
    (project / "b2.deck.yaml").write_text(textwrap.dedent(DECK))
    result = build_deck(
        project / "b2.deck.yaml", theme_path=project / "testtheme.yaml", out=project / "b2.pptx"
    )
    data = json.loads(result.manifest.read_text())
    animations = data["slides"][0]["animations"]
    assert animations and animations[0]["kind"] == "click_sequence"
    assert len(animations[0]["steps"]) == 2


def test_a_body_key_fails_the_build_pointing_at_the_new_shape(project):
    (project / "bad.deck.yaml").write_text(
        "theme: testtheme\nout: bad.pptx\n---\ntitle: T\nbody:\n  type: nonesuch\n"
    )
    from pptxkit.errors import SpecError

    with pytest.raises(SpecError, match=r"'body' is gone.*'place:'"):
        build_deck(
            project / "bad.deck.yaml",
            theme_path=project / "testtheme.yaml",
            out=project / "bad.pptx",
        )


TWO_UP = """
    theme: testtheme
    out: out/TwoUp.pptx
    ---
    kicker: Task 6
    title: Side by side
    subtitle: Two placements
    place:
      - at: {cols: left-half}
        bullets: {items: [First point, Second point]}
      - at: {cols: right-half}
        stats:
          items:
            - {value: "1", label: one}
"""


def test_two_placements_on_one_slide_get_distinct_boxes(project):
    """The defect this engine exists to fix: both otherwise get the whole body band."""
    (project / "two.deck.yaml").write_text(textwrap.dedent(TWO_UP))
    result = build_deck(
        project / "two.deck.yaml", theme_path=project / "testtheme.yaml", out=project / "two.pptx"
    )
    shapes = json.loads(result.manifest.read_text())["slides"][0]["shapes"]
    left = box_of([s for s in shapes if "First point" in (s.get("text") or "")][0])
    right = box_of([s for s in shapes if "one" in (s.get("text") or "")][0])
    assert right[0] >= left[0] + left[2]


def test_the_chrome_lines_are_recorded_in_order_from_one_box(project):
    (project / "chrome.deck.yaml").write_text(textwrap.dedent(TWO_UP))
    result = build_deck(
        project / "chrome.deck.yaml",
        theme_path=project / "testtheme.yaml",
        out=project / "chrome.pptx",
    )
    shapes = json.loads(result.manifest.read_text())["slides"][0]["shapes"]
    chrome = [s for s in shapes if s["text"] in ("Task 6", "Side by side", "Two placements")]
    assert [s["text"] for s in chrome] == ["Task 6", "Side by side", "Two placements"]
    assert len({box_of(s) for s in chrome}) == 1


def test_the_content_band_starts_below_the_chrome_box(project):
    (project / "under.deck.yaml").write_text(textwrap.dedent(TWO_UP))
    result = build_deck(
        project / "under.deck.yaml",
        theme_path=project / "testtheme.yaml",
        out=project / "under.pptx",
    )
    shapes = json.loads(result.manifest.read_text())["slides"][0]["shapes"]
    chrome = box_of([s for s in shapes if s.get("text") == "Task 6"][0])
    body = box_of([s for s in shapes if "First point" in (s.get("text") or "")][0])
    assert body[1] >= chrome[1] + chrome[3]
