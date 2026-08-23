import json
from dataclasses import dataclass

import pytest
from pf_core.exceptions import InvalidInputError, PreconditionError
from pptx import Presentation
from pptx.util import Inches

from pptxkit.compile.record import Box, box_of
from pptxkit.compile.manifest import ManifestRecorder
from pptxkit.errors import SpecError
from pptxkit.theme.model import Rect


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide):
    return slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(0.5))


def test_records_a_shape_with_its_box_in_inches(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page", section="One")
    m.record(_box(slide), text="Hello", font_pt=32, fg="2D0937", bg="FFFFFF")
    rec = m.to_dict()["slides"][0]["shapes"][0]
    assert rec["box"] == {"x": 1.0, "y": 2.0, "w": 3.0, "h": 0.5}
    assert rec["text"] == "Hello"
    assert rec["font_pt"] == 32
    assert rec["fg"] == "2D0937"
    # "native" is the default, so it is omitted; the image case is pinned below.
    assert "rendered" not in rec


def test_texts_collects_only_native_text(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record(_box(slide), text="native text")
    m.record(_box(slide), text="inside a panel", rendered="image")
    assert m.slides[0].texts() == ["native text"]


def test_panel_text_is_still_recorded_for_reference(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record(_box(slide), text="inside a panel", rendered="image")
    assert m.to_dict()["slides"][0]["shapes"][0]["rendered"] == "image"


def test_animation_steps_name_the_shapes_they_reveal(slide):
    """A shape id says nothing to a reader and is not unique on a slide."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.p1.callouts"
    first, second = m.record(_box(slide)), m.record(_box(slide))
    m.record_animation("click_sequence", [[first.shape_id], [second.shape_id]])
    assert m.to_dict()["slides"][0]["animations"] == [
        {"kind": "click_sequence", "steps": [["s1.p1.callouts#1"], ["s1.p1.callouts#2"]]}
    ]


def test_a_reveal_carrying_a_motion_role_is_named_like_any_other(slide):
    """A component may report (id, role); the role picked an OOXML preset already."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.p1.rule"
    rec = m.record(_box(slide))
    m.record_animation("click_sequence", [[(rec.shape_id, "line")]])
    assert m.slides[0].animations[0]["steps"] == [["s1.p1.rule#1"]]


def test_a_reveal_target_that_was_never_recorded_keeps_its_id(slide):
    """Better a number that says which shape than a name invented for it."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record_animation("click_sequence", [[99]])
    assert m.slides[0].animations[0]["steps"] == [["shape 99"]]


def test_recording_without_beginning_a_slide_is_an_error(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    with pytest.raises(PreconditionError, match="begin_slide"):
        m.record(_box(slide), text="x")


def test_recording_an_animation_without_beginning_a_slide_is_an_error(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    with pytest.raises(PreconditionError, match="begin_slide"):
        m.record_animation("click_sequence", [[2, 3]])


def test_marking_a_backdrop_without_beginning_a_slide_is_an_error(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    with pytest.raises(PreconditionError, match="begin_slide"):
        m.mark_backdrop()


def test_slides_keep_spec_order(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    for i in (1, 2, 3):
        m.begin_slide(i, background="page")
    assert [s["index"] for s in m.to_dict()["slides"]] == [1, 2, 3]


def test_write_produces_readable_json(tmp_path, slide):
    m = ManifestRecorder(deck="d.pptx", theme="base")
    m.begin_slide(1, background="inverse")
    m.record(_box(slide), text="Demo")
    out = tmp_path / "d.manifest.json"
    m.write(out)
    data = json.loads(out.read_text())
    assert data["theme"] == "base"
    assert data["slides"][0]["shapes"][0]["text"] == "Demo"


def test_slide_dimensions_and_theme_path_default_to_empty():
    m = ManifestRecorder(deck="d.pptx", theme="t")
    data = m.to_dict()
    assert data["canvas"] == {"w": 0.0, "h": 0.0, "unit": "in"}
    assert data["theme_path"] == ""


def test_slide_dimensions_and_theme_path_are_recorded():
    m = ManifestRecorder(
        deck="d.pptx",
        theme="t",
        slide_w=13.333,
        slide_h=7.5,
        theme_path="templates/brand.theme.yaml",
    )
    data = m.to_dict()
    assert data["canvas"] == {"w": 13.333, "h": 7.5, "unit": "in"}
    assert data["theme_path"] == "templates/brand.theme.yaml"


def test_lines_are_recorded_individually(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record(_box(slide), lines=["•  alpha", "•  beta"])
    assert m.to_dict()["slides"][0]["shapes"][0]["lines"] == ["•  alpha", "•  beta"]


def test_texts_flattens_lines(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record(_box(slide), lines=["alpha", "beta"])
    m.record(_box(slide), text="gamma")
    assert m.slides[0].texts() == ["alpha", "beta", "gamma"]


def test_text_defaults_to_the_joined_lines(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    rec = m.record(_box(slide), lines=["alpha", "beta"])
    assert rec.text == "alpha beta"


def test_an_explicit_text_wins_over_the_join(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    rec = m.record(_box(slide), text="explicit", lines=["a", "b"])
    assert rec.text == "explicit"


def test_lines_default_to_empty(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    assert m.record(_box(slide), text="solo").lines == []


def test_image_rendered_lines_are_excluded_from_texts(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record(_box(slide), lines=["in a panel"], rendered="image")
    assert m.slides[0].texts() == []


def test_an_unknown_rendered_value_is_rejected(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    with pytest.raises(InvalidInputError, match="rendered must be"):
        m.record(_box(slide), text="x", rendered="Image")


def test_rendered_accepts_both_valid_values(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    assert m.record(_box(slide), text="a", rendered="native").rendered == "native"
    assert m.record(_box(slide), text="b", rendered="image").rendered == "image"


def test_the_slide_records_the_background_it_was_painted_on(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="inverse")
    record = m.to_dict()["slides"][0]
    assert record["background"] == "inverse"
    assert "layout" not in record


def test_a_slide_records_the_templates_own_backdrop():
    """Nothing is drawn for it, so the flag is all QA gets to tell a photograph from paint."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.mark_backdrop()
    m.begin_slide(2, background="page")
    records = m.to_dict()["slides"]
    assert records[0]["backdrop"] is True
    # Omitted rather than written false: a default-valued key is not recorded.
    assert "backdrop" not in records[1]


@dataclass(frozen=True)
class _NotAShape:
    """A table cell as the recorder sees it — mirrors ``table._CellBox``. No ``_element``:
    renaming would raise rather than be skipped."""

    shape_id: int = 2
    name: str = "Table 1 r1c1"
    left: int = 914400
    top: int = 914400
    width: int = 914400
    height: int = 914400


def test_an_origin_names_the_shape_in_the_package_too(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.hero.card"
    shape = _box(slide)
    rec = m.record(shape)
    assert rec.name == "s1.hero.card#1"
    # The package name is the half that survives a hand-edit in PowerPoint.
    assert shape.name == "s1.hero.card#1"


def test_parts_are_numbered_from_one_within_each_origin(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.p1.card"
    m.record(_box(slide))
    m.record(_box(slide))
    m.origin = "s1.p2.card"
    m.record(_box(slide))
    assert [s.name for s in m.slides[0].shapes] == [
        "s1.p1.card#1",
        "s1.p1.card#2",
        "s1.p2.card#1",
    ]


def test_parts_sharing_one_frame_leave_the_package_name_at_the_origin(slide):
    """Chrome's stacked lines are three paragraphs in one textbox — the condition has to be
    built, or the collapsed name is asserted against a case that was never constructed."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.chrome"
    shape = _box(slide)
    names = [m.record(shape, part=p).name for p in ("kicker", "title", "subtitle")]
    assert names == ["s1.chrome.kicker", "s1.chrome.title", "s1.chrome.subtitle"]
    # One shape cannot carry three names, so the package keeps the shared origin.
    assert shape.name == "s1.chrome"


def test_parts_on_separate_shapes_each_keep_their_own_package_name(slide):
    """A chrome field with its own `at:` is its own textbox, and must say which it is —
    collapsing the names strands an author who moved one by hand."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.chrome"
    shapes = [_box(slide) for _ in range(3)]
    for shape, part in zip(shapes, ("kicker", "title", "subtitle"), strict=True):
        m.record(shape, part=part)
    assert [s.name for s in shapes] == [
        "s1.chrome.kicker",
        "s1.chrome.title",
        "s1.chrome.subtitle",
    ]


def test_one_frame_per_origin_is_tracked_separately(slide):
    """The shared-frame memory is per origin, or the second placement inherits the first."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    first, second = _box(slide), _box(slide)
    m.origin = "s1.p1.card"
    m.record(first, part="head")
    m.origin = "s1.p2.card"
    m.record(second, part="head")
    assert [first.name, second.name] == ["s1.p1.card.head", "s1.p2.card.head"]


def test_a_shape_recorded_outside_any_origin_keeps_its_own_name(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    shape = _box(slide)
    shape.name = "Untouched 7"
    assert m.record(shape).name == "Untouched 7"


def test_a_cell_that_is_not_a_shape_is_recorded_without_being_renamed():
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.p1.table"
    assert m.record(_NotAShape(), part="r1c1").name == "s1.p1.table.r1c1"


def test_a_box_is_recorded_keyed_and_rounded_to_the_thousandth(slide):
    """A 13.333in canvas divides out of EMU as 13.33299978127734."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    wide = slide.shapes.add_textbox(
        Inches(0.7999792213473316),
        Inches(0.375),
        Inches(11.733039151356081),
        Inches(1.04998687664042),
    )
    assert m.record(wide).box == Box(x=0.8, y=0.375, w=11.733, h=1.05)


def test_a_font_size_is_recorded_rounded_to_the_hundredth(slide):
    """The EMU round trip records 13pt as 12.99975; nobody reads that as 13."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    assert m.record(_box(slide), text="x", font_pt=12.99975).font_pt == 13.0


def test_keys_still_at_their_default_are_not_written(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record(_box(slide), text="Hello")
    written = m.to_dict()["slides"][0]["shapes"][0]
    assert set(written) == {"shape_id", "name", "box", "text"}


def test_a_value_that_moves_off_its_default_is_written(slide):
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record(_box(slide), text="Hello", rendered="image", fg="FFFFFF")
    written = m.to_dict()["slides"][0]["shapes"][0]
    assert written["rendered"] == "image"
    assert written["fg"] == "FFFFFF"
    assert "bg" not in written


def test_box_of_reads_a_keyed_box_back_as_left_top_width_height():
    assert box_of({"box": {"x": 1.0, "y": 2.0, "w": 3.0, "h": 0.5}}) == (1.0, 2.0, 3.0, 0.5)


def test_box_of_returns_none_for_a_record_with_no_box():
    assert box_of({"name": "s1.p1.rule#1"}) is None


def test_a_manifest_predating_keyed_boxes_says_so_instead_of_failing_obscurely():
    """`qa` on a deck built before this change would otherwise raise TypeError from
    indexing a list with a string."""
    # A positional box, as every manifest written before the cutover carries it.
    stale = {"name": "TextBox 2", "box": [1.0, 2.0, 3.0, 0.5]}
    with pytest.raises(SpecError, match="rebuild the deck"):
        box_of(stale)


def test_the_canvas_is_recorded_rounded_like_every_other_inch():
    """A theme's slide width arrives as the EMU round trip left it."""
    m = ManifestRecorder(deck="d.pptx", theme="t", slide_w=13.33299978127734, slide_h=7.5)
    assert m.to_dict()["canvas"] == {"w": 13.333, "h": 7.5, "unit": "in"}


class _RecycledIdShape:
    """A shape sharing one `_element` address with another — a reused lxml proxy."""

    def __init__(self, element, shape_id: int) -> None:
        self._element = element
        self.shape_id = shape_id
        self.name = ""
        self.left = self.top = 914400  # EMU; record reads geometry before naming
        self.width = self.height = 914400


def test_distinct_shapes_are_not_merged_by_a_recycled_element_address():
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.chrome"

    shared_address = object()  # one object => one id, standing in for reuse
    kicker = _RecycledIdShape(shared_address, shape_id=2)
    title = _RecycledIdShape(shared_address, shape_id=3)

    m.record(kicker, part="kicker")
    m.record(title, part="title")

    assert [kicker.name, title.name] == ["s1.chrome.kicker", "s1.chrome.title"], (
        f"distinct shapes collapsed onto one name: {[kicker.name, title.name]}"
    )


def test_a_genuinely_shared_frame_still_shares_one_name():
    """Paragraphs in one frame are one shape, so the second part reuses the origin."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.chrome"

    frame = _RecycledIdShape(object(), shape_id=7)
    m.record(frame, part="kicker")
    first = frame.name
    m.record(frame, part="title")

    assert first == "s1.chrome.kicker"
    assert frame.name == "s1.chrome", f"shared frame was renamed per part: {frame.name}"


def test_animation_name_map_matches_the_package_name():
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.origin = "s1.chrome"

    shared_address = object()
    kicker = _RecycledIdShape(shared_address, shape_id=2)
    title = _RecycledIdShape(shared_address, shape_id=3)
    m.record(kicker, part="kicker")
    m.record(title, part="title")

    m.record_animation("appear", [[2], [3]])
    steps = m.to_dict()["slides"][0]["animations"][0]["steps"]
    assert steps == [["s1.chrome.kicker"], ["s1.chrome.title"]], (
        f"animation steps disagree with the package names: {steps}"
    )


def test_a_recorded_placement_survives_json(slide):
    """`_slim`'s list-of-dataclass branch turns the nested `Box` into a dict; without it
    every build dies at `write` while every `to_dict()` test stays green."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    m.record_placement("s1.hero.card", "card", Rect(1.0, 2.0, 3.0, 4.0))

    written = json.loads(json.dumps(m.to_dict()))
    placements = written["slides"][0]["placements"]
    assert placements == [
        {
            "origin": "s1.hero.card",
            "component": "card",
            "box": {"x": 1.0, "y": 2.0, "w": 3.0, "h": 4.0},
        }
    ]
    # box_of reads a placement unchanged, which is why the field is named `box`.
    assert box_of(placements[0]) == (1.0, 2.0, 3.0, 4.0)


def test_a_slide_with_no_placements_records_no_key(slide):
    """Chrome-only slides exist; an empty list must drop out rather than ship as `[]`."""
    m = ManifestRecorder(deck="d.pptx", theme="t")
    m.begin_slide(1, background="page")
    assert "placements" not in m.to_dict()["slides"][0]


def test_recording_a_placement_before_a_slide_refuses():
    m = ManifestRecorder(deck="d.pptx", theme="t")
    with pytest.raises(PreconditionError):
        m.record_placement("s1.p1.card", "card", Rect(0.0, 0.0, 1.0, 1.0))
