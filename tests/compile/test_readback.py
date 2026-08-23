"""Reading a hand-edited deck back against the build that made it."""

from __future__ import annotations

import json

import pytest
from pptx.util import Inches

from pptxkit.compile.readback import read_back, render_drift
from pptxkit.errors import SpecError


@pytest.fixture
def built(project):
    """A real build: a deck, and the manifest that describes it."""
    from pptxkit.compile import build_deck

    (project / "d.deck.yaml").write_text(
        "theme: testtheme\nout: out/D.pptx\n---\ntitle: A title\n"
        "place:\n  - at: {cols: full}\n    bullets: {items: [alpha, beta]}\n"
    )
    return build_deck(project / "d.deck.yaml", theme_path=project / "testtheme.yaml")


def _edit(deck, change):
    from pptx import Presentation

    prs = Presentation(str(deck))
    change(prs.slides[0])
    prs.save(str(deck))


def _named(slide, name):
    return [sh for sh in slide.shapes if sh.name == name][0]


def test_an_untouched_deck_reads_back_as_unedited(built):
    drift = read_back(built.deck)
    assert drift.edited is False
    assert drift.changes == ()


def test_a_moved_shape_is_reported_with_both_rectangles(built):
    _edit(
        built.deck,
        lambda s: setattr(
            _named(s, "s1.chrome.title"), "left", _named(s, "s1.chrome.title").left + Inches(1)
        ),
    )
    drift = read_back(built.deck)
    assert drift.edited is True
    moved = drift.of("moved")
    assert [c.shape for c in moved] == ["s1.chrome.title"]
    assert "→" in moved[0].detail


def test_a_nudge_below_the_tolerance_is_not_a_move(built):
    """The manifest rounds to a thousandth; float dust is not an edit."""
    _edit(
        built.deck,
        lambda s: setattr(
            _named(s, "s1.chrome.title"), "left", _named(s, "s1.chrome.title").left + 900
        ),
    )
    assert read_back(built.deck).of("moved") == []


def test_retyped_text_is_reported_against_what_was_built(built):
    def retype(slide):
        _named(slide, "s1.chrome.title").text_frame.paragraphs[0].runs[0].text = "Other"

    _edit(built.deck, retype)
    retyped = read_back(built.deck).of("retyped")
    assert [c.shape for c in retyped] == ["s1.chrome.title"]
    assert "'A title' → 'Other'" in retyped[0].detail


def test_a_shape_added_by_hand_says_no_placement_made_it(built):
    _edit(built.deck, lambda s: s.shapes.add_textbox(Inches(1), Inches(6), Inches(2), Inches(0.4)))
    added = read_back(built.deck).of("added")
    assert len(added) == 1
    assert "added by hand" in added[0].detail


def test_a_shape_deleted_by_hand_is_reported_gone(built):
    def drop(slide):
        shape = _named(slide, "s1.p1.bullets#1")
        shape._element.getparent().remove(shape._element)

    _edit(built.deck, drop)
    assert [c.shape for c in read_back(built.deck).of("gone")] == ["s1.p1.bullets#1"]


def test_one_frame_answers_for_every_line_it_carries(project):
    """Chrome records one manifest row per line but draws one shape named `sN.chrome`, so
    matching from the manifest side reports every line missing."""
    from pptxkit.compile import build_deck

    (project / "c.deck.yaml").write_text(
        "theme: testtheme\nout: out/C.pptx\n---\nkicker: K\ntitle: T\nsubtitle: S\n"
    )
    made = build_deck(project / "c.deck.yaml", theme_path=project / "testtheme.yaml")
    drift = read_back(made.deck)
    assert drift.of("gone") == []
    assert drift.of("added") == []


def test_a_deck_with_no_manifest_says_how_to_get_one(built):
    built.manifest.unlink()
    with pytest.raises(SpecError, match="manifest not found"):
        read_back(built.deck)


def test_a_manifest_recording_no_deck_hash_cannot_call_it_edited(built):
    """An older manifest has nothing to compare, and guessing 'edited' would cry wolf."""
    data = json.loads(built.manifest.read_text())
    del data["deck_hash"]
    built.manifest.write_text(json.dumps(data))
    _edit(built.deck, lambda s: s.shapes.add_textbox(Inches(1), Inches(6), Inches(2), Inches(0.4)))
    assert read_back(built.deck).edited is False


def test_the_report_names_the_spec_to_carry_the_edit_back_into(built):
    _edit(
        built.deck,
        lambda s: setattr(
            _named(s, "s1.chrome.title"), "left", _named(s, "s1.chrome.title").left + Inches(1)
        ),
    )
    report = render_drift(read_back(built.deck))
    assert "d.deck.yaml" in report
    assert "**moved** `s1.chrome.title`" in report


def test_an_unedited_deck_reports_that_there_is_nothing_to_carry_back(built):
    assert "nothing to carry back" in render_drift(read_back(built.deck))
