from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from pptxkit.errors import ThemeError
from pptxkit.theme.media import resolve_media


@pytest.fixture(autouse=True)
def _isolated_cache(tmp_path, monkeypatch):
    """Redirect the module's cache dir under tmp_path so tests never write into the repo."""
    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path / "cache"))


@pytest.fixture
def template_with_media(tmp_path) -> Path:
    """A .pptx that genuinely embeds one picture, so it has a real ppt/media/ entry."""
    seed = tmp_path / "seed.png"
    Image.new("RGB", (8, 8), "red").save(seed)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(seed), Inches(0), Inches(0), Inches(1), Inches(1))

    template = tmp_path / "assets" / "t.pptx"
    template.parent.mkdir()
    prs.save(str(template))
    return template


def _embedded_media_name(template: Path) -> str:
    """The filename python-pptx assigned to the one embedded picture."""
    with zipfile.ZipFile(template) as archive:
        names = [n.rsplit("/", 1)[-1] for n in archive.namelist() if n.startswith("ppt/media/")]
    assert len(names) == 1
    return names[0]


def test_loose_file_beside_the_template_wins(template_with_media):
    name = _embedded_media_name(template_with_media)
    loose = template_with_media.parent / name
    loose.write_bytes(b"override-bytes")

    result = resolve_media(name, template=template_with_media)

    assert result == loose
    assert result.read_bytes() == b"override-bytes"


def test_media_absent_loosely_is_extracted_from_the_template(template_with_media):
    name = _embedded_media_name(template_with_media)

    result = resolve_media(name, template=template_with_media)

    assert result.is_file()
    assert result.stat().st_size > 0


def test_the_extracted_copy_lands_under_pptxkit_cache_dir(
    template_with_media, tmp_path, monkeypatch
):
    """The env var is the only way to move extracted art off the default relative
    ``.pptxkit-cache`` — an operator who redirects it must not still get the default."""
    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path / "elsewhere"))
    name = _embedded_media_name(template_with_media)

    result = resolve_media(name, template=template_with_media)

    assert result.parent.parent == tmp_path / "elsewhere" / "media"
    assert result.name == name
    assert result.is_file()


def test_a_second_call_reuses_the_cache(template_with_media):
    name = _embedded_media_name(template_with_media)
    first = resolve_media(name, template=template_with_media)
    sentinel = b"cached-sentinel"
    first.write_bytes(sentinel)

    second = resolve_media(name, template=template_with_media)

    assert second == first
    assert second.read_bytes() == sentinel


def test_name_in_neither_place_raises_theme_error(template_with_media):
    with pytest.raises(ThemeError, match="was not found in .*, nor inside "):
        resolve_media("nonexistent.png", template=template_with_media)


def test_a_name_with_nowhere_to_look_raises_theme_error():
    """Theme.template is optional, so None reaches here; it must not be an AttributeError."""
    with pytest.raises(ThemeError, match="the theme names no template to fall back on"):
        resolve_media("logo.png")


def test_a_climbing_name_is_refused_before_anything_is_written(tmp_path, monkeypatch):
    """The archive really carries a traversing member, so the extract this refuses is
    the one that would write ``payload`` outside the cache."""
    evil = tmp_path / "evil.pptx"
    with zipfile.ZipFile(evil, "w") as archive:
        archive.writestr(zipfile.ZipInfo("ppt/media/../../escaped.txt"), b"payload")
    with zipfile.ZipFile(evil) as archive:
        assert "ppt/media/../../escaped.txt" in archive.namelist()

    monkeypatch.setenv("PPTXKIT_CACHE_DIR", str(tmp_path / "cache" / "deep" / "deeper"))
    before = {p for p in tmp_path.rglob("*") if p.is_file()}

    with pytest.raises(ThemeError, match="climbs out"):
        resolve_media("../../escaped.txt", template=evil)

    assert {p for p in tmp_path.rglob("*") if p.is_file()} == before


def test_a_climbing_name_does_not_read_a_file_beside_the_template(template_with_media):
    """The guard is in resolve_media, not the template branch: a ``..`` name must be
    refused before the search roots can hand back what it points at."""
    secret = template_with_media.parent.parent / "secret.txt"
    secret.write_bytes(b"not-yours")

    with pytest.raises(ThemeError, match="climbs out"):
        resolve_media("../secret.txt", template=template_with_media, roots=())


def test_a_name_in_a_subdirectory_of_a_root_still_resolves(tmp_path):
    """Only ``..`` is refused — a name may still carry directories."""
    photo = tmp_path / "art" / "photo.png"
    photo.parent.mkdir()
    photo.write_bytes(b"art-bytes")

    result = resolve_media("art/photo.png", roots=[tmp_path])

    assert result == photo


def _rewrite_media_member(template: Path, member_name: str, new_bytes: bytes) -> Path:
    """Replace one ``ppt/media/*`` member's bytes **at the same path**: editing in place is the
    case a filename-keyed cache cannot survive, and writing a sibling would prove nothing."""
    out = template
    with zipfile.ZipFile(template) as src:
        members = [
            (
                item,
                new_bytes
                if item.filename == f"ppt/media/{member_name}"
                else src.read(item.filename),
            )
            for item in src.infolist()
        ]
    with zipfile.ZipFile(out, "w") as dst:
        for item, data in members:
            dst.writestr(item, data)
    return out


def test_a_changed_template_yields_the_new_bytes_not_the_stale_cache(template_with_media):
    name = _embedded_media_name(template_with_media)
    first = resolve_media(name, template=template_with_media)
    original_bytes = first.read_bytes()

    edited_template = _rewrite_media_member(template_with_media, name, b"brand-new-art-bytes")
    assert edited_template == template_with_media  # edited in place, same path and stem
    second = resolve_media(name, template=edited_template)

    assert second.read_bytes() == b"brand-new-art-bytes"
    assert second.read_bytes() != original_bytes
