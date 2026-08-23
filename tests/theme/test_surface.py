"""What the template already paints behind a slide. Every resolution here is invisible to a corpus
build: most sample templates reference a fill style that resolves to white, which is also what a
slide shows when nothing resolves at all."""

from __future__ import annotations

from lxml import etree
from PIL import Image
from pptx import Presentation

from pptxkit.layouts.resolve import pick_compose_layout
from pptxkit.theme.surface import inherited_surface

_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Three distinguishable fills, so an index off by one names a different colour rather
# than the same white every entry of a stock template's list resolves to.
_STYLES = """
<a:bgFillStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:solidFill><a:srgbClr val="110000"/></a:solidFill>
  <a:solidFill><a:srgbClr val="220000"/></a:solidFill>
  <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
</a:bgFillStyleLst>
"""


def _template(tmp_path, bg_xml: str, *, styles: str = _STYLES):
    """A .pptx whose master carries ``bg_xml`` and whose theme carries ``styles``."""
    prs = Presentation()
    master = prs.slide_masters[0]
    csld = master._element.find(f"{_P}cSld")
    existing = csld.find(f"{_P}bg")
    if existing is not None:
        csld.remove(existing)
    csld.insert(0, etree.fromstring(bg_xml))

    part = master.part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    root = etree.fromstring(part.blob)
    fmt = root.find(f".//{_A}fmtScheme")
    old = fmt.find(f"{_A}bgFillStyleLst")
    fmt.replace(old, etree.fromstring(styles))
    part._blob = etree.tostring(root)

    path = tmp_path / "surfaced.pptx"
    prs.save(str(path))
    return path


def _resolved(tmp_path, bg_xml: str, **kwargs):
    path = _template(tmp_path, bg_xml, **kwargs)
    return inherited_surface(pick_compose_layout(Presentation(str(path))))


def test_a_master_that_paints_nothing_leaves_no_surface(tmp_path):
    prs = Presentation()
    csld = prs.slide_masters[0]._element.find(f"{_P}cSld")
    bg = csld.find(f"{_P}bg")
    if bg is not None:
        csld.remove(bg)
    path = tmp_path / "bare.pptx"
    prs.save(str(path))
    assert inherited_surface(pick_compose_layout(Presentation(str(path)))) is None


def test_a_direct_solid_fill_is_the_surface(tmp_path):
    surface = _resolved(
        tmp_path,
        f"""
        <p:bg xmlns:p="{_P[1:-1]}" xmlns:a="{_A[1:-1]}">
          <p:bgPr><a:solidFill><a:srgbClr val="C0FFEE"/></a:solidFill></p:bgPr>
        </p:bg>""",
    )
    assert surface.flat == "C0FFEE"


def test_a_bgref_selects_the_style_its_index_names(tmp_path):
    """``idx="1002"`` is the *second* background fill style, not the second overall."""
    surface = _resolved(
        tmp_path,
        f"""
        <p:bg xmlns:p="{_P[1:-1]}" xmlns:a="{_A[1:-1]}">
          <p:bgRef idx="1002"><a:srgbClr val="FFFFFF"/></p:bgRef>
        </p:bg>""",
    )
    assert surface.flat == "220000"


def test_a_bgrefs_own_colour_stands_in_for_phclr(tmp_path):
    surface = _resolved(
        tmp_path,
        f"""
        <p:bg xmlns:p="{_P[1:-1]}" xmlns:a="{_A[1:-1]}">
          <p:bgRef idx="1003"><a:srgbClr val="ABCDEF"/></p:bgRef>
        </p:bg>""",
    )
    assert surface.flat == "ABCDEF"


def test_a_shade_on_the_fills_colour_darkens_the_surface(tmp_path):
    """50% shade of FF8040 is 7F4020 — a transform dropped leaves the colour untouched."""
    surface = _resolved(
        tmp_path,
        f"""
        <p:bg xmlns:p="{_P[1:-1]}" xmlns:a="{_A[1:-1]}">
          <p:bgPr><a:solidFill>
            <a:srgbClr val="FF8040"><a:shade val="50000"/></a:srgbClr>
          </a:solidFill></p:bgPr>
        </p:bg>""",
    )
    assert surface.flat == "804020"


def test_a_gradient_surface_reports_every_stop_and_is_not_flat(tmp_path):
    surface = _resolved(
        tmp_path,
        f"""
        <p:bg xmlns:p="{_P[1:-1]}" xmlns:a="{_A[1:-1]}">
          <p:bgPr><a:gradFill><a:gsLst>
            <a:gs pos="0"><a:srgbClr val="000080"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="80C0FF"/></a:gs>
          </a:gsLst></a:gradFill></p:bgPr>
        </p:bg>""",
    )
    assert surface.fills == ("000080", "80C0FF")
    assert surface.flat is None


def test_a_scheme_colour_resolves_through_the_masters_clrmap(tmp_path):
    """``bg1`` is an alias the master remaps; reading it as ``lt1`` is a coin flip."""
    path = _template(
        tmp_path,
        f"""
        <p:bg xmlns:p="{_P[1:-1]}" xmlns:a="{_A[1:-1]}">
          <p:bgPr><a:solidFill><a:schemeClr val="bg1"/></a:solidFill></p:bgPr>
        </p:bg>""",
    )
    prs = Presentation(str(path))
    master = prs.slide_masters[0]
    master._element.find(f"{_P}clrMap").set("bg1", "dk2")
    remapped = tmp_path / "remapped.pptx"
    prs.save(str(remapped))

    scheme_dk2 = "1F497D"  # the stock theme's dk2; lt1, which bg1 usually names, is white
    surface = inherited_surface(pick_compose_layout(Presentation(str(remapped))))
    assert surface.flat == scheme_dk2


def test_a_picture_surface_names_the_media_it_stretches(tmp_path):
    source = tmp_path / "art.png"
    Image.new("RGB", (32, 18), (10, 20, 30)).save(source)

    prs = Presentation()
    master = prs.slide_masters[0]
    _, rid = master.part.get_or_add_image_part(str(source))

    csld = master._element.find(f"{_P}cSld")
    bg = csld.find(f"{_P}bg")
    if bg is not None:
        csld.remove(bg)
    csld.insert(
        0,
        etree.fromstring(f"""
        <p:bg xmlns:p="{_P[1:-1]}" xmlns:a="{_A[1:-1]}"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:bgPr><a:blipFill><a:blip r:embed="{rid}"/>
            <a:stretch><a:fillRect/></a:stretch></a:blipFill></p:bgPr>
        </p:bg>"""),
    )
    path = tmp_path / "pictured.pptx"
    prs.save(str(path))

    surface = inherited_surface(pick_compose_layout(Presentation(str(path))))
    assert surface.media.startswith("image")
    assert surface.flat is None
