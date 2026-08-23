import textwrap

import pytest

from pptxkit.errors import ThemeError
from pptxkit.theme import load_theme


def _write(tmp_path, template, body: str):
    (tmp_path / "assets").mkdir(exist_ok=True)
    dest = tmp_path / "assets" / "t.pptx"
    dest.write_bytes(template.read_bytes())
    path = tmp_path / "t.yaml"
    path.write_text(textwrap.dedent(body))
    return path


BASE = """
    name: testtheme
    template: assets/t.pptx
    bind:
      page: lt1
      ink: dk1
    type:
      ramp:
        title: {pt: 32, bold: true}
        body: {pt: 13.5}
      min_pt: 10.5
    scale:
      margin: {top: 4%, right: 4.5751%, bottom: 6.6667%, left: 4.6501%}
      columns: 12
      gutter: 1.35%
      body_top: 22.6667%
"""


def test_fonts_come_from_the_template(tmp_path, synthetic_template):
    theme = load_theme(_write(tmp_path, synthetic_template, BASE))
    assert theme.face == "Calibri"  # the stock Office theme's latin typeface
    assert theme.mono == "Courier New"  # default when unset


def test_slide_size_comes_from_the_template(tmp_path, synthetic_template):
    theme = load_theme(_write(tmp_path, synthetic_template, BASE))
    assert theme.grid.slide_w == pytest.approx(13.333, abs=1e-3)
    assert theme.grid.slide_h == pytest.approx(7.5, abs=1e-3)


def test_type_ramp_is_parsed_into_styles(tmp_path, synthetic_template):
    theme = load_theme(_write(tmp_path, synthetic_template, BASE))
    assert theme.style("title").size == pytest.approx(32)
    assert theme.style("title").bold is True
    assert theme.style("body").bold is False


def test_line_weight_defaults_when_the_type_block_omits_it(tmp_path, synthetic_template):
    theme = load_theme(_write(tmp_path, synthetic_template, BASE))
    assert theme.line_weight == pytest.approx(2.25)  # the 0.30 rung at 7.5in


def test_line_weight_is_read_from_the_type_block_in_points(tmp_path, synthetic_template):
    body = BASE.replace("      min_pt: 10.5\n", "      min_pt: 10.5\n      line_weight_pt: 3\n")
    theme = load_theme(_write(tmp_path, synthetic_template, body))
    assert theme.line_weight == pytest.approx(3.0)


def test_missing_template_is_rejected(tmp_path, synthetic_template):
    body = BASE.replace("template: assets/t.pptx", "template: assets/nope.pptx")
    with pytest.raises(ThemeError, match="template not found"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_template_that_is_not_a_readable_pptx_is_a_theme_error(tmp_path, synthetic_template):
    """Otherwise it reaches the CLI as a python-pptx traceback, naming neither the template nor
    the theme that points at it."""
    path = _write(tmp_path, synthetic_template, BASE)
    (tmp_path / "assets" / "t.pptx").write_bytes(b"PK\x03\x04 truncated")

    with pytest.raises(ThemeError, match=r"template .*t\.pptx is not a readable \.pptx"):
        load_theme(path)


def test_missing_theme_file_is_rejected(tmp_path):
    with pytest.raises(ThemeError, match="theme file not found"):
        load_theme(tmp_path / "absent.yaml")


def test_a_reserve_poly_is_read_as_canvas_fractions(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve:
      - name: logo-wedge
        poly: [{x: 100%, y: 72.27%}, {x: 100%, y: 100%}, {x: 82.5%, y: 100%}]
"""
    )
    theme = load_theme(_write(tmp_path, synthetic_template, body))
    assert len(theme.reserve) == 1
    assert theme.reserve[0].poly == ((1.0, 0.7227), (1.0, 1.0), (0.825, 1.0))


def test_a_reserve_poly_scales_onto_the_canvas_it_was_loaded_against(tmp_path, synthetic_template):
    """Fractions, not inches: the wedge covers the bottom-right corner at any size."""
    body = (
        BASE
        + """
    reserve:
      - name: logo-wedge
        poly: [{x: 100%, y: 72.27%}, {x: 100%, y: 100%}, {x: 82.5%, y: 100%}]
"""
    )
    theme = load_theme(_write(tmp_path, synthetic_template, body))
    box = theme.reserve[0].rect(theme.scale)
    assert box.right == pytest.approx(theme.scale.slide_w)
    assert box.bottom == pytest.approx(theme.scale.slide_h)


def test_a_reserve_entry_without_a_poly_is_rejected(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve:
      - name: bad
"""
    )
    with pytest.raises(ThemeError, match="needs a 'poly'"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_reserve_entry_that_is_not_a_mapping_is_rejected(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve: [logo-wedge]
"""
    )
    with pytest.raises(ThemeError, match=r"each 'reserve' entry is a mapping"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_reserve_block_that_is_not_a_list_is_rejected(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve: logo-wedge
"""
    )
    with pytest.raises(ThemeError, match=r"'reserve' is a list of regions"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_poly_point_missing_its_second_number_is_rejected(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve:
      - name: wedge
        poly: [{x: 100%, y: 72.27%}, {x: 100%}]
"""
    )
    with pytest.raises(
        ThemeError,
        match=r"reserved region 'wedge': every 'poly' point is an \{x, y\} "
        r"mapping in percents of the canvas, got \{'x': '100%'\}",
    ):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_flat_list_of_numbers_is_not_a_poly(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve:
      - name: wedge
        poly: [[100%, 72.27%], [82.5%, 100%]]
"""
    )
    with pytest.raises(ThemeError, match=r"a 'poly' point is keyed — write \{x: 78%, y: 0%\}"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_non_numeric_poly_point_is_rejected(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve:
      - name: wedge
        poly: [{x: left, y: top}, {x: right, y: bottom}]
"""
    )
    with pytest.raises(ThemeError, match=r"poly.x is a percent of the canvas, got 'left'"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_the_old_safe_zones_key_is_rejected_by_name(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    safe_zones:
      - name: wedge
        poly: [{x: 1333.3%, y: 542%}, {x: 1333.3%, y: 750%}, {x: 1100%, y: 750%}]
"""
    )
    with pytest.raises(ThemeError, match="'safe_zones' is gone — use 'reserve'"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_theme_hash_is_stable_and_content_sensitive(tmp_path, synthetic_template):
    a = load_theme(_write(tmp_path, synthetic_template, BASE))
    b = load_theme(_write(tmp_path, synthetic_template, BASE))
    c = load_theme(_write(tmp_path, synthetic_template, BASE.replace("pt: 32", "pt: 37.5")))
    assert a.hash == b.hash
    assert a.hash != c.hash


def test_malformed_yaml_is_rejected_as_a_theme_error(tmp_path, synthetic_template):
    with pytest.raises(ThemeError, match="invalid YAML"):
        load_theme(_write(tmp_path, synthetic_template, "name: t\ntemplate: [unclosed\n"))


def test_non_mapping_theme_file_is_rejected(tmp_path, synthetic_template):
    with pytest.raises(ThemeError, match="mapping at its top level"):
        load_theme(_write(tmp_path, synthetic_template, "- just\n- a list\n"))


def test_an_applies_to_on_a_reserved_region_names_its_replacement(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve:
      - name: wedge
        poly: [{x: 100%, y: 72.27%}, {x: 100%, y: 100%}, {x: 82.5%, y: 100%}]
        applies_to: [content]
"""
    )
    with pytest.raises(
        ThemeError,
        match=r"reserved region 'wedge': 'applies_to' is gone — a slide has "
        r"no layout to scope a region to; every region applies to "
        r"every slide",
    ):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_an_unknown_reserve_key_lists_what_is_accepted(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    reserve:
      - name: wedge
        poly: [{x: 100%, y: 72.27%}, {x: 100%, y: 100%}, {x: 82.5%, y: 100%}]
        colour: orange
"""
    )
    with pytest.raises(
        ThemeError,
        match=r"reserved region 'wedge': unknown key 'colour'; "
        r"known keys: name, poly",
    ):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_theme_hash_changes_when_the_template_changes(tmp_path, synthetic_template):
    from pptx import Presentation

    path = _write(tmp_path, synthetic_template, BASE)
    before = load_theme(path).hash
    prs = Presentation(str(tmp_path / "assets" / "t.pptx"))
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(tmp_path / "assets" / "t.pptx"))
    assert load_theme(path).hash != before


FACES = BASE.replace(
    "      min_pt: 10.5",
    "      face: Aptos\n      heading_face: Aptos Display\n      min_pt: 10.5",
).replace(
    "        title: {pt: 32, bold: true}",
    "        title: {pt: 32, bold: true, face: heading}\n"
    "        kicker: {pt: 12, face: Courier New}",
)


def test_a_declared_face_beats_the_templates_font_scheme(tmp_path, synthetic_template):
    """A template's fontScheme routinely lags its real typeface, so an explicit face wins."""
    theme = load_theme(_write(tmp_path, synthetic_template, FACES))
    assert theme.face == "Aptos"  # not the template's Calibri
    assert theme.heading_face == "Aptos Display"


def test_a_ramp_rung_renders_in_the_face_it_asks_for(tmp_path, synthetic_template):
    theme = load_theme(_write(tmp_path, synthetic_template, FACES))
    assert theme.font_for(theme.style("title")) == "Aptos Display"  # face: heading
    assert theme.font_for(theme.style("body")) == "Aptos"  # names none -> body face
    assert theme.font_for(theme.style("kicker")) == "Courier New"  # a literal typeface


def test_the_faces_fall_back_to_the_template_when_undeclared(tmp_path, synthetic_template):
    """major is the display face, minor the body face — both Calibri in stock Office."""
    theme = load_theme(_write(tmp_path, synthetic_template, BASE))
    assert theme.face == "Calibri"
    assert theme.heading_face == "Calibri"
    assert theme.font_for(theme.style("title")) == "Calibri"


def test_a_theme_may_omit_its_template_entirely(tmp_path):
    """The design system stands alone: a theme file needs no brand asset to load."""
    path = tmp_path / "bare.yaml"
    path.write_text("name: bare\n")
    theme = load_theme(path)
    assert theme.template is None
    assert theme.palette.accents  # the built-in ramp, not an empty one
    assert theme.style("body").size > 0


def test_a_templateless_theme_still_gets_the_whole_type_ramp(tmp_path):
    """A theme restates only what it moves, so an omitted ramp is the built-in one."""
    path = tmp_path / "bare.yaml"
    path.write_text("name: bare\n")
    ramp = load_theme(path).ramp
    assert {
        "kicker",
        "caption",
        "body",
        "lead",
        "head",
        "stat",
        "subtitle",
        "title",
        "display",
        "hero",
    } <= set(ramp)


def test_a_declared_rung_overrides_only_itself(tmp_path):
    path = tmp_path / "bare.yaml"
    path.write_text("name: bare\ntype:\n  ramp:\n    body: {pt: 22.5}\n")
    theme = load_theme(path)
    assert theme.style("body").size == pytest.approx(3.0 * theme.grid.slide_h)
    assert "title" in theme.ramp  # the rest of the ramp survives


def test_binding_without_a_template_is_rejected(tmp_path):
    path = tmp_path / "bare.yaml"
    path.write_text("name: bare\nbind:\n  accent-1: accent1\n")
    with pytest.raises(ThemeError, match="declares 'bind:' but no 'template:'"):
        load_theme(path)


def test_marks_without_a_template_are_rejected(tmp_path):
    """Mark media resolves beside the template or out of it, so a themeless theme has
    nowhere to look — say so at load, not with an AttributeError mid-build."""
    path = tmp_path / "bare.yaml"
    path.write_text("name: bare\nmarks:\n  inverse: {media: art.jpg}\n")
    with pytest.raises(ThemeError, match="declares 'marks:' but no 'template:'"):
        load_theme(path)


def test_a_mark_named_after_a_painted_backdrop_is_kept(tmp_path, synthetic_template):
    body = BASE + "    marks:\n      inverse: {media: art.jpg}\n"
    assert load_theme(_write(tmp_path, synthetic_template, body)).marks == {
        "inverse": {"media": "art.jpg"}
    }


def test_a_mark_naming_no_painted_backdrop_is_rejected(tmp_path, synthetic_template):
    """A mark nothing lays down is the silent drop this loader exists to reject."""
    body = BASE + "    marks:\n      wordmark: {media: logo.png, left: 10.55, top: 6.35}\n"
    with pytest.raises(ThemeError, match=r"mark 'wordmark' names no painted backdrop"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_chrome_treatment_declared_in_inches_is_rejected(tmp_path, synthetic_template):
    """Inch values written where the loader expects canvas percents leave the canvas."""
    body = (
        BASE
        + """
    chrome:
      title: {at: {box: {x: 62%, y: 60%, w: 1130%, h: 125%}}}
"""
    )
    with pytest.raises(ThemeError, match=r"box .* leaves the canvas"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_chrome_treatment_reaches_the_theme(tmp_path, synthetic_template):
    body = (
        BASE
        + """
    chrome:
      title: {at: {box: {x: 0%, y: 5%, w: 100%, h: 14%}}, align: center, anchor: bottom}
"""
    )
    theme = load_theme(_write(tmp_path, synthetic_template, body))
    assert theme.chrome["title"].align == "center"
    assert theme.chrome["title"].anchoring == "bottom"
    assert theme.chrome["title"].at == {"box": pytest.approx((0.0, 0.05, 1.0, 0.14))}


def test_a_theme_declaring_no_chrome_leaves_every_field_at_its_default(
    tmp_path, synthetic_template
):
    theme = load_theme(_write(tmp_path, synthetic_template, BASE))
    assert theme.chrome == {}


def test_a_chrome_field_the_vocabulary_does_not_name_is_rejected(tmp_path, synthetic_template):
    body = BASE + "    chrome:\n      eyebrow: {align: center}\n"
    with pytest.raises(ThemeError, match="unknown chrome field 'eyebrow'"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_chrome_align_outside_the_vocabulary_is_rejected(tmp_path, synthetic_template):
    body = BASE + "    chrome:\n      title: {align: justify}\n"
    with pytest.raises(ThemeError, match="align must be one of left, center, right"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_mark_whose_value_is_not_a_mapping_is_rejected_at_load(tmp_path, synthetic_template):
    """Caught at load, not mid-build from the compositor on whichever slide reaches it first."""
    body = BASE + "    marks:\n      inverse: art.jpg\n"
    with pytest.raises(ThemeError, match=r"mark 'inverse' needs a mapping with a 'media:'"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_mark_without_media_is_rejected_at_load(tmp_path, synthetic_template):
    body = BASE + "    marks:\n      inverse: {opacity: 0.5}\n"
    with pytest.raises(ThemeError, match=r"mark 'inverse' needs a mapping with a 'media:'"):
        load_theme(_write(tmp_path, synthetic_template, body))


def test_a_ramp_entry_naming_only_a_size_keeps_its_rungs_weight_and_face(tmp_path):
    """`title: {pt: 34}` resizes the title; it must not quietly un-bold it or strip the
    heading face — which is exactly what the packaged base ramp did to every deck."""
    path = tmp_path / "t.yaml"
    path.write_text("name: t\ntype:\n  ramp:\n    title: {pt: 34}\n    body: {pt: 14}\n")
    theme = load_theme(path)
    assert theme.style("title").bold is True
    assert theme.style("title").face == theme.heading_face
    assert theme.style("body").bold is False


def test_an_explicit_bold_false_still_wins(tmp_path):
    path = tmp_path / "t.yaml"
    path.write_text("name: t\ntype:\n  ramp:\n    title: {pt: 34, bold: false}\n")
    theme = load_theme(path)
    assert theme.style("title").bold is False


def test_a_bare_theme_name_loads_the_packaged_builtin(tmp_path, monkeypatch):
    """`pptxkit.load_theme("base")` is the advertised way in; a path was the only one."""
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path / "no-such-dir"))
    assert load_theme("base").name == "base"


def test_a_bare_name_prefers_the_theme_directory_over_the_packaged_builtin(tmp_path, monkeypatch):
    (tmp_path / "base.theme.yaml").write_text("name: local-override\n")
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path))
    assert load_theme("base").name == "local-override"


def test_an_unknown_name_names_the_directory_it_searched_and_the_remedy(tmp_path, monkeypatch):
    """'theme file not found: acme' named no directory, no env var and no way out."""
    monkeypatch.setenv("PPTXKIT_THEME_DIR", str(tmp_path))
    with pytest.raises(ThemeError) as excinfo:
        load_theme("acme")
    message = str(excinfo.value)
    assert "unknown theme 'acme'" in message
    assert str(tmp_path) in message
    assert "PPTXKIT_THEME_DIR" in message
    assert "pptxkit conform <brand>.pptx --adopt acme" in message
    assert "packaged: base" in message


def test_a_path_shaped_reference_is_reported_as_a_path_not_an_unknown_name(tmp_path):
    """A mistyped path must not be answered with 'onboard a brand template'."""
    with pytest.raises(ThemeError) as excinfo:
        load_theme(tmp_path / "absent.yaml")
    message = str(excinfo.value)
    assert "theme file not found" in message
    assert "unknown theme" not in message
