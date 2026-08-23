"""The OOXML a masked picture writes. Nothing reads back the preset, the adjust value or the
crop, and no renderer this project can drive agrees with PowerPoint about all three."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxkit.errors import LayoutError
from pptxkit.imagery.draw import place_picture
from pptxkit.imagery.fit import fit_image
from pptxkit.theme.model import Rect

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
EMU = 914400


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _place(slide, path, box, **kwargs):
    """A 2:1 source covering ``box``, masked as asked."""
    fit = fit_image(source_aspect=2.0, box=box, fit=kwargs.pop("fit", "cover"))
    return place_picture(slide, str(path), fit, **kwargs)


def _geom(pic):
    return pic._element.spPr.find(A + "prstGeom")


def test_a_circle_mask_writes_the_ellipse_preset(slide, white_wide):
    """``prst`` is what clips the picture; ``ellipse`` inscribed in a square is a circle."""
    pic = _place(slide, white_wide, Rect(1.0, 1.0, 3.0, 3.0), mask="circle")
    assert _geom(pic).get("prst") == "ellipse"


def test_a_rounded_mask_writes_its_radius_as_the_preset_adjust_value(slide, white_wide):
    """OOXML states ``roundRect``'s corner as ``adj`` in 100000ths of the short side."""
    pic = _place(slide, white_wide, Rect(1.0, 1.0, 4.0, 2.0), mask="rounded", radius=0.25)
    guides = _geom(pic).find(A + "avLst").findall(A + "gd")
    assert [(g.get("name"), g.get("fmla")) for g in guides] == [("adj", "val 25000")]


def test_a_circle_mask_leaves_no_adjust_value_behind(slide, white_wide):
    """``ellipse`` takes no adjust; one inherited from the preset before it would stick."""
    pic = _place(slide, white_wide, Rect(1.0, 1.0, 3.0, 3.0), mask="circle")
    assert _geom(pic).find(A + "avLst").findall(A + "gd") == []


def test_a_masked_picture_is_cropped_rather_than_squashed(slide, white_wide):
    """A 2:1 source in a square box loses a quarter off each side, not half its width."""
    pic = _place(slide, white_wide, Rect(1.0, 1.0, 3.0, 3.0), mask="circle")
    src = pic._element.find(P + "blipFill").find(A + "srcRect")
    assert src.attrib == {"l": "25000", "r": "25000"}
    assert (pic.width, pic.height) == (Inches(3.0), Inches(3.0))


def test_an_unknown_mask_names_the_ones_that_exist(slide, white_wide):
    with pytest.raises(LayoutError, match="none, circle, rounded"):
        _place(slide, white_wide, Rect(1.0, 1.0, 3.0, 3.0), mask="hexagon")
