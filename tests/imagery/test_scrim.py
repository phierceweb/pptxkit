"""Which end of a gradient scrim is opaque, and what a malformed one is told. A built deck cannot
tell a ``top`` gradient from a ``bottom`` one — the stop alphas below are where that is written."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxkit.errors import LayoutError
from pptxkit.imagery.draw import paint_scrim
from pptxkit.imagery.scrim import Scrim, gradient_fraction, scrim_spec
from pptxkit.theme.model import Rect

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _stops(shape):
    """Each gradient stop as ``(position, alpha)``, in document order."""
    fill = shape._element.spPr.find(A + "gradFill")
    return [
        (stop.get("pos"), stop.find(A + "srgbClr").find(A + "alpha").get("val"))
        for stop in fill.find(A + "gsLst")
    ]


def _lin(shape):
    return shape._element.spPr.find(A + "gradFill").find(A + "lin")


def test_a_top_gradient_is_opaque_at_the_first_stop(slide):
    """With the fill running downward, stop 0 is the top edge."""
    shape = paint_scrim(slide, Rect(0, 0, 4, 3), Scrim("101010", "FFFFFF", 0.8, "top"))
    assert _stops(shape) == [("0", "80000"), ("100000", "0")]


def test_a_bottom_gradient_is_opaque_at_the_last_stop(slide):
    shape = paint_scrim(slide, Rect(0, 0, 4, 3), Scrim("101010", "FFFFFF", 0.8, "bottom"))
    assert _stops(shape) == [("0", "0"), ("100000", "80000")]


def test_the_gradient_runs_down_the_shape(slide):
    """5400000 sixty-thousandths of a degree is 90 — straight down. Sideways is wrong."""
    shape = paint_scrim(slide, Rect(0, 0, 4, 3), Scrim("101010", "FFFFFF", 0.8, "top"))
    assert _lin(shape).get("ang") == "5400000"


def test_a_flat_scrim_is_a_solid_fill_at_its_opacity(slide):
    shape = paint_scrim(slide, Rect(0, 0, 4, 3), Scrim("101010", "FFFFFF", 0.4))
    solid = shape._element.spPr.find(A + "solidFill")
    assert solid.find(A + "srgbClr").get("val") == "101010"
    assert solid.find(A + "srgbClr").find(A + "alpha").get("val") == "40000"


def test_a_fully_transparent_scrim_draws_nothing(slide):
    """An invisible rectangle would still sit above the picture and take its clicks."""
    assert paint_scrim(slide, Rect(0, 0, 4, 3), Scrim("101010", "FFFFFF", 0.0)) is None
    assert len(slide.shapes) == 0


def test_a_bottom_gradient_is_weakest_at_a_bands_top_edge():
    """The text nearest the clear end decides how much scrim the band really gets."""
    assert gradient_fraction("bottom", band_top=0.7, band_bottom=0.95) == 0.7


def test_a_top_gradient_is_weakest_at_a_bands_bottom_edge():
    assert gradient_fraction("top", band_top=0.05, band_bottom=0.3) == pytest.approx(0.7)


def test_a_flat_scrim_reaches_a_band_in_full():
    assert gradient_fraction("none", band_top=0.7, band_bottom=0.95) == 1.0


def test_an_unknown_gradient_names_the_ones_that_exist():
    with pytest.raises(LayoutError, match="none, top, bottom"):
        gradient_fraction("sideways", band_top=0.0, band_bottom=1.0)


def test_a_scrim_that_is_not_a_mapping_is_refused():
    with pytest.raises(LayoutError, match="'scrim' must be a mapping"):
        scrim_spec("heavy", default_pair="inverse", where="slide 1")


def test_a_scrim_key_that_does_not_exist_names_the_ones_that_do():
    with pytest.raises(LayoutError, match="scrim has no key 'colour'"):
        scrim_spec({"colour": "black"}, default_pair="inverse", where="slide 1")


def test_a_scrim_opacity_outside_zero_to_one_is_refused():
    with pytest.raises(LayoutError, match="fraction 0..1"):
        scrim_spec({"opacity": 55}, default_pair="inverse", where="slide 1")


def test_a_bare_true_means_solve_the_opacity_from_the_pixels():
    assert scrim_spec(True, default_pair="inverse", where="slide 1").opacity is None


def test_the_word_auto_means_the_same():
    assert scrim_spec({"opacity": "auto"}, default_pair="inverse", where="slide 1").opacity is None
