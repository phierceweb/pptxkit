"""Prove the variance guard can still tell a designed template from an empty one: hand
`tests/test_templates.py`'s own assertions a blank `Presentation()` and fail if they accept it.

Three of them read `Theme.palette`, which pptxkit fills from its own defaults, and so passed.
This needs no brand template, so unlike that module it runs everywhere, always."""

from __future__ import annotations

import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxkit.conform.derive import derive
from pptxkit.theme.stock import is_stock_accent


@pytest.fixture(scope="module")
def blank(tmp_path_factory):
    """Stock Office, 16:9, and nothing else — the same file `synthetic_template` makes."""
    path = tmp_path_factory.mktemp("blank") / "blank.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    prs.save(str(path))
    return path


@pytest.fixture(scope="module")
def blank_bind(blank):
    return derive(blank).get("bind", {})


def test_a_blank_office_file_binds_no_brand_accent(blank_bind):
    """Stock accents say nothing about a brand, so `derive` rejects all six. Read `bind` here, never
    a loaded `Theme` — `palette.accents` answers with pptxkit's own four defaults for this file."""
    bound = [role for role in blank_bind if role.startswith("accent-")]
    assert bound == [], (
        f"a blank Office file bound {bound} — either derive stopped rejecting stock "
        f"accents, or the stock table no longer covers this Office generation"
    )


def test_every_accent_in_a_blank_file_is_one_microsoft_ships(blank):
    """Why the assertion above holds, stated over the file rather than the code path,
    so a `derive` that silently stopped filtering cannot make both agree."""
    from pptxkit.theme.clrscheme import parse_color_scheme, read_theme_xml

    scheme = parse_color_scheme(read_theme_xml(Presentation(str(blank)).slide_masters[0]))
    accents = {slot: scheme[slot] for slot in scheme if slot.startswith("accent")}
    assert accents, "the stock scheme declares no accents — this file changed shape"
    assert [s for s, v in sorted(accents.items()) if not is_stock_accent(v)] == []


def test_the_generated_sample_does_pass_that_same_gate(tmp_path):
    """The negative control: if a blank file and a designed one both fail, the gate is broken rather
    than discriminating."""
    from pptxkit.conform.sample import write_sample

    bind = derive(write_sample(tmp_path / "sample.pptx")).get("bind", {})
    bound = [role for role in bind if role.startswith("accent-")]
    assert bound == [f"accent-{i}" for i in range(1, 7)]


def test_a_loaded_theme_still_reports_accents_for_a_template_that_bound_none(blank):
    """Asserts the `Theme.palette` fallback EXISTS — correct behaviour for rendering, and wrong as
    evidence about a template."""
    from pptxkit.theme import load_theme

    theme = load_theme(None)
    assert theme.palette.accents, (
        "the default palette no longer carries accents — if that is deliberate, the "
        "corpus assertions may read palette.accents again"
    )


# --- the loop closed: the real assertions, handed the failure case ------------------


@pytest.fixture(scope="module")
def blank_built(blank, tmp_path_factory):
    """What `conform` would hand the corpus assertions for a blank template. Assembled rather than
    built: the assertions below read only the derived theme and the template beside it."""
    import yaml

    out = tmp_path_factory.mktemp("blank-conform")
    theme = derive(blank)
    theme["template"] = os.path.relpath(blank, out)
    theme_path = out / "blank.theme.yaml"
    theme_path.write_text(yaml.safe_dump(theme, sort_keys=False))

    class _Built:
        template = blank.name
        theme = theme_path

    return _Built()


def test_a_blank_office_file_offers_no_brand_accent_to_bind(blank_built):
    """What makes the accent assertions in `test_templates.py` non-vacuous: they read the
    template's own scheme, so a file carrying only stock colours has nothing for them to examine.
    Revert either to reading `Theme.palette` — which pptxkit fills from its own defaults — and this
    reddens, because the blank file would suddenly appear to own brand accents."""
    import tests.test_templates as corpus

    assert corpus._brand_accents_in_scheme(blank_built) == {}
    assert corpus._bound_accents(blank_built) == {}


def test_those_same_assertions_accept_the_generated_sample(tmp_path):
    """The other half: a designed template must offer brand accents and pass both assertions, or
    the check above is merely blind rather than discriminating."""
    import yaml

    import tests.test_templates as corpus
    from pptxkit.conform.sample import write_sample

    sample = write_sample(tmp_path / "sample.pptx")
    theme = derive(sample)
    theme["template"] = sample.name
    theme_path = tmp_path / "sample.theme.yaml"
    theme_path.write_text(yaml.safe_dump(theme, sort_keys=False))

    class _Built:
        template = sample.name
        theme = theme_path

    built = _Built()
    corpus.test_every_brand_accent_the_template_owns_is_bound(built)
    corpus.test_no_unedited_microsoft_accent_is_adopted_as_a_brand_colour(built)
