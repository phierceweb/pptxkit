"""The capability demo: the exercise catalogue projected onto a named theme."""

from __future__ import annotations

import pytest
from pptx import Presentation

from pptxkit.conform import demo as demo_mod
from pptxkit.conform.demo import demo
from pptxkit.errors import ThemeError


def test_an_unknown_theme_names_where_it_looked(tmp_path):
    """A successful build never reaches this, and nothing else guards it."""
    with pytest.raises(ThemeError, match="no theme named 'nonesuch'"):
        demo("nonesuch", tmp_path)


def test_the_default_catalogue_is_the_exercise_registry(tmp_path, monkeypatch):
    """The point of generating this deck: a capability added to the registry reaches
    it with nobody editing anything here."""
    monkeypatch.setattr(
        demo_mod, "EXERCISE", {"one": {"title": "First"}, "two": {"title": "Second"}}
    )
    deck = demo("base", tmp_path)
    assert len(Presentation(str(deck)).slides) == 2


def test_one_slide_per_exercise_in_catalogue_order(tmp_path):
    deck = demo(
        "base",
        tmp_path,
        exercises={
            "a": {"title": "Alpha"},
            "b": {"title": "Beta"},
            "c": {"title": "Gamma"},
        },
    )
    prs = Presentation(str(deck))
    titles = [
        "".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        for slide in prs.slides
    ]
    assert titles == ["Alpha", "Beta", "Gamma"]


def test_the_deck_is_named_for_the_theme_it_was_built_against(tmp_path):
    deck = demo("base", tmp_path, exercises={"a": {"title": "Alpha"}})
    assert deck.name == "base capabilities.pptx"


def test_the_words_are_written_beside_the_deck(tmp_path):
    """The catalogue is the whole registry; the content view is how you read it without opening
    PowerPoint."""
    deck = demo("base", tmp_path, exercises={"a": {"title": "Alpha"}})
    assert "### Alpha" in deck.with_suffix(".content.md").read_text()


def test_the_spec_that_built_it_is_kept(tmp_path):
    """Generated, so it is the one artifact saying what this deck was."""
    demo("base", tmp_path, exercises={"a": {"title": "Alpha"}})
    assert (tmp_path / "base capabilities.deck.yaml").is_file()
