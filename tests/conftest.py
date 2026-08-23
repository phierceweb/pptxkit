"""Shared fixtures. Tests never touch the real brand template."""

from __future__ import annotations

import json
import pathlib
import textwrap

import pytest
from pptx import Presentation
from pptx.util import Inches

import pptxkit.components  # noqa: F401 — registers the built-in components
from pptxkit.compile.manifest import ManifestRecorder
from pptxkit.layouts.components import registered_components
from pptxkit.layouts.place import content_rect
from pptxkit.layouts.registry import SlideCtx
from pptxkit.spec.model import Background, SlideSpec
from pptxkit.theme import Grid, Scale
from pptxkit.theme.chartstyle import ChartStyle
from pptxkit.theme.defaults import DEFAULT_PAIRS
from pptxkit.theme.model import Theme, TypeStyle
from pptxkit.theme.palette import build_palette


LEGACY_GLYPHS = (
    "arrow-down",
    "arrow-left",
    "arrow-right",
    "arrow-up",
    "bell",
    "bolt",
    "calendar",
    "chart-bar",
    "chart-line",
    "chart-pie",
    "check",
    "circle",
    "clock",
    "close",
    "cloud",
    "diamond",
    "document",
    "download",
    "eye",
    "flag",
    "folder",
    "gear",
    "globe",
    "grid",
    "heart",
    "info",
    "layers",
    "lightbulb",
    "list",
    "lock",
    "mail",
    "minus",
    "pin",
    "plus",
    "ring",
    "search",
    "shield",
    "square",
    "star",
    "target",
    "triangle",
    "upload",
    "user",
    "users",
    "warning",
)
"""Glyph names decks in the wild are written against. Every one has to keep resolving.

They are the library's oldest public surface and the only names guaranteed to be
spelled pptxkit's way rather than Material's, so they are what a rename or a
re-vendoring breaks first.
"""


@pytest.fixture
def legacy_glyphs() -> tuple[str, ...]:
    return LEGACY_GLYPHS


@pytest.fixture
def synthetic_template(tmp_path):
    """A 16:9 .pptx with the stock Office theme — stands in for a brand template."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    path = tmp_path / "synthetic.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def wide_template(tmp_path):
    """A 26.666x15in .pptx — exactly twice the standard canvas on both axes."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(26.666), Inches(15.0)
    path = tmp_path / "wide.pptx"
    prs.save(str(path))
    return path


@pytest.fixture(autouse=True)
def _isolated_registries():
    """Registrations made in a test never leak. ``_LOADED_EXTENSIONS`` is half of the same
    state — roll back only the registry and an extension reloads with its components gone."""
    from pptxkit.layouts import components as components_mod
    from pptxkit.layouts import registry as registry_mod

    saved = dict(components_mod._REGISTRY)
    loaded = set(registry_mod._LOADED_EXTENSIONS)
    try:
        yield
    finally:
        components_mod._REGISTRY.clear()
        components_mod._REGISTRY.update(saved)
        registry_mod._LOADED_EXTENSIONS.clear()
        registry_mod._LOADED_EXTENSIONS.update(loaded)


SCALE = Scale(13.333, 7.5)

RAMP = {
    "kicker": TypeStyle(14 / 7.5, SCALE, bold=True),
    "caption": TypeStyle(13 / 7.5, SCALE),
    "body": TypeStyle(11 / 7.5, SCALE),
    "lead": TypeStyle(19 / 7.5, SCALE),
    "head": TypeStyle(18 / 7.5, SCALE, bold=True),
    "stat": TypeStyle(36 / 7.5, SCALE, bold=True),
    "subtitle": TypeStyle(15 / 7.5, SCALE),
    "title": TypeStyle(32 / 7.5, SCALE, bold=True),
    "display": TypeStyle(46 / 7.5, SCALE, bold=True),
    "hero": TypeStyle(56 / 7.5, SCALE, bold=True),
}

CHART = ChartStyle()


FIXTURE_ROLES = {
    "page": "FFFFFF",
    "ink": "2D0937",
    "muted": "573C65",
    "line": "EDEDED",
    "surface": "F5F6F8",
    "surface-ink": "2D0937",
    "inverse": "2D0937",
    "inverse-ink": "FFFFFF",
    "accent-1": "27B94C",
    "accent-2": "18CEDA",
    "accent-3": "4DB6AC",
    "accent-4": "A78BD0",
}
PALETTE = build_palette(FIXTURE_ROLES, pairs=DEFAULT_PAIRS)


@pytest.fixture
def theme():
    return Theme(
        name="t",
        template=pathlib.Path("x.pptx"),
        drop_template_slides=False,
        palette=PALETTE,
        scale=SCALE,
        face="Calibri",
        mono="Consolas",
        ramp=RAMP,
        min_pt=10.5,
        grid=Grid(
            scale=SCALE,
            top_frac=0.30 / 7.5,
            right_frac=0.61 / 13.333,
            bottom_frac=0.5 / 7.5,
            left_frac=0.62 / 13.333,
            columns=12,
            rows=12,
            gutter_frac=0.18 / 13.333,
            body_top_frac=1.7 / 7.5,
        ),
        marks={},
        line_weight=2.25,
        chart=CHART,
    )


@pytest.fixture
def ctx_factory(theme):
    """A ctx positioned on one placement. ``content`` maps chrome fields and at most
    one component key."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    def make(
        content,
        *,
        section=None,
        sections=("One", "Two", "Three"),
        theme_override=None,
        animate=None,
        background="page",
        base=None,
    ):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        manifest = ManifestRecorder(deck="d", theme="t")
        manifest.begin_slide(1, background=background, section=section)
        name = next((k for k in content if k in registered_components()), None)
        spec = SlideSpec(
            index=1,
            section=section,
            animate=animate,
            background=Background(kind=background),
            title=content.get("title"),
            kicker=content.get("kicker"),
            subtitle=content.get("subtitle"),
        )
        live = theme_override or theme
        return SlideCtx(
            slide=slide,
            theme=live,
            spec=spec,
            manifest=manifest,
            sections=sections,
            component=name,
            body=(content.get(name) or {}) if name else {},
            rect=content_rect(grid=live.grid),
            # The deck spec's own directory. A real build always sets it, so a
            # path resolved against `media_roots` is untestable while it is None.
            base=base,
        )

    return make


_QA_THEME = """
    name: testtheme
    template: assets/t.pptx
    bind: {page: lt1, ink: dk1, inverse: dk1, line: lt2}
    type:
      min_pt: 10.5
      ramp:
        kicker: {pt: 14, bold: true}
        caption: {pt: 12}
        body: {pt: 13.5}
        lead: {pt: 19}
        head: {pt: 18, bold: true}
        stat: {pt: 30, bold: true}
        subtitle: {pt: 15}
        title: {pt: 32, bold: true}
        display: {pt: 46, bold: true}
        hero: {pt: 52, bold: true}
    scale:
      margin: {top: 4%, right: 4.5751%, bottom: 6.6667%, left: 4.6501%}
      columns: 12
      gutter: 1.35%
      body_top: 22.6667%
    reserve:
      - name: logo-wedge
        poly: [{x: 100%, y: 72.2667%}, {x: 100%, y: 100%}, {x: 82.5%, y: 100%}]
"""


@pytest.fixture
def theme_file(tmp_path, synthetic_template):
    """A loadable theme YAML with its template beside it."""
    (tmp_path / "assets").mkdir(exist_ok=True)
    (tmp_path / "assets" / "t.pptx").write_bytes(synthetic_template.read_bytes())
    path = tmp_path / "testtheme.yaml"
    path.write_text(textwrap.dedent(_QA_THEME))
    return path


def save_blank_deck(path):
    """A real, openable one-slide .pptx at ``path`` — QA's package check reads the saved file,
    so a zero-byte placeholder is a finding of its own."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    return path


def _qa_deck(tmp_path, theme_file, *, box):
    deck = save_blank_deck(tmp_path / "d.pptx")
    manifest = tmp_path / "d.manifest.json"
    manifest.write_text(
        json.dumps(
            {
                "deck": "d.pptx",
                "theme": "testtheme",
                "theme_path": str(theme_file),
                "slide_w": 13.333,
                "slide_h": 7.5,
                "slides": [
                    {
                        "index": 1,
                        "background": "page",
                        "animations": [],
                        "shapes": [
                            {
                                "shape_id": 2,
                                "name": "Box",
                                "box": dict(zip("xywh", box, strict=True)),
                                "text": "hello",
                                "lines": ["hello"],
                                "font_pt": 13.5,
                                "fg": "2D0937",
                                "bg": "FFFFFF",
                                "rendered": "native",
                            }
                        ],
                    }
                ],
            }
        )
    )
    return deck, manifest


@pytest.fixture
def clean_manifest_deck(tmp_path, theme_file):
    return _qa_deck(tmp_path, theme_file, box=[1.0, 1.0, 3.0, 1.0])


@pytest.fixture
def dirty_manifest_deck(tmp_path, theme_file):
    return _qa_deck(tmp_path, theme_file, box=[12.0, 1.0, 3.0, 1.0])


#: A checkout has bin/, CLAUDE.md and the gitignored content directories; an sdist
#: ships the package, the suite, the docs and the examples and nothing else. Tests
#: that assert on the *repository's* layout have nothing to check in the second case.
IS_CHECKOUT = (pathlib.Path(__file__).resolve().parents[1] / "bin").is_dir()


def pytest_configure(config):
    config.addinivalue_line(
        "markers", "repo: asserts on the repository layout; skipped outside a checkout"
    )


def pytest_collection_modifyitems(config, items):
    if IS_CHECKOUT:
        return
    skip = pytest.mark.skip(reason="not a checkout — running from an sdist")
    for item in items:
        if "repo" in item.keywords:
            item.add_marker(skip)
